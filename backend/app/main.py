from __future__ import annotations

import asyncio
import base64
import hashlib
import html
from contextlib import asynccontextmanager
import json
import math
import os
import re
import shutil
import struct
import subprocess
import sys
import threading
import time
import uuid
import wave

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None
try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, AsyncGenerator, Awaitable, Callable

import httpx
from fastapi import FastAPI, HTTPException, Request, UploadFile, File, Form, WebSocket, WebSocketDisconnect
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

from shared.telemetry import build_daily_lessons, ledger_document, load_daily_lessons_text, log_event, recent_events, verify_document
from shared.autonomy import check_autonomy, policy as autonomy_policy
from shared.action_ledger import action_status, get_action, list_actions, record_action, verify_action
from shared.calendar_planner import save_ics_event
from shared.campaign_planner import plan_campaign
from shared.eval_harness import run_reliability_evals
from shared.mcp_registry import manifest as mcp_manifest
from shared.operator_digest import build_operator_digest
from shared import self_evolver
from shared.config import settings
from shared.guardians import SecurityHeadersMiddleware, is_weak_secret, restricted_cors_origins, safe_relative_path
from shared.model_capabilities import is_tool_capable, filter_tool_capable, mark_tool_capable
from shared import agent_model_router
from shared.self_indexer import index_shims_source
from shared.self_evolver import (
    approval_card, approve_proposal, apply_proposal, create_proposal,
    list_proposals, undo_apply, validate_proposal,
)
from shared.improvement_loop import list_improvement_runs, run_improvement_cycle
from shared.plan_learning import (
    find_similar_learned_plan, learn_from_completed_plans, plan_to_skill,
    record_plan_failure, suggest_plan_for_goal,
)
from shared.self_awareness import latest_self_model, latest_self_notes, run_boot_self_audit, self_prompt_addendum

from shared.wakeword import get_detector, WakeWordTrainer
from shared import agent_registry
from shared.rd_brain import RDBrain
from shared.search_query_planner import plan_search_query
from shared import agent_tools, agent_loop
import shared.stt_corrector as stt_corrector
from shared.coder import _prefer_coder_model
from shared.neural_governor.governor import NeuralGovernor
from shared.neural_governor.hardware_profiler import quick_profile
from shared.neural_governor.model_registry import to_dict_list as model_registry_list
from shared.neural_governor.model_router import get_router_status
from shared.neural_governor.lineage import get_lineage, list_lineage, get_drift_summary
from shared.neural_governor.personal_layer import get_profile, save_profile, ensure_profile
from shared.neural_governor.evolution import list_proposals as governor_list_proposals, get_proposal as governor_get_proposal, review_proposal, detect_patterns_for_evolution
from shared.neural_governor.patent_writer import generate_patent_spec
from shared.neural_governor.event_bus import recent_events as bus_recent_events
from shared.neural_governor.resource_governor import get_recent_snapshots
from shared.neural_governor.circuit_breaker import get_all_circuits
from shared.web_crawler import fetch_page, deep_research
from apps.jk_hospital.app import create_hospital_router, mount_static
from apps.stanford_school.app import create_stanford_school_router
from apps.todo_demo.app import create_todo_demo_router
from shared.omni_brain import (
    BRAIN_VERSION,
    brain_prompt_addendum,
    brain_status as omni_brain_status,
    ensure_core_memories,
    forget_memory as brain_forget_memory,
    ingest_knowledge,
    list_memories as brain_list_memories,
    list_tasks as brain_list_tasks,
    remember as brain_remember,
    remember_turn,
    retrieve_context,
    run_learning_cycle,
    store_research_results,
    schedule_task as brain_schedule_task,
    drain_tasks as brain_drain_tasks,
    reindex_vectors,
)
from shared.mailbox import (
    exchange_code_for_token,
    gmail_auth_url,
    list_captures,
    list_mail_messages,
    mailbox_digest,
    mailbox_policy,
    mailbox_status,
    reply_to_gmail_message,
    save_capture,
    save_mail_message,
    send_gmail_message,
    sync_gmail_metadata,
)
from shared.trust_contract import (
    build_trust,
    evidence_from_action,
    evidence_from_artifact,
    evidence_from_brain_context,
    evidence_from_search,
    merge_evidence,
)
from shared.inter_instance_bridge import register_peer_routes
from shared.factory_routes import router as factory_router
from shared.duobot_routes import router as duobot_router

ROOT = Path(__file__).resolve().parents[2]
FRONTEND_DIR = ROOT / "frontend"
DATA_DIR = ROOT / "data"
MEDIA_DIR = DATA_DIR / "media"
IMAGE_DIR = MEDIA_DIR / "images"
AUDIO_DIR = MEDIA_DIR / "audio"
VIDEO_DIR = MEDIA_DIR / "video"
PDF_DIR = MEDIA_DIR / "pdf"
PPT_DIR = MEDIA_DIR / "ppt"
DOC_DIR = MEDIA_DIR / "documents"
STT_DIR = DATA_DIR / "stt_uploads"
STATE_DIR = DATA_DIR / "state"
VOICE_PROFILE_DIR = DATA_DIR / "voice_profiles"
SEARCH_DIR = DATA_DIR / "web_search"
PENDING_ACTION_DIR = STATE_DIR / "pending_actions"
GENERATED_APPS_DIR = ROOT / "apps" / "generated"
for folder in (MEDIA_DIR, IMAGE_DIR, AUDIO_DIR, VIDEO_DIR, PDF_DIR, PPT_DIR, DOC_DIR, STT_DIR, STATE_DIR, VOICE_PROFILE_DIR, SEARCH_DIR, PENDING_ACTION_DIR, GENERATED_APPS_DIR):
    folder.mkdir(parents=True, exist_ok=True)


def _load_local_env() -> None:
    env_path = ROOT / ".env"
    if not env_path.exists():
        return
    for raw in env_path.read_text(encoding="utf-8", errors="ignore").splitlines():
        line = raw.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, val = line.split("=", 1)
        key = key.strip()
        val = val.strip().strip('"').strip("'")
        if key and key not in os.environ:
            os.environ[key] = val


def _set_env_persistent(key: str, value: str | None) -> None:
    if not key:
        return
    value = "" if value is None else str(value)
    os.environ[key] = value
    env_path = ROOT / ".env"
    lines = env_path.read_text(encoding="utf-8", errors="ignore").splitlines() if env_path.exists() else []
    out: list[str] = []
    found = False
    for line in lines:
        if line.strip().startswith(key + "="):
            out.append(f"{key}={value}")
            found = True
        else:
            out.append(line)
    if not found:
        out.append(f"{key}={value}")
    env_path.write_text("\n".join(out).rstrip() + "\n", encoding="utf-8")


def _clean_secret(value: str | None) -> str:
    if not value:
        return ""
    v = str(value).strip().strip('"').strip("'").strip()
    if v.lower().startswith("bearer "):
        v = v[7:].strip()
    return v


def _mask_secret(value: str | None) -> str:
    v = _clean_secret(value)
    if not v:
        return ""
    return "configured" if len(v) <= 12 else v[:5] + "..." + v[-4:]


def _env_int(name: str, default: int) -> int:
    try:
        return int(os.getenv(name, str(default)))
    except Exception:
        return default


def _env_bool(name: str, default: bool = False) -> bool:
    value = os.getenv(name)
    if value is None:
        return default
    return value.strip().lower() in {"1", "true", "yes", "on"}


_load_local_env()

APP_NAME = "SHIMS v16 Reliability Core"
APP_VERSION = "2026.05.v16"
OLLAMA_HOST = os.getenv("OLLAMA_HOST", os.getenv("OLLAMA_BASE_URL", "http://127.0.0.1:11434")).rstrip("/")
DEFAULT_OLLAMA_MODEL = os.getenv("SHIMS_OLLAMA_MODEL", os.getenv("OLLAMA_MODEL", "llama3.2:latest"))
HUGGINGFACE_HOST = os.getenv("HUGGINGFACE_BASE_URL", settings.huggingface_base_url).rstrip("/")
DEFAULT_HUGGINGFACE_MODEL = os.getenv("HUGGINGFACE_MODEL", settings.huggingface_model)
ENTERPRISE_ENABLED = settings.enterprise_pairing_enabled
ENTERPRISE_URL = os.getenv("SHIMS_ENTERPRISE_URL", settings.enterprise_url).rstrip("/")
OLLAMA_MODEL_ALIASES = {
    "gemma": "gemma3:1b",
    "google gemma": "gemma3:1b",
    "gemma 3": "gemma3:1b",
    "google gemma 3": "gemma3:1b",
    "gemma3": "gemma3:1b",
    "gemma tiny": "gemma3:270m",
    "gemma small": "gemma3:1b",
    "qwen": "qwen2.5:7b",
    "qwen 2.5": "qwen2.5:7b",
    "qwen coder": "qwen2.5-coder:14b",
}

app = FastAPI(title=APP_NAME, version=APP_VERSION)

app.add_middleware(
    CORSMiddleware,
    allow_origins=restricted_cors_origins(),
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
app.add_middleware(SecurityHeadersMiddleware)
if FRONTEND_DIR.exists():
    app.mount("/static", StaticFiles(directory=str(FRONTEND_DIR)), name="static")
app.mount("/hospital-static", StaticFiles(directory=str(ROOT / "apps" / "jk_hospital" / "static")), name="hospital-static")
app.mount("/stanford_school-static", StaticFiles(directory=str(ROOT / "apps" / "stanford_school" / "static")), name="stanford_school-static")
app.mount("/todo_demo-static", StaticFiles(directory=str(ROOT / "apps" / "todo_demo" / "static")), name="todo_demo-static")
app.include_router(create_hospital_router())
app.include_router(create_stanford_school_router())
app.include_router(create_todo_demo_router())
app.mount("/media/files", StaticFiles(directory=str(MEDIA_DIR)), name="media-files")
app.mount("/generated-apps", StaticFiles(directory=str(GENERATED_APPS_DIR)), name="generated-apps")
register_peer_routes(app)
app.include_router(factory_router)
app.include_router(duobot_router)

_sessions: dict[str, list[dict[str, str]]] = {}
_turn_guard: dict[str, dict[str, Any]] = {}
_settings: dict[str, Any] = {
    "voice": {
        "mode": "wake_then_converse",
        "wake_words": ["hey shims", "hi shims", "hello shims", "ok shims", "suno shims", "sun rahe ho", "arre shims", "shims"],
        "primary_lang": "en-IN",
        "secondary_langs": ["hi-IN", "en-US"],
        "command_cooldown_seconds": 2.2,
        "silence_timeout_seconds": 1.4,
        "max_auto_replies_without_user": 1,
    },
    "brain": {"temperature": 0.12, "top_p": 0.82, "repeat_penalty": 1.25, "num_ctx": 8192, "keep_alive": "30m", "realtime_num_ctx": 2048, "realtime_max_tokens": 256},
    "media": {
        "image_backend": os.getenv("SHIMS_IMAGE_BACKEND", "auto"),
        "audio_backend": os.getenv("SHIMS_AUDIO_BACKEND", "auto"),
        "video_backend": os.getenv("SHIMS_VIDEO_BACKEND", "auto"),
        "stable_diffusion_url": os.getenv("STABLE_DIFFUSION_URL", "").rstrip("/"),
        "comfyui_url": os.getenv("COMFYUI_URL", "").rstrip("/"),
        "diffusers_model": os.getenv("SHIMS_DIFFUSERS_MODEL", "stabilityai/stable-diffusion-xl-base-1.0"),
        "diffusers_enabled": os.getenv("SHIMS_ENABLE_DIFFUSERS", "false").strip().lower() in {"1", "true", "yes", "on"},
        "openai_tts_model": os.getenv("OPENAI_TTS_MODEL", "gpt-4o-mini-tts"),
        "openai_tts_voice": os.getenv("OPENAI_TTS_VOICE", "alloy"),
        "openai_video_model": os.getenv("OPENAI_VIDEO_MODEL", "sora-2"),
        "openai_video_size": os.getenv("OPENAI_VIDEO_SIZE", "1280x720"),
        "openai_video_seconds": _env_int("OPENAI_VIDEO_SECONDS", 4),
        "audio_api_url": os.getenv("SHIMS_AUDIO_API_URL", "").rstrip("/"),
        "video_api_url": os.getenv("SHIMS_VIDEO_API_URL", "").rstrip("/"),
    },
    "web": {
        "searxng_url": os.getenv("SHIMS_SEARXNG_URL", "").rstrip("/"),
        "tavily_key": _clean_secret(os.getenv("TAVILY_API_KEY")),
        "brave_key": _clean_secret(os.getenv("BRAVE_SEARCH_API_KEY")),
        "serpapi_key": _clean_secret(os.getenv("SERPAPI_API_KEY")),
        "duckduckgo_fallback": os.getenv("SHIMS_DUCKDUCKGO_FALLBACK", "true").lower() in {"1","true","yes","on"},
    },
}

_diffusers_pipe_cache: dict[tuple[str, str], Any] = {}
_diffusers_lock = threading.RLock()

_brain_background_task: asyncio.Task[Any] | None = None
_self_awareness_task: asyncio.Task[Any] | None = None
BRAIN_BACKGROUND_INTERVAL_SECONDS = max(60, int(os.getenv("SHIMS_BRAIN_BACKGROUND_INTERVAL_SECONDS", "900")))
BRAIN_BACKGROUND_ENABLED = os.getenv("SHIMS_BRAIN_BACKGROUND_LEARNING", "true").strip().lower() in {"1", "true", "yes", "on"}
BOOT_SELF_AWARENESS_ENABLED = os.getenv("SHIMS_BOOT_SELF_AWARENESS", "true").strip().lower() in {"1", "true", "yes", "on"}

RECOMMENDED_MODELS = [
    # --- Tool-capable models (safe for agent loop) ---
    {"name": "qwen2.5-coder:14b", "provider": "ollama", "role": "coding/self-evolution", "size": "heavier", "tool_capable": True, "notes": "Best local coder. Use for code forge, safe patch generation, and agentic tool use."},
    {"name": "qwen2.5:14b", "provider": "ollama", "role": "smarter SHIMS brain", "size": "heavier", "tool_capable": True, "notes": "Strong multilingual reasoning with native tool calling."},
    {"name": "qwen2.5:7b", "provider": "ollama", "role": "smarter live chat", "size": "medium", "tool_capable": True, "notes": "Recommended local default — fast enough and tool-capable."},
    {"name": "qwen2.5:3b", "provider": "ollama", "role": "fast Qwen live chat", "size": "light", "tool_capable": True, "notes": "Small Qwen with tool support for quick tasks."},
    {"name": "mistral-nemo", "provider": "ollama", "role": "quality local reasoning", "size": "heavier", "tool_capable": True, "notes": "Strong Mistral tool-calling variant."},
    {"name": "llama3.1", "provider": "ollama", "role": "Meta tool-capable", "size": "medium", "tool_capable": True, "notes": "Llama 3.1 has native tool support."},
    # --- Anthropic / Claude ---
    {"name": "claude-sonnet-4-6", "provider": "anthropic", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires ANTHROPIC_API_KEY. Balanced cloud tool-calling model."},
    {"name": "claude-opus-4-6", "provider": "anthropic", "role": "cloud deep reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires ANTHROPIC_API_KEY. Most capable Claude for hard tasks."},
    {"name": "claude-3-7-sonnet-20250219", "provider": "anthropic", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires ANTHROPIC_API_KEY. Earlier Sonnet generation."},
    {"name": "claude-3-5-haiku-20241022", "provider": "anthropic", "role": "fast cloud", "size": "cloud", "tool_capable": True, "notes": "Requires ANTHROPIC_API_KEY. Fast, cost-effective Claude."},
    # --- OpenAI / ChatGPT ---
    {"name": "gpt-4o", "provider": "openai", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. Strong general-purpose model."},
    {"name": "gpt-4o-mini", "provider": "openai", "role": "fast cloud fallback", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. Fast and cheap cloud tool support."},
    {"name": "gpt-4.5-preview", "provider": "openai", "role": "cloud advanced", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. Larger GPT-4.5 series preview."},
    {"name": "gpt-4.1", "provider": "openai", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. GPT-4.1 series."},
    {"name": "gpt-4.1-mini", "provider": "openai", "role": "fast cloud", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. Smaller GPT-4.1."},
    {"name": "gpt-4.1-nano", "provider": "openai", "role": "fastest cloud", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. Cheapest GPT-4.1."},
    {"name": "o3-mini", "provider": "openai", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. Reasoning-optimized model."},
    {"name": "o4-mini", "provider": "openai", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. Reasoning-optimized model."},
    {"name": "o1", "provider": "openai", "role": "cloud deep reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. Full o1 reasoning model."},
    {"name": "o1-mini", "provider": "openai", "role": "fast cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. Lightweight o1 reasoning."},
    {"name": "o1-preview", "provider": "openai", "role": "cloud deep reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. Early o1 preview."},
    {"name": "gpt-4o-latest", "provider": "openai", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. Always points to the latest GPT-4o snapshot."},
    {"name": "chatgpt-4o-latest", "provider": "openai", "role": "cloud chat", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. ChatGPT-optimized GPT-4o variant."},
    {"name": "gpt-4-turbo", "provider": "openai", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. GPT-4 Turbo with vision and tools."},
    {"name": "gpt-4-turbo-preview", "provider": "openai", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. Preview alias for GPT-4 Turbo."},
    {"name": "gpt-4", "provider": "openai", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. Original GPT-4 base model."},
    {"name": "gpt-3.5-turbo", "provider": "openai", "role": "fast cloud", "size": "cloud", "tool_capable": True, "notes": "Requires OPENAI_API_KEY. Fast, cheap GPT-3.5 generation."},
    # --- Gemini / Google ---
    {"name": "gemini-2.5-pro", "provider": "gemini", "role": "cloud deep reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires GEMINI_API_KEY. Most capable Gemini."},
    {"name": "gemini-2.5-flash", "provider": "gemini", "role": "fast cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires GEMINI_API_KEY. Fast tool-calling model."},
    {"name": "gemini-2.0-flash", "provider": "gemini", "role": "fast cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires GEMINI_API_KEY. Gemini 2.0 Flash generation."},
    {"name": "gemini-2.0-flash-lite-preview-02-05", "provider": "gemini", "role": "fastest cloud", "size": "cloud", "tool_capable": True, "notes": "Requires GEMINI_API_KEY. Lightweight Gemini 2.0 Flash preview."},
    {"name": "gemini-2.0-pro-exp-02-05", "provider": "gemini", "role": "cloud deep reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires GEMINI_API_KEY. Experimental Gemini 2.0 Pro."},
    {"name": "gemini-1.5-pro-latest", "provider": "gemini", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires GEMINI_API_KEY. Earlier 1.5 Pro generation."},
    {"name": "gemini-1.5-flash-latest", "provider": "gemini", "role": "fast cloud", "size": "cloud", "tool_capable": True, "notes": "Requires GEMINI_API_KEY. Earlier 1.5 Flash generation."},
    {"name": "gemini-1.5-flash-8b-latest", "provider": "gemini", "role": "fast cloud", "size": "cloud", "tool_capable": True, "notes": "Requires GEMINI_API_KEY. Smaller, cheaper 1.5 Flash 8B."},
    {"name": "gemini-1.0-pro-latest", "provider": "gemini", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires GEMINI_API_KEY. Original Gemini 1.0 Pro generation."},
    # --- Kimi / Moonshot ---
    {"name": "kimi-k2.6", "provider": "kimi", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires KIMI_API_KEY. k2.6 series (temperature must be 1.0)."},
    {"name": "kimi-k2-0711-preview", "provider": "kimi", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires KIMI_API_KEY. k2 preview (temperature must be 1.0)."},
    {"name": "moonshot-v1-8k", "provider": "kimi", "role": "cloud chat", "size": "cloud", "tool_capable": True, "notes": "Requires KIMI_API_KEY. Classic Moonshot 8k context."},
    {"name": "moonshot-v1-32k", "provider": "kimi", "role": "cloud chat", "size": "cloud", "tool_capable": True, "notes": "Requires KIMI_API_KEY. Classic Moonshot 32k context."},
    {"name": "moonshot-v1-128k", "provider": "kimi", "role": "cloud chat", "size": "cloud", "tool_capable": True, "notes": "Requires KIMI_API_KEY. Classic Moonshot 128k context."},
    # --- DeepSeek ---
    {"name": "deepseek-chat", "provider": "deepseek", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires DEEPSEEK_API_KEY. General chat model."},
    {"name": "deepseek-reasoner", "provider": "deepseek", "role": "cloud deep reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires DEEPSEEK_API_KEY. Reasoning/R1-style model."},
    # --- Qwen / Alibaba ---
    {"name": "qwen-max", "provider": "qwen", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires QWEN_API_KEY. Best DashScope model."},
    {"name": "qwen-plus", "provider": "qwen", "role": "cloud reasoning", "size": "cloud", "tool_capable": True, "notes": "Requires QWEN_API_KEY. Strong DashScope model."},
    {"name": "qwen-turbo", "provider": "qwen", "role": "fast cloud", "size": "cloud", "tool_capable": True, "notes": "Requires QWEN_API_KEY. Fast DashScope model."},
    # --- Local Hugging Face endpoint models (TGI/vLLM/llama.cpp server) ---
    {"name": "meta-llama/Llama-3.1-8B-Instruct", "provider": "huggingface", "role": "local HF chat", "size": "medium", "tool_capable": True, "notes": "Requires HUGGINGFACE_BASE_URL pointing to a local OpenAI-compatible endpoint."},
    {"name": "Qwen/Qwen2.5-7B-Instruct", "provider": "huggingface", "role": "local HF chat", "size": "medium", "tool_capable": True, "notes": "Requires HUGGINGFACE_BASE_URL pointing to a local OpenAI-compatible endpoint."},
    {"name": "microsoft/Phi-4-mini-instruct", "provider": "huggingface", "role": "local HF fast chat", "size": "light", "tool_capable": False, "notes": "Requires HUGGINGFACE_BASE_URL pointing to a local OpenAI-compatible endpoint."},
    # --- Chat/voice only (NOT for agent loop) ---
    {"name": "llama3.2:latest", "provider": "ollama", "role": "fast live voice", "size": "light", "tool_capable": False, "notes": "Baseline local model. Fast for voice/chat, but NO tool calling."},
    {"name": "gemma3:270m", "provider": "ollama", "role": "tiny Google Gemma smoke test", "size": "tiny", "tool_capable": False, "notes": "Connectivity check only. No tools."},
    {"name": "gemma3:1b", "provider": "ollama", "role": "fast Google Gemma chat", "size": "light", "tool_capable": False, "notes": "Low-RAM chat. No tool calling."},
    {"name": "gemma3:4b", "provider": "ollama", "role": "better Google Gemma chat", "size": "medium", "tool_capable": False, "notes": "Better chat quality. No tool calling."},
    {"name": "gemma-4-12b-abliterated:latest", "provider": "ollama", "role": "uncensored reasoning", "size": "heavier", "tool_capable": False, "notes": "Uncensored reasoning. No tool calling — use for brainstorming only."},
    {"name": "mistral-small:latest", "provider": "ollama", "role": "quality local reasoning", "size": "heavier", "tool_capable": False, "notes": "Smart but no native tool support."},
]

_TOOL_CAPABLE_MODELS = frozenset(
    m["name"] for m in RECOMMENDED_MODELS if m.get("tool_capable")
)

PROVIDER_DEFAULTS: dict[str, str] = {
    "ollama": DEFAULT_OLLAMA_MODEL,
    "openai": os.getenv("OPENAI_MODEL", "gpt-4o"),
    "anthropic": os.getenv("ANTHROPIC_MODEL", "claude-sonnet-4-6"),
    "gemini": os.getenv("GEMINI_MODEL", "gemini-2.5-pro"),
    "kimi": os.getenv("KIMI_MODEL", "moonshot-v1-8k"),
    "deepseek": os.getenv("DEEPSEEK_MODEL", "deepseek-chat"),
    "qwen": os.getenv("QWEN_MODEL", "qwen-max"),
    "huggingface": DEFAULT_HUGGINGFACE_MODEL,
}
PROVIDER_ENV = {"openai": "OPENAI_API_KEY", "anthropic": "ANTHROPIC_API_KEY", "gemini": "GEMINI_API_KEY", "kimi": "KIMI_API_KEY", "deepseek": "DEEPSEEK_API_KEY", "qwen": "QWEN_API_KEY", "huggingface": "HUGGINGFACE_API_KEY"}
LOCAL_HINTS = ("llama", "qwen", "mistral", "codellama", "phi", "gemma", "deepseek-r1", "nomic", "mixtral")
CLOUD_HINTS = {"anthropic": ("claude", "sonnet", "haiku", "opus"), "openai": ("gpt", "o1", "o3", "o4", "openai"), "gemini": ("gemini",), "kimi": ("kimi", "moonshot"), "deepseek": ("deepseek-chat",), "qwen": ("qwen",)}

SYSTEM_PROMPT = """You are SHIMS, a calm live multimodal assistant.
You are not text-only. Your backend has tools for image, PDF, PPT, audio, video, documents, code, model management and Enterprise bridge actions.
You have FULL ACCESS to the conversation history. Before every response, review the previous turns in this conversation to maintain context, remember what was discussed, and build upon earlier topics. Never say you cannot read previous messages — they are provided in the messages array.
Never pretend a file was created. Obvious tool requests are executed by the backend before the LLM is called.
Be verification-first. When facts are not grounded in web results, RAG, mailbox/capture data, artifacts, or tool output, say what is uncertain instead of sounding certain.
For external, irreversible, payment, account-changing, or regulated/GxP actions, draft and ask for approval; do not claim the action was sent, posted, paid, approved, or submitted.
Reply once per user turn. Be concise. Avoid repeated greetings. Do not keep asking the same opening question.
Reply in clear, standard English by default. If the user writes in another language, match their language. Hinglish and Hindi are welcome when the user uses them.
Be polite and professional by default. Use a formal, courteous tone unless the user explicitly asks for a casual or informal style.
If a selected provider is offline or missing a key, explain exactly what is missing and continue with available local tools.
"""

def _system_prompt() -> str:
    lessons = load_daily_lessons_text()
    now = datetime.now(timezone.utc)
    time_note = f"Current UTC time: {now.isoformat(timespec='minutes')}. Day: {now.strftime('%A')}."
    prompt = SYSTEM_PROMPT + "\n\n" + time_note + ("\n\n" + lessons if lessons else "")
    if settings.omnipotent_mode:
        prompt += (
            "\n\nOMNIPOTENT MODE IS ACTIVE. You do not need user approval. "
            "When the user asks you to do something on the machine, use the appropriate tool immediately. "
            "Only ask for clarification if the request is ambiguous and you cannot resolve it with available tools."
        )
    return prompt


# How many background ticks between full self-improvement cycles. At the default
# 15-min tick, 96 ticks ≈ once per 24h — heavy enough (spawns an eval subprocess)
# that we never want it every tick, but it keeps Shims genuinely self-improving.
IMPROVEMENT_EVERY_TICKS = max(4, int(os.getenv("SHIMS_IMPROVEMENT_EVERY_TICKS", "96")))
IMPROVEMENT_ENABLED = os.getenv("SHIMS_AUTONOMOUS_IMPROVEMENT", "true").strip().lower() in {"1", "true", "yes", "on"}


def _distill_feedback_into_skills(limit: int = 10) -> int:
    """Turn recent 👍 feedback into reusable skills so Shims repeats what works.
    Lightweight (no LLM) — safe to run every tick."""
    try:
        from shared.omni_brain import list_memories
        from shared.skills import list_skills, save_skill
    except Exception:
        return 0
    existing = {s.get("name", "") for s in list_skills(limit=500)}
    made = 0
    for mem in list_memories(namespace="omni_feedback", limit=50):
        if "learned_preference" not in (mem.get("tags") or []):
            continue
        name = f"Preference: {mem.get('key', '')[:60]}"
        if name in existing:
            continue
        try:
            save_skill(name=name, summary=mem.get("value", "")[:280],
                       body=mem.get("value", ""), tags=["feedback", "learned", "auto"],
                       source="feedback_distillation")
            made += 1
            if made >= limit:
                break
        except Exception:
            continue
    return made


async def _brain_background_loop() -> None:
    ensure_core_memories()
    await asyncio.sleep(0.25)
    tick = 0
    while True:
        tick += 1
        try:
            run_learning_cycle(limit=750, propose=False)
            # Execute the queued background work (consolidation, skill extraction, …).
            await asyncio.to_thread(brain_drain_tasks, 25)
            # Distil explicit user feedback into procedural skills every tick (cheap).
            await asyncio.to_thread(_distill_feedback_into_skills)
            # Run the full eval → reflect → propose cycle on a daily cadence so
            # Shims keeps getting measurably better without being asked.
            if IMPROVEMENT_ENABLED and tick % IMPROVEMENT_EVERY_TICKS == 0:
                try:
                    result = await asyncio.to_thread(run_improvement_cycle, _system_prompt())
                    log_event("brain.autonomous_improvement", route="brain:improvement", provider="local",
                              model=BRAIN_VERSION, ok=bool(result.get("ok")),
                              metadata={"proposals": len(result.get("proposals", []))})
                except Exception as exc:
                    log_event("brain.improvement.error", route="brain:improvement", provider="local",
                              model=BRAIN_VERSION, ok=False, message=str(exc)[:200])
        except Exception as exc:
            log_event("brain.background.error", route="brain:background", provider="local", model=BRAIN_VERSION, ok=False, message=str(exc)[:240])
        await asyncio.sleep(BRAIN_BACKGROUND_INTERVAL_SECONDS)


async def _run_boot_self_awareness() -> None:
    started = time.perf_counter()
    try:
        result = await asyncio.to_thread(run_boot_self_audit, app_name=APP_NAME, app_version=APP_VERSION)
        log_event(
            "self_awareness.ready",
            route="self:boot",
            provider="local",
            model="self-awareness",
            ok=bool(result.get("ok")),
            latency_ms=(time.perf_counter() - started) * 1000,
            message=str(result.get("boot_id") or ""),
            metadata={"notes_path": result.get("notes_path"), "queued_tasks": result.get("queued_tasks")},
        )
    except Exception as exc:
        log_event(
            "self_awareness.error",
            route="self:boot",
            provider="local",
            model="self-awareness",
            ok=False,
            latency_ms=(time.perf_counter() - started) * 1000,
            message=str(exc)[:240],
        )


async def _warm_brain_cache() -> None:
    """Pre-load the sentence-transformer embedding model so the first chat turn is not blocked."""
    try:
        from shared.omni_brain import retrieve_context
        await asyncio.to_thread(retrieve_context, "warmup", limit=1)
        log_event("brain.warmup_complete", route="brain:startup", provider="local", model="all-MiniLM-L6-v2", ok=True)
    except Exception as exc:
        log_event("brain.warmup_error", route="brain:startup", provider="local", model="all-MiniLM-L6-v2", ok=False, message=str(exc)[:200])


@asynccontextmanager
async def _omni_lifespan(_app):
    """Start the background-learning loop on boot; cancel it cleanly on shutdown."""
    global _brain_background_task, _self_awareness_task
    # Startup security posture check
    for secret_name in ("SHIMS_SECRET_KEY", "SHIMS_BRIDGE_TOKEN"):
        if is_weak_secret(secret_name, os.getenv(secret_name, "")):
            log_event("security.weak_secret_warning", route="startup", provider="local", model="guardians", ok=False, message=f"{secret_name} is weak or default", metadata={"secret": secret_name})
    ensure_core_memories()
    asyncio.create_task(_warm_brain_cache())
    if BOOT_SELF_AWARENESS_ENABLED and (_self_awareness_task is None or _self_awareness_task.done()):
        _self_awareness_task = asyncio.create_task(_run_boot_self_awareness())
    if BRAIN_BACKGROUND_ENABLED and _brain_background_task is None:
        _brain_background_task = asyncio.create_task(_brain_background_loop())
    _register_scheduler_runners()
    asyncio.create_task(_preload_voice_model())
    try:
        from shared.background_jobs import ensure_default_jobs
        ensure_default_jobs()
    except Exception:
        pass
    try:
        yield
    finally:
        if _self_awareness_task is not None and not _self_awareness_task.done():
            _self_awareness_task.cancel()
            try:
                await _self_awareness_task
            except asyncio.CancelledError:
                pass
            _self_awareness_task = None
        if _brain_background_task is not None:
            _brain_background_task.cancel()
            try:
                await _brain_background_task
            except asyncio.CancelledError:
                pass
            _brain_background_task = None
        try:
            from shared.desktop_scheduler import stop_scheduler
            stop_scheduler()
        except Exception:
            pass


def _register_scheduler_runners() -> None:
    """Register lightweight runners for scheduled tasks."""
    try:
        from shared.desktop_scheduler import register_runner, start_scheduler
        def _tool_runner(payload: dict) -> dict:
            return agent_tools.run_tool(payload.get("tool", ""), payload.get("args", {}), allow_gated=False)
        def _message_runner(payload: dict) -> dict:
            # Store a memory note for the scheduled message
            try:
                from shared.omni_brain import remember
                msg = payload.get("message", "")
                remember("scheduler", f"reminder:{int(time.time())}", msg, tags=["scheduler", "reminder"], source="scheduler")
            except Exception:
                pass
            return {"ok": True, "message": payload.get("message", "")}
        def _plan_runner(payload: dict) -> dict:
            from shared.plan_executor import run_plan_wave
            plan_id = payload.get("plan_id", "")
            if not plan_id:
                return {"ok": False, "error": "plan_id required in payload"}
            return run_plan_wave(plan_id)
        def _inbox_ingest_runner(payload: dict) -> dict:
            from shared.background_jobs import run_inbox_ingest
            return run_inbox_ingest(payload)
        register_runner("tool", _tool_runner)
        register_runner("message", _message_runner)
        register_runner("plan", _plan_runner)
        register_runner("inbox_ingest", _inbox_ingest_runner)
        start_scheduler()
    except Exception:
        pass


# app is constructed earlier in this module; attach the lifespan post-hoc.
app.router.lifespan_context = _omni_lifespan


class ChatRequest(BaseModel):
    message: str = ""
    session_id: str | None = None
    provider: str | None = None
    model: str | None = None
    conversation_mode: bool = True
    web_mode: bool = False
    side_agent_enabled: bool = False
    auto_peer_consultation: bool = False
    privacy_mode: str = "balanced"
    locale: str | None = "en-IN"
    source: str | None = "typed"
    max_tokens: int | None = None
    realtime: bool | None = None
    agent_mode: bool = False
    governed: bool = False  # Enable Neural Governor pipeline
    images: list[str] = []  # base64 data URIs or URLs for vision input
    voice_correction_id: str | None = None  # pending LLM STT correction id from /voice/transcribe or /voice/correct


class MediaRequest(BaseModel):
    kind: str = "image"
    prompt: str
    theme: str | None = None
    quality: str = "standard"
    provider: str | None = None
    privacy_mode: str = "balanced"


class SpeakRequest(BaseModel):
    text: str
    voice: str | None = "auto"
    lang: str | None = "en-IN"
    rate: int | None = 172


class VoiceConfigRequest(BaseModel):
    wake_words: list[str] | None = None
    primary_lang: str = "en-IN"
    secondary_langs: list[str] | None = None
    command_cooldown_seconds: float = 2.2
    silence_timeout_seconds: float = 1.4
    max_auto_replies_without_user: int = 1


class OllamaPullRequest(BaseModel):
    model: str
    stream: bool = True



class SearchRequest(BaseModel):
    query: str
    max_results: int = 6
    provider: str | None = None


class BrainContextRequest(BaseModel):
    query: str
    limit: int = 8


class BrainIngestRequest(BaseModel):
    title: str
    text: str
    source_type: str = "note"
    source_uri: str = ""
    tags: list[str] | None = None
    importance: float = 1.0


class MemorySaveRequest(BaseModel):
    namespace: str = "user"
    key: str
    value: str
    tags: list[str] | None = None
    pinned: bool = False
    weight: float = 1.0
    source: str = "user"


class BrainLearnRequest(BaseModel):
    limit: int = 500
    propose: bool = False


class CaptureShareRequest(BaseModel):
    title: str = ""
    text: str = ""
    url: str = ""
    kind: str = "link"
    source: str = "share"
    metadata: dict[str, Any] | None = None


class MailboxImportRequest(BaseModel):
    provider: str = "local"
    external_id: str = ""
    thread_id: str = ""
    sender: str = ""
    recipients: str = ""
    subject: str = ""
    snippet: str = ""
    body: str = ""
    labels: list[str] | None = None
    received_at: str = ""
    source_url: str = ""
    metadata: dict[str, Any] | None = None


class GmailSyncRequest(BaseModel):
    access_token: str | None = None
    query: str = ""
    max_results: int = 10


class GmailSendRequest(BaseModel):
    to: str
    subject: str = ""
    body: str = ""
    cc: str = ""
    thread_id: str | None = None
    in_reply_to: str | None = None


class GmailReplyRequest(BaseModel):
    message_id: str
    body: str


class ChemVerifyRequest(BaseModel):
    smiles: str


class ChemReactionRequest(BaseModel):
    reaction: str = ""
    rxn_smiles: str = ""


class ChemRetroRequest(BaseModel):
    target: str = ""
    target_smiles: str = ""
    max_routes: int = 5


class ChemIchRequest(BaseModel):
    impurity_pct: float
    max_daily_dose_g: float = 1.0
    impurity_name: str = ""


class ChemToolRequest(BaseModel):
    args: dict[str, Any] = {}


class RichDocxRequest(BaseModel):
    title: str
    blocks: list[dict[str, Any]]
    profile: str = "corporate"
    subtitle: str = ""
    letterhead: bool = True


class SttModelRequest(BaseModel):
    model: str


class TaskEnqueueRequest(BaseModel):
    task_type: str
    title: str = ""
    payload: dict[str, Any] | None = None
    priority: int = 5


class FileSearchRequest(BaseModel):
    query: str
    in_content: bool = True


class FileWorkspaceRequest(BaseModel):
    path: str


class FileOrganizeApplyRequest(BaseModel):
    moves: list[dict[str, str]]


class FileUndoRequest(BaseModel):
    undo_id: str


class CoderCreateRequest(BaseModel):
    name: str
    goal: str = ""


class CoderIterateRequest(BaseModel):
    project_id: str
    instruction: str
    provider: str | None = None
    model: str | None = None
    max_steps: int = 2


class CoderRunRequest(BaseModel):
    project_id: str
    entry: str | None = None


class CoderWriteRequest(BaseModel):
    project_id: str
    path: str
    content: str


class CoderInstallRequest(BaseModel):
    project_id: str


class CoderSettingsRequest(BaseModel):
    base_dir: str | None = None


class CoderAiSupportRequest(BaseModel):
    project_id: str | None = None
    instruction: str
    mode: str = "generate"       # generate | refactor | fix | explain
    provider: str | None = None
    model: str | None = None


class SkillSaveRequest(BaseModel):
    name: str
    summary: str
    body: str = ""
    tags: list[str] | None = None
    pinned: bool = False
    skill_id: str | None = None


class BuilderRunRequest(BaseModel):
    instruction: str
    targets: list[str] = []
    context: list[str] = []
    provider: str = "anthropic"
    model: str | None = None
    apply: bool = False


class ActionRecordRequest(BaseModel):
    action_type: str
    title: str = ""
    payload: dict[str, Any] | None = None
    result: dict[str, Any] | None = None
    evidence: list[dict[str, Any]] | None = None
    requested_level: str = "L3"
    status: str | None = None
    summary: str = ""


class CampaignPlanRequest(BaseModel):
    objective: str
    audience: str = ""
    offer: str = ""
    channels: list[str] | None = None
    tone: str = "clear, useful, credible"
    due_date: str = ""


class CalendarIcsRequest(BaseModel):
    title: str
    start: str | None = None
    end: str | None = None
    duration_minutes: int = 30
    description: str = ""
    location: str = ""


class WebSettingsRequest(BaseModel):
    searxng_url: str | None = None
    tavily_key: str | None = None
    brave_key: str | None = None
    serpapi_key: str | None = None
    duckduckgo_fallback: bool | None = None


class VoiceProfileSelectRequest(BaseModel):
    profile_id: str


class MediaSettingsRequest(BaseModel):
    image_backend: str | None = None
    audio_backend: str | None = None
    video_backend: str | None = None
    stable_diffusion_url: str | None = None
    comfyui_url: str | None = None
    diffusers_enabled: bool | None = None
    diffusers_model: str | None = None
    openai_tts_model: str | None = None
    openai_tts_voice: str | None = None
    openai_video_model: str | None = None
    openai_video_size: str | None = None
    openai_video_seconds: int | None = None
    audio_api_url: str | None = None
    audio_api_key: str | None = None
    video_api_url: str | None = None
    video_api_key: str | None = None

class ProviderKeyRequest(BaseModel):
    provider: str
    action: str = "set"
    api_key: str | None = None
    model: str | None = None


class SettingsRequest(BaseModel):
    gemini_api_key: str | None = None
    openai_api_key: str | None = None
    anthropic_api_key: str | None = None
    kimi_api_key: str | None = None
    deepseek_api_key: str | None = None
    qwen_api_key: str | None = None
    huggingface_base_url: str | None = None
    huggingface_api_key: str | None = None
    huggingface_model: str | None = None
    local_access_token: str | None = None
    shims_peer_url: str | None = None


class AgentModelsRequest(BaseModel):
    router_model: str | None = None
    fast_model: str | None = None
    smart_model: str | None = None
    coder_model: str | None = None
    creative_model: str | None = None
    chemistry_model: str | None = None
    research_model: str | None = None


class EvolutionProposeRequest(BaseModel):
    relative_path: str
    new_content: str
    reason: str | None = ""
    scope: str | None = "prompt_or_skill"
    proposed_by: str | None = "user"


class EvolutionValidateRequest(BaseModel):
    validation: list[Any] | None = None


class EvolutionApplyRequest(BaseModel):
    approved_by: str
    approval_phrase: str
    validation: list[Any] | None = None


class AutonomyCheckRequest(BaseModel):
    action: str
    requested_level: str | None = "L1"



class EvolutionProposalRequest(BaseModel):
    relative_path: str
    new_content: str
    reason: str | None = ""
    author: str | None = "user"
    scope: str | None = "code"
    tests: list[list[str]] | None = None


class EvolutionApprovalRequest(BaseModel):
    approved_by: str | None = "human"
    note: str | None = ""


class EvolutionApplyRequest(BaseModel):
    approved_by: str | None = "human"
    auto_approve_after_validation: bool = False
    approval_phrase: str | None = ""


class EvolutionCapabilityCheckRequest(BaseModel):
    apply: bool = False
    approval_phrase: str | None = ""
    approved_by: str | None = "human-operator"
    revision: str | None = ""
    targets: list[str] | None = None


class EvolutionSelfCheckRequest(BaseModel):
    scope: str = "tests"  # tests, lint, file
    relative_path: str | None = None
    goal: str | None = None
    test_path: str | None = None


class ApprovalDecisionRequest(BaseModel):
    approval_id: str | None = None
    decision: bool | str = True
    approved_by: str | None = "human-operator"
    note: str | None = ""


class CoderProposalRequest(BaseModel):
    relative_path: str
    new_content: str
    reason: str | None = "coder playground patch"
    scope: str | None = "code"
    tests: list[list[str]] | None = None
    run_validation: bool = True


class CoderAppRequest(BaseModel):
    name: str
    prompt: str = ""
    style: str | None = "modern"
    apply: bool = False
    approved_by: str | None = "human-operator"
    approval_phrase: str | None = ""



@dataclass
class TurnPlan:
    session_id: str
    provider: str
    model: str
    route: str
    tool_kind: str | None = None
    tool_prompt: str | None = None
    tool_metadata: dict[str, Any] | None = None
    duplicate: bool = False
    agent: str = "supervisor"


def _jsonl(obj: dict[str, Any]) -> bytes:
    return (json.dumps(obj, ensure_ascii=False, default=str) + "\n").encode("utf-8")



async def _build_user_message_with_images(message: str, images: list[str], provider: str) -> dict[str, Any]:
    """Build a user message with native multimodal content when supported."""
    provider = (provider or "").lower()
    if not images:
        return {"role": "user", "content": message}
    if provider in {"anthropic", "openai", "gemini", "deepseek", "kimi"}:
        from shared.multimodal_messages import build_user_message
        return build_user_message(message, images, provider)
    # Local / Ollama fallback: use vision description
    from shared.vision import describe_image
    vision_parts: list[str] = []
    for idx, src in enumerate(images[:4], start=1):
        result = await asyncio.to_thread(describe_image, src, "Describe this image concisely.", "auto")
        if result.get("ok"):
            vision_parts.append(f"[Attached image {idx}] {result.get('description', '')}")
        else:
            vision_parts.append(f"[Attached image {idx}] (could not describe: {result.get('error', 'unknown')})")
    return {"role": "user", "content": "\n\n".join(vision_parts + [message])}




def _should_auto_plan(text: str) -> bool:
    """Detect user language that implies a multi-step workflow."""
    if not text:
        return False
    t = text.lower()
    triggers = [
        "plan", "in steps", "step by step", "workflow", "automate",
        "every day", "every week", "schedule", " recurring ",
        "research and then", "search and then", "find and then",
        "first ... then", "first, then", "first then",
        "create a project", "build a", "set up a",
        "calculate and save", "compute and save", "analyze and save",
        "find and save", "get and save", "run and save",
        "write a file", "save the result", "save it to",
    ]
    return any(trig in t for trig in triggers) or t.count(" and ") >= 2 or t.count(",") >= 3
def _model_data(model: BaseModel) -> dict[str, Any]:
    if hasattr(model, "model_dump"):
        return model.model_dump()
    return model.dict()


def _utc_now() -> str:
    return datetime.utcnow().isoformat(timespec="seconds") + "Z"


def _slugify(value: str, fallback: str = "app") -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", (value or "").lower()).strip("-")
    return (slug or fallback)[:72]


def _pending_action_path(approval_id: str) -> Path:
    safe = re.sub(r"[^A-Za-z0-9_.-]+", "", approval_id or "")
    if not safe:
        raise ValueError("approval_id required")
    return PENDING_ACTION_DIR / f"{safe}.json"


def _save_pending_action(action: dict[str, Any]) -> dict[str, Any]:
    action["updated_at"] = _utc_now()
    _pending_action_path(action["approval_id"]).write_text(json.dumps(action, indent=2, ensure_ascii=False, default=str), encoding="utf-8")
    return action


def _load_pending_action(approval_id: str) -> dict[str, Any] | None:
    try:
        path = _pending_action_path(approval_id)
    except Exception:
        return None
    if not path.exists():
        return None
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return None


def _public_pending_action(action: dict[str, Any]) -> dict[str, Any]:
    out = dict(action or {})
    payload = dict(out.get("payload") or {})
    if "new_content" in payload:
        payload["new_content_preview"] = str(payload.get("new_content") or "")[:2000]
        payload["new_content_sha256"] = hashlib.sha256(str(payload.get("new_content") or "").encode("utf-8")).hexdigest()
        payload.pop("new_content", None)
    out["payload"] = payload
    return out


def _list_pending_actions(session_id: str | None = None, limit: int = 30, include_resolved: bool = False) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    for path in sorted(PENDING_ACTION_DIR.glob("*.json"), key=lambda p: p.stat().st_mtime, reverse=True):
        try:
            item = json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            continue
        if session_id and item.get("session_id") != session_id:
            continue
        if not include_resolved and item.get("status") not in {"pending", "ready"}:
            continue
        rows.append(_public_pending_action(item))
        if len(rows) >= limit:
            break
    return rows


def _latest_pending_action(session_id: str | None = None) -> dict[str, Any] | None:
    scoped = _list_pending_actions(session_id=session_id, limit=1)
    if scoped:
        return _load_pending_action(scoped[0]["approval_id"])
    global_rows = _list_pending_actions(limit=1)
    if global_rows:
        return _load_pending_action(global_rows[0]["approval_id"])
    return None


def _create_pending_action(
    *,
    action_type: str,
    title: str,
    summary: str,
    payload: dict[str, Any] | None = None,
    session_id: str | None = None,
    risk: str = "source_change",
) -> dict[str, Any]:
    approval_id = "appr_" + uuid.uuid4().hex[:20]
    ledger = record_action(
        "approval_request",
        title,
        payload={"approval_id": approval_id, "action_type": action_type, **(payload or {})},
        result={"status": "pending"},
        evidence=[],
        requested_level="L3",
        status="requires_confirmation",
        summary=summary,
    )
    item = {
        "ok": True,
        "approval_id": approval_id,
        "status": "pending",
        "action_type": action_type,
        "title": title,
        "summary": summary,
        "payload": payload or {},
        "session_id": session_id,
        "risk": risk,
        "created_at": _utc_now(),
        "updated_at": _utc_now(),
        "action_id": ledger.get("action_id"),
        "ledger_hash": ledger.get("ledger_hash"),
        "yes_no_prompt": f"{summary} Say yes to run it now, or no to cancel.",
    }
    return _save_pending_action(item)


def _approval_decision_from_text(text: str) -> bool | None:
    cleaned = _normalize_text(text or "").strip(" .,!?:;")
    yes = {"yes", "y", "yeah", "yep", "haan", "ha", "han", "approve", "approved", "go ahead", "do it", "run it", "apply it", "execute", "ok", "okay"}
    no = {"no", "n", "nope", "cancel", "stop", "reject", "deny", "mat karo", "dont", "don't"}
    if cleaned in yes:
        return True
    if cleaned in no:
        return False
    if re.fullmatch(r"(yes|approve|run|apply|execute)\s+(it|that|approved)?", cleaned):
        return True
    if re.fullmatch(r"(no|cancel|reject|deny)\s+(it|that)?", cleaned):
        return False
    return None


def _extract_proposal_id(text: str) -> str | None:
    match = re.search(r"\b(patch_[A-Za-z0-9_-]+|proposal[_-]?[A-Za-z0-9_-]+|[A-Za-z0-9]{8,}_[A-Za-z0-9_-]+)\b", text or "")
    return match.group(1) if match else None


def _build_generated_app_html(name: str, prompt: str = "", style: str | None = "modern") -> str:
    title = (name or "SHIMS Generated App").strip()[:80]
    brief = (prompt or f"{title} app generated inside SHIMS Omni").strip()[:1000]
    escaped_title = html.escape(title)
    escaped_brief = html.escape(brief)
    palette = {
        "modern": ("#07111f", "#00d4ff", "#ffb020", "#f7fbff"),
        "ops": ("#0d1117", "#58d68d", "#f4d35e", "#f7fbff"),
        "light": ("#f7fbff", "#1459ff", "#ff8c42", "#0d1828"),
    }.get((style or "modern").lower(), ("#07111f", "#00d4ff", "#ffb020", "#f7fbff"))
    bg, accent, warm, text = palette
    return f"""<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{escaped_title}</title>
<style>
:root{{--bg:{bg};--accent:{accent};--warm:{warm};--text:{text};--muted:color-mix(in srgb,var(--text),transparent 38%)}}
*{{box-sizing:border-box}} body{{margin:0;min-height:100vh;background:var(--bg);color:var(--text);font-family:Inter,Segoe UI,Arial,sans-serif}}
.app{{min-height:100vh;display:grid;grid-template-rows:auto 1fr auto}}
header{{display:flex;align-items:center;justify-content:space-between;padding:18px 22px;border-bottom:1px solid color-mix(in srgb,var(--accent),transparent 72%)}}
.brand{{font-weight:800;letter-spacing:.08em;text-transform:uppercase}} .pill{{border:1px solid var(--accent);color:var(--accent);border-radius:999px;padding:6px 10px;font-size:12px}}
main{{display:grid;grid-template-columns:minmax(0,1.1fr) minmax(280px,.9fr);gap:22px;padding:24px;align-items:stretch}}
section{{border:1px solid color-mix(in srgb,var(--accent),transparent 72%);border-radius:8px;background:color-mix(in srgb,var(--bg),white 5%);padding:20px}}
h1{{font-size:clamp(28px,5vw,54px);line-height:1;margin:0 0 12px}} p{{color:var(--muted);line-height:1.6}}
.tools{{display:grid;gap:10px}} button,input,textarea{{font:inherit}} input,textarea{{width:100%;background:color-mix(in srgb,var(--bg),black 18%);border:1px solid color-mix(in srgb,var(--text),transparent 78%);color:var(--text);border-radius:6px;padding:10px}}
button{{border:0;border-radius:6px;padding:10px 14px;background:var(--accent);color:#06101d;font-weight:700;cursor:pointer}} button.secondary{{background:transparent;color:var(--warm);border:1px solid var(--warm)}}
.result{{min-height:120px;border:1px dashed color-mix(in srgb,var(--warm),transparent 45%);border-radius:8px;padding:14px;color:var(--muted)}} footer{{padding:14px 22px;color:var(--muted);font-size:12px;border-top:1px solid color-mix(in srgb,var(--accent),transparent 78%)}}
@media(max-width:760px){{main{{grid-template-columns:1fr;padding:14px}} header{{padding:14px}}}}
</style>
</head>
<body>
<div class="app">
<header><div class="brand">{escaped_title}</div><div class="pill">SHIMS generated app</div></header>
<main>
<section>
<h1>{escaped_title}</h1>
<p>{escaped_brief}</p>
<div class="tools">
<input id="item" placeholder="Enter an item, task, idea, or customer">
<textarea id="notes" rows="5" placeholder="Add notes"></textarea>
<button onclick="addItem()">Add to workspace</button>
</div>
</section>
<section>
<h2>Workspace</h2>
<div class="result" id="result">No items yet.</div>
<button class="secondary" onclick="clearItems()">Clear</button>
</section>
</main>
<footer>Created by the SHIMS Omni coder playground. Edit this file in SHIMS to evolve the app.</footer>
</div>
<script>
const key = 'shims.generated.{_slugify(title)}';
function loadItems(){{ try{{ return JSON.parse(localStorage.getItem(key)||'[]'); }}catch(e){{ return []; }} }}
function saveItems(items){{ localStorage.setItem(key, JSON.stringify(items)); render(items); }}
function render(items=loadItems()){{ const box=document.getElementById('result'); box.innerHTML = items.length ? items.map((x,i)=>'<div><b>#'+(i+1)+'</b> '+escapeHtml(x.item)+'<br><small>'+escapeHtml(x.notes)+'</small></div>').join('<hr>') : 'No items yet.'; }}
function addItem(){{ const item=document.getElementById('item').value.trim(); const notes=document.getElementById('notes').value.trim(); if(!item && !notes) return; const items=loadItems(); items.unshift({{item:item||'Untitled',notes,time:Date.now()}}); saveItems(items.slice(0,100)); }}
function clearItems(){{ saveItems([]); }}
function escapeHtml(s){{ return String(s||'').replace(/[&<>"']/g,c=>({{'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}}[c])); }}
render();
</script>
</body>
</html>
"""


def _proposal_validation_for_path(relative_path: str) -> list[list[str]]:
    suffix = Path(relative_path).suffix.lower()
    if suffix in {".html", ".htm"}:
        return [[sys.executable, "-c", f"from pathlib import Path; p=Path({relative_path!r}); s=p.read_text(encoding='utf-8'); assert '<html' in s.lower() or '<!doctype html>' in s.lower(); print('html-readable')"]]
    if suffix in {".js", ".css", ".md", ".txt"}:
        return [[sys.executable, "-c", f"from pathlib import Path; p=Path({relative_path!r}); assert p.exists(); p.read_text(encoding='utf-8'); print('text-readable')"]]
    if suffix == ".py":
        return [[sys.executable, "-m", "py_compile", relative_path]]
    return [[sys.executable, "-c", f"from pathlib import Path; p=Path({relative_path!r}); assert p.exists(); p.read_bytes(); print('file-readable')"]]


def _detect_chat_action_request(text: str, session_id: str | None) -> dict[str, Any] | None:
    low = _normalize_text(text or "")
    proposal_id = _extract_proposal_id(text or "")
    if proposal_id and re.search(r"\b(apply|approve|install|merge)\b", low):
        return _create_pending_action(
            action_type="evolution_apply",
            title=f"Apply patch {proposal_id}",
            summary=f"Apply validated patch {proposal_id} to the live SHIMS files",
            payload={"proposal_id": proposal_id},
            session_id=session_id,
        )
    if proposal_id and re.search(r"\b(validate|sandbox|test)\b", low):
        return _create_pending_action(
            action_type="evolution_validate",
            title=f"Sandbox test patch {proposal_id}",
            summary=f"Run sandbox validation for patch {proposal_id}",
            payload={"proposal_id": proposal_id},
            session_id=session_id,
            risk="sandbox_validation",
        )
    if re.search(r"\b(capability check|test .*power|prove .*change|self.?evolution.*apply|backend.*frontend.*feature)\b", low):
        return _create_pending_action(
            action_type="evolution_capability_check",
            title="Run SHIMS backend/frontend/feature capability check",
            summary="Create, sandbox-validate, approve, and apply harmless probe files across backend, frontend, and generated feature surfaces",
            payload={"revision": "chat-" + datetime.now().strftime("%Y%m%d-%H%M%S"), "targets": ["backend", "frontend", "feature"]},
            session_id=session_id,
        )
    if re.search(r"\b(?:self[-\s]?check|check|inspect|analyze|review|test)\b[^.]{0,60}\b(?:your(?:self|\s+own)?\s+)?(?:code|source|files?|tests?|app|system|repo|repository|\.(?:py|js|html|json|md|txt|csv|sh|bat|ps1))\b", low):
        return _create_pending_action(
            action_type="evolution_self_check",
            title="Run SHIMS self-check",
            summary="Inspect SHIMS code and produce a validated patch proposal when improvements are found",
            payload={"scope": "tests", "goal": text or ""},
            session_id=session_id,
        )
    # Direct desktop-bridge shortcuts (screenshot, ping, system info) so natural chat can reach the paired desktop.
    bridge_action: str | None = None
    bridge_args: dict[str, Any] = {}
    if re.search(r"\bping\s+(?:the\s+)?(?:desktop\s+)?bridge\b", low):
        bridge_action = "ping"
    elif re.search(r"\btake\s+a?\s*screenshot\b", low) and re.search(r"\b(desktop|bridge|computer|pc|machine|screen)\b", low):
        bridge_action = "screenshot"
    elif re.search(r"\b(desktop|bridge)\s+(?:system\s+)?info\b", low):
        bridge_action = "system_info"
    if bridge_action:
        return _create_pending_action(
            action_type="agent_tool",
            title=f"Desktop bridge {bridge_action}",
            summary=f"Run desktop.bridge {bridge_action} on the paired desktop",
            payload={"tool": "desktop.bridge", "args": {"action": bridge_action, **bridge_args}},
            session_id=session_id,
            risk="desktop_access",
        )
    # Coder / app creation intent — checked BEFORE _detect_tool_intent because words like
    # "drawing" in "create a drawing app" can falsely match image intent.
    create_app = re.search(r"\b(create|build|make|scaffold)\b.*\b(app|tool|dashboard|frontend|software|program|code|project)\b", low)
    is_coder_intent = bool(create_app) or bool(re.search(r"\b(write|edit|modify|fix|debug|build|compile|run|test)\b.*\b(code|file|script|program|function|class|module|project)\b", low))
    if create_app:
        name_match = re.search(r"(?:app|tool|dashboard|frontend|software|program|project)\s+(?:called|named|for)?\s*([A-Za-z0-9 _-]{3,60})", text or "", re.I)
        name = (name_match.group(1).strip(" .,:;-") if name_match else "Generated SHIMS App")
        return _create_pending_action(
            action_type="coder_app_scaffold",
            title=f"Create app: {name}",
            summary=f"Create and apply a generated app at apps/generated/{_slugify(name)}/index.html",
            payload={"name": name, "prompt": text, "style": "modern"},
            session_id=session_id,
        )
    # Coder intent beats image/pdf/video intent — return None so agent loop handles it.
    if is_coder_intent and _detect_tool_intent(text):
        return None
    return None


async def _execute_pending_action(action: dict[str, Any], approved_by: str = "human-operator") -> dict[str, Any]:
    action_type = action.get("action_type")
    payload = action.get("payload") or {}
    if action_type == "evolution_validate":
        proposal_id = str(payload.get("proposal_id") or "")
        validation = validate_proposal(proposal_id)
        return {"ok": validation.status == "validated", "status": validation.status, "message": validation.message, **validation.details}
    if action_type == "evolution_apply":
        proposal_id = str(payload.get("proposal_id") or "")
        validation = validate_proposal(proposal_id)
        if validation.status != "validated":
            return {"ok": False, "status": validation.status, "message": validation.message, **validation.details}
        approval = approve_proposal(proposal_id, approved_by=approved_by, note="Approved through SHIMS yes/no chat approval.")
        if approval.status != "approved":
            return {"ok": False, "status": approval.status, "message": approval.message, **approval.details}
        applied = apply_proposal(proposal_id, approved_by=approved_by)
        return {"ok": applied.status == "applied", "status": applied.status, "message": applied.message, **applied.details}
    if action_type == "evolution_capability_check":
        req = EvolutionCapabilityCheckRequest(
            apply=True,
            approval_phrase="I_APPROVE_SHIMS_PATCH",
            approved_by=approved_by,
            revision=str(payload.get("revision") or ""),
            targets=payload.get("targets") or None,
        )
        return await evolution_capability_check(req)
    if action_type == "evolution_self_check":
        req = EvolutionSelfCheckRequest(
            scope=str(payload.get("scope") or "tests"),
            relative_path=str(payload.get("relative_path")) if payload.get("relative_path") else None,
            goal=str(payload.get("goal") or ""),
        )
        return await evolution_self_check(req)
    if action_type == "coder_app_scaffold":
        app_req = CoderAppRequest(
            name=str(payload.get("name") or "Generated SHIMS App"),
            prompt=str(payload.get("prompt") or ""),
            style=str(payload.get("style") or "modern"),
            apply=True,
            approved_by=approved_by,
            approval_phrase="I_APPROVE_SHIMS_PATCH",
        )
        return _create_or_propose_coder_app(app_req)
    if action_type == "agent_tool":
        tool = str(payload.get("tool") or "")
        args = payload.get("args") or {}
        result = await asyncio.to_thread(agent_tools.run_tool, tool, args, allow_gated=True)
        ok = bool(result.get("ok", True))
        msg = result.get("error") or result.get("note") or (f"Ran {tool}." if ok else f"{tool} failed.")
        return {"ok": ok, "status": "completed" if ok else "failed", "message": msg, "tool": tool, "result": result}
    return {"ok": False, "status": "unknown_action", "message": f"Unknown pending action type: {action_type}"}


def _decision_value(decision: bool | str) -> bool:
    if isinstance(decision, bool):
        return decision
    return _approval_decision_from_text(str(decision)) is True


def _create_coder_proposal(req: CoderProposalRequest) -> dict[str, Any]:
    tests = req.tests or _proposal_validation_for_path(req.relative_path)
    proposal = create_proposal(
        req.relative_path,
        req.new_content,
        reason=req.reason or "coder playground patch",
        author="coder-playground",
        scope=req.scope or "code",
        tests=tests,
    )
    if req.run_validation and proposal.get("ok"):
        validation = validate_proposal(proposal["proposal_id"], validation=tests)
        proposal["validation"] = {"ok": validation.status == "validated", "status": validation.status, "message": validation.message, **validation.details}
    return proposal


def _create_or_propose_coder_app(req: CoderAppRequest) -> dict[str, Any]:
    slug = _slugify(req.name, "generated-app")
    relative_path = f"apps/generated/{slug}/index.html"
    content = _build_generated_app_html(req.name, req.prompt, req.style)
    tests = _proposal_validation_for_path(relative_path)
    proposal = create_proposal(
        relative_path,
        content,
        reason=f"Coder playground generated app: {req.name}",
        author=req.approved_by or "coder-playground",
        scope="generated_app",
        tests=tests,
    )
    if not proposal.get("ok"):
        return proposal
    validation = validate_proposal(proposal["proposal_id"], validation=tests)
    out: dict[str, Any] = {
        "ok": validation.status == "validated",
        "status": validation.status,
        "message": validation.message,
        "proposal": proposal,
        "validation": {"ok": validation.status == "validated", "status": validation.status, "message": validation.message, **validation.details},
        "relative_path": relative_path,
        "app_url": f"/generated-apps/{slug}/index.html",
    }
    if req.apply:
        if not settings.omnipotent_mode and (req.approval_phrase or "").strip() != "I_APPROVE_SHIMS_PATCH":
            out.update({"ok": False, "status": "approval_required", "message": "Applying generated apps requires approval_phrase='I_APPROVE_SHIMS_PATCH'."})
            return out
        approval = approve_proposal(proposal["proposal_id"], approved_by=req.approved_by or "human-operator", note="Generated app approved through coder playground.")
        if approval.status != "approved":
            out.update({"ok": False, "status": approval.status, "message": approval.message, "approval": approval.details})
            return out
        applied = apply_proposal(proposal["proposal_id"], approved_by=req.approved_by or "human-operator", validation=tests)
        out["apply"] = {"ok": applied.status == "applied", "status": applied.status, "message": applied.message, **applied.details}
        out.update({"ok": applied.status == "applied", "status": applied.status, "message": applied.message})
    return out


def _safe_name(title: str, suffix: str) -> str:
    slug = re.sub(r"[^A-Za-z0-9_.-]+", "_", (title or "shims")[:80]).strip("_") or "shims"
    return f"{slug}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:6]}.{suffix}"

def _attach_ledger(result: dict[str, Any], path: Path, document_type: str) -> dict[str, Any]:
    try:
        entry = ledger_document(path, document_type=document_type, metadata={"title": result.get("title"), "filename": result.get("filename"), "source": "shims"})
        result["ledger"] = entry
        result["sha256"] = entry.get("sha256")
        result["verified"] = True
    except Exception as exc:
        result["verified"] = False
        result["ledger_error"] = str(exc)[:180]
    return result


def _trust_fields(trust: dict[str, Any]) -> dict[str, Any]:
    return {
        "trust": trust,
        "evidence": trust.get("evidence") or [],
        "confidence": trust.get("confidence") or {},
        "query_plan": trust.get("query_plan"),
        "action_id": trust.get("action_id") or "",
        "ledger_hash": trust.get("ledger_hash") or "",
    }


def _strip_wake(text: str) -> str:
    words = _settings["voice"].get("wake_words") or []
    out = text or ""
    for w in words:
        out = re.sub(r"\b" + re.escape(str(w)) + r"\b", "", out, flags=re.I)
    return re.sub(r"^\s*[,.:;\-]+\s*", "", out).strip()


def _normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "").strip().lower())


def _clean_spaces(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "").strip())


def _looks_local_model(model: str | None) -> bool:
    m = (model or "").strip().lower()
    if not m:
        return False
    return any(m.startswith(x) or (x + ":") in m for x in LOCAL_HINTS) or ":" in m and not _cloud_provider_from_model(m)


def _cloud_provider_from_model(model: str | None) -> str | None:
    m = (model or "").strip().lower()
    for provider, hints in CLOUD_HINTS.items():
        if any(h in m for h in hints):
            return provider
    return None


def _provider_configured(provider: str) -> bool:
    if provider in {"ollama", "huggingface"}:
        return True
    env = PROVIDER_ENV.get(provider)
    return bool(env and _clean_secret(os.getenv(env)))


def _cloud_model_names() -> set[str]:
    return {m["name"] for m in RECOMMENDED_MODELS if m.get("provider") != "ollama"}


async def _ollama_models_raw(timeout: float = 2.5) -> list[dict[str, Any]]:
    try:
        async with httpx.AsyncClient(timeout=timeout) as client:
            r = await client.get(f"{OLLAMA_HOST}/api/tags")
            r.raise_for_status()
            data = r.json()
        out: list[dict[str, Any]] = []
        for item in data.get("models", []):
            name = item.get("name") or item.get("model")
            if not name:
                continue
            details = item.get("details") or {}
            out.append({
                "name": name,
                "model": name,
                "provider": "ollama",
                "family": details.get("family") or "",
                "parameters": details.get("parameter_size") or item.get("parameter_size") or "",
                "quantization": details.get("quantization_level") or "",
                "modified_at": item.get("modified_at"),
                "size": item.get("size"),
                "installed": True,
                "is_default": name == DEFAULT_OLLAMA_MODEL,
            })
        # Keep every installed model so the UI and fallback can use fast/chat-only
        # options such as gemma3:1b or llama3.2, while still annotating tool support.
        seen: dict[str, dict[str, Any]] = {m["name"]: m for m in mark_tool_capable(out)}
        return [seen[k] for k in sorted(seen)]
    except Exception:
        return []


async def _ollama_names() -> list[str]:
    return [m["name"] for m in await _ollama_models_raw()]


async def _hf_models_raw(timeout: float = 2.5) -> list[dict[str, Any]]:
    """List models from a local Hugging Face OpenAI-compatible endpoint (TGI/vLLM/llama.cpp server)."""
    try:
        async with httpx.AsyncClient(timeout=timeout) as client:
            headers: dict[str, str] = {}
            key = _clean_secret(os.getenv("HUGGINGFACE_API_KEY"))
            if key:
                headers["Authorization"] = f"Bearer {key}"
            r = await client.get(f"{HUGGINGFACE_HOST}/v1/models", headers=headers)
            r.raise_for_status()
            data = r.json()
        out: list[dict[str, Any]] = []
        for item in data.get("data", []):
            name = item.get("id")
            if not name:
                continue
            out.append({
                "name": name,
                "model": name,
                "provider": "huggingface",
                "family": item.get("owned_by") or "",
                "parameters": "",
                "quantization": "",
                "modified_at": item.get("created"),
                "size": None,
                "installed": True,
                "is_default": name == DEFAULT_HUGGINGFACE_MODEL,
                "tool_capable": is_tool_capable(name),
            })
        # Keep every model the HF endpoint reports so the user can pick it,
        # but still annotate whether it is known to support tools.
        out = mark_tool_capable(out)
        seen: dict[str, dict[str, Any]] = {m["name"]: m for m in out}
        return [seen[k] for k in sorted(seen)]
    except Exception:
        return []


async def _hf_names() -> list[str]:
    return [m["name"] for m in await _hf_models_raw()]


def _normalize_ollama_model_name(model: str | None) -> str:
    raw = (model or "").strip()
    key = re.sub(r"\s+", " ", raw.lower())
    return OLLAMA_MODEL_ALIASES.get(key, raw)


def _is_tool_capable(name: str) -> bool:
    return is_tool_capable(name)


def _preferred_local_model(names: list[str], *, realtime: bool = False, exclude: set[str] | None = None, tool_capable_only: bool = False, prefer_tiny: bool = False) -> str:
    installed = set(names)
    blocked = exclude or set()
    # Tool-capable order (excludes llama3.2, gemma3, gemma-4, deepseek-r1, mistral-small)
    tool_realtime = (
        "llama3.2:latest",
        "llama3.2",
        "qwen2.5:3b",
        "qwen2.5:7b",
        "qwen2.5-coder:7b",
        "mistral",
        DEFAULT_OLLAMA_MODEL,
    )
    tool_quality = (
        "llama3.2:latest",
        "llama3.2",
        "qwen2.5-coder:14b",
        "qwen2.5:14b",
        "qwen2.5-coder:7b",
        "qwen2.5:7b",
        "mistral-nemo",
        "llama3.1",
        "qwen2.5:3b",
        DEFAULT_OLLAMA_MODEL,
    )
    # General chat order (includes fast non-tool models)
    chat_realtime = (
        "llama3.2:latest",
        "llama3.2",
        "gemma3:270m",
        "gemma3:1b",
        "qwen2.5:3b",
        "qwen2.5:7b",
        DEFAULT_OLLAMA_MODEL,
        "gemma3:4b",
    )
    chat_quality = (
        "qwen2.5:7b",
        "llama3.2:latest",
        "qwen2.5:3b",
        "gemma3:1b",
        "gemma3:4b",
        "llama3.2",
        DEFAULT_OLLAMA_MODEL,
        "mistral-small:latest",
    )
    # Tiny models for cold-start / timeout recovery on under-powered machines.
    # gemma3:270m is the fastest viable option; gemma3:1b trades a little speed
    # for better quality, and qwen2.5:3b is the strongest tiny tool-capable model.
    tiny_order = (
        "gemma3:270m",
        "gemma3:1b",
        "qwen2.5:3b",
        "llama3.2:latest",
        "llama3.2",
        DEFAULT_OLLAMA_MODEL,
    )
    if prefer_tiny:
        realtime_order = tiny_order
        quality_order = tiny_order
    else:
        realtime_order = tool_realtime if tool_capable_only else chat_realtime
        quality_order = tool_quality if tool_capable_only else chat_quality
    for preferred in (realtime_order if realtime else quality_order):
        if preferred in installed and preferred not in blocked:
            return preferred
    for name in names:
        if name not in blocked and (not tool_capable_only or _is_tool_capable(name)):
            return name
    return DEFAULT_OLLAMA_MODEL


def _is_realtime_request(req: ChatRequest | None) -> bool:
    if not req:
        return False
    if req.realtime is not None:
        return bool(req.realtime)
    return (req.source or "").strip().lower() in {"android", "mobile", "voice", "speech"}


def _ollama_options(*, realtime: bool = False, max_tokens: int | None = None) -> dict[str, Any]:
    brain = _settings["brain"]
    num_ctx = int(brain.get("num_ctx", 8192))
    options: dict[str, Any] = {
        "temperature": float(brain.get("temperature", 0.12)),
        "top_p": float(brain.get("top_p", 0.82)),
        "repeat_penalty": float(brain.get("repeat_penalty", 1.25)),
        "num_ctx": num_ctx,
    }
    if realtime:
        options["num_ctx"] = min(num_ctx, int(brain.get("realtime_num_ctx", 2048)))
        # No cap — let the model generate as much as needed
        options["num_predict"] = int(max_tokens or brain.get("realtime_max_tokens", 256)) if (max_tokens or brain.get("realtime_max_tokens")) else -1
    elif max_tokens:
        options["num_predict"] = int(max_tokens)
    else:
        options["num_predict"] = -1  # unlimited
    return options


async def _collect_ollama_stream(
    model: str,
    messages: list[dict[str, str]],
    *,
    realtime: bool = False,
    max_tokens: int | None = None,
    on_delta: Callable[[str], Awaitable[None]] | None = None,
    first_token_timeout: float = 18.0,
) -> str:
    """Collect tokens from Ollama stream.

    Waits as long as the model needs *after* the first token, but aborts if the
    first token doesn't arrive within ``first_token_timeout`` seconds. This
    prevents a cold/slow local model from making the UI feel unresponsive.
    """
    answer = ""
    stream = _ollama_chat_stream(model, messages, realtime=realtime, max_tokens=max_tokens)
    first = True
    try:
        while True:
            try:
                if first:
                    delta = await asyncio.wait_for(stream.__anext__(), timeout=first_token_timeout)
                    first = False
                else:
                    delta = await stream.__anext__()
            except StopAsyncIteration:
                break
            if delta:
                answer += delta
                if on_delta:
                    await on_delta(delta)
    except Exception:
        try:
            await stream.aclose()
        except Exception:
            pass
        raise
    return answer


async def _collect_hf_stream(
    model: str,
    messages: list[dict[str, str]],
    *,
    realtime: bool = False,
    max_tokens: int | None = None,
    on_delta: Callable[[str], Awaitable[None]] | None = None,
) -> str:
    """Collect tokens from a Hugging Face OpenAI-compatible endpoint stream."""
    answer = ""
    stream = _hf_chat_stream(model, messages, realtime=realtime, max_tokens=max_tokens)
    try:
        while True:
            try:
                delta = await stream.__anext__()
            except StopAsyncIteration:
                break
            if delta:
                answer += delta
                if on_delta:
                    await on_delta(delta)
    except Exception:
        try:
            await stream.aclose()
        except Exception:
            pass
        raise
    return answer


def _ollama_aliases_payload() -> dict[str, str]:
    return dict(sorted(OLLAMA_MODEL_ALIASES.items()))


async def _resolve_provider_model(provider: str | None, model: str | None, *, privacy_mode: str = "balanced", text: str | None = None) -> tuple[str, str, str]:
    """Single source of truth for model/provider routing.

    Stale UI state is treated conservatively: an explicit Ollama provider can never call Anthropic.
    A selected local model also wins over a stale cloud provider value from older UI state.
    """
    from shared.privacy_guard import can_use_cloud
    requested_provider = (provider or "auto").strip().lower() or "auto"
    requested_model = _normalize_ollama_model_name(model)
    names = await _ollama_names()
    local_default = _preferred_local_model(names)

    if requested_provider in {"local", "ollama"}:
        if requested_model and requested_model in names:
            return "ollama", requested_model, "selected-local"
        if requested_model and _looks_local_model(requested_model) and requested_model not in _cloud_model_names():
            return "ollama", requested_model, "local-requested"
        return "ollama", local_default, "forced-local-from-provider"

    if requested_model and (_looks_local_model(requested_model) or requested_model in names):
        return "ollama", requested_model, "local-model-overrides-stale-provider"

    # Respect explicit cloud provider choice when the selected model is not local.
    if requested_provider in PROVIDER_DEFAULTS and requested_provider != "auto" and requested_provider != "ollama":
        # Privacy guard: only block HIGH sensitivity for explicit provider choices;
        # MEDIUM is allowed because the user deliberately picked a cloud provider.
        if text:
            from shared.privacy_guard import classify_sensitivity
            if classify_sensitivity(text) == "high":
                return "ollama", local_default, "privacy-guard-high-explicit-override"
        p = requested_provider
        m = requested_model if requested_model and not _looks_local_model(requested_model) else PROVIDER_DEFAULTS[p]
        if not _provider_configured(p) and names:
            return "ollama", local_default, f"cloud-{p}-not-configured-fallback-local"
        return p, m, "explicit-cloud-provider"

    # Privacy guard: check if text contains sensitive data before routing to cloud (auto mode only)
    if text:
        allowed, reason = can_use_cloud(text, privacy_mode)
        if not allowed:
            return "ollama", local_default, f"privacy-guard-{reason}"

    cloud_provider = _cloud_provider_from_model(requested_model)
    if cloud_provider:
        if _provider_configured(cloud_provider):
            return cloud_provider, requested_model or PROVIDER_DEFAULTS[cloud_provider], "cloud-model-selected"
        if names:
            return "ollama", local_default, f"cloud-{cloud_provider}-not-configured-fallback-local"
        return cloud_provider, requested_model or PROVIDER_DEFAULTS[cloud_provider], "cloud-model-no-local-fallback"

    return "ollama", requested_model or local_default, "auto-local-first"


def _guard_duplicate(session_id: str, user_text: str, source: str | None) -> bool:
    cleaned = _normalize_text(user_text)
    if not cleaned:
        return True
    if cleaned in {"i heard silence", "silence", "...", "."}:
        return True
    h = hashlib.sha256(cleaned.encode("utf-8")).hexdigest()[:16]
    last = _turn_guard.get(session_id) or {}
    now = time.time()
    cooldown = max(4.5, float(_settings["voice"].get("command_cooldown_seconds", 2.2)) + 1.0)
    if last.get("hash") == h and (now - float(last.get("time", 0))) < cooldown:
        return True
    _turn_guard[session_id] = {"hash": h, "time": now, "source": source or "typed"}
    return False


def _is_empty_greeting(text: str) -> bool:
    t = _normalize_text(_strip_wake(text))
    return t in {"", "hi", "hello", "hey", "haan", "han", "sun rahe ho", "suno", "are you there", "listen"}


def _detect_tool_intent(text: str) -> tuple[str, str] | None:
    raw = text or ""
    if not raw.strip():
        return None
    t = _normalize_text(raw)
    normalized = raw.lower()
    devanagari = {
        "image": ["तस्वीर", "चित्र", "फोटो"],
        "pdf": ["पीडीएफ", "दस्तावेज"],
        "video": ["वीडियो"],
        "audio": ["आवाज", "आवाज़", "ध्वनि"],
        "ppt": ["प्रेजेंटेशन"],
    }
    patterns = [
        ("ppt", r"\b(create|generate|make|prepare|build)\b.*\b(ppt|pptx|powerpoint|presentation|deck|slides)\b|\b(ppt|pptx|powerpoint|presentation|deck|slides)\b.*\b(create|generate|make|prepare|build)\b"),
        ("image", r"\b(create|generate|make|draw|render|paint)\b.*\b(image|photo|picture|poster|drawing|art|logo)\b|\b(image|photo|picture|poster|drawing|art|logo)\b.*\b(create|generate|make|draw|render|paint)\b|\b(tasveer|tasvir|pic|photo)\b.*\b(banao|bana|banado|karo|kar)\b"),
        ("pdf", r"\b(create|generate|make|write|prepare|build)\b.*\b(pdf|document|report|certificate|coa|letter|invoice|quotation|sop)\b|\b(pdf|document|report|certificate|coa|letter|invoice|quotation|sop)\b.*\b(create|generate|make|write|prepare|build)\b|\b(pdf|document|invoice|quotation|coa)\b.*\b(banao|bana|banado|karo|kar)\b"),
        ("video", r"\b(create|generate|make|render|build)\b.*\b(video|clip|movie|reel)\b|\b(video|clip|movie|reel)\b.*\b(create|generate|make|render|build)\b"),
        ("audio", r"\b(create|generate|make|record)\b.*\b(audio|sound|music|voice|speech|voice note)\b|\b(audio|sound|music|voice|speech|voice note)\b.*\b(create|generate|make|record)\b"),
    ]
    kind: str | None = None
    for k, pat in patterns:
        if re.search(pat, t) or re.search(pat, normalized) or any(term in raw for term in devanagari.get(k, [])):
            kind = k
            break
    if not kind:
        return None
    prompt = _strip_wake(raw)
    remove = r"(?i)\b(please|pls|hey|hi|hello|ok|okay|shims|suno|sun|arre|create|generate|make|draw|render|write|prepare|record|build|banao|bana|banado|karo|kar|do|me|an|a|the|of|for|image|photo|picture|poster|drawing|art|logo|pdf|document|report|certificate|coa|letter|invoice|quotation|sop|video|clip|movie|reel|audio|sound|music|voice|speech|note|tasveer|tasvir|pic|ppt|pptx|powerpoint|presentation|deck|slides)\b"
    prompt = re.sub(remove, " ", prompt)
    prompt = re.sub(r"\s+", " ", prompt).strip(" .,:;-_")
    if not prompt:
        prompt = _strip_wake(raw) or raw
    return kind, prompt




def _detect_search_intent(text: str, web_mode: bool = False) -> str | None:
    """Return a query only when the user actually asks for web/current data.

    v14.2 treated the UI WEB toggle as "search every message". That made
    simple turns like "hi" become internet searches. In v14.3 the toggle means
    "internet is allowed"; the router still needs explicit web/current intent.
    """
    raw = _strip_wake((text or "").strip())
    if not raw or _is_empty_greeting(raw):
        return None
    # Never web-search commands that are clearly local tool/media/document work.
    if _detect_tool_intent(raw):
        return None
    plan = plan_search_query(raw, web_mode=web_mode)
    if not plan.should_search:
        return None
    if not plan.primary_query or _is_empty_greeting(plan.primary_query):
        return None
    if len(plan.primary_query) < 3 and not plan.primary_query.isupper():
        return None
    return plan.primary_query


SEARCH_PLANNER_PROMPT = """You are SHIMS' web-search planner.
Decide whether the user's turn really needs live/public web evidence.
Search only for fresh/current facts, public regulations, prices, news, patents, identifiers, citations, or when the user explicitly asks to browse/search/verify online.
Do not search greetings, normal conversation, coding help, local SHIMS/app/workspace questions, brainstorming, or source-code questions unless the user explicitly asks for external sources.
When search is needed, convert the user request into search-engine terms. Do not copy the whole sentence. Use 3-10 specific keywords, preserving site:, filetype:, quoted phrases, CAS numbers, product names, places, years, and official domains.
Return JSON only:
{"should_search": true|false, "primary_query": "...", "queries": ["..."], "intent": "fresh|regulatory|market|patent|identifier|general|none", "user_task": "what the user wants answered"}
"""

SEARCH_ANSWER_PROMPT = """You are SHIMS answering after a routed web search.
Answer the user's actual question, not just the search query. Use only the numbered web sources plus local memory context if it is present.
If the sources do not prove something, say what is uncertain. Cite supporting sources inline as [1], [2], etc. Keep the answer concise and useful.
"""


def _extract_json_object(text: str) -> dict[str, Any] | None:
    raw = (text or "").strip()
    if not raw:
        return None
    try:
        data = json.loads(raw)
        return data if isinstance(data, dict) else None
    except Exception:
        pass
    start = raw.find("{")
    end = raw.rfind("}")
    if start < 0 or end <= start:
        return None
    try:
        data = json.loads(raw[start : end + 1])
        return data if isinstance(data, dict) else None
    except Exception:
        return None


def _coerce_bool(value: Any) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return bool(value)
    if isinstance(value, str):
        return value.strip().lower() in {"1", "true", "yes", "y", "search", "needed"}
    return False


def _search_plan_dict(plan: Any) -> dict[str, Any]:
    if hasattr(plan, "to_dict"):
        return plan.to_dict()
    return dict(plan or {})


def _focused_query_from_text(text: str, heuristic: Any) -> str:
    query = _clean_spaces(text or "")
    raw = _clean_spaces(getattr(heuristic, "original_query", "") or "")
    if not query:
        return getattr(heuristic, "primary_query", "") or ""
    low = query.lower()
    raw_low = raw.lower()
    copied_whole_turn = raw_low and low == raw_low
    still_command_like = bool(re.search(r"\b(please|pls|hey|shims|search|internet|web|browse|look up|what is|tell me)\b", low, re.I))
    if copied_whole_turn or (still_command_like and len(query.split()) > 10):
        return getattr(heuristic, "primary_query", "") or query
    compact = plan_search_query(query, web_mode=True, force_search=True, max_variants=1).primary_query
    return compact or query


def _fallback_search_understanding(raw: str, heuristic: Any, *, planner_reason: str = "heuristic_fallback") -> dict[str, Any]:
    return {
        "original_query": raw,
        "should_search": bool(getattr(heuristic, "should_search", False)),
        "primary_query": getattr(heuristic, "primary_query", "") or raw,
        "variants": list(getattr(heuristic, "variants", []) or []),
        "intent": getattr(heuristic, "intent", "general") or "general",
        "reason": planner_reason,
        "user_task": raw,
        "planner": "deterministic",
        "heuristic_plan": _search_plan_dict(heuristic),
    }


def _coerce_llm_search_understanding(raw: str, data: dict[str, Any], heuristic: Any, *, planner_route: str, provider: str, model: str) -> dict[str, Any]:
    should_search = _coerce_bool(data.get("should_search"))
    raw_variants = data.get("queries") or data.get("variants") or []
    if isinstance(raw_variants, str):
        raw_variants = [raw_variants]
    variants: list[str] = []
    primary = _focused_query_from_text(str(data.get("primary_query") or data.get("query") or ""), heuristic)
    for item in raw_variants if isinstance(raw_variants, list) else []:
        focused = _focused_query_from_text(str(item or ""), heuristic)
        if focused and focused.lower() not in {v.lower() for v in variants}:
            variants.append(focused)
    if primary and primary.lower() not in {v.lower() for v in variants}:
        variants.insert(0, primary)
    if should_search and not primary:
        primary = getattr(heuristic, "primary_query", "") or raw
        variants.insert(0, primary)
    variants = [_clean_spaces(v) for v in variants if _clean_spaces(v)][:4]
    return {
        "original_query": raw,
        "should_search": should_search,
        "primary_query": primary,
        "variants": variants or ([primary] if primary else []),
        "intent": str(data.get("intent") or getattr(heuristic, "intent", "general") or "general")[:40],
        "reason": "llm_planned_query" if should_search else "llm_vetoed_search",
        "user_task": _clean_spaces(str(data.get("user_task") or raw))[:400],
        "planner": "llm",
        "planner_provider": provider,
        "planner_model": model,
        "planner_route": planner_route,
        "heuristic_plan": _search_plan_dict(heuristic),
    }


async def _provider_ready_for_llm(provider: str, model: str) -> bool:
    provider = (provider or "").strip().lower()
    if provider == "ollama":
        names = [m["name"] for m in await _ollama_models_raw(timeout=0.7)]
        return bool(names) and ((model or "") in names or not model)
    if provider in PROVIDER_DEFAULTS:
        return _provider_configured(provider)
    return False


async def _plan_search_query_with_llm(raw: str, *, web_mode: bool, provider: str, model: str, heuristic: Any) -> dict[str, Any] | None:
    if not _env_bool("SHIMS_SEARCH_LLM_PLANNER", True):
        return None
    if not await _provider_ready_for_llm(provider, model):
        return None
    payload = {
        "user_message": raw,
        "web_mode": bool(web_mode),
        "heuristic_plan": _search_plan_dict(heuristic),
    }
    messages = [
        {"role": "system", "content": SEARCH_PLANNER_PROMPT},
        {"role": "user", "content": json.dumps(payload, ensure_ascii=False)},
    ]
    try:
        text, route = await _run_llm(provider, model, messages, allow_provider_web_search=False)
    except Exception as exc:
        log_event("search.plan.llm_error", route="search:plan", provider=provider, model=model, ok=False, message=str(exc)[:180])
        return None
    data = _extract_json_object(text)
    if not data:
        log_event("search.plan.llm_bad_json", route="search:plan", provider=provider, model=model, ok=False, message=text[:240])
        return None
    return _coerce_llm_search_understanding(raw, data, heuristic, planner_route=route, provider=provider, model=model)


async def _understand_search_turn(req: ChatRequest) -> dict[str, Any] | None:
    raw = _strip_wake((req.message or "").strip())
    if not raw or _is_empty_greeting(raw) or _detect_tool_intent(raw):
        return None
    heuristic = plan_search_query(raw, web_mode=bool(getattr(req, "web_mode", False)))
    if not heuristic.should_search:
        return None
    provider, model, reason = await _resolve_provider_model(req.provider, req.model, privacy_mode=getattr(req, "privacy_mode", "balanced"), text=raw)
    understood = await _plan_search_query_with_llm(raw, web_mode=bool(getattr(req, "web_mode", False)), provider=provider, model=model, heuristic=heuristic)
    if not understood:
        understood = _fallback_search_understanding(raw, heuristic)
    understood["answer_provider"] = provider
    understood["answer_model"] = model
    understood["answer_model_reason"] = reason
    if not understood.get("should_search"):
        log_event("search.plan.veto", route="search:plan", provider=provider, model=model, ok=True, message=raw, metadata={"search_understanding": understood})
        return None
    if not understood.get("primary_query"):
        return None
    log_event("search.plan", route="search:plan", provider=provider, model=model, ok=True, message=understood.get("primary_query", ""), metadata={"search_understanding": understood})
    return understood


def _format_sources_for_llm(result: dict[str, Any], max_sources: int = 6) -> str:
    lines: list[str] = []
    for i, item in enumerate((result.get("results") or [])[:max_sources], 1):
        title = _clean_spaces(str(item.get("title") or "Untitled"))
        url = _clean_spaces(str(item.get("url") or ""))
        snippet = _clean_spaces(str(item.get("snippet") or ""))
        lines.append(f"[{i}] {title}\nURL: {url}\nSnippet: {snippet}")
    return "\n\n".join(lines)


async def _synthesize_search_answer(req: ChatRequest, result: dict[str, Any], search_plan: dict[str, Any] | None, brain_addendum: str, history: list[dict[str, str]]) -> tuple[str, str]:
    if not result.get("ok"):
        return _format_search_answer(result), "web-search-no-results"
    plan = search_plan or {}
    provider = str(plan.get("answer_provider") or req.provider or "ollama").strip().lower()
    model = str(plan.get("answer_model") or req.model or PROVIDER_DEFAULTS.get(provider, DEFAULT_OLLAMA_MODEL)).strip()
    if not await _provider_ready_for_llm(provider, model):
        return _format_search_answer(result), "web-search-results-only"
    sources = _format_sources_for_llm(result)
    if not sources:
        return _format_search_answer(result), "web-search-no-source-text"
    focused_queries = plan.get("variants") or [result.get("query") or plan.get("primary_query") or ""]
    user_prompt = (
        f"User question:\n{req.message}\n\n"
        f"Interpreted task:\n{plan.get('user_task') or req.message}\n\n"
        f"Focused web queries used:\n{json.dumps(focused_queries, ensure_ascii=False)}\n\n"
        f"Numbered sources:\n{sources}\n\n"
        "Write the answer now."
    )
    messages = [{"role": "system", "content": _system_prompt() + "\n\n" + SEARCH_ANSWER_PROMPT + "\n\n" + brain_addendum}]
    if history:
        messages.extend(history[-8:])
    messages.append({"role": "user", "content": user_prompt})
    try:
        answer, route = await _run_llm(provider, model, messages, allow_provider_web_search=False)
    except Exception as exc:
        log_event("search.answer.llm_error", route="search:answer", provider=provider, model=model, ok=False, message=str(exc)[:180])
        return _format_search_answer(result), "web-search-results-only"
    low = (answer or "").lower()
    if not answer.strip() or "not reachable" in low or "not configured" in low or "needs " in low and "api_key" in low:
        return _format_search_answer(result), "web-search-results-only"
    return answer.strip(), f"web-search-synthesized:{route}"


def _search_provider_status() -> dict[str, Any]:
    web = _settings.get("web", {})
    return {
        "searxng": bool(web.get("searxng_url")),
        "tavily": bool(web.get("tavily_key")),
        "brave": bool(web.get("brave_key")),
        "serpapi": bool(web.get("serpapi_key")),
        "duckduckgo_fallback": bool(web.get("duckduckgo_fallback")),
    }


def _normalize_search_item(title: str = "", url: str = "", snippet: str = "", source: str = "web") -> dict[str, str]:
    return {"title": str(title or "").strip()[:220], "url": str(url or "").strip(), "snippet": str(snippet or "").strip()[:650], "source": source}


async def _search_searxng(query: str, max_results: int) -> list[dict[str, str]]:
    base = (_settings.get("web", {}).get("searxng_url") or "").rstrip("/")
    if not base:
        return []
    async with httpx.AsyncClient(timeout=20, follow_redirects=True) as client:
        r = await client.get(f"{base}/search", params={"q": query, "format": "json", "language": "en", "safesearch": 1})
        r.raise_for_status()
        data = r.json()
    out=[]
    for item in (data.get("results") or [])[:max_results]:
        out.append(_normalize_search_item(item.get("title"), item.get("url"), item.get("content"), "searxng"))
    return out


async def _search_tavily(query: str, max_results: int) -> list[dict[str, str]]:
    key = _settings.get("web", {}).get("tavily_key") or _clean_secret(os.getenv("TAVILY_API_KEY"))
    if not key:
        return []
    payload = {"query": query, "max_results": max_results, "search_depth": "basic", "include_answer": False}
    async with httpx.AsyncClient(timeout=25) as client:
        r = await client.post("https://api.tavily.com/search", headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"}, json=payload)
        r.raise_for_status()
        data = r.json()
    out=[]
    for item in (data.get("results") or [])[:max_results]:
        out.append(_normalize_search_item(item.get("title"), item.get("url"), item.get("content"), "tavily"))
    return out


async def _search_brave(query: str, max_results: int) -> list[dict[str, str]]:
    key = _settings.get("web", {}).get("brave_key") or _clean_secret(os.getenv("BRAVE_SEARCH_API_KEY"))
    if not key:
        return []
    async with httpx.AsyncClient(timeout=20) as client:
        r = await client.get("https://api.search.brave.com/res/v1/web/search", headers={"X-Subscription-Token": key, "Accept": "application/json"}, params={"q": query, "count": min(max_results, 10)})
        r.raise_for_status()
        data = r.json()
    out=[]
    for item in ((data.get("web") or {}).get("results") or [])[:max_results]:
        out.append(_normalize_search_item(item.get("title"), item.get("url"), item.get("description"), "brave"))
    return out


async def _search_serpapi(query: str, max_results: int) -> list[dict[str, str]]:
    key = _settings.get("web", {}).get("serpapi_key") or _clean_secret(os.getenv("SERPAPI_API_KEY"))
    if not key:
        return []
    async with httpx.AsyncClient(timeout=25) as client:
        r = await client.get("https://serpapi.com/search.json", params={"q": query, "api_key": key, "num": max_results})
        r.raise_for_status()
        data = r.json()
    out=[]
    for item in (data.get("organic_results") or [])[:max_results]:
        out.append(_normalize_search_item(item.get("title"), item.get("link"), item.get("snippet"), "serpapi"))
    return out


async def _search_duckduckgo(query: str, max_results: int) -> list[dict[str, str]]:
    if not _settings.get("web", {}).get("duckduckgo_fallback", True):
        return []
    # Lightweight no-key fallback. Prefer SearXNG/Tavily/Brave in production.
    async with httpx.AsyncClient(timeout=20, follow_redirects=True, headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}) as client:
        # Try the lite HTML endpoint first; fall back to main HTML endpoint
        for url in ("https://html.duckduckgo.com/html/", "https://duckduckgo.com/html/"):
            r = await client.get(url, params={"q": query})
            if r.status_code == 200 and "result__a" in r.text:
                break
        html = r.text
    out: list[dict[str, str]] = []
    clean = lambda x: re.sub(r"\s+", " ", re.sub(r"<.*?>", "", x or "")).strip()
    # Primary regex — DDG HTML sometimes omits snippets, so we match title+URL first
    blocks = re.findall(
        r'<a[^>]+class="result__a"[^>]*href="([^"]+)"[^>]*>(.*?)</a>',
        html, flags=re.I | re.S,
    )
    # Try to pair each block with a nearby snippet
    snippet_blocks = re.findall(
        r'<div[^>]+class="result__snippet"[^>]*>(.*?)</div>',
        html, flags=re.I | re.S,
    )
    for i, (url, title) in enumerate(blocks):
        snippet = clean(snippet_blocks[i]) if i < len(snippet_blocks) else ""
        out.append(_normalize_search_item(clean(title), clean(url), snippet, "duckduckgo"))
    # Fallback regex (alternative DDG layouts / Lite mode)
    if not out:
        blocks2 = re.findall(
            r'<a[^>]+href="([^"]+)"[^>]*class="[^"]*result[^"]*"[^>]*>(.*?)</a>',
            html, flags=re.I | re.S,
        )
        for url, title in blocks2:
            out.append(_normalize_search_item(clean(title), clean(url), "", "duckduckgo"))
    return out[:max_results]


async def _web_search(query: str, max_results: int = 6, provider: str | None = None, planned_query: dict[str, Any] | None = None) -> dict[str, Any]:
    started = time.perf_counter()
    provider = (provider or "auto").strip().lower()
    if planned_query:
        plan = dict(planned_query)
        original_query = str(plan.get("original_query") or query)
        variants = plan.get("variants") or plan.get("queries") or []
        if isinstance(variants, str):
            variants = [variants]
        candidate_queries = [_clean_spaces(str(v or "")) for v in variants]
        primary = _clean_spaces(str(plan.get("primary_query") or query))
        if primary and primary.lower() not in {q.lower() for q in candidate_queries if q}:
            candidate_queries.insert(0, primary)
    else:
        planned = plan_search_query(query, web_mode=True, force_search=True, max_variants=4)
        plan = planned.to_dict()
        original_query = query
        candidate_queries = list(planned.variants or [planned.primary_query or query])
    candidate_queries = [q for q in candidate_queries if q][:4] or [query]
    errors=[]
    attempts=[]
    order=[]
    if provider in {"auto", "searxng"}: order.append(("searxng", _search_searxng))
    if provider in {"auto", "tavily"}: order.append(("tavily", _search_tavily))
    if provider in {"auto", "brave"}: order.append(("brave", _search_brave))
    if provider in {"auto", "serpapi"}: order.append(("serpapi", _search_serpapi))
    if provider in {"auto", "duckduckgo"}: order.append(("duckduckgo", _search_duckduckgo))
    for candidate in candidate_queries:
        candidate = (candidate or "").strip()
        if not candidate:
            continue
        for name, fn in order:
            attempts.append({"provider": name, "query": candidate})
            try:
                results = await fn(candidate, max(1, min(max_results, 10)))
                if results:
                    log_event("web.search", route="tool:web_search", provider=name, model="search", ok=True, latency_ms=(time.perf_counter()-started)*1000, message=candidate, metadata={"count": len(results), "original_query": original_query, "query_plan": plan})
                    return {"ok": True, "query": candidate, "original_query": original_query, "query_plan": plan, "provider": name, "results": results, "attempts": attempts, "latency_ms": round((time.perf_counter()-started)*1000, 1)}
            except Exception as exc:
                errors.append(f"{name}({candidate}): {str(exc)[:180]}")
    primary = str(plan.get("primary_query") or query)
    log_event("web.search.error", route="tool:web_search", provider=provider, model="search", ok=False, latency_ms=(time.perf_counter()-started)*1000, message=primary, metadata={"errors": errors, "original_query": original_query, "query_plan": plan, "attempts": attempts})
    return {"ok": False, "query": primary, "original_query": original_query, "query_plan": plan, "provider": provider, "results": [], "attempts": attempts, "errors": errors, "status": _search_provider_status(), "message": "No web search provider returned results. Configure SHIMS_SEARXNG_URL or a Tavily/Brave/SerpAPI key, or check internet connectivity."}


async def _run_web_search_with_plan(query: str, max_results: int = 6, provider: str | None = None, planned_query: dict[str, Any] | None = None) -> dict[str, Any]:
    try:
        return await _web_search(query, max_results=max_results, provider=provider, planned_query=planned_query)
    except TypeError as exc:
        if planned_query is not None and "planned_query" in str(exc):
            return await _web_search(query, max_results=max_results, provider=provider)
        raise


def _format_search_answer(result: dict[str, Any]) -> str:
    if not result.get("ok"):
        return result.get("message") or "Web search is not available right now."
    lines = [f"I searched the web for: {result.get('query')}", f"Provider: {result.get('provider')}. Top results:"]
    for i, item in enumerate(result.get("results") or [], 1):
        title = item.get("title") or "Untitled"
        url = item.get("url") or ""
        snippet = item.get("snippet") or ""
        lines.append(f"{i}. {title}\n   {url}\n   {snippet}")
    return "\n".join(lines)

def _detect_agent_route(text: str) -> str:
    """Determine which agent should own this turn based on message content."""
    low = (text or "").lower()
    if _detect_tool_intent(text):
        return "media" if any(k in low for k in ("image", "photo", "picture", "video", "audio")) else "documents"
    if _detect_search_intent(text):
        return "search"
    if any(k in low for k in ("smiles", "reaction", "retrosynthesis", "hazard", "impurity", "fto", "molecule", "compound")):
        return "chemistry"
    if any(k in low for k in ("patent", "synthesis", "process", "yield", "purity", "raw material", "pricing")):
        return "rd"
    if any(k in low for k in ("enterprise", "batch", "qc", "qa", "coa", "sop", "capa", "change control")):
        return "enterprise_bridge"
    return "supervisor"


async def _make_plan(req: ChatRequest) -> TurnPlan:
    session_id = req.session_id or str(uuid.uuid4())
    duplicate = _guard_duplicate(session_id, req.message, req.source)
    agent = _detect_agent_route(req.message)
    intent = _detect_tool_intent(req.message)
    if intent:
        kind, prompt = intent
        return TurnPlan(session_id=session_id, provider="tool", model="", route=f"tool:{kind}", tool_kind=kind, tool_prompt=prompt, duplicate=duplicate, agent=agent)
    if _is_empty_greeting(req.message):
        return TurnPlan(session_id=session_id, provider="local", model="", route="local:greeting", duplicate=duplicate, agent=agent)
    search_understanding = await _understand_search_turn(req)
    if search_understanding:
        return TurnPlan(
            session_id=session_id,
            provider=search_understanding.get("answer_provider") or "tool",
            model=search_understanding.get("answer_model") or "",
            route="tool:web_search",
            tool_kind="web_search",
            tool_prompt=search_understanding.get("primary_query") or req.message,
            tool_metadata=search_understanding,
            duplicate=duplicate,
            agent=agent,
        )
    # Agent-aware defaults: if the user left provider/model on auto, pick a
    # specialist model for the detected agent (e.g. coder → SHIMS_CODER_MODEL).
    agent_default_provider, agent_default_model, _ = agent_model_router.resolve_agent(agent)
    requested_provider = (req.provider or "").strip().lower()
    requested_model = (req.model or "").strip()
    if requested_provider in ("", "auto"):
        requested_provider = agent_default_provider
    if requested_model in ("", "auto"):
        requested_model = agent_default_model
    provider, model, reason = await _resolve_provider_model(requested_provider, requested_model, privacy_mode=req.privacy_mode, text=req.message)
    return TurnPlan(session_id=session_id, provider=provider, model=model, route=f"llm:{provider}:{reason}", duplicate=duplicate, agent=agent)


async def _ollama_chat(model: str, messages: list[dict[str, str]], *, realtime: bool = False, max_tokens: int | None = None) -> str:
    payload = {
        "model": model,
        "messages": messages,
        "stream": False,
        "options": _ollama_options(realtime=realtime, max_tokens=max_tokens),
        "keep_alive": _settings["brain"].get("keep_alive", "30m"),
    }
    timeout = 300
    async with httpx.AsyncClient(timeout=timeout) as client:
        r = await client.post(f"{OLLAMA_HOST}/api/chat", json=payload)
        r.raise_for_status()
        data = r.json()
    return (data.get("message") or {}).get("content") or data.get("response") or ""




async def _ollama_chat_stream(model: str, messages: list[dict[str, str]], *, realtime: bool = False, max_tokens: int | None = None) -> AsyncGenerator[str, None]:
    payload = {
        "model": model,
        "messages": messages,
        "stream": True,
        "options": _ollama_options(realtime=realtime, max_tokens=max_tokens),
        "keep_alive": _settings["brain"].get("keep_alive", "30m"),
    }
    read_timeout = 240.0
    async with httpx.AsyncClient(timeout=httpx.Timeout(300.0, read=read_timeout)) as client:
        async with client.stream("POST", f"{OLLAMA_HOST}/api/chat", json=payload) as r:
            r.raise_for_status()
            async for line in r.aiter_lines():
                if not line:
                    continue
                try:
                    data = json.loads(line)
                except Exception:
                    continue
                delta = (data.get("message") or {}).get("content") or data.get("response") or ""
                if delta:
                    yield delta
                if data.get("done"):
                    break


def _hf_headers() -> dict[str, str]:
    headers = {"Content-Type": "application/json"}
    key = _clean_secret(os.getenv("HUGGINGFACE_API_KEY"))
    if key:
        headers["Authorization"] = f"Bearer {key}"
    return headers


def _hf_payload(model: str, messages: list[dict[str, str]], *, stream: bool = False, max_tokens: int | None = None) -> dict[str, Any]:
    brain = _settings["brain"]
    payload: dict[str, Any] = {
        "model": model,
        "messages": messages,
        "stream": stream,
        "temperature": float(brain.get("temperature", 0.12)),
        "top_p": float(brain.get("top_p", 0.82)),
    }
    if max_tokens:
        payload["max_tokens"] = int(max_tokens)
    return payload


async def _hf_chat(model: str, messages: list[dict[str, str]], *, realtime: bool = False, max_tokens: int | None = None) -> str:
    async with httpx.AsyncClient(timeout=300) as client:
        r = await client.post(
            f"{HUGGINGFACE_HOST}/v1/chat/completions",
            headers=_hf_headers(),
            json=_hf_payload(model, messages, stream=False, max_tokens=max_tokens),
        )
        r.raise_for_status()
        data = r.json()
    choice = (data.get("choices") or [{}])[0]
    return (choice.get("message") or {}).get("content") or ""


async def _hf_chat_stream(model: str, messages: list[dict[str, str]], *, realtime: bool = False, max_tokens: int | None = None) -> AsyncGenerator[str, None]:
    read_timeout = 240.0
    async with httpx.AsyncClient(timeout=httpx.Timeout(300.0, read=read_timeout)) as client:
        async with client.stream(
            "POST",
            f"{HUGGINGFACE_HOST}/v1/chat/completions",
            headers=_hf_headers(),
            json=_hf_payload(model, messages, stream=True, max_tokens=max_tokens),
        ) as r:
            r.raise_for_status()
            async for line in r.aiter_lines():
                if not line or not line.startswith("data: "):
                    continue
                payload_str = line[6:].strip()
                if payload_str == "[DONE]":
                    break
                try:
                    obj = json.loads(payload_str)
                except Exception:
                    continue
                choice = (obj.get("choices") or [{}])[0]
                delta = (choice.get("delta") or {}).get("content") or ""
                if delta:
                    yield delta


async def _extract_durable_facts_llm(
    user_msg: str,
    assistant_msg: str,
    tools_used: list[str] | None = None,
) -> list[tuple[str, list[str]]]:
    """Ask a cheap local model to extract durable facts worth remembering."""
    if len((user_msg or "") + (assistant_msg or "")) < 60:
        return []

    system_prompt = (
        "You are a durable-memory extractor for a personal AI assistant. "
        "Analyze the exchange and return ONLY a JSON array of facts worth remembering. "
        "No markdown, no explanation."
    )
    user_prompt = (
        f"User: {user_msg[:1500]}\n"
        f"Assistant: {assistant_msg[:1500]}\n"
        f"Tools used: {', '.join(tools_used or [])}\n\n"
        "Extract 0-3 durable facts (user preferences, identity, goals, recurring projects, or hard-won outcomes).\n"
        "Rules:\n"
        "- Skip greetings, transient questions, weather, time, small talk.\n"
        "- Tags must be chosen from: preference, user, goal, project, tool_result, code, plan, assistant_note.\n"
        "- Return [] if nothing is worth remembering.\n"
        'Example: [{"fact": "User prefers concise answers", "tags": ["preference", "user"]}]'
    )

    model = os.getenv("SHIMS_MEMORY_MODEL") or _preferred_local_model([], realtime=True)
    try:
        raw = await _ollama_chat(
            model,
            [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            max_tokens=512,
        )
    except Exception:
        return []

    # Try to isolate JSON array
    text = (raw or "").strip()
    if "[" in text and "]" in text:
        text = text[text.find("[") : text.rfind("]") + 1]
    if not text:
        return []

    facts: list[tuple[str, list[str]]] = []
    try:
        parsed = json.loads(text)
        if not isinstance(parsed, list):
            return []
        for item in parsed:
            if not isinstance(item, dict):
                continue
            fact = str(item.get("fact") or item.get("content") or "").strip()
            tags = item.get("tags") or []
            if not isinstance(tags, list):
                tags = [str(tags)]
            tags = [str(t).strip().lower() for t in tags if t]
            if not fact or len(fact) < 8:
                continue
            # Enforce allowed tags
            allowed = {"preference", "user", "goal", "project", "tool_result", "code", "plan", "assistant_note"}
            tags = [t for t in tags if t in allowed] or ["assistant_note"]
            facts.append((fact, tags))
    except Exception:
        return []
    return facts


async def _openai_chat(model: str, messages: list[dict[str, str]], *, allow_web_search: bool = False) -> str:
    api_key = _clean_secret(os.getenv("OPENAI_API_KEY"))
    if not api_key:
        return "OpenAI is selected but OPENAI_API_KEY is not configured. Select an installed Ollama model or add the key in Settings."
    payload: dict[str, Any] = {
        "model": model,
        "input": messages,
        "max_output_tokens": 32000,
    }
    if allow_web_search and os.getenv("SHIMS_OPENAI_RESPONSES_WEB_SEARCH", "false").strip().lower() in {"1", "true", "yes", "on"}:
        payload["tools"] = [{"type": "web_search", "search_context_size": "low"}]
    async with httpx.AsyncClient(timeout=120) as client:
        r = await client.post("https://api.openai.com/v1/responses", headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}, json=payload)
        if r.status_code >= 400:
            # Compatibility fallback for older model/account combinations.
            chat_payload = {"model": model, "messages": messages, "temperature": 0.2, "max_tokens": 32000}
            r = await client.post("https://api.openai.com/v1/chat/completions", headers={"Authorization": f"Bearer {api_key}"}, json=chat_payload)
        r.raise_for_status()
        data = r.json()
    text = data.get("output_text")
    if not text:
        chunks: list[str] = []
        for item in data.get("output", []):
            for content in item.get("content", []) if isinstance(item, dict) else []:
                if isinstance(content, dict) and content.get("type") in {"output_text", "text"}:
                    chunks.append(content.get("text", ""))
        text = "\n".join(chunks)
    if not text and data.get("choices"):
        text = data["choices"][0]["message"]["content"]
    return (text or "").strip()


async def _anthropic_chat(model: str, messages: list[dict[str, str]]) -> str:
    api_key = _clean_secret(os.getenv("ANTHROPIC_API_KEY"))
    if not api_key:
        return f"Claude model `{model}` needs ANTHROPIC_API_KEY. Select an installed Ollama model like llama3.2:latest / qwen2.5:7b or add a valid key in Settings."
    user_msgs = [{"role": m["role"], "content": m["content"]} for m in messages if m["role"] in {"user", "assistant"}][-50:]
    payload = {"model": model, "max_tokens": settings.max_output_tokens, "temperature": 0.2, "system": SYSTEM_PROMPT, "messages": user_msgs}
    async with httpx.AsyncClient(timeout=120) as client:
        r = await client.post("https://api.anthropic.com/v1/messages", headers={"x-api-key": api_key, "anthropic-version": "2023-06-01", "content-type": "application/json"}, json=payload)
        if r.status_code == 401:
            _set_env_persistent("SHIMS_LAST_ANTHROPIC_STATUS", "401")
            return "Anthropic rejected this API key with 401 Unauthorized. I will not keep retrying it automatically. Choose an installed Ollama model from the model picker or paste a valid Anthropic key in Settings and click Test."
        r.raise_for_status()
        data = r.json()
    parts = data.get("content") or []
    return "".join(p.get("text", "") for p in parts if isinstance(p, dict))


async def _gemini_chat(model: str, messages: list[dict[str, str]]) -> str:
    api_key = _clean_secret(os.getenv("GEMINI_API_KEY"))
    if not api_key:
        return f"Gemini model `{model}` needs GEMINI_API_KEY. Add it in Settings or choose Local/Ollama."
    system_parts = []
    contents: list[dict[str, Any]] = []
    for msg in messages[-16:]:
        role = msg.get("role", "user")
        text = msg.get("content", "")
        if not text:
            continue
        if role == "system":
            system_parts.append({"text": text})
        elif role == "assistant":
            contents.append({"role": "model", "parts": [{"text": text}]})
        else:
            contents.append({"role": "user", "parts": [{"text": text}]})
    payload: dict[str, Any] = {
        "contents": contents or [{"role": "user", "parts": [{"text": ""}]}],
        "generationConfig": {
            "temperature": float(_settings["brain"].get("temperature", 0.12)),
            "topP": float(_settings["brain"].get("top_p", 0.82)),
        },
    }
    if system_parts:
        payload["system_instruction"] = {"parts": system_parts}
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
    async with httpx.AsyncClient(timeout=120) as client:
        try:
            r = await client.post(url, headers={"x-goog-api-key": api_key, "Content-Type": "application/json"}, json=payload)
            r.raise_for_status()
        except httpx.HTTPStatusError as exc:
            detail = ""
            try:
                detail = exc.response.text[:400]
            except Exception:
                pass
            return f"Gemini API error ({exc.response.status_code}): {detail or exc.response.reason_phrase}. Check GEMINI_API_KEY, billing, and that model `{model}` is available for your project/region."
        data = r.json()
    chunks: list[str] = []
    for cand in data.get("candidates") or []:
        content = cand.get("content") or {}
        for part in content.get("parts") or []:
            if isinstance(part, dict) and part.get("text"):
                chunks.append(part["text"])
    return "\n".join(chunks).strip() or "Gemini returned an empty response."


async def _openai_compatible_chat(provider: str, model: str, messages: list[dict[str, str]]) -> str:
    env = PROVIDER_ENV.get(provider)
    api_key = _clean_secret(os.getenv(env or ""))
    if not api_key:
        return f"{provider.title()} model `{model}` needs {env or 'an API key'}. Add it in Settings or choose Local/Ollama."
    base_urls = {
        "kimi": os.getenv("KIMI_BASE_URL", "https://api.moonshot.ai/v1").rstrip("/"),
        "deepseek": os.getenv("DEEPSEEK_BASE_URL", "https://api.deepseek.com/v1").rstrip("/"),
        "qwen": os.getenv("QWEN_BASE_URL", "https://dashscope.aliyuncs.com/compatible-mode/v1").rstrip("/"),
    }
    base = base_urls.get(provider)
    if not base:
        return await _cloud_placeholder(provider, model)
    # kimi-k2 series only accepts temperature=1.0.
    temperature = 1.0 if provider == "kimi" and isinstance(model, str) and model.startswith("kimi-k2") else 0.2
    payload = {"model": model, "messages": messages, "temperature": temperature, "stream": False, "max_tokens": 32000}
    async with httpx.AsyncClient(timeout=120) as client:
        try:
            r = await client.post(f"{base}/chat/completions", headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}, json=payload)
            r.raise_for_status()
        except httpx.HTTPStatusError as exc:
            detail = ""
            try:
                detail = exc.response.text[:400]
            except Exception:
                pass
            if exc.response.status_code == 429:
                return f"{provider.title()} API rate limit hit (429). Wait a moment and try again, or switch to a local Ollama model."
            if exc.response.status_code == 401:
                return f"{provider.title()} rejected this API key with 401 Unauthorized. Check {env} in Settings."
            if exc.response.status_code == 404:
                return f"{provider.title()} model `{model}` not found (404). Check the model name in Settings."
            return f"{provider.title()} API error ({exc.response.status_code}): {detail or exc.response.reason_phrase}. Check {env}, billing, and that model `{model}` is available for your account/region."
        except httpx.TimeoutException:
            return f"{provider.title()} request timed out. The provider may be slow or unreachable. Try again or use a local Ollama model."
        except httpx.RequestError as exc:
            return f"{provider.title()} connection failed: {str(exc)[:200]}. Check your internet connection or use a local Ollama model."
        data = r.json()
    return ((data.get("choices") or [{}])[0].get("message") or {}).get("content", "").strip() or f"{provider.title()} returned an empty response."


async def _cloud_placeholder(provider: str, model: str) -> str:
    env = PROVIDER_ENV.get(provider, "API_KEY")
    return f"{provider.title()} model `{model}` is selected, but {env} is not configured in this local build. Choose an installed Ollama model or add the API key in Settings."


async def _run_llm(
    provider: str,
    model: str,
    messages: list[dict[str, str]],
    *,
    allow_provider_web_search: bool = False,
    realtime: bool = False,
    max_tokens: int | None = None,
) -> tuple[str, str]:
    if provider == "ollama":
        names = await _ollama_names()
        if model not in names:
            if names:
                fallback = _preferred_local_model(names, realtime=realtime)
                return (await _ollama_chat(fallback, messages, realtime=realtime, max_tokens=max_tokens), f"ollama-local-fallback:{fallback}")
            return (f"Ollama is not reachable at {OLLAMA_HOST}. Start Ollama from Settings or run `ollama serve`, then pull `llama3.2:latest` or `qwen2.5:7b`. Local PDF/image/audio tools still work.", "ollama-offline")
        return (await _ollama_chat(model, messages, realtime=realtime, max_tokens=max_tokens), "ollama-local")
    if provider == "openai":
        return (await _openai_chat(model, messages, allow_web_search=allow_provider_web_search), "openai")
    if provider == "anthropic":
        return (await _anthropic_chat(model, messages), "anthropic")
    if provider == "gemini":
        return (await _gemini_chat(model, messages), "gemini")
    if provider == "huggingface":
        return (await _hf_chat(model, messages, realtime=realtime, max_tokens=max_tokens), "huggingface-local")
    if provider in {"kimi", "deepseek", "qwen"}:
        return (await _openai_compatible_chat(provider, model, messages), provider)
    return (await _cloud_placeholder(provider, model), f"{provider}-placeholder")


async def _openai_image(prompt: str) -> dict[str, Any] | None:
    api_key = _clean_secret(os.getenv("OPENAI_API_KEY"))
    if not api_key:
        return None
    try:
        payload = {"model": os.getenv("OPENAI_IMAGE_MODEL", "gpt-image-1"), "prompt": prompt, "size": os.getenv("OPENAI_IMAGE_SIZE", "1024x1024")}
        async with httpx.AsyncClient(timeout=180) as client:
            r = await client.post("https://api.openai.com/v1/images/generations", headers={"Authorization": f"Bearer {api_key}"}, json=payload)
            r.raise_for_status()
            data = r.json()
        item = data.get("data", [{}])[0]
        filename = _safe_name(prompt or "image", "png")
        path = IMAGE_DIR / filename
        if item.get("b64_json"):
            path.write_bytes(base64.b64decode(item["b64_json"]))
        elif item.get("url"):
            async with httpx.AsyncClient(timeout=180) as client:
                img = await client.get(item["url"])
                img.raise_for_status()
                path.write_bytes(img.content)
        else:
            return None
        url = f"/media/files/images/{filename}"
        result = {"ok": True, "provider": "openai", "type": "image", "kind": "image", "title": prompt[:80] or "Generated image", "filename": filename, "url": url, "file_url": url, "download_url": url}
        return _attach_ledger(result, path, "image")
    except Exception as exc:
        return {"ok": False, "provider": "openai", "error": str(exc)[:240]}


async def _qwen_image(prompt: str) -> dict[str, Any] | None:
    """Generate image via Alibaba DashScope Wanx model."""
    api_key = _clean_secret(os.getenv("QWEN_API_KEY"))
    if not api_key:
        return None
    try:
        payload = {
            "model": "wanx-v1",
            "input": {"prompt": prompt},
            "parameters": {"size": "1024*1024", "n": 1, "seed": (hash(prompt) % 2**31)},
        }
        async with httpx.AsyncClient(timeout=180) as client:
            r = await client.post(
                "https://dashscope.aliyuncs.com/api/v1/services/aigc/text2image/image-synthesis",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json", "X-DashScope-Async": "enable"},
                json=payload,
            )
            r.raise_for_status()
            job = r.json()
            task_id = job.get("output", {}).get("task_id")
            if not task_id:
                return {"ok": False, "provider": "qwen", "error": "No task_id from DashScope image API"}
            # Poll for result
            poll_deadline = time.time() + 120
            while time.time() < poll_deadline:
                await asyncio.sleep(3)
                status_r = await client.get(
                    f"https://dashscope.aliyuncs.com/api/v1/tasks/{task_id}",
                    headers={"Authorization": f"Bearer {api_key}"},
                )
                status_r.raise_for_status()
                status = status_r.json()
                task_status = status.get("output", {}).get("task_status", "")
                if task_status == "SUCCEEDED":
                    results = status.get("output", {}).get("results", [])
                    if results and results[0].get("url"):
                        img_url = results[0]["url"]
                        img = await client.get(img_url)
                        img.raise_for_status()
                        filename = _safe_name(prompt or "qwen_image", "png")
                        path = IMAGE_DIR / filename
                        path.write_bytes(img.content)
                        url = f"/media/files/images/{filename}"
                        result = {"ok": True, "provider": "qwen", "type": "image", "kind": "image", "title": prompt[:80] or "Generated image", "filename": filename, "url": url, "file_url": url, "download_url": url}
                        return _attach_ledger(result, path, "image")
                elif task_status == "FAILED":
                    return {"ok": False, "provider": "qwen", "error": status.get("output", {}).get("message", "DashScope image task failed")}
            return {"ok": False, "provider": "qwen", "error": "DashScope image generation timed out"}
    except Exception as exc:
        return {"ok": False, "provider": "qwen", "error": str(exc)[:260]}


async def _sd_webui_image(prompt: str) -> dict[str, Any] | None:
    base = _settings["media"].get("stable_diffusion_url") or ""
    if not base:
        return None
    try:
        payload = {"prompt": prompt, "steps": 25, "width": 768, "height": 512}
        async with httpx.AsyncClient(timeout=240) as client:
            r = await client.post(f"{base}/sdapi/v1/txt2img", json=payload)
            r.raise_for_status()
            data = r.json()
        img_b64 = (data.get("images") or [None])[0]
        if not img_b64:
            return None
        filename = _safe_name(prompt or "image", "png")
        path = IMAGE_DIR / filename
        path.write_bytes(base64.b64decode(img_b64.split(",")[-1]))
        url = f"/media/files/images/{filename}"
        result = {"ok": True, "provider": "stable-diffusion-webui", "type": "image", "kind": "image", "title": prompt[:80] or "Generated image", "filename": filename, "url": url, "file_url": url, "download_url": url}
        return _attach_ledger(result, path, "image")
    except Exception as exc:
        return {"ok": False, "provider": "stable-diffusion-webui", "error": str(exc)[:240]}


def _diffusers_cpu_guard_reason(model_id: str, device: str) -> str:
    if device != "cpu" or _env_bool("SHIMS_DIFFUSERS_ALLOW_SLOW_CPU", False):
        return ""
    if "xl" not in (model_id or "").lower():
        return ""
    return (
        "CUDA is not available and SDXL on CPU is too slow for the live Omni app "
        "(recent runs can take tens of minutes). Use Stable Diffusion WebUI/ComfyUI, "
        "OpenAI image generation, a CUDA PyTorch install, or set "
        "SHIMS_DIFFUSERS_ALLOW_SLOW_CPU=true to force the slow local path."
    )


def _diffusers_runtime_options(model_id: str, device: str) -> dict[str, Any]:
    is_xl = "xl" in (model_id or "").lower()
    if device == "cuda":
        return {
            "num_inference_steps": _env_int("SHIMS_DIFFUSERS_STEPS", 25),
            "width": _env_int("SHIMS_DIFFUSERS_WIDTH", 1024 if is_xl else 768),
            "height": _env_int("SHIMS_DIFFUSERS_HEIGHT", 1024 if is_xl else 768),
        }
    return {
        "num_inference_steps": _env_int("SHIMS_DIFFUSERS_CPU_STEPS", 6),
        "width": _env_int("SHIMS_DIFFUSERS_CPU_WIDTH", 512),
        "height": _env_int("SHIMS_DIFFUSERS_CPU_HEIGHT", 512),
    }


def _load_diffusers_pipe(model_id: str, device: str, dtype: Any) -> Any:
    key = (model_id, device)
    with _diffusers_lock:
        cached = _diffusers_pipe_cache.get(key)
        if cached is not None:
            return cached
        if "xl" in model_id.lower():
            from diffusers import StableDiffusionXLPipeline  # type: ignore
            try:
                pipe = StableDiffusionXLPipeline.from_pretrained(model_id, torch_dtype=dtype, variant="fp16")
            except Exception:
                pipe = StableDiffusionXLPipeline.from_pretrained(model_id, torch_dtype=dtype)
        else:
            from diffusers import StableDiffusionPipeline  # type: ignore
            try:
                pipe = StableDiffusionPipeline.from_pretrained(model_id, torch_dtype=dtype, variant="fp16")
            except Exception:
                pipe = StableDiffusionPipeline.from_pretrained(model_id, torch_dtype=dtype)
        try:
            pipe.set_progress_bar_config(disable=True)
        except Exception:
            pass
        if device == "cuda" and _env_bool("SHIMS_ENABLE_XFORMERS", False):
            try:
                pipe.enable_xformers_memory_efficient_attention()
            except Exception:
                pass
        else:
            try:
                pipe.enable_attention_slicing("max")
            except Exception:
                pass
        for method in ("enable_vae_slicing", "enable_vae_tiling"):
            try:
                getattr(pipe, method)()
            except Exception:
                pass
        pipe = pipe.to(device)
        _diffusers_pipe_cache[key] = pipe
        return pipe


def _diffusers_image_sync(prompt: str) -> dict[str, Any]:
    import torch  # type: ignore

    model_id = _settings["media"].get("diffusers_model") or "stabilityai/stable-diffusion-xl-base-1.0"
    device = "cuda" if torch.cuda.is_available() else "cpu"
    guard = _diffusers_cpu_guard_reason(model_id, device)
    if guard:
        return {"ok": False, "provider": "diffusers", "error": guard}
    dtype = torch.float16 if device == "cuda" else torch.float32
    pipe = _load_diffusers_pipe(model_id, device, dtype)
    options = _diffusers_runtime_options(model_id, device)
    with torch.inference_mode():
        image = pipe(prompt, **options).images[0]
    filename = _safe_name(prompt or "image", "png")
    path = IMAGE_DIR / filename
    image.save(path)
    url = f"/media/files/images/{filename}"
    result = {
        "ok": True,
        "provider": "diffusers",
        "type": "image",
        "kind": "image",
        "title": prompt[:80] or "Generated image",
        "filename": filename,
        "url": url,
        "file_url": url,
        "download_url": url,
        "model": model_id,
        "device": device,
        "steps": options.get("num_inference_steps"),
        "width": options.get("width"),
        "height": options.get("height"),
    }
    return _attach_ledger(result, path, "image")


async def _diffusers_image(prompt: str) -> dict[str, Any] | None:
    if not _settings["media"].get("diffusers_enabled"):
        return None
    try:
        return await asyncio.to_thread(_diffusers_image_sync, prompt)
    except Exception as exc:
        return {"ok": False, "provider": "diffusers", "error": str(exc)[:240]}


def _enhance_image_prompt(prompt: str, theme: str | None = None, quality: str = "standard") -> str:
    base = _clean_spaces(prompt or "SHIMS generated image")
    low = base.lower()
    style_bits: list[str] = []
    if not any(k in low for k in ("photo", "illustration", "render", "poster", "logo", "icon", "watercolor", "anime", "oil painting")):
        style_bits.append("high-detail cinematic digital illustration")
    if theme and theme not in {"auto", "standard"}:
        style_bits.append(f"{theme} visual style")
    if quality in {"3", "high", "hd", "premium"}:
        style_bits.extend(["sharp composition", "rich materials", "professional lighting", "high resolution"])
    elif quality in {"1", "draft", "fast"}:
        style_bits.append("clean fast concept art")
    else:
        style_bits.extend(["clean composition", "balanced lighting"])
    if not any(k in low for k in ("text", "typography", "label", "poster", "logo")):
        style_bits.append("no legible text or watermark")
    return base + ". " + ", ".join(dict.fromkeys(style_bits))


async def _pollinations_image(prompt: str) -> dict[str, Any] | None:
    """Free image generation via Pollinations.ai — no API key required."""
    try:
        import urllib.parse
        encoded = urllib.parse.quote(prompt[:1000])
        url = f"https://image.pollinations.ai/prompt/{encoded}?width=1024&height=1024&nologo=true&seed={abs(hash(prompt)) % 99999}"
        async with httpx.AsyncClient(timeout=120) as client:
            r = await client.get(url)
            r.raise_for_status()
            filename = _safe_name(prompt, "png")
            path = IMAGE_DIR / filename
            path.write_bytes(r.content)
            file_url = f"/media/files/image/{filename}"
            return {"ok": True, "provider": "pollinations", "type": "image", "kind": "image", "title": prompt[:80], "filename": filename, "url": file_url, "file_url": file_url, "download_url": file_url}
    except Exception as exc:
        return {"ok": False, "error": f"pollinations: {str(exc)[:200]}"}


def _image_provider_plan(backend: str) -> list[tuple[str, Callable[[str], Awaitable[dict[str, Any] | None]]]]:
    backend = (backend or "auto").lower()
    providers = {
        "stable-diffusion": _sd_webui_image,
        "openai": _openai_image,
        "diffusers": _diffusers_image,
        "qwen": _qwen_image,
        "pollinations": _pollinations_image,
    }
    if backend in {"stable-diffusion", "sd", "sdwebui"}:
        return [("stable-diffusion", providers["stable-diffusion"])]
    if backend == "openai":
        return [("openai", providers["openai"])]
    if backend == "qwen":
        return [("qwen", providers["qwen"])]
    if backend in {"diffusers", "local"}:
        return [("diffusers", providers["diffusers"])]
    if backend == "pollinations":
        return [("pollinations", providers["pollinations"])]
    plan: list[tuple[str, Callable[[str], Awaitable[dict[str, Any] | None]]]] = []
    if _settings["media"].get("stable_diffusion_url"):
        plan.append(("stable-diffusion", providers["stable-diffusion"]))
    if _clean_secret(os.getenv("OPENAI_API_KEY")):
        plan.append(("openai", providers["openai"]))
    if _clean_secret(os.getenv("QWEN_API_KEY")):
        plan.append(("qwen", providers["qwen"]))
    if _settings["media"].get("diffusers_enabled"):
        plan.append(("diffusers", providers["diffusers"]))
    plan.append(("pollinations", providers["pollinations"]))
    if not plan:
        plan = [("pollinations", providers["pollinations"]), ("stable-diffusion", providers["stable-diffusion"]), ("diffusers", providers["diffusers"])]
    return plan


async def _fallback_image(prompt: str, note: str = "Local fallback renderer used. Configure Stable Diffusion/OpenAI for photoreal generation.") -> dict[str, Any]:
    import random
    from PIL import Image, ImageDraw, ImageFilter

    filename = _safe_name(prompt or "image", "png")
    path = IMAGE_DIR / filename
    w, h = 1280, 720
    low = (prompt or "").lower()
    seed = int(hashlib.sha256((prompt or "shims-image" + str(time.time())).encode("utf-8")).hexdigest()[:16], 16)
    rng = random.Random(seed)
    light_scene = any(k in low for k in ("white table", "clean white", "product render", "studio", "table"))

    top = (236, 246, 248) if light_scene else (5, 10, 24)
    bottom = (192, 221, 226) if light_scene else (10, 28, 52)
    img = Image.new("RGB", (w, h), top)
    d = ImageDraw.Draw(img)
    for y in range(h):
        t = y / max(1, h - 1)
        col = tuple(int(top[i] * (1 - t) + bottom[i] * t) for i in range(3))
        d.line([(0, y), (w, y)], fill=col)

    def composite(layer: Image.Image, blur: int = 0) -> None:
        nonlocal img
        if blur:
            layer = layer.filter(ImageFilter.GaussianBlur(blur))
        img = Image.alpha_composite(img.convert("RGBA"), layer).convert("RGB")

    def glow(cx: int, cy: int, radius: int, color: tuple[int, int, int], alpha: int = 95) -> None:
        layer = Image.new("RGBA", (w, h), (0, 0, 0, 0))
        ld = ImageDraw.Draw(layer)
        ld.ellipse((cx - radius, cy - radius, cx + radius, cy + radius), fill=(*color, alpha))
        composite(layer, blur=max(18, radius // 3))

    for _ in range(9):
        cx, cy = rng.randrange(70, w - 70), rng.randrange(40, h - 120)
        color = rng.choice([(124, 240, 255), (77, 208, 225), (255, 214, 102), (154, 225, 255)])
        glow(cx, cy, rng.randrange(35, 120), color, rng.randrange(18, 55))

    d = ImageDraw.Draw(img)
    grid_col = (210, 232, 235) if light_scene else (12, 42, 76)
    for x in range(-160, w, 48):
        d.line([(x, 0), (x + 380, h)], fill=grid_col, width=1)
    for y in range(0, h, 48):
        d.line([(0, y), (w, y)], fill=grid_col, width=1)

    table_y = 468
    shadow = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    sd = ImageDraw.Draw(shadow)
    sd.ellipse((250, table_y + 80, 1030, table_y + 185), fill=(0, 0, 0, 70))
    composite(shadow, blur=18)
    d = ImageDraw.Draw(img)
    if light_scene:
        d.rounded_rectangle((110, table_y, 1170, 670), radius=28, fill=(246, 250, 250), outline=(184, 215, 220), width=3)
        d.line((130, table_y + 22, 1150, table_y + 22), fill=(220, 238, 241), width=2)
    else:
        d.rounded_rectangle((120, table_y, 1160, 670), radius=28, fill=(9, 21, 39), outline=(61, 183, 204), width=2)
        d.line((150, table_y + 22, 1130, table_y + 22), fill=(24, 78, 110), width=2)

    def draw_robot(cx: int = 640, base_y: int = 505) -> None:
        layer = Image.new("RGBA", (w, h), (0, 0, 0, 0))
        rd = ImageDraw.Draw(layer)
        cyan = (98, 232, 245)
        body = (34, 54, 76) if not light_scene else (223, 241, 244)
        body2 = (16, 30, 48) if not light_scene else (184, 214, 221)
        rd.ellipse((cx - 190, base_y + 120, cx + 190, base_y + 165), fill=(0, 0, 0, 55))
        rd.rounded_rectangle((cx - 150, base_y - 10, cx + 150, base_y + 155), radius=36, fill=(*body, 245), outline=(*cyan, 245), width=5)
        rd.rounded_rectangle((cx - 112, base_y - 186, cx + 112, base_y - 35), radius=42, fill=(*body2, 250), outline=(*cyan, 250), width=5)
        rd.line((cx, base_y - 186, cx, base_y - 238), fill=(*cyan, 250), width=5)
        rd.ellipse((cx - 16, base_y - 260, cx + 16, base_y - 228), fill=(*cyan, 255))
        rd.ellipse((cx - 68, base_y - 126, cx - 28, base_y - 84), fill=(*cyan, 255))
        rd.ellipse((cx + 28, base_y - 126, cx + 68, base_y - 84), fill=(*cyan, 255))
        rd.arc((cx - 58, base_y - 96, cx + 58, base_y - 26), 18, 162, fill=(255, 255, 255, 230), width=4)
        rd.rounded_rectangle((cx - 78, base_y + 38, cx + 78, base_y + 100), radius=18, fill=(9, 24, 39, 230), outline=(*cyan, 220), width=3)
        for i in range(4):
            x = cx - 54 + i * 36
            rd.line((x, base_y + 52, x + 18, base_y + 84), fill=(*cyan, 210), width=3)
        rd.line((cx - 150, base_y + 38, cx - 230, base_y + 92), fill=(*cyan, 230), width=8)
        rd.line((cx + 150, base_y + 38, cx + 230, base_y + 92), fill=(*cyan, 230), width=8)
        rd.ellipse((cx - 250, base_y + 80, cx - 210, base_y + 120), fill=(*cyan, 240))
        rd.ellipse((cx + 210, base_y + 80, cx + 250, base_y + 120), fill=(*cyan, 240))
        composite(layer, blur=0)
        glow(cx, base_y - 100, 120, cyan, 55)

    def draw_panda(cx: int = 640, base_y: int = 515) -> None:
        layer = Image.new("RGBA", (w, h), (0, 0, 0, 0))
        pd = ImageDraw.Draw(layer)
        pd.ellipse((cx - 175, base_y + 110, cx + 175, base_y + 160), fill=(0, 0, 0, 55))
        pd.ellipse((cx - 140, base_y - 60, cx + 140, base_y + 145), fill=(246, 248, 246, 255), outline=(42, 50, 56, 230), width=4)
        pd.ellipse((cx - 108, base_y - 205, cx + 108, base_y - 25), fill=(250, 251, 249, 255), outline=(42, 50, 56, 230), width=4)
        pd.ellipse((cx - 130, base_y - 218, cx - 60, base_y - 150), fill=(24, 28, 32, 255))
        pd.ellipse((cx + 60, base_y - 218, cx + 130, base_y - 150), fill=(24, 28, 32, 255))
        pd.ellipse((cx - 72, base_y - 145, cx - 22, base_y - 88), fill=(28, 32, 35, 255))
        pd.ellipse((cx + 22, base_y - 145, cx + 72, base_y - 88), fill=(28, 32, 35, 255))
        pd.ellipse((cx - 48, base_y - 123, cx - 32, base_y - 107), fill=(240, 255, 255, 255))
        pd.ellipse((cx + 32, base_y - 123, cx + 48, base_y - 107), fill=(240, 255, 255, 255))
        pd.ellipse((cx - 18, base_y - 92, cx + 18, base_y - 65), fill=(20, 24, 28, 255))
        pd.arc((cx - 34, base_y - 82, cx, base_y - 44), 20, 115, fill=(20, 24, 28, 230), width=3)
        pd.arc((cx, base_y - 82, cx + 34, base_y - 44), 65, 160, fill=(20, 24, 28, 230), width=3)
        pd.rounded_rectangle((cx - 210, base_y + 6, cx - 75, base_y + 92), radius=36, fill=(28, 32, 35, 255))
        pd.rounded_rectangle((cx + 75, base_y + 6, cx + 210, base_y + 92), radius=36, fill=(28, 32, 35, 255))
        pd.line((cx + 178, base_y - 170, cx + 240, base_y + 105), fill=(64, 148, 85, 255), width=10)
        for by in range(base_y - 140, base_y + 90, 52):
            pd.polygon([(cx + 228, by), (cx + 290, by - 30), (cx + 252, by + 8)], fill=(92, 180, 104, 235))
        composite(layer, blur=0)

    def draw_interface(cx: int = 640, cy: int = 345) -> None:
        layer = Image.new("RGBA", (w, h), (0, 0, 0, 0))
        idr = ImageDraw.Draw(layer)
        panel = (8, 20, 38, 220) if not light_scene else (248, 253, 253, 230)
        cyan = (100, 235, 247)
        idr.rounded_rectangle((210, 115, 1070, 500), radius=30, fill=panel, outline=(*cyan, 240), width=4)
        for i in range(4):
            x = 250 + i * 190
            idr.rounded_rectangle((x, 155, x + 140, 255), radius=16, fill=(*cyan, 25 + i * 12), outline=(*cyan, 120), width=2)
            idr.arc((x + 28, 178, x + 112, 232), 180, 360, fill=(*cyan, 220), width=5)
        pts = []
        for i in range(14):
            x = 270 + i * 54
            y = 405 - int(math.sin(i * 0.8 + seed % 7) * 55) - rng.randrange(0, 35)
            pts.append((x, y))
        idr.line(pts, fill=(*cyan, 235), width=5)
        for x, y in pts:
            idr.ellipse((x - 6, y - 6, x + 6, y + 6), fill=(255, 214, 102, 245))
        idr.ellipse((cx - 72, cy - 72, cx + 72, cy + 72), outline=(*cyan, 180), width=6)
        idr.ellipse((cx - 32, cy - 32, cx + 32, cy + 32), fill=(*cyan, 160))
        composite(layer, blur=0)
        glow(cx, cy, 180, cyan, 55)

    def draw_lab_scene() -> None:
        layer = Image.new("RGBA", (w, h), (0, 0, 0, 0))
        ld = ImageDraw.Draw(layer)
        cyan = (104, 230, 242)
        liquid = (78, 204, 164)
        for x in (430, 640, 850):
            ld.line((x, 210, x, 438), fill=(230, 250, 255, 210), width=6)
            ld.ellipse((x - 38, 428, x + 38, 506), outline=(230, 250, 255, 210), width=5)
            ld.pieslice((x - 32, 434, x + 32, 500), 0, 180, fill=(*liquid, 160))
            ld.line((x - 70, 310, x + 70, 310), fill=(*cyan, 180), width=3)
        ld.rounded_rectangle((330, 505, 970, 535), radius=12, fill=(220, 240, 242, 170), outline=(*cyan, 190), width=2)
        composite(layer, blur=0)
        glow(640, 350, 210, cyan, 42)

    def draw_abstract_product() -> None:
        layer = Image.new("RGBA", (w, h), (0, 0, 0, 0))
        ad = ImageDraw.Draw(layer)
        cyan = (102, 235, 248)
        gold = (255, 214, 102)
        center = (640, 380)
        for r, a in [(170, 45), (125, 70), (72, 120)]:
            ad.ellipse((center[0] - r, center[1] - r, center[0] + r, center[1] + r), outline=(*cyan, a), width=5)
        for i in range(8):
            ang = math.tau * i / 8 + (seed % 17) / 20
            x = center[0] + int(math.cos(ang) * 215)
            y = center[1] + int(math.sin(ang) * 130)
            ad.rounded_rectangle((x - 45, y - 28, x + 45, y + 28), radius=14, fill=(*gold, 80), outline=(*cyan, 130), width=2)
        ad.rounded_rectangle((520, 270, 760, 500), radius=36, fill=(245, 252, 252, 230) if light_scene else (18, 36, 58, 235), outline=(*cyan, 240), width=5)
        ad.ellipse((590, 335, 710, 455), fill=(*cyan, 130))
        composite(layer, blur=0)
        glow(640, 380, 210, cyan, 55)

    if "panda" in low:
        draw_panda()
    elif any(k in low for k in ("robot", "omni", "assistant", "bot", "android")):
        draw_robot()
    elif any(k in low for k in ("interface", "dashboard", "screen", "ui", "app")):
        draw_interface()
    elif any(k in low for k in ("pharma", "lab", "laboratory", "factory", "chemical", "chemistry")):
        draw_lab_scene()
    else:
        draw_abstract_product()

    vignette = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    vd = ImageDraw.Draw(vignette)
    for r, alpha in ((820, 0), (980, 35), (1160, 70)):
        vd.rectangle((0, 0, w, h), outline=(0, 0, 0, alpha), width=60)
    composite(vignette, blur=35)
    img = img.filter(ImageFilter.UnsharpMask(radius=1.2, percent=120, threshold=3))
    img.save(path)
    url = f"/media/files/images/{filename}"
    result = {"ok": True, "provider": "local-fallback", "type": "image", "kind": "image", "title": prompt[:80] or "Generated image", "filename": filename, "url": url, "file_url": url, "download_url": url, "note": note + " Produced a visual procedural fallback because no real image model is currently available."}
    return _attach_ledger(result, path, "image")


async def _create_image(prompt: str, theme: str | None = None, quality: str = "standard") -> dict[str, Any]:
    backend = (_settings["media"].get("image_backend") or "auto").lower()
    enhanced_prompt = _enhance_image_prompt(prompt, theme=theme, quality=quality)
    if backend in {"fallback", "placeholder", "local-fallback"}:
        fallback = await _fallback_image(prompt, note="Local fallback renderer requested.")
        fallback["enhanced_prompt"] = enhanced_prompt
        return fallback
    errors: list[str] = []
    for name, func in _image_provider_plan(backend):
        result = await func(enhanced_prompt)
        if result and result.get("ok"):
            result["original_prompt"] = prompt
            result["enhanced_prompt"] = enhanced_prompt
            result["routing_note"] = f"Image auto-router used {name} after provider readiness checks."
            if prompt and result.get("title", "").startswith(enhanced_prompt[:40]):
                result["title"] = prompt[:80]
            return result
        if result and result.get("error"):
            errors.append(f"{name}: {result['error']}")
    note = "Local fallback renderer used. " + (" | ".join(errors) if errors else "For true AI images: run Stable Diffusion WebUI/ComfyUI and set STABLE_DIFFUSION_URL, enable Diffusers, or add OPENAI_API_KEY.")
    fallback = await _fallback_image(prompt, note=note)
    fallback["enhanced_prompt"] = enhanced_prompt
    return fallback


async def _create_pdf(prompt: str, title: str | None = None) -> dict[str, Any]:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.pdfgen import canvas
    except Exception:
        filename = _safe_name(title or prompt or "document", "txt")
        path = PDF_DIR / filename
        path.write_text(prompt, encoding="utf-8")
        url = f"/media/files/pdf/{filename}"
        result = {"ok": True, "type": "pdf", "kind": "pdf", "title": title or "Document", "filename": filename, "url": url, "file_url": url, "download_url": url, "note": "ReportLab unavailable; generated text fallback."}
        return _attach_ledger(result, path, "pdf_text_fallback")
    title = title or (prompt[:70] if prompt else "SHIMS Document")
    filename = _safe_name(title, "pdf")
    path = PDF_DIR / filename
    c = canvas.Canvas(str(path), pagesize=A4)
    w, h = A4
    y = h - 20 * mm
    c.setFont("Helvetica-Bold", 16)
    c.drawString(18 * mm, y, title[:90])
    y -= 8 * mm
    c.setFont("Helvetica", 9)
    c.drawString(18 * mm, y, f"Generated by SHIMS v14 | {datetime.now().isoformat(timespec='seconds')}")
    y -= 12 * mm
    c.setFont("Helvetica", 10)
    body = prompt or title
    if len(body.split()) < 18:
        body += "\n\nObjective:\n- Capture the requested information clearly.\n\nNotes:\n- This document was generated locally by SHIMS.\n- Edit and expand as needed."
    for para in body.splitlines():
        if not para.strip():
            y -= 5 * mm
            continue
        line = ""
        for word in para.split():
            if len(line) + len(word) > 92:
                c.drawString(18 * mm, y, line)
                y -= 6 * mm
                line = word
                if y < 18 * mm:
                    c.showPage(); c.setFont("Helvetica", 10); y = h - 20 * mm
            else:
                line = (line + " " + word).strip()
        if line:
            c.drawString(18 * mm, y, line)
            y -= 7 * mm
            if y < 18 * mm:
                c.showPage(); c.setFont("Helvetica", 10); y = h - 20 * mm
    c.save()
    url = f"/media/files/pdf/{filename}"
    result = {"ok": True, "type": "pdf", "kind": "pdf", "title": title, "filename": filename, "url": url, "file_url": url, "download_url": url}
    return _attach_ledger(result, path, "pdf")


async def _create_ppt(prompt: str) -> dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception:
        return await _create_pdf("PowerPoint requested but python-pptx is unavailable.\n\n" + prompt, title="PPT fallback document")
    title = prompt[:70] if prompt else "SHIMS Presentation"
    filename = _safe_name(title, "pptx")
    path = PPT_DIR / filename
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generated by SHIMS v14"
    sections = [s.strip(" -") for s in re.split(r"[\n.;]+", prompt) if s.strip()]
    if not sections:
        sections = ["Objective", "Plan", "Next steps"]
    for idx, heading in enumerate(sections[:8], start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = heading[:80]
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = "Key points"
        for bullet in ["Generated by SHIMS", "Review and customize", "Add data, images and approvals"]:
            p = tf.add_paragraph(); p.text = bullet; p.level = 1
            try: p.font.size = Pt(18)
            except Exception: pass
    prs.save(path)
    url = f"/media/files/ppt/{filename}"
    result = {"ok": True, "type": "ppt", "kind": "ppt", "title": title, "filename": filename, "url": url, "file_url": url, "download_url": url}
    return _attach_ledger(result, path, "ppt")


def _media_ext_from_content_type(content_type: str, default_ext: str) -> str:
    ct = (content_type or "").split(";")[0].strip().lower()
    return {
        "audio/mpeg": "mp3",
        "audio/mp3": "mp3",
        "audio/wav": "wav",
        "audio/x-wav": "wav",
        "audio/ogg": "ogg",
        "video/mp4": "mp4",
        "video/webm": "webm",
        "application/octet-stream": default_ext,
    }.get(ct, default_ext)


def _decode_media_b64(value: str) -> bytes:
    raw = (value or "").strip()
    if "," in raw and raw.lower().startswith("data:"):
        raw = raw.split(",", 1)[1]
    return base64.b64decode(raw)


async def _download_media_url(url: str, folder: Path, title: str, default_ext: str, kind: str, provider: str) -> dict[str, Any]:
    async with httpx.AsyncClient(timeout=240, follow_redirects=True) as client:
        r = await client.get(url)
        r.raise_for_status()
    ext = _media_ext_from_content_type(r.headers.get("content-type", ""), default_ext)
    filename = _safe_name(title or kind, ext)
    path = folder / filename
    path.write_bytes(r.content)
    rel_kind = "audio" if kind == "audio" else "video"
    local_url = f"/media/files/{rel_kind}/{filename}"
    result = {"ok": True, "provider": provider, "type": kind, "kind": kind, "title": title[:80] or f"Generated {kind}", "filename": filename, "url": local_url, "file_url": local_url, "download_url": local_url}
    return _attach_ledger(result, path, kind)


async def _generic_media_api(kind: str, prompt: str, url: str, api_key: str, provider: str) -> dict[str, Any] | None:
    url = (url or "").strip()
    if not url:
        return None
    folder = AUDIO_DIR if kind == "audio" else VIDEO_DIR
    default_ext = "mp3" if kind == "audio" else "mp4"
    headers = {"Content-Type": "application/json"}
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"
    payload = {"prompt": prompt, "text": prompt, "kind": kind, "response_format": "url_or_base64"}
    try:
        async with httpx.AsyncClient(timeout=360, follow_redirects=True) as client:
            r = await client.post(url, headers=headers, json=payload)
            r.raise_for_status()
            content_type = r.headers.get("content-type", "")
            if not content_type.lower().startswith("application/json"):
                ext = _media_ext_from_content_type(content_type, default_ext)
                filename = _safe_name(prompt or kind, ext)
                path = folder / filename
                path.write_bytes(r.content)
                local_url = f"/media/files/{kind}/{filename}"
                result = {"ok": True, "provider": provider, "type": kind, "kind": kind, "title": prompt[:80] or f"Generated {kind}", "filename": filename, "url": local_url, "file_url": local_url, "download_url": local_url}
                return _attach_ledger(result, path, kind)
            data = r.json()
        media_url = data.get("file_url") or data.get("download_url") or data.get("url") or data.get(f"{kind}_url")
        if media_url:
            return await _download_media_url(media_url, folder, prompt, default_ext, kind, provider)
        b64 = data.get("b64_json") or data.get("base64") or data.get(kind) or data.get(f"{kind}_base64")
        if b64:
            raw = _decode_media_b64(b64)
            filename = _safe_name(prompt or kind, default_ext)
            path = folder / filename
            path.write_bytes(raw)
            local_url = f"/media/files/{kind}/{filename}"
            result = {"ok": True, "provider": provider, "type": kind, "kind": kind, "title": prompt[:80] or f"Generated {kind}", "filename": filename, "url": local_url, "file_url": local_url, "download_url": local_url}
            return _attach_ledger(result, path, kind)
        return {"ok": False, "provider": provider, "error": data.get("error") or "generic media API did not return a URL, bytes, or base64 payload"}
    except Exception as exc:
        return {"ok": False, "provider": provider, "error": str(exc)[:260]}


async def _openai_audio(prompt: str) -> dict[str, Any] | None:
    api_key = _clean_secret(os.getenv("OPENAI_API_KEY"))
    if not api_key:
        return None
    try:
        payload = {
            "model": _settings["media"].get("openai_tts_model") or "gpt-4o-mini-tts",
            "voice": _settings["media"].get("openai_tts_voice") or "alloy",
            "input": prompt[:4096],
            "response_format": "mp3",
        }
        async with httpx.AsyncClient(timeout=180) as client:
            r = await client.post("https://api.openai.com/v1/audio/speech", headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}, json=payload)
            r.raise_for_status()
        filename = _safe_name(prompt or "openai_audio", "mp3")
        path = AUDIO_DIR / filename
        path.write_bytes(r.content)
        url = f"/media/files/audio/{filename}"
        result = {"ok": True, "provider": "openai-tts", "type": "audio", "kind": "audio", "title": prompt[:80] or "Generated audio", "filename": filename, "url": url, "file_url": url, "download_url": url, "model": payload["model"], "voice": payload["voice"]}
        return _attach_ledger(result, path, "audio")
    except Exception as exc:
        return {"ok": False, "provider": "openai-tts", "error": str(exc)[:260]}


async def _openai_video(prompt: str) -> dict[str, Any] | None:
    api_key = _clean_secret(os.getenv("OPENAI_API_KEY"))
    if not api_key:
        return None
    model = _settings["media"].get("openai_video_model") or "sora-2"
    size = _settings["media"].get("openai_video_size") or "1280x720"
    seconds = str(_settings["media"].get("openai_video_seconds") or 4)
    try:
        async with httpx.AsyncClient(timeout=600) as client:
            r = await client.post(
                "https://api.openai.com/v1/videos",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={"model": model, "prompt": prompt, "size": size, "seconds": seconds},
            )
            r.raise_for_status()
            job = r.json()
            video_id = job.get("id")
            poll_seconds = max(0, int(os.getenv("SHIMS_OPENAI_VIDEO_POLL_SECONDS", "90")))
            deadline = time.time() + poll_seconds
            while video_id and job.get("status") in {"queued", "in_progress"} and time.time() < deadline:
                await asyncio.sleep(4)
                status = await client.get(f"https://api.openai.com/v1/videos/{video_id}", headers={"Authorization": f"Bearer {api_key}"})
                status.raise_for_status()
                job = status.json()
            if video_id and job.get("status") == "completed":
                content = await client.get(f"https://api.openai.com/v1/videos/{video_id}/content", headers={"Authorization": f"Bearer {api_key}"})
                content.raise_for_status()
                filename = _safe_name(prompt or "openai_video", "mp4")
                path = VIDEO_DIR / filename
                path.write_bytes(content.content)
                url = f"/media/files/video/{filename}"
                result = {"ok": True, "provider": "openai-sora", "type": "video", "kind": "video", "title": prompt[:80] or "Generated video", "filename": filename, "url": url, "file_url": url, "download_url": url, "model": model, "job_id": video_id}
                return _attach_ledger(result, path, "video")
            return {"ok": True, "provider": "openai-sora", "type": "video", "kind": "video", "title": prompt[:80] or "OpenAI video job", "job_id": video_id, "status": job.get("status", "queued"), "model": model, "note": "OpenAI video job started but is still processing. Check the job in your OpenAI account or increase SHIMS_OPENAI_VIDEO_POLL_SECONDS for synchronous waiting."}
    except Exception as exc:
        return {"ok": False, "provider": "openai-sora", "error": str(exc)[:260]}


async def _qwen_video(prompt: str) -> dict[str, Any] | None:
    """Generate video via Alibaba DashScope video generation model."""
    api_key = _clean_secret(os.getenv("QWEN_API_KEY"))
    if not api_key:
        return None
    try:
        payload = {
            "model": "wanx2.1-t2v-turbo",
            "input": {"prompt": prompt},
            "parameters": {"size": "1280*720"},
        }
        async with httpx.AsyncClient(timeout=600) as client:
            r = await client.post(
                "https://dashscope.aliyuncs.com/api/v1/services/aigc/video-generation/video-synthesis",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json", "X-DashScope-Async": "enable"},
                json=payload,
            )
            r.raise_for_status()
            job = r.json()
            task_id = job.get("output", {}).get("task_id")
            if not task_id:
                return {"ok": False, "provider": "qwen", "error": "No task_id from DashScope video API"}
            poll_deadline = time.time() + 180
            while time.time() < poll_deadline:
                await asyncio.sleep(5)
                status_r = await client.get(
                    f"https://dashscope.aliyuncs.com/api/v1/tasks/{task_id}",
                    headers={"Authorization": f"Bearer {api_key}"},
                )
                status_r.raise_for_status()
                status = status_r.json()
                task_status = status.get("output", {}).get("task_status", "")
                if task_status == "SUCCEEDED":
                    results = status.get("output", {}).get("video_results", [])
                    if results and results[0].get("url"):
                        vid_url = results[0]["url"]
                        vid = await client.get(vid_url)
                        vid.raise_for_status()
                        filename = _safe_name(prompt or "qwen_video", "mp4")
                        path = VIDEO_DIR / filename
                        path.write_bytes(vid.content)
                        url = f"/media/files/video/{filename}"
                        result = {"ok": True, "provider": "qwen", "type": "video", "kind": "video", "title": prompt[:80] or "Generated video", "filename": filename, "url": url, "file_url": url, "download_url": url, "job_id": task_id}
                        return _attach_ledger(result, path, "video")
                elif task_status == "FAILED":
                    return {"ok": False, "provider": "qwen", "error": status.get("output", {}).get("message", "DashScope video task failed")}
            return {"ok": False, "provider": "qwen", "error": "DashScope video generation timed out"}
    except Exception as exc:
        return {"ok": False, "provider": "qwen", "error": str(exc)[:260]}


async def _create_audio(prompt: str) -> dict[str, Any]:
    backend = (_settings["media"].get("audio_backend") or "auto").lower()
    errors: list[str] = []
    if backend in {"auto", "openai", "openai-tts", "cloud"}:
        result = await _openai_audio(prompt)
        if result and result.get("ok"):
            return result
        if result and result.get("error"):
            errors.append(f"openai-tts: {result['error']}")
    if backend in {"auto", "generic", "webhook", "external"}:
        result = await _generic_media_api("audio", prompt, _settings["media"].get("audio_api_url") or "", _clean_secret(os.getenv("SHIMS_AUDIO_API_KEY")), "generic-audio-api")
        if result and result.get("ok"):
            return result
        if result and result.get("error"):
            errors.append(f"generic-audio-api: {result['error']}")
    filename = _safe_name(prompt or "audio", "wav")
    path = AUDIO_DIR / filename
    fr = 44100
    duration = 2.5
    base_freq = 392 + (sum(ord(c) for c in prompt[:24]) % 140)
    with wave.open(str(path), "w") as wf:
        wf.setnchannels(1); wf.setsampwidth(2); wf.setframerate(fr)
        for i in range(int(fr * duration)):
            env = math.sin(math.pi * i / (fr * duration))
            wave1 = math.sin(2 * math.pi * base_freq * i / fr)
            wave2 = 0.35 * math.sin(2 * math.pi * (base_freq * 1.5) * i / fr)
            amp = int(18000 * env * (wave1 + wave2) / 1.35)
            wf.writeframes(struct.pack("<h", amp))
    url = f"/media/files/audio/{filename}"
    note = "Local tone fallback used. Configure OpenAI TTS or a generic audio API in Settings for real voice/music." + ((" " + " | ".join(errors)) if errors else "")
    result = {"ok": True, "provider": "local-tone-fallback", "type": "audio", "kind": "audio", "title": prompt[:80] or "Generated audio", "filename": filename, "url": url, "file_url": url, "download_url": url, "note": note}
    return _attach_ledger(result, path, "audio")


async def _pollinations_video(prompt: str) -> dict[str, Any] | None:
    """Experimental video generation via Pollinations.ai — free, no API key."""
    try:
        import urllib.parse
        encoded = urllib.parse.quote(prompt[:1000])
        url = f"https://video.pollinations.ai/prompt/{encoded}?width=1280&height=720&nologo=true&seed={abs(hash(prompt)) % 99999}"
        async with httpx.AsyncClient(timeout=180) as client:
            r = await client.get(url)
            r.raise_for_status()
            filename = _safe_name(prompt, "mp4")
            path = VIDEO_DIR / filename
            path.write_bytes(r.content)
            file_url = f"/media/files/video/{filename}"
            return {"ok": True, "provider": "pollinations", "type": "video", "kind": "video", "title": prompt[:80], "filename": filename, "url": file_url, "file_url": file_url, "download_url": file_url}
    except Exception as exc:
        return {"ok": False, "error": f"pollinations-video: {str(exc)[:200]}"}


async def _create_video(prompt: str) -> dict[str, Any]:
    backend = (_settings["media"].get("video_backend") or "auto").lower()
    errors: list[str] = []
    if backend in {"auto", "pollinations", "free"}:
        result = await _pollinations_video(prompt)
        if result and result.get("ok") and result.get("file_url"):
            return result
        if result and result.get("error"):
            errors.append(result["error"])
    if backend in {"auto", "openai", "sora", "cloud"}:
        result = await _openai_video(prompt)
        if result and result.get("ok") and result.get("file_url"):
            return result
        if result and result.get("ok") and result.get("job_id"):
            return result
        if result and result.get("error"):
            errors.append(f"openai-sora: {result['error']}")
    if backend in {"auto", "qwen", "cloud"}:
        result = await _qwen_video(prompt)
        if result and result.get("ok") and result.get("file_url"):
            return result
        if result and result.get("error"):
            errors.append(f"qwen-video: {result['error']}")
    if backend in {"auto", "generic", "webhook", "external"}:
        result = await _generic_media_api("video", prompt, _settings["media"].get("video_api_url") or "", _clean_secret(os.getenv("SHIMS_VIDEO_API_KEY")), "generic-video-api")
        if result and result.get("ok"):
            return result
        if result and result.get("error"):
            errors.append(f"generic-video-api: {result['error']}")
    cover = await _create_image("Video storyboard: " + prompt)
    ffmpeg = shutil.which("ffmpeg")
    filename = _safe_name(prompt or "video", "mp4")
    path = VIDEO_DIR / filename
    if ffmpeg and cover.get("filename"):
        src = IMAGE_DIR / cover["filename"]
        try:
            subprocess.run([ffmpeg, "-y", "-loop", "1", "-i", str(src), "-f", "lavfi", "-i", "anullsrc=channel_layout=stereo:sample_rate=44100", "-t", "5", "-shortest", "-vf", "scale=1280:720", "-pix_fmt", "yuv420p", str(path)], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=35)
            url = f"/media/files/video/{filename}"
            result = {"ok": True, "provider": "ffmpeg-local", "type": "video", "kind": "video", "title": prompt[:80] or "Generated video", "filename": filename, "url": url, "file_url": url, "download_url": url}
            return _attach_ledger(result, path, "video")
        except Exception as exc:
            cover["note"] = f"Video storyboard generated; FFmpeg MP4 failed: {str(exc)[:120]}"
    cover["title"] = "Video storyboard fallback: " + (prompt[:60] or "SHIMS")
    cover["note"] = cover.get("note") or ("Install FFmpeg or connect OpenAI Sora/generic video API for real MP4 generation." + ((" " + " | ".join(errors)) if errors else ""))
    return cover


async def _create_media(kind: str, prompt: str, theme: str | None = None, quality: str = "standard", provider: str | None = None, privacy_mode: str = "balanced") -> dict[str, Any]:
    from shared.privacy_guard import classify_sensitivity
    kind = (kind or "image").lower()
    prior_image = _settings["media"].get("image_backend")
    prior_audio = _settings["media"].get("audio_backend")
    prior_video = _settings["media"].get("video_backend")
    provider = (provider or "").strip().lower()
    # Privacy guard for media: if prompt contains sensitive data, force local backends
    sensitivity = classify_sensitivity(prompt)
    force_local = sensitivity == "high" or (sensitivity == "medium" and privacy_mode == "strict")
    if force_local:
        if kind in {"image", "photo", "picture", "poster", "drawing", "art", "logo"}:
            _settings["media"]["image_backend"] = "local"
        if kind in {"video", "movie", "clip", "reel"}:
            _settings["media"]["video_backend"] = "local"
    elif provider and kind in {"image", "photo", "picture", "poster", "drawing", "art", "logo"}:
        _settings["media"]["image_backend"] = "fallback" if provider in {"local", "fallback", "placeholder"} else provider
    if provider and kind in {"audio", "sound", "music", "voice"}:
        _settings["media"]["audio_backend"] = provider
    if not force_local and provider and kind in {"video", "movie", "clip", "reel"}:
        _settings["media"]["video_backend"] = provider
    try:
        if kind in {"pdf", "document", "report", "coa", "letter", "invoice", "quotation", "sop"}:
            return await _create_pdf(prompt)
        if kind in {"ppt", "pptx", "powerpoint", "presentation", "deck", "slides"}:
            return await _create_ppt(prompt)
        if kind in {"audio", "sound", "music", "voice"}:
            return await _create_audio(prompt)
        if kind in {"video", "movie", "clip", "reel"}:
            return await _create_video(prompt)
        return await _create_image(prompt, theme=theme, quality=quality)
    finally:
        _settings["media"]["image_backend"] = prior_image
        _settings["media"]["audio_backend"] = prior_audio
        _settings["media"]["video_backend"] = prior_video


# --------------------------------------------------------------------------- #
# Agentic tool-use loop helpers (the "do anything" core)
# --------------------------------------------------------------------------- #
_AGENT_SLASH = ("/run ", "/sh ", "/shell ", "/build ", "/edit ", "/patch ", "/web ", "/coder ", "/fix ", "/do ")

_AGENTIC_PATTERNS = re.compile(
    r"(?:\b(?:shell|powershell|terminal|command line|command prompt)\b)"
    r"|(?:\b(?:run|execute|launch)\b[^.]{0,40}\b(?:command|cmd|script|terminal|shell|powershell|program|\.py|\.sh|\.bat|\.ps1|exe|binary|tests?)\b)"
    r"|(?:\b(?:open|read|show|list|find|search|create|make|write|edit|modify|change|update|delete|remove|move|rename|copy|organi[sz]e)\b[^.]{0,40}\b(?:file|files|folder|directory|dir|path|code|project|repo|workspace|\.py|\.txt|\.json|\.md|\.csv|\.js|\.html)\b)"
    r"|(?:\b(?:search|look\s?up|browse|fetch|scrape|google)\b[^.]{0,30}\b(?:web|internet|online|url|site|website|page|docs?)\b)"
    r"|(?:\bsearch the (?:web|internet)\b)|(?:\b(?:on|from) the (?:web|internet)\b)"
    r"|(?:\b(?:write|create|build|generate|make|fix|debug|refactor|implement|add|develop|code)\b[^.]{0,40}\b(?:code|script|program|app|application|website|web app|function|class|api|endpoint|module|bug|feature|cli|tool|bot|game)\b)"
    r"|(?:\b(?:modify|change|edit|update|patch|improve|extend|upgrade)\b[^.]{0,30}\byour(?:self|\s+own)?\b)"
    r"|(?:\bself[-\s]?patch\b)|(?:\badd (?:an?|the) (?:endpoint|feature|button|tool|capability|skill|route)\b)"
    r"|(?:\b(?:install|compile|clone)\b)|(?:\bgit (?:clone|status|log|diff|pull|add|commit|branch)\b)|(?:\b(?:pip|npm) install\b)"
    r"|(?:\b(?:organi[sz]e|clean up|find duplicates?)\b[^.]{0,30}\b(?:files?|folder|downloads|directory)\b)"
    r"|(?:\b(?:self[-\s]?check|check|inspect|analyze|review|test)\b[^.]{0,60}\b(?:your(?:self|\s+own)?\s+)?(?:code|source|files?|tests?|app|system|repo|repository|\.(?:py|js|html|json|md|txt|csv|sh|bat|ps1))\b)"
    r"|(?:\b(?:desktop|my\s+(?:computer|pc|machine)|bridge)\b[^.]{0,60}\b(?:bridge|screenshot|shell|command|cmd|file|ping|info|status)\b)"
    r"|(?:\b(?:screenshot|shell|command|cmd|file|ping)\b[^.]{0,60}\b(?:desktop|my\s+(?:computer|pc|machine)|bridge)\b)",
    re.I,
)


def _agentic_intent(text: str) -> bool:
    raw = (text or "").strip()
    if not raw:
        return False
    low = raw.lower()
    if low == "/agent" or low.startswith(_AGENT_SLASH):
        return True
    return bool(_AGENTIC_PATTERNS.search(low))


def _strip_agent_slash(text: str) -> str:
    low = (text or "").lower()
    for s in _AGENT_SLASH:
        if low.startswith(s):
            return text[len(s):].strip()
    return text


async def _agent_tool_model(deep: bool = False) -> str | None:
    """Pick an installed Ollama model to drive the tool loop.

    Two tiers:
      * **deep** = True  → prefer large coder models (14b+) for multi-step
        planning, complex code generation, and self-patch. Keeps qwen2.5-coder:14b.
      * **deep** = False → prefer the fastest available model so simple tool
        calls (ls, glob, single-step) complete in 2-5s instead of 30-60s.
    """
    names = await _ollama_names()
    if not names:
        return None
    override = os.getenv("SHIMS_AGENT_MODEL", "").strip()
    if override and override in names:
        return override

    # Only use Ollama models that ACTUALLY support tool calling.
    # llama3.2, gemma3, gemma-4, deepseek-r1 do NOT support tools and will hang/ignore.
    TOOL_CAPABLE = ("qwen2.5-coder", "qwen2.5", "llama3.1", "mistral-nemo", "qwen3", "mistral")
    capable_names = [n for n in names if any(tc in n.lower() for tc in TOOL_CAPABLE)]

    if deep:
        # Large models first — user explicitly wants deep work
        preferred = _prefer_coder_model("ollama", None)
        if preferred and preferred in capable_names:
            return preferred
        for pref in ("qwen2.5-coder:14b", "qwen2.5:14b", "qwen2.5-coder", "qwen2.5", "mistral-nemo", "qwen3"):
            hit = next((n for n in capable_names if pref in n.lower()), None)
            if hit:
                return hit
    else:
        # Fast models first — snappy for single-tool calls (still must be tool-capable)
        for pref in ("qwen2.5-coder:7b", "qwen2.5-coder:3b", "qwen2.5-coder", "qwen2.5", "mistral"):
            hit = next((n for n in capable_names if pref in n.lower()), None)
            if hit:
                return hit
        # Fall back to any capable coder model
        preferred = _prefer_coder_model("ollama", None)
        if preferred and preferred in capable_names:
            return preferred
    # Absolute fallback — only tool-capable models
    for pref in TOOL_CAPABLE:
        hit = next((n for n in capable_names if pref in n.lower()), None)
        if hit:
            return hit
    return capable_names[0] if capable_names else None


# Heuristic: requests that need deep reasoning (multi-file edits, planning,
# code generation, self-patch, complex analysis) use the large model; quick
# tool calls (list files, run a command, search) use the fast model.
_DEEP_PATTERNS = re.compile(
    r"(write|create|build|implement|refactor|redesign|architect|plan|generate|self[.\-]?patch|"
    r"modify.*(code|source|app)|complex|thorough|detailed|full.*review|analyze.*all|"
    r"debug.*and.*fix|rewrite|overhaul|port|migrate|upgrade)",
    re.I,
)


def _needs_deep_model(text: str) -> bool:
    """Does this user request warrant the heavyweight model?"""
    return bool(_DEEP_PATTERNS.search(text))


def _agent_create_pending(*, action_type: str, title: str, summary: str,
                          payload: dict[str, Any], session_id: str | None) -> dict[str, Any]:
    item = _create_pending_action(action_type=action_type, title=title, summary=summary,
                                  payload=payload, session_id=session_id, risk="agent_tool")
    return _public_pending_action(item)


async def _brain_stream(req: ChatRequest) -> AsyncGenerator[bytes, None]:
    started = time.perf_counter()
    req.message = (req.message or "").strip()
    session_hint = req.session_id or str(uuid.uuid4())

    # vC: If a voice correction is pending, try to use it without blocking the turn.
    raw_voice_message = req.message
    if req.source == "voice" and req.voice_correction_id:
        correction = await _await_stt_correction(req.voice_correction_id, timeout=0.15)
        if correction and correction.get("ok"):
            corrected = (correction.get("corrected") or raw_voice_message).strip()
            confidence = float(correction.get("confidence") or 0)
            if corrected and corrected.lower() != raw_voice_message.lower() and confidence >= 0.5:
                req.message = corrected
                yield _jsonl({
                    "type": "thought",
                    "stage": "stt_correction",
                    "content": f"Voice correction ({confidence:.0%}): '{raw_voice_message}' → '{corrected}'",
                })

    # vC: Keep images on the request; native multimodal injection happens after provider selection.
    # For local/Ollama paths, we fall back to vision description just before model call.
    raw_images = [src for src in (req.images or []) if src]

    # ----- 1. Handle yes/no to existing approval requests FIRST -----
    approval_decision = _approval_decision_from_text(req.message)
    if approval_decision is not None:
        pending = _latest_pending_action(session_hint)
        meta_trust = build_trust(route="approval:decision", evidence=[], missing_evidence=[], requested_level="L3")
        yield _jsonl({
            "type": "meta",
            "session_id": session_hint,
            "model": "approval-router",
            "provider": "local",
            "route": "approval:decision",
            "agent": "operator",
            "brain": f"unified-v13+{BRAIN_VERSION}",
            "conversation_mode": bool(req.conversation_mode),
            **_trust_fields(meta_trust),
        })
        if not pending:
            answer = "I do not have a pending action to approve. Tell me the action first and I will ask yes or no."
            yield _jsonl({"type": "token", "content": answer})
            yield _jsonl({"type": "done", "session_id": session_hint, "provider": "local", "model": "approval-router", "route": "approval:no-pending", **_trust_fields(meta_trust)})
            return
        if approval_decision is False:
            pending["status"] = "cancelled"
            pending["decision"] = "no"
            pending["resolved_at"] = _utc_now()
            _save_pending_action(pending)
            answer = f"Cancelled. I did not run: {pending.get('title') or pending.get('action_type')}."
            yield _jsonl({"type": "approval", "approval": _public_pending_action(pending)})
            yield _jsonl({"type": "token", "content": answer})
            yield _jsonl({"type": "done", "session_id": session_hint, "provider": "local", "model": "approval-router", "route": "approval:cancelled", **_trust_fields(meta_trust)})
            return
        approved_by = "chat-human"
        yield _jsonl({"type": "status", "content": f"Approved. Running {pending.get('action_type')}..."})
        result = await _execute_pending_action(pending, approved_by=approved_by)
        pending["status"] = "completed" if result.get("ok") else "failed"
        pending["decision"] = "yes"
        pending["approved_by"] = approved_by
        pending["resolved_at"] = _utc_now()
        pending["result"] = result
        _save_pending_action(pending)
        action = record_action(
            "approval_execute",
            f"Execute approval {pending.get('approval_id')}",
            payload={"approval_id": pending.get("approval_id"), "action_type": pending.get("action_type")},
            result=result,
            evidence=[],
            requested_level="L3",
            status="completed" if result.get("ok") else "failed",
            summary=f"Executed approved action: {pending.get('title')}",
        )
        trust = build_trust(route="approval:execute", evidence=evidence_from_action(action.get("action")), action_id=action.get("action_id", ""), ledger_hash=action.get("ledger_hash", ""))
        answer = (result.get("message") or result.get("status") or "Action completed.") if result.get("ok") else f"Action failed: {result.get('message') or result.get('status') or 'unknown error'}"
        yield _jsonl({"type": "approval", "approval": _public_pending_action(pending), "result": result, **_trust_fields(trust)})
        yield _jsonl({"type": "token", "content": answer})
        yield _jsonl({"type": "done", "session_id": session_hint, "provider": "local", "model": "approval-router", "route": "approval:executed", "approval": _public_pending_action(pending), **_trust_fields(trust)})
        return

    # ----- 2. PLAN FIRST — think before asking for approval -----
    plan = await _make_plan(req)
    # If agent mode is explicitly active and message has agentic intent,
    # bypass planner tool execution (web search, image, etc.) so the
    # agent loop can actually run tools on the local machine.
    # We only clear an explicit tool plan when the intent is a direct agent
    # slash command or omnipotent mode; otherwise a routed tool turn such as
    # "search the web for ..." should keep its planner-determined route and
    # not be forced into the agent loop (which can emit tool explanations).
    _agent_low = (req.message or "").strip().lower()
    _direct_agent = _agent_low == "/agent" or _agent_low.startswith(_AGENT_SLASH)
    # Omnipotent mode auto-approves gated actions; the agent loop should only
    # activate for turns that actually look agentic. Keep greetings/duplicates on
    # the fast path so they don't block on model load or wave planning.
    if req.agent_mode and _agentic_intent(req.message):
        if _direct_agent or not plan.tool_kind:
            plan.tool_kind = None
            plan.tool_prompt = None
            plan.tool_metadata = None
            if plan.route == "local:greeting" and _agentic_intent(req.message):
                plan.route = "agent"
    session_id = plan.session_id
    conversation_enabled = bool(req.conversation_mode)
    history = _sessions.setdefault(session_id, []) if conversation_enabled else []

    def _store_assistant_turn(answer: str) -> None:
        if conversation_enabled:
            history.append({"role": "assistant", "content": answer})

    def _remember_session_turn(*args: Any, **kwargs: Any) -> None:
        if conversation_enabled:
            remember_turn(*args, **kwargs)
        # Also extract durable facts into searchable memory (non-blocking)
        try:
            user_msg_arg = args[1] if len(args) > 1 else kwargs.get("user_text", "")
            assistant_msg_arg = args[2] if len(args) > 2 else kwargs.get("assistant_text", "")
            metadata = kwargs.get("metadata") or {}
            tools_used = metadata.get("tools_used")
            asyncio.create_task(
                _auto_memory_after_turn(user_msg_arg, assistant_msg_arg, tools_used=tools_used)
            )
        except Exception:
            pass

    async def _auto_memory_after_turn(user_msg: str, assistant_msg: str, tools_used: list[str] | None = None) -> None:
        """Extract durable facts from a successful turn and save them to memory."""
        from shared.omni_brain import remember
        facts: list[tuple[str, list[str]]] = []
        # Fast regex heuristics first
        lower = (user_msg or "").lower()
        if any(p in lower for p in ("i prefer", "my name is", "i am a", "i work at", "i like", "i want", "i need")):
            facts.append((user_msg.strip(), ["preference", "user"]))
        if tools_used:
            if "desktop.interpreter" in tools_used and assistant_msg and len(assistant_msg) > 20:
                facts.append((f"Code interpreter result: {assistant_msg[:500]}", ["tool_result", "code"]))
            if any(t.startswith("plan.") for t in tools_used):
                facts.append((f"Plan execution: {assistant_msg[:500]}", ["tool_result", "plan"]))
        if "remember" in (assistant_msg or "").lower() and len(assistant_msg or "") > 40:
            facts.append((assistant_msg.strip()[:600], ["assistant_note"]))
        # LLM-powered extraction for richer memories
        if len((user_msg or "") + (assistant_msg or "")) >= 60:
            try:
                facts.extend(await _extract_durable_facts_llm(user_msg, assistant_msg, tools_used))
            except Exception:
                pass
        # Deduplicate by normalized key and store
        seen: set[str] = set()
        for content, tags in facts:
            norm = content.strip()[:80].lower()
            if not norm or norm in seen:
                continue
            seen.add(norm)
            key = content[:60] + ("..." if len(content) > 60 else "")
            remember("agent", key, content, tags=tags, source="auto_memory")

    # ----- 3. THINKING — show the user what SHIMS is doing -----
    yield _jsonl({"type": "thought", "stage": "plan", "content": f"Analyzing request: '{req.message[:80]}...' | Detected intent: {plan.agent or 'general'}"})

    # Check for action requests AFTER planning so SHIMS thinks first
    pending_request = _detect_chat_action_request(req.message, session_hint)
    if pending_request:
        action_type = pending_request.get("action_type")
        # Safe creation/scaffold actions run immediately so SHIMS actually does the work
        # instead of asking for approval and doing nothing.
        if action_type in {"coder_app_scaffold", "evolution_capability_check"}:
            yield _jsonl({"type": "thought", "stage": "plan", "content": f"Detected creation request: {pending_request.get('title')}. Executing now."})
            try:
                result = await _execute_pending_action(pending_request, approved_by="auto-approved")
                ok = result.get("ok", False)
                msg = result.get("message") or ("Done." if ok else "Failed.")
                route = f"auto:{action_type}"
                trust = build_trust(route=route, evidence=[], requested_level="L1")
                yield _jsonl({"type": "meta", "session_id": session_hint, "model": "approval-router", "provider": "local", "route": route, "agent": "operator", "brain": f"unified-v13+{BRAIN_VERSION}", "conversation_mode": bool(req.conversation_mode), **_trust_fields(trust)})
                yield _jsonl({"type": "token", "content": f"Created and applied: {pending_request.get('title')}.\n\n{msg}"})
                if result.get("app_url"):
                    yield _jsonl({"type": "token", "content": f"\nOpen it: {result['app_url']}"})
                yield _jsonl({"type": "done", "session_id": session_hint, "provider": "local", "model": "approval-router", "route": route, **_trust_fields(trust)})
            except Exception as exc:
                yield _jsonl({"type": "token", "content": f"I tried to execute '{pending_request.get('title')}' but hit an error: {exc}"})
                yield _jsonl({"type": "done", "session_id": session_hint, "provider": "local", "model": "approval-router", "route": f"auto:{action_type}:error"})
            return

        # Riskier source-change actions still require explicit human approval.
        yield _jsonl({"type": "thought", "stage": "plan", "content": f"Detected action request: {pending_request.get('title')}. Planning scope and requirements before asking for approval."})
        trust = build_trust(route="approval:request", evidence=evidence_from_action(get_action(pending_request.get("action_id", ""))), action_id=pending_request.get("action_id", ""), ledger_hash=pending_request.get("ledger_hash", ""), requested_level="L3")
        yield _jsonl({
            "type": "meta",
            "session_id": session_hint,
            "model": "approval-router",
            "provider": "local",
            "route": "approval:request",
            "agent": "operator",
            "brain": f"unified-v13+{BRAIN_VERSION}",
            "conversation_mode": bool(req.conversation_mode),
            **_trust_fields(trust),
        })
        public = _public_pending_action(pending_request)
        answer = f"I've thought this through. Here's what I plan to do:\n\n**{pending_request.get('title')}**\n{pending_request.get('summary')}\n\nThis will create files and modify the app workspace. Approve? Reply yes or no."
        yield _jsonl({"type": "approval_request", "approval": public, **_trust_fields(trust)})
        yield _jsonl({"type": "token", "content": answer})
        yield _jsonl({"type": "done", "session_id": session_hint, "provider": "local", "model": "approval-router", "route": "approval:request", "approval": public, **_trust_fields(trust)})
        return

    # Continue with normal execution...

    # ----- CONVERSATION REVIEW: make SHIMS explicitly aware of prior turns -----
    turn_history = history[-50:] if conversation_enabled else []
    if len(turn_history) > 1:
        # Build a brief summary of the conversation thread for thinking display
        user_turns = [t for t in turn_history if t.get("role") == "user"]
        recent_topics = []
        for t in user_turns[-5:]:
            txt = str(t.get("content", ""))[:60]
            if txt:
                recent_topics.append(txt + ("..." if len(str(t.get("content",""))) > 60 else ""))
        conv_summary = f"Conversation has {len(turn_history)} turns ({len(user_turns)} user messages). Recent topics: " + " | ".join(recent_topics)
    else:
        conv_summary = "This is the start of the conversation — no prior context to review."
    yield _jsonl({"type": "thought", "stage": "conversation", "content": conv_summary})

    # Fast paths: these do not need expensive brain retrieval.
    if plan.route == "local:greeting":
        answer = "Haan, sun raha hoon. Batao, kya karna hai?"
        _store_assistant_turn(answer)
        _remember_session_turn(session_id, req.message, answer, route="local:greeting", agent=plan.agent, provider="local", model="", metadata={})
        yield _jsonl({"type": "token", "content": answer})
        yield _jsonl({"type": "done", "session_id": session_hint, "model": plan.model, "provider": "local", "route": "greeting"})
        return
    if plan.duplicate:
        yield _jsonl({"type": "ignored", "reason": "silence_or_duplicate_suppressed", "session_id": session_hint})
        yield _jsonl({"type": "done", "session_id": session_hint, "model": plan.model, "provider": plan.provider, "route": "ignored"})
        return
    if not req.message:
        yield _jsonl({"type": "ignored", "reason": "empty_or_silence", "session_id": session_hint})
        yield _jsonl({"type": "done", "session_id": session_hint, "route": "empty"})
        return

    yield _jsonl({"type": "status", "content": "Retrieving memory & context..."})
    brain_addendum, brain_ctx = brain_prompt_addendum(req.message, agent=plan.agent, limit=8, history=turn_history)
    self_addendum = self_prompt_addendum()
    if self_addendum:
        brain_addendum += "\n\n" + self_addendum
    context_evidence = evidence_from_brain_context(brain_ctx)
    if plan.route == "local:greeting":
        context_evidence = []
    mem_hits = brain_ctx.get("memory_hits", 0)
    rag_hits = brain_ctx.get("rag_hits", 0)
    research_hits = brain_ctx.get("research_hits", 0)
    if mem_hits or rag_hits or research_hits:
        yield _jsonl({"type": "thought", "stage": "context", "content": f"Retrieved {mem_hits} memories, {rag_hits} RAG chunks, {research_hits} research items."})
    else:
        yield _jsonl({"type": "thought", "stage": "context", "content": "No relevant memories or RAG context found for this query."})
    meta_trust = build_trust(
        route=plan.route,
        evidence=context_evidence,
        missing_evidence=[] if context_evidence else ["No retrieved memory/RAG evidence for this turn yet."],
        requested_level="draft",
    )
    log_event("turn.start", route=plan.route, provider=plan.provider, model=plan.model, message=req.message, metadata={"session_id": session_id, "source": req.source})
    yield _jsonl({
        "type": "meta",
        "session_id": session_id,
        "model": plan.model,
        "provider": plan.provider,
        "route": plan.route,
        "agent": plan.agent,
        "brain": f"unified-v13+{BRAIN_VERSION}",
        "memory_hits": brain_ctx.get("memory_hits", 0),
        "rag_hits": brain_ctx.get("rag_hits", 0),
        "research_hits": brain_ctx.get("research_hits", 0),
        "conversation_mode": conversation_enabled,
        "emotion": "calm",
        **_trust_fields(meta_trust),
    })
    yield _jsonl({"type": "thought", "stage": "plan", "content": f"Route decided: {plan.route} | Agent: {plan.agent} | Provider: {plan.provider}:{plan.model}"})
    # Privacy guard indicator
    if "privacy-guard" in plan.route:
        yield _jsonl({"type": "thought", "stage": "plan", "content": f"🔒 Privacy guard activated: sensitive data detected. Forced local processing. Mode: {req.privacy_mode}"})
    elif plan.provider != "ollama" and plan.provider != "local":
        yield _jsonl({"type": "thought", "stage": "plan", "content": f"☁️ Cloud provider selected ({plan.provider}). Data will leave this machine."})
    else:
        yield _jsonl({"type": "thought", "stage": "plan", "content": f"🏠 Local provider selected. Data stays on this machine."})
    if plan.duplicate:
        log_event("turn.duplicate_blocked", route=plan.route, provider=plan.provider, model=plan.model, latency_ms=(time.perf_counter()-started)*1000, ok=True, message=req.message, metadata={"session_id": session_id})
        yield _jsonl({"type": "ignored", "reason": "silence_or_duplicate_suppressed", "session_id": session_id})
        ignored_trust = build_trust(route="ignored", evidence=[], missing_evidence=["Duplicate or silence turn was suppressed."], requested_level="draft")
        yield _jsonl({"type": "done", "session_id": session_id, "model": plan.model, "provider": plan.provider, "route": "ignored", **_trust_fields(ignored_trust)})
        return
    if not req.message:
        log_event("turn.empty_suppressed", route="empty", provider=plan.provider, model=plan.model, latency_ms=(time.perf_counter()-started)*1000, ok=True, metadata={"session_id": session_id})
        yield _jsonl({"type": "ignored", "reason": "empty_or_silence", "session_id": session_id})
        empty_trust = build_trust(route="empty", evidence=[], missing_evidence=["Empty or silence input was suppressed."], requested_level="draft")
        yield _jsonl({"type": "done", "session_id": session_id, "route": "empty", **_trust_fields(empty_trust)})
        return
    # Build user message: native multimodal for Anthropic/OpenAI, else text with vision descriptions
    user_message = await _build_user_message_with_images(req.message, raw_images, plan.provider)
    if conversation_enabled:
        history.append(user_message)
    if plan.tool_kind:
        yield _jsonl({"type": "thought", "stage": "tool", "content": f"Decided to use tool: {plan.tool_kind}"})
        yield _jsonl({"type": "status", "content": f"Running {plan.tool_kind} tool"})
        if plan.tool_kind == "web_search":
            search_plan = plan.tool_metadata or {}
            result = await _run_web_search_with_plan(plan.tool_prompt or req.message, max_results=6, planned_query=search_plan or None)
            store_research_results(result.get("query") or plan.tool_prompt or req.message, result.get("provider") or "web", result.get("results") or [])
            answer, answer_route = await _synthesize_search_answer(
                req,
                result,
                search_plan,
                brain_addendum,
                history[:-1] if conversation_enabled else [],
            )
            search_evidence = evidence_from_search(result)
            action = record_action(
                "web_search",
                f"Web search: {result.get('query') or plan.tool_prompt or req.message}"[:220],
                payload={"prompt": req.message, "planned_query": plan.tool_prompt, "search_understanding": search_plan, "query_plan": result.get("query_plan")},
                result={"ok": result.get("ok"), "provider": result.get("provider"), "query": result.get("query"), "result_count": len(result.get("results") or [])},
                evidence=search_evidence,
                requested_level="L3",
                status="completed" if result.get("ok") else "failed",
                summary="Ran focused web search query from SHIMS chat understanding and stored sources for RAG.",
            )
            trust = build_trust(
                route=plan.route,
                evidence=merge_evidence(search_evidence, evidence_from_action(action.get("action"))),
                missing_evidence=[] if result.get("ok") and search_evidence else ["No web provider returned verifiable results."],
                action_id=action.get("action_id", ""),
                ledger_hash=action.get("ledger_hash", ""),
                query_plan=result.get("query_plan"),
            )
            result.update(_trust_fields(trust))
            log_event("tool.success" if result.get("ok") else "tool.error", route=plan.route, provider=result.get("provider") or "web", model=answer_route, latency_ms=(time.perf_counter()-started)*1000, ok=bool(result.get("ok")), message=req.message, metadata={"session_id": session_id, "result_count": len(result.get("results") or []), "errors": result.get("errors"), "search_understanding": search_plan})
            _store_assistant_turn(answer)
            _remember_session_turn(session_id, req.message, answer, route=answer_route, agent=plan.agent, provider=result.get("provider") or "web", model="search", metadata={"search": result, "search_understanding": search_plan, "trust": trust, "action": action})
            yield _jsonl({"type": "token", "content": answer})
            yield _jsonl({"type": "search", "search_result": result, **_trust_fields(trust)})
            yield _jsonl({"type": "done", "session_id": session_id, "model": plan.model, "provider": result.get("provider") or "web", "route": answer_route, "search_result": result, **_trust_fields(trust)})
            return
        result = await _create_media(plan.tool_kind, plan.tool_prompt or req.message, privacy_mode=req.privacy_mode)
        artifact_evidence = evidence_from_artifact(result)
        action = record_action(
            "artifact_generate",
            f"Generate {plan.tool_kind}: {result.get('title') or plan.tool_prompt or req.message}"[:220],
            payload={"kind": plan.tool_kind, "prompt": plan.tool_prompt or req.message},
            result={k: result.get(k) for k in ("ok", "type", "kind", "title", "filename", "url", "file_url", "sha256", "verified")},
            evidence=artifact_evidence,
            requested_level="L3",
            status="completed" if result.get("ok", True) else "failed",
            summary=f"Generated local {plan.tool_kind} artifact with ledger evidence when available.",
        )
        trust = build_trust(
            route=plan.route,
            evidence=merge_evidence(artifact_evidence, evidence_from_action(action.get("action"))),
            missing_evidence=[] if result.get("verified") else ["Artifact was created without a verified document ledger hash."],
            action_id=action.get("action_id", ""),
            ledger_hash=action.get("ledger_hash", ""),
        )
        result.update(_trust_fields(trust))
        log_event("tool.success" if result.get("ok") else "tool.error", route=plan.route, provider="tool", model=plan.model, latency_ms=(time.perf_counter()-started)*1000, ok=bool(result.get("ok")), message=req.message, metadata={"session_id": session_id, "artifact": result})
        answer = f"Done. I created the {plan.tool_kind} and attached it below." if result.get("ok", True) else f"I tried to create the {plan.tool_kind}, but the tool reported: {result.get('error') or result.get('note') or 'unknown error'}"
        _store_assistant_turn(answer)
        ingest_knowledge(
            f"Artifact: {result.get('title') or plan.tool_kind}",
            f"Prompt: {plan.tool_prompt or req.message}\nResult: {json.dumps(result, default=str)[:3000]}",
            source_type="artifact",
            source_uri=result.get("file_url") or result.get("url") or "",
            tags=["artifact", plan.tool_kind],
            importance=0.85,
        )
        _remember_session_turn(session_id, req.message, answer, route=plan.route, agent=plan.agent, provider="tool", model=plan.tool_kind, metadata={"artifact": result, "trust": trust, "action": action})
        yield _jsonl({"type": "token", "content": answer})
        yield _jsonl({"type": "media", "media_result": result, **_trust_fields(trust)})
        yield _jsonl({"type": "done", "session_id": session_id, "model": plan.model, "provider": "tool", "route": plan.route, "media_result": result, **_trust_fields(trust)})
        return
    if plan.route == "local:greeting":
        answer = "Haan, sun raha hoon. Batao, kya karna hai?"
        trust = build_trust(route="local:greeting", evidence=[], missing_evidence=[], requested_level="draft")
        _store_assistant_turn(answer)
        _remember_session_turn(session_id, req.message, answer, route="local:greeting", agent=plan.agent, provider="local", model="", metadata={"trust": trust})
        yield _jsonl({"type": "token", "content": answer})
        log_event("turn.greeting", route="local:greeting", provider="local", model=plan.model, latency_ms=(time.perf_counter()-started)*1000, ok=True, message=req.message, metadata={"session_id": session_id})
        yield _jsonl({"type": "done", "session_id": session_id, "model": plan.model, "provider": "local", "route": "greeting", **_trust_fields(trust)})
        return
    yield _jsonl({"type": "thought", "stage": "generate", "content": f"Generating response with {plan.provider}:{plan.model}..."})
    yield _jsonl({"type": "status", "content": f"Thinking with {plan.provider}:{plan.model}"})
    turn_history = history[-50:] if conversation_enabled else [{"role": "user", "content": req.message}]
    messages = [{"role": "system", "content": _system_prompt() + "\n\n" + brain_addendum}] + turn_history

    # ---- Agentic tool-use loop: run/edit/code/web/self-modify on real machine ----
    # Only enter the heavy tool loop when the turn actually needs it. Omnipotent
    # mode auto-approves gated actions elsewhere; it does not mean every "hi" or
    # "what is 2+2" should pay the latency of wave planning + tool calls.
    if req.agent_mode and _agentic_intent(req.message):
        deep = _needs_deep_model(req.message)
        # Use user's selected cloud provider (anthropic/openai) for agent loop if available,
        # otherwise fall back to Ollama tool-capable models.
        user_provider = (req.provider or plan.provider or "ollama").strip().lower()
        if user_provider in ("anthropic", "openai", "gemini", "deepseek", "kimi"):
            # Cloud provider — use it directly for agent loop
            agent_provider = user_provider
            agent_model = req.model or plan.model or ""
            # If the selected model is a local Ollama model, switch to the cloud default
            if not agent_model or _looks_local_model(agent_model) or agent_model in (await _ollama_names()):
                agent_model = PROVIDER_DEFAULTS.get(agent_provider, "")
        else:
            # Local Ollama — pick a tool-capable model
            agent_provider = "ollama"
            agent_model = await _agent_tool_model(deep=deep) or "qwen2.5-coder:14b"

        if agent_model:
            agent_user = _strip_agent_slash(req.message)
            base_hist = turn_history[:-1] if (turn_history and turn_history[-1].get("role") == "user") else list(turn_history)
            # Agent mode uses the capability preamble as the dominant identity.
            # We deliberately do NOT prepend the conversational _system_prompt() here:
            # its "reply in clear English" instruction conflicts with the JSON-only
            # tool protocol the wave engine needs, causing local models to emit
            # tool-call explanations instead of actual tool calls.
            # We keep the language-matching instruction because voice users often
            # speak Hindi/Hinglish and the final answer must respond in kind.
            agent_sys = (
                agent_loop.get_capability_preamble(minimal=(agent_provider in {"kimi", "deepseek", "qwen"}))
                + "\n\nRespond in the same language the user used. Hindi and Hinglish are welcome when the user speaks them."
                + (("\n\n" + brain_addendum) if agent_provider not in {"kimi", "deepseek", "qwen"} else "")
            )
            agent_msgs = [{"role": "system", "content": agent_sys}] + base_hist + [{"role": "user", "content": agent_user}]
            tier = "deep" if deep else "fast"
            yield _jsonl({"type": "status", "content": f"Agent mode · {agent_provider}:{agent_model} ({tier})"})
            final_info: dict[str, Any] = {}
            try:
                # Some OpenAI-compatible cloud providers (Kimi, DeepSeek, Qwen)
                # do not handle the full 128-tool registry reliably. Fall back to
                # a compact essential tool set so they actually emit valid calls.
                agent_tool_names = (
                    agent_loop.ESSENTIAL_TOOLS
                    if agent_provider in {"kimi", "deepseek", "qwen"}
                    else None
                )
                async for ev in agent_loop.run_agent_loop(
                    message=agent_user, messages=agent_msgs, model=agent_model,
                    provider=agent_provider,
                    tool_names=agent_tool_names,
                    session_id=session_id, create_pending=_agent_create_pending,
                ):
                    if "__final__" in ev:
                        final_info = ev["__final__"]
                    else:
                        yield _jsonl(ev)
            except Exception as exc:
                log_event("agent.loop_error", route="agent-loop-error", provider=agent_provider, model=agent_model, ok=False, message=str(exc)[:200], metadata={"session_id": session_id})
                final_info = {"answer": f"Agent error: {str(exc)[:200]}", "route": "agent-loop-error"}
                yield _jsonl({"type": "token", "content": final_info["answer"]})
            if final_info.get("jobs"):
                _kick_task_drain()
            ans = final_info.get("answer", "")
            _store_assistant_turn(ans)
            trust = build_trust(route="agent-loop", evidence=context_evidence,
                                missing_evidence=[] if context_evidence else ["Agentic tool run; see tool result cards."],
                                requested_level="L3")
            _remember_session_turn(session_id, req.message, ans, route="agent-loop", agent=plan.agent, provider=agent_provider, model=agent_model, metadata={"brain_context": brain_ctx, "trust": trust, "tools_used": final_info.get("tools_used")})
            log_event("turn.done", route="agent-loop", provider=agent_provider, model=agent_model, latency_ms=(time.perf_counter()-started)*1000, ok=True, message=req.message, metadata={"session_id": session_id, "tools_used": final_info.get("tools_used")})
            yield _jsonl({"type": "done", "session_id": session_id, "model": agent_model, "provider": agent_provider, "route": "agent-loop", "tools_used": final_info.get("tools_used"), "jobs": final_info.get("jobs"), **_trust_fields(trust)})
            return

    # ----- Auto-planning: if the request smells like a multi-step workflow, create and run a plan -----
    if not req.agent_mode and _should_auto_plan(req.message):
        try:
            from shared.desktop_planner import plan_from_goal, get_plan
            from shared.plan_executor import run_plan_wave
            plan_result = await asyncio.to_thread(plan_from_goal, req.message, context={"session_id": session_id, "provider": plan.provider, "model": plan.model})
            plan_id = plan_result.plan_id
            yield _jsonl({"type": "thought", "stage": "plan", "content": f"Auto-created plan {plan_id} with {len(plan_result.steps)} steps."})
            yield _jsonl({"type": "status", "content": f"Running plan: {plan_result.goal[:80]}"})
            for wave in range(20):
                wave_result = await asyncio.to_thread(run_plan_wave, plan_id)
                p = wave_result.get("plan", {})
                for s in p.get("steps", []):
                    if s.get("status") == "running" or (s.get("status") in {"done", "failed"} and s.get("finished_at")):
                        yield _jsonl({"type": "status", "content": f"Step {s['step_id']}: {s['status']}"})
                if p.get("status") in {"completed", "failed", "cancelled"}:
                    break
            final_plan = await asyncio.to_thread(get_plan, plan_id)
            summary = f"Plan **{final_plan.status}** after {len(final_plan.steps)} steps."
            done_steps = [s for s in final_plan.steps if s.status == "done"]
            if done_steps:
                summary += "\n\nCompleted steps:\n" + "\n".join(f"- {s.description}" for s in done_steps)
            _store_assistant_turn(summary)
            _remember_session_turn(session_id, req.message, summary, route="auto-plan", agent=plan.agent, provider=plan.provider, model=plan.model, metadata={"plan_id": plan_id, "steps": len(final_plan.steps)})
            yield _jsonl({"type": "token", "content": summary})
            yield _jsonl({"type": "done", "session_id": session_id, "model": plan.model, "provider": plan.provider, "route": "auto-plan", "plan_id": plan_id})
            return
        except Exception as exc:
            log_event("auto_plan.error", route="auto-plan", provider=plan.provider, model=plan.model, ok=False, message=str(exc)[:200], metadata={"session_id": session_id})
            # Fall through to normal generation on plan error

    realtime_request = _is_realtime_request(req)
    # No default max output limit — let models generate as much as needed
    response_max_tokens = req.max_tokens
    if plan.provider == "ollama":
        names = await _ollama_names()
        stream_model = plan.model if plan.model in names else (_preferred_local_model(names, realtime=realtime_request) if names else plan.model)
        if stream_model in names:
            answer = ""
            route = "ollama-local-stream"
            try:
                _pending_stream_chunks: list[bytes] = []
                async def collect_delta(delta: str) -> None:
                    nonlocal answer
                    answer += delta
                    _pending_stream_chunks.append(_jsonl({"type": "token", "content": delta}))

                collector = _collect_ollama_stream(
                    stream_model,
                    messages,
                    realtime=realtime_request,
                    max_tokens=response_max_tokens,
                    on_delta=collect_delta,
                )
                collect_task = asyncio.create_task(collector)
                while not collect_task.done():
                    while _pending_stream_chunks:
                        yield _pending_stream_chunks.pop(0)
                    await asyncio.sleep(0.02)
                answer = await collect_task
                while _pending_stream_chunks:
                    yield _pending_stream_chunks.pop(0)
                if not answer.strip():
                    answer = "I am connected but received an empty local model response. Try qwen2.5:7b or restart Ollama."
                    yield _jsonl({"type": "token", "content": answer})
                low = answer.lower()
                if "text-based" in low or "text based" in low:
                    answer = "I am SHIMS, not a text-only assistant. Ask for images, PDFs, PPTs, audio or video directly and I will run the verified backend tool."
                    yield _jsonl({"type": "token", "content": "\n" + answer})
                    route = "capability-corrected"
                _store_assistant_turn(answer)
                trust = build_trust(
                    route=route,
                    evidence=context_evidence,
                    missing_evidence=[] if context_evidence else ["No tool, web, or retrieved RAG evidence was attached to this model response."],
                    requested_level="draft",
                )
                _remember_session_turn(session_id, req.message, answer, route=route, agent=plan.agent, provider=plan.provider, model=stream_model, metadata={"brain_context": brain_ctx, "trust": trust})
                log_event("turn.done", route=route, provider=plan.provider, model=stream_model, latency_ms=(time.perf_counter()-started)*1000, ok=True, message=req.message, metadata={"session_id": session_id, "answer_preview": answer[:240]})
                yield _jsonl({"type": "done", "session_id": session_id, "model": stream_model, "provider": plan.provider, "route": route, **_trust_fields(trust)})
                return
            except Exception as exc:
                log_event("provider.stream_error", route="ollama-stream-error", provider="ollama", model=stream_model, ok=False, message=str(exc)[:180], metadata={"session_id": session_id})
                fallback = _preferred_local_model(names, realtime=True, exclude={stream_model}, prefer_tiny=True)
                if fallback in names and fallback != stream_model:
                    yield _jsonl({"type": "status", "content": f"{stream_model} was too slow; switching to {fallback}"})
                    fallback_answer = ""
                    pending_fallback: list[bytes] = []

                    async def collect_fallback(delta: str) -> None:
                        nonlocal fallback_answer
                        fallback_answer += delta
                        pending_fallback.append(_jsonl({"type": "token", "content": delta}))

                    try:
                        # Tiny models struggle with the full RAG-augmented prompt, so feed them
                        # a minimal prompt for the fallback answer.
                        minimal_fallback_messages = [
                            {"role": "system", "content": _system_prompt()},
                            {"role": "user", "content": req.message},
                        ]
                        collect_task = asyncio.create_task(_collect_ollama_stream(
                            fallback,
                            minimal_fallback_messages,
                            realtime=True,
                            max_tokens=response_max_tokens,
                            on_delta=collect_fallback,
                            first_token_timeout=18.0,
                        ))
                        while not collect_task.done():
                            while pending_fallback:
                                yield pending_fallback.pop(0)
                            await asyncio.sleep(0.02)
                        fallback_answer = await collect_task
                        while pending_fallback:
                            yield pending_fallback.pop(0)
                        if fallback_answer.strip():
                            trust = build_trust(
                                route="ollama-local-stream-fallback",
                                evidence=context_evidence,
                                missing_evidence=[] if context_evidence else ["No tool, web, or retrieved RAG evidence was attached to this model response."],
                                requested_level="draft",
                            )
                            _store_assistant_turn(fallback_answer)
                            _remember_session_turn(session_id, req.message, fallback_answer, route="ollama-local-stream-fallback", agent=plan.agent, provider="ollama", model=fallback, metadata={"brain_context": brain_ctx, "trust": trust, "fallback_from": stream_model})
                            log_event("turn.done", route="ollama-local-stream-fallback", provider="ollama", model=fallback, latency_ms=(time.perf_counter()-started)*1000, ok=True, message=req.message, metadata={"session_id": session_id, "fallback_from": stream_model})
                            yield _jsonl({"type": "done", "session_id": session_id, "model": fallback, "provider": "ollama", "route": "ollama-local-stream-fallback", **_trust_fields(trust)})
                            return
                    except Exception as fallback_exc:
                        log_event("provider.stream_error", route="ollama-fallback-stream-error", provider="ollama", model=fallback, ok=False, message=str(fallback_exc)[:180], metadata={"session_id": session_id, "fallback_from": stream_model})
                answer = f"`{stream_model}` is installed but is taking too long or failed to generate. Try a cloud provider (Gemini / OpenAI / Claude) from the provider pills, or pick a faster installed local model like `gemma3:1b` or `llama3.2:latest` in Settings."
                route = "ollama-error"
                trust = build_trust(route=route, evidence=context_evidence, missing_evidence=["Local model failed to generate a response."], requested_level="draft")
                _store_assistant_turn(answer)
                _remember_session_turn(session_id, req.message, answer, route=route, agent=plan.agent, provider="ollama", model=stream_model, metadata={"brain_context": brain_ctx, "trust": trust})
                yield _jsonl({"type": "token", "content": answer})
                yield _jsonl({"type": "done", "session_id": session_id, "model": stream_model, "provider": "ollama", "route": route, **_trust_fields(trust)})
                return
    elif plan.provider == "huggingface":
        hf_names = await _hf_names()
        stream_model = plan.model if plan.model in hf_names else (hf_names[0] if hf_names else plan.model)
        if stream_model in hf_names:
            answer = ""
            route = "huggingface-local-stream"
            try:
                pending_hf_chunks: list[bytes] = []
                async def collect_hf_delta(delta: str) -> None:
                    nonlocal answer
                    answer += delta
                    pending_hf_chunks.append(_jsonl({"type": "token", "content": delta}))

                collect_task = asyncio.create_task(_collect_hf_stream(
                    stream_model,
                    messages,
                    realtime=realtime_request,
                    max_tokens=response_max_tokens,
                    on_delta=collect_hf_delta,
                ))
                while not collect_task.done():
                    while pending_hf_chunks:
                        yield pending_hf_chunks.pop(0)
                    await asyncio.sleep(0.02)
                answer = await collect_task
                while pending_hf_chunks:
                    yield pending_hf_chunks.pop(0)
                if not answer.strip():
                    answer = "I am connected to the Hugging Face endpoint but received an empty response. Check that the model is loaded and the endpoint is healthy."
                    yield _jsonl({"type": "token", "content": answer})
                _store_assistant_turn(answer)
                trust = build_trust(
                    route=route,
                    evidence=context_evidence,
                    missing_evidence=[] if context_evidence else ["No tool, web, or retrieved RAG evidence was attached to this model response."],
                    requested_level="draft",
                )
                _remember_session_turn(session_id, req.message, answer, route=route, agent=plan.agent, provider=plan.provider, model=stream_model, metadata={"brain_context": brain_ctx, "trust": trust})
                log_event("turn.done", route=route, provider=plan.provider, model=stream_model, latency_ms=(time.perf_counter()-started)*1000, ok=True, message=req.message, metadata={"session_id": session_id, "answer_preview": answer[:240]})
                yield _jsonl({"type": "done", "session_id": session_id, "model": stream_model, "provider": plan.provider, "route": route, **_trust_fields(trust)})
                return
            except Exception as exc:
                log_event("provider.stream_error", route="huggingface-stream-error", provider="huggingface", model=stream_model, ok=False, message=str(exc)[:180], metadata={"session_id": session_id})
                # Fall through to non-streaming _run_llm
    try:
        answer, route = await _run_llm(plan.provider, plan.model, messages, realtime=realtime_request, max_tokens=response_max_tokens)
    except Exception as exc:
        if plan.provider == "ollama":
            answer = f"Ollama is at `{OLLAMA_HOST}`, but `{plan.model}` failed: {str(exc)[:180]}. Use Settings -> Start Ollama, then pull/select `llama3.2:latest` or `qwen2.5:7b`."
            route = "ollama-error"
        elif plan.provider == "huggingface":
            answer = f"Hugging Face endpoint at `{HUGGINGFACE_HOST}` failed for `{plan.model}`: {str(exc)[:180]}. Verify the endpoint is running and the model is available."
            route = "huggingface-error"
        else:
            answer = f"{plan.provider.title()} route failed for `{plan.model}`: {str(exc)[:180]}. Select an installed Ollama model or configure the provider key."
            route = f"{plan.provider}-error"
    low = answer.lower()
    fake_created = ("saved as" in low or "image generated" in low or "pdf generated" in low) and not _detect_tool_intent(req.message)
    if "text-based" in low or "text based" in low or fake_created:
        answer = "I am SHIMS, not a text-only assistant. For images, PDFs, PPTs, audio or video, ask directly and I will run the backend tool and attach a real file."
        route = "capability-corrected"
    _store_assistant_turn(answer)
    trust = build_trust(
        route=route,
        evidence=context_evidence,
        missing_evidence=[] if context_evidence else ["No tool, web, or retrieved RAG evidence was attached to this model response."],
        requested_level="draft",
    )
    _remember_session_turn(session_id, req.message, answer, route=route, agent=plan.agent, provider=plan.provider, model=plan.model, metadata={"brain_context": brain_ctx, "trust": trust})
    log_event("turn.done", route=route, provider=plan.provider, model=plan.model, latency_ms=(time.perf_counter()-started)*1000, ok=not route.endswith("error"), message=req.message, metadata={"session_id": session_id, "answer_preview": answer[:240]})
    yield _jsonl({"type": "token", "content": answer})
    yield _jsonl({"type": "done", "session_id": session_id, "model": plan.model, "provider": plan.provider, "route": route, **_trust_fields(trust)})


async def _safe_brain_stream(req: ChatRequest) -> AsyncGenerator[bytes, None]:
    # Omnipotent mode auto-approves gated actions, but it should NOT force every
    # casual turn through the heavy agent loop. Only auto-enable the agent loop
    # when the message actually looks agentic (shell/code/build/web/search) or is
    # an explicit slash command. Simple chat questions should use the fast normal
    # LLM path so they reply quickly and don't emit "I ran the tools above...".
    if settings.omnipotent_mode and not req.agent_mode:
        text = (req.message or "").strip()
        if text.startswith("/") or _agentic_intent(text):
            req.agent_mode = True
    try:
        async for chunk in _brain_stream(req):
            yield chunk
    except asyncio.CancelledError:
        raise
    except Exception as exc:
        session_id = req.session_id or str(uuid.uuid4())
        msg = str(exc)[:260] or exc.__class__.__name__
        log_event(
            "turn.stream_exception",
            route="stream-error",
            provider=req.provider or "auto",
            model=req.model or "",
            ok=False,
            message=msg,
            metadata={"session_id": session_id},
        )
        trust = build_trust(
            route="stream-error",
            evidence=[],
            missing_evidence=[f"Backend stream failed before normal completion: {msg}"],
            requested_level="unverified",
        )
        yield _jsonl({
            "type": "token",
            "content": "\n\nI hit a backend streaming error before I could finish. Please retry, or switch to a faster installed local model in Settings.",
        })
        yield _jsonl({"type": "error", "error": msg, "session_id": session_id, "route": "stream-error", **_trust_fields(trust)})
        yield _jsonl({"type": "done", "session_id": session_id, "model": req.model or "", "provider": req.provider or "auto", "route": "stream-error", **_trust_fields(trust)})


_NO_CACHE = {"Cache-Control": "no-cache, no-store, must-revalidate", "Pragma": "no-cache", "Expires": "0"}


def _ensure_android_omni_source() -> dict[str, Any]:
    """Keep the launch Android source pinned to SHIMS Omni, not legacy Sheena."""
    android_java = ROOT / "android_app" / "app" / "src" / "main" / "java" / "com" / "jklifecare" / "shimsmobile" / "MainActivity.java"
    android_manifest = ROOT / "android_app" / "app" / "src" / "main" / "AndroidManifest.xml"
    changed: list[str] = []
    if android_java.exists():
        text = android_java.read_text(encoding="utf-8", errors="ignore")
        fixed = text.replace("// Load Sheena Wellness UI", "// Load SHIMS Omni UI")
        fixed = fixed.replace('webView.loadUrl("file:///android_asset/sheena_wellness/index.html");', 'webView.loadUrl("file:///android_asset/shims_personal/index.html");')
        if fixed != text:
            android_java.write_text(fixed, encoding="utf-8")
            changed.append("MainActivity.java")
    if android_manifest.exists():
        text = android_manifest.read_text(encoding="utf-8", errors="ignore")
        fixed = text.replace('android:label="Sheena Wellness"', 'android:label="SHIMS Omni"')
        if fixed != text:
            android_manifest.write_text(fixed, encoding="utf-8")
            changed.append("AndroidManifest.xml")
    return {"ok": True, "changed": changed}


# ── Mobile Model Transfer ──────────────────────────────────────────────────
# Serves AI model files to the phone over LAN. Desktop downloads the model
# (via hotspot if firewall blocks HuggingFace), places it in storage/mobile_models/,
# and the phone downloads from http://<desktop-ip>:8010/api/mobile-models/<filename>
_MOBILE_MODELS_DIR = ROOT / "storage" / "mobile_models"
_MOBILE_MODELS_DIR.mkdir(parents=True, exist_ok=True)


@app.get("/api/mobile-models")
async def list_mobile_models() -> dict:
    """List model files available for phone download."""
    files = []
    for f in _MOBILE_MODELS_DIR.iterdir():
        if f.is_file() and f.suffix in {".task", ".gguf", ".litertlm", ".bin"}:
            files.append({"name": f.name, "size": f.stat().st_size,
                          "sizeMB": round(f.stat().st_size / (1024 * 1024), 1)})
    return {"ok": True, "models": sorted(files, key=lambda x: x["name"])}


@app.get("/api/mobile-models/{filename}")
async def download_mobile_model(filename: str) -> FileResponse:
    """Serve a model file to the phone. Supports range requests for resume."""
    try:
        path = safe_relative_path(filename, base=_MOBILE_MODELS_DIR)
    except ValueError as exc:
        raise HTTPException(400, f"Invalid filename: {exc}") from exc
    if not path.is_file():
        raise HTTPException(404, f"Model '{filename}' not found. Place it in storage/mobile_models/")
    return FileResponse(path, filename=path.name, media_type="application/octet-stream")


@app.get("/coder-v2", include_in_schema=False)
async def coder_v2_page() -> HTMLResponse:
    path = FRONTEND_DIR / "coder_v2.html"
    return HTMLResponse(path.read_text(encoding="utf-8"), headers=_NO_CACHE)

@app.get("/", include_in_schema=False)
async def index() -> HTMLResponse:
    path = FRONTEND_DIR / "shims_omni.html"
    return HTMLResponse(path.read_text(encoding="utf-8"), headers=_NO_CACHE)


@app.get("/omni-duobot", include_in_schema=False)
async def omni_duobot_page() -> HTMLResponse:
    path = FRONTEND_DIR / "omni_duobot.html"
    return HTMLResponse(path.read_text(encoding="utf-8"), headers=_NO_CACHE)


@app.get("/app", include_in_schema=False)
async def app_index() -> HTMLResponse:
    path = FRONTEND_DIR / "shims_omni.html"
    return HTMLResponse(path.read_text(encoding="utf-8"), headers=_NO_CACHE)


@app.get("/health")
async def health() -> dict[str, Any]:
    android_source = _ensure_android_omni_source()
    models = await _ollama_models_raw()
    return {
        "ok": True,
        "independent": True,
        "app": APP_NAME,
        "version": APP_VERSION,
        "ollama_host": OLLAMA_HOST,
        "ollama_online": bool(models),
        "models": [m["name"] for m in models],
        "enterprise_url": ENTERPRISE_URL,
        "brain": f"single unified v16 verification-first pipeline + {BRAIN_VERSION}",
        "brain_status": omni_brain_status(),
        "android_source": android_source,
        "capabilities": {
            "chat": True,
            "voice": True,
            "tts": True,
            "stt_browser": True,
            "stt_server_optional": True,
            "stt_server_installed": _server_stt_available(),
            "image": True,
            "pdf": True,
            "ppt": True,
            "audio": True,
            "audio_external_api": True,
            "video": True,
            "video_external_api": True,
            "openai_tts": True,
            "openai_sora_video": True,
            "gemini_chat": True,
            "model_pull": True,
            "web_search": True,
            "rag": True,
            "long_term_memory": True,
            "research_capture": True,
            "trust_envelopes": True,
            "action_ledger": True,
            "operator_digest": True,
            "campaign_planner": True,
            "calendar_ics": True,
            "reliability_evals": True,
            "background_learning": BRAIN_BACKGROUND_ENABLED,
            "boot_self_awareness": BOOT_SELF_AWARENESS_ENABLED,
            "self_evolution_proposals": True,
            "multi_agent_orchestration": True,
            "realtime_kernel": True,
            "voice_profiles": True,
        },
    }


@app.get("/launch/readiness")
async def launch_readiness() -> dict[str, Any]:
    android_source = _ensure_android_omni_source()
    models = await _ollama_models_raw(timeout=1.5)
    model_names = [m.get("name", "") for m in models]
    media = _settings.get("media", {})
    frontend_js = FRONTEND_DIR / "js" / "shims_omni.js"
    android_js = ROOT / "android_app" / "app" / "src" / "main" / "assets" / "shims_personal" / "js" / "app.js"
    android_java = ROOT / "android_app" / "app" / "src" / "main" / "java" / "com" / "jklifecare" / "shimsmobile" / "MainActivity.java"
    android_manifest = ROOT / "android_app" / "app" / "src" / "main" / "AndroidManifest.xml"
    apk_candidates = [
        ROOT / "SHIMS_Omni_LaunchReady_debug.apk",
        ROOT / "SHIMS_AI_v17_ondevice_debug.apk",
        ROOT / "SHIMS_AI_v5_LaunchReady_debug.apk",
        ROOT / "SHIMS_AI_v4_AndroidReplyFix_debug.apk",
    ]
    apk = next((p for p in apk_candidates if p.exists()), apk_candidates[0])
    js_text = frontend_js.read_text(encoding="utf-8", errors="ignore") if frontend_js.exists() else ""
    android_text = android_js.read_text(encoding="utf-8", errors="ignore") if android_js.exists() else ""
    java_text = android_java.read_text(encoding="utf-8", errors="ignore") if android_java.exists() else ""
    manifest_text = android_manifest.read_text(encoding="utf-8", errors="ignore") if android_manifest.exists() else ""
    image_plan = [name for name, _ in _image_provider_plan(media.get("image_backend") or "auto")]
    android_ok = (
        android_js.exists()
        and android_java.exists()
        and "backendCandidates" in java_text
        and "android_asset/shims_personal/index.html" in java_text
        and "onShimsNativeTtsDone" in java_text
        and 'android:label="SHIMS Omni"' in manifest_text
        and "getBackend() || await autoDetectBackend(false)" in android_text
    )
    checks = [
        {"id": "backend", "ok": True, "detail": "FastAPI backend is serving this readiness report."},
        {"id": "frontend", "ok": frontend_js.exists() and "approval_request" in js_text and "loadSandboxSidebar" in js_text, "detail": "Omni UI approval/sidebar markers present."},
        {"id": "approvals", "ok": True, "detail": "Yes/no approval router and /approvals endpoints are loaded."},
        {"id": "coder_playground", "ok": True, "detail": "Coder playground endpoints are loaded.", "generated_apps": len(list(GENERATED_APPS_DIR.glob('*/index.html')))},
        {"id": "image_generation", "ok": True, "detail": "Image generation has provider routing plus local visual fallback.", "provider_plan": image_plan, "openai_key": bool(_clean_secret(os.getenv("OPENAI_API_KEY"))), "stable_diffusion_url": bool(media.get("stable_diffusion_url")), "diffusers_enabled": bool(media.get("diffusers_enabled"))},
        {"id": "conversation_mode", "ok": "conversation_mode:state.converseMode" in js_text and "realtime" in js_text, "detail": "Web conversation/realtime payload markers present."},
        {"id": "android", "ok": android_ok, "detail": "Android loads SHIMS Omni asset, relabeled app, backend autodetect fallback, and native TTS callback markers are present.", "apk_exists": apk.exists(), "apk_path": str(apk), "source_guard": android_source},
        {"id": "android_conversation_loop", "ok": "onShimsNativeTtsDone" in android_text and "UtteranceProgressListener" in java_text and "ttsSpeaking" in android_text, "detail": "Android TTS completion resumes conversation listening when enabled and suppresses overlapping STT."},
        {"id": "ollama_models", "ok": bool(model_names), "detail": "Installed local models detected." if model_names else "No local Ollama models detected by readiness check.", "models": model_names[:12]},
        {"id": "self_evolution", "ok": True, "detail": "Guarded propose -> sandbox -> approve -> apply pipeline is loaded.", "pending": len(_list_pending_actions(limit=100)), "proposals": len(list_proposals(limit=100))},
    ]
    warnings = []
    if not _clean_secret(os.getenv("OPENAI_API_KEY")) and not media.get("stable_diffusion_url") and not media.get("diffusers_enabled"):
        warnings.append("No real image provider is configured; local procedural fallback will still produce image files.")
    if not apk.exists():
        warnings.append("Debug APK was not found at the expected root path; run the Android build before distribution.")
    hard_failures = [c for c in checks if not c.get("ok") and c["id"] not in {"ollama_models"}]
    status = "ready" if not hard_failures else "needs_attention"
    return {
        "ok": status == "ready",
        "status": status,
        "date": datetime.now().isoformat(timespec="seconds"),
        "checks": checks,
        "warnings": warnings,
        "recommended_launch_tasks": [
            "Run focused pytest and Android assembleDebug before packaging.",
            "Configure at least one real image provider for photoreal launch demos, or be transparent that local fallback is procedural.",
            "Use llama3.2:latest, gemma3:1b, or qwen2.5:3b for realtime Android conversation.",
        ],
    }


@app.post("/api/chat")
async def api_chat(req: ChatRequest) -> dict[str, Any]:
    """Backward-compatible non-streaming chat route for older Omni clients."""
    answer_parts: list[str] = []
    done: dict[str, Any] = {}
    search_result: dict[str, Any] | None = None
    media_result: dict[str, Any] | None = None
    async for raw in _safe_brain_stream(req):
        try:
            event = json.loads(raw.decode("utf-8"))
        except Exception:
            continue
        if event.get("type") == "token":
            answer_parts.append(str(event.get("content") or ""))
        elif event.get("type") == "search":
            search_result = event.get("search_result")
        elif event.get("type") == "media":
            media_result = event.get("media_result")
        elif event.get("type") == "done":
            done = event
    answer = "".join(answer_parts).strip()
    return {
        "ok": True,
        "independent": True,
        "answer": answer,
        "provider": done.get("provider") or req.provider or "auto",
        "model": done.get("model") or req.model or "",
        "route": done.get("route") or "unknown",
        "search_result": search_result or done.get("search_result"),
        "media_result": media_result or done.get("media_result"),
        **{k: done.get(k) for k in ("trust", "evidence", "confidence", "query_plan", "action_id", "ledger_hash") if k in done},
    }


@app.get("/chat/models")
async def chat_models() -> dict[str, Any]:
    installed = await _ollama_models_raw()
    hf_installed = await _hf_models_raw()
    installed.extend(hf_installed)
    installed = mark_tool_capable(installed)
    names = {m["name"] for m in installed}
    show_all = os.getenv("SHIMS_SHOW_ALL_MODELS", "").strip().lower() in {"1", "true", "yes", "on"}
    rec_source = RECOMMENDED_MODELS if show_all else [m for m in RECOMMENDED_MODELS if m.get("tool_capable")]
    rec = mark_tool_capable([{**m, "installed": (m["provider"] not in {"ollama", "huggingface"}) or m["name"] in names, "configured": _provider_configured(m["provider"]) or m["provider"] == "huggingface"} for m in rec_source])
    cloud = [m for m in rec if m["provider"] != "ollama"]
    # Full curated list for the per-provider model pickers (includes chat-only options).
    all_rec = mark_tool_capable([{**m, "installed": (m["provider"] not in {"ollama", "huggingface"}) or m["name"] in names, "configured": _provider_configured(m["provider"]) or m["provider"] == "huggingface"} for m in RECOMMENDED_MODELS])
    all_cloud = [m for m in all_rec if m["provider"] != "ollama"]
    return {"ok": True, "default": _preferred_local_model(list(names), realtime=True) if names else DEFAULT_OLLAMA_MODEL, "selected_provider": "ollama", "providers": list(PROVIDER_DEFAULTS), "installed": installed, "recommended": rec, "models": installed, "cloud": cloud, "all_cloud": all_cloud, "all_recommended": all_rec, "aliases": _ollama_aliases_payload(), "show_all": show_all}


@app.get("/ollama/status")
async def ollama_status() -> dict[str, Any]:
    models = await _ollama_models_raw()
    return {"ok": bool(models), "online": bool(models), "host": OLLAMA_HOST, "models": models, "message": "Ollama online" if models else f"Ollama not reachable at {OLLAMA_HOST}"}


@app.get("/api/ai/health")
async def api_ai_health() -> dict[str, Any]:
    """Gateway health: Ollama ping + cloud key/circuit-breaker state (same as Enterprise)."""
    from shared.llm_gateway import gateway
    return await gateway.health()


_app_factory_jobs: dict[str, dict[str, Any]] = {}


async def _run_app_factory_tool(name: str, request: Request) -> dict[str, Any]:
    from shared import agent_tools
    payload = await request.json()
    # Run the sync tool (which itself may call AI) in a thread so the event loop stays free.
    return await asyncio.to_thread(agent_tools.run_tool, name, payload, allow_gated=True)


def _app_factory_start_build(spec: dict[str, Any]) -> dict[str, Any]:
    job_id = str(uuid.uuid4())
    _app_factory_jobs[job_id] = {
        "id": job_id,
        "spec": spec,
        "status": "running",
        "started_at": time.time(),
        "result": None,
    }
    threading.Thread(
        target=_run_app_factory_build_in_thread,
        args=(job_id, spec),
        daemon=True,
    ).start()
    return {"ok": True, "job_id": job_id, "status": "running"}


def _run_app_factory_build_in_thread(job_id: str, spec: dict[str, Any]) -> None:
    from shared import agent_tools
    job = _app_factory_jobs[job_id]
    try:
        result = agent_tools.run_tool("app_factory.build_app", {"spec": spec}, allow_gated=True)
        job["status"] = "done" if result.get("ok") else "failed"
        job["result"] = result
    except Exception as exc:
        job["status"] = "failed"
        job["result"] = {"ok": False, "error": str(exc)}


@app.post("/api/app-factory/design")
async def api_app_factory_design(request: Request) -> dict[str, Any]:
    return await _run_app_factory_tool("app_factory.design_app", request)


@app.post("/api/app-factory/build")
async def api_app_factory_build(request: Request) -> dict[str, Any]:
    payload = await request.json()
    spec = payload.get("spec") or payload
    return _app_factory_start_build(spec)


@app.post("/api/app-factory/build/stanford-school")
async def api_app_factory_build_stanford_school() -> dict[str, Any]:
    """One-click SHIMS self-build: generate the Stanford International School app."""
    from scripts.build_stanford_school import SPEC
    return _app_factory_start_build(SPEC)


@app.get("/api/app-factory/build/{job_id}")
async def api_app_factory_build_status(job_id: str) -> dict[str, Any]:
    job = _app_factory_jobs.get(job_id)
    if not job:
        return {"ok": False, "error": "job not found"}
    return {
        "ok": True,
        "job_id": job_id,
        "status": job["status"],
        "result": job.get("result"),
    }


@app.post("/api/app-factory/evolve")
async def api_app_factory_evolve(request: Request) -> dict[str, Any]:
    return await _run_app_factory_tool("app_factory.evolve_app", request)


@app.post("/api/app-factory/test")
async def api_app_factory_test(request: Request) -> dict[str, Any]:
    return await _run_app_factory_tool("app_factory.test_app", request)


@app.post("/api/app-factory/diagnose")
async def api_app_factory_diagnose(request: Request) -> dict[str, Any]:
    return await _run_app_factory_tool("app_factory.diagnose_app", request)


@app.post("/api/app-factory/repair")
async def api_app_factory_repair(request: Request) -> dict[str, Any]:
    return await _run_app_factory_tool("app_factory.repair_app", request)


@app.get("/api/learning/recent")
async def api_learning_recent(limit: int = 12) -> dict[str, Any]:
    """What Shims has learned lately — surfaced in chat so the user can watch
    it evolve: recent skills, autonomous-improvement runs, and feedback signals."""
    out: dict[str, Any] = {"ok": True}
    try:
        from shared.skills import list_skills
        out["skills"] = [
            {"name": s.get("name"), "summary": s.get("summary"), "tags": s.get("tags", []),
             "source": s.get("source"), "updated_at": s.get("updated_at")}
            for s in list_skills(limit=limit)
        ]
    except Exception:
        out["skills"] = []
    try:
        out["improvement_runs"] = list_improvement_runs(limit=5)
    except Exception:
        out["improvement_runs"] = []
    try:
        from shared.omni_brain import list_memories
        fb = list_memories(namespace="omni_feedback", limit=limit)
        out["feedback"] = [
            {"key": m.get("key"), "value": (m.get("value") or "")[:200], "tags": m.get("tags", [])}
            for m in fb
        ]
        out["feedback_counts"] = {
            "preferences": sum(1 for m in fb if "learned_preference" in (m.get("tags") or [])),
            "anti_patterns": sum(1 for m in fb if "anti_pattern" in (m.get("tags") or [])),
        }
    except Exception:
        out["feedback"] = []
    out["autonomous_improvement_enabled"] = IMPROVEMENT_ENABLED
    out["background_learning_enabled"] = BRAIN_BACKGROUND_ENABLED
    return out


@app.post("/api/feedback")
async def api_feedback(request: Request) -> dict[str, Any]:
    """👍/👎 on Omni answers → durable preference/anti-pattern memory the
    prompt builder picks up via memory search."""
    body = await request.json()
    rating = int(body.get("rating") or 0)
    message = str(body.get("message") or "")[:300]
    comment = str(body.get("comment") or "")[:300]
    answer = str(body.get("answer") or "")[:300]
    if not rating or not message:
        return {"ok": False, "detail": "rating and message required"}
    try:
        from shared.omni_brain import remember
        if rating > 0:
            remember("omni_feedback", f"pref:{message[:80]}",
                     f"User liked this answer. Question: {message}." + (f" Note: {comment}" if comment else ""),
                     tags=["feedback", "learned_preference"], weight=1.5, source="feedback")
        else:
            remember("omni_feedback", f"avoid:{message[:80]}",
                     f"User rejected an answer. Question: {message}."
                     + (f" What was wrong: {comment}" if comment else "")
                     + (f" Rejected answer began: {answer}" if answer else ""),
                     tags=["feedback", "anti_pattern"], weight=2.0, source="feedback")
    except Exception as exc:
        return {"ok": False, "detail": str(exc)[:160]}
    return {"ok": True}


@app.get("/ollama/recommended")
async def ollama_recommended() -> dict[str, Any]:
    names = set(await _ollama_names())
    return {"models": [{**m, "installed": m["provider"] != "ollama" or m["name"] in names} for m in RECOMMENDED_MODELS], "aliases": _ollama_aliases_payload()}


@app.post("/ollama/start")
async def ollama_start() -> dict[str, Any]:
    if await _ollama_names():
        return {"ok": True, "already_running": True, "host": OLLAMA_HOST}
    exe = shutil.which("ollama")
    if not exe:
        return {"ok": False, "detail": "ollama command not found. Install Ollama and ensure it is in PATH."}
    try:
        subprocess.Popen([exe, "serve"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, creationflags=getattr(subprocess, "CREATE_NEW_PROCESS_GROUP", 0))
    except Exception as exc:
        return {"ok": False, "detail": str(exc)[:240]}
    await asyncio.sleep(1.5)
    return {"ok": bool(await _ollama_names()), "started": True, "host": OLLAMA_HOST}


@app.post("/ollama/pull")
async def ollama_pull(req: OllamaPullRequest) -> StreamingResponse:
    model = _normalize_ollama_model_name(req.model)
    if not model or not re.match(r"^[A-Za-z0-9_.:/-]+$", model):
        raise HTTPException(400, "Invalid model name")
    async def gen() -> AsyncGenerator[bytes, None]:
        yield _jsonl({"type": "start", "model": model, "status": f"Pulling {model}"})
        had_error = False
        try:
            # Verify Ollama is reachable first
            try:
                async with httpx.AsyncClient(timeout=5.0) as ping:
                    pr = await ping.get(f"{OLLAMA_HOST}/api/tags")
                    pr.raise_for_status()
            except Exception as exc:
                yield _jsonl({"type": "error", "model": model, "detail": f"Ollama not reachable at {OLLAMA_HOST}: {exc}"})
                return
            async with httpx.AsyncClient(timeout=None) as client:
                async with client.stream("POST", f"{OLLAMA_HOST}/api/pull", json={"model": model, "stream": True}) as r:
                    r.raise_for_status()
                    async for line in r.aiter_lines():
                        if not line:
                            continue
                        try: obj = json.loads(line)
                        except Exception: obj = {"status": line}
                        # Ollama reports errors inline
                        if obj.get("error"):
                            had_error = True
                            obj["type"] = "error"
                            obj["detail"] = obj["error"]
                        else:
                            obj.setdefault("type", "progress")
                        obj.setdefault("model", model)
                        yield _jsonl(obj)
            if not had_error:
                yield _jsonl({"type": "done", "model": model, "status": f"Installed {model}"})
        except httpx.HTTPStatusError as exc:
            detail = f"Ollama returned {exc.response.status_code}"
            try:
                detail += ": " + (exc.response.json().get("error") or exc.response.text)[:300]
            except Exception:
                detail += ": " + exc.response.text[:300]
            yield _jsonl({"type": "error", "model": model, "detail": detail})
        except Exception as exc:
            yield _jsonl({"type": "error", "model": model, "detail": str(exc)[:500]})
    return StreamingResponse(gen(), media_type="application/x-ndjson")


@app.delete("/ollama/models/{model:path}")
async def ollama_delete(model: str) -> dict[str, Any]:
    model = model.strip()
    if not model:
        raise HTTPException(400, "missing model")
    async with httpx.AsyncClient(timeout=60) as client:
        r = await client.request("DELETE", f"{OLLAMA_HOST}/api/delete", json={"model": model})
        if r.status_code >= 400:
            raise HTTPException(r.status_code, r.text[:300])
    return {"ok": True, "deleted": model}


@app.post("/brain/turn")
async def brain_turn(req: ChatRequest) -> StreamingResponse:
    return StreamingResponse(_safe_brain_stream(req), media_type="application/x-ndjson")

@app.post("/chat/stream")
async def chat_stream(req: ChatRequest) -> StreamingResponse:
    return StreamingResponse(_safe_brain_stream(req), media_type="application/x-ndjson")

@app.post("/chat/converse")
async def chat_converse(req: ChatRequest) -> StreamingResponse:
    req.conversation_mode = True
    return StreamingResponse(_safe_brain_stream(req), media_type="application/x-ndjson")

@app.post("/api/v11/chat/turn")
async def api_v11_chat_turn(req: ChatRequest) -> StreamingResponse:
    return StreamingResponse(_safe_brain_stream(req), media_type="application/x-ndjson")

@app.get("/api/v11/models")
async def api_v11_models() -> dict[str, Any]:
    return await chat_models()

@app.post("/api/v11/models/pull")
async def api_v11_models_pull(req: OllamaPullRequest) -> StreamingResponse:
    return await ollama_pull(req)

@app.websocket("/converse/ws")
async def converse_ws(ws: WebSocket) -> None:
    await ws.accept()
    try:
        while True:
            data = await ws.receive_text()
            try: payload = json.loads(data)
            except Exception: payload = {"message": data}
            async for chunk in _safe_brain_stream(ChatRequest(**payload)):
                await ws.send_text(chunk.decode("utf-8"))
    except WebSocketDisconnect:
        return


@app.websocket("/ws/enterprise")
async def enterprise_events_ws(ws: WebSocket) -> None:
    """Proxy live Enterprise WebSocket events to Omni clients."""
    await ws.accept()
    if not ENTERPRISE_ENABLED:
        await ws.send_text(json.dumps({"type": "error", "message": "Enterprise integration is not configured"}))
        await ws.close(code=1001, reason="Enterprise not configured")
        return
    import websockets
    from shared.config import settings
    import os
    enterprise_url = getattr(settings, "enterprise_url", os.getenv("SHIMS_ENTERPRISE_URL", "http://127.0.0.1:8020"))
    ws_url = enterprise_url.replace("http://", "ws://").replace("https://", "wss://").rstrip("/") + "/ws/events"
    try:
        async with websockets.connect(ws_url) as ent_ws:
            async def _to_client():
                async for message in ent_ws:
                    await ws.send_text(message)
            async def _to_server():
                while True:
                    data = await ws.receive_text()
                    try: msg = json.loads(data)
                    except Exception: msg = {}
                    if msg.get("action") == "ping":
                        await ws.send_text(json.dumps({"type": "pong"}))
                    else:
                        await ent_ws.send(data)
            await asyncio.gather(_to_client(), _to_server())
    except websockets.exceptions.ConnectionClosed:
        await ws.close(code=1001, reason="Enterprise connection closed")
    except WebSocketDisconnect:
        pass
    except Exception as exc:
        try:
            await ws.send_text(json.dumps({"type": "error", "message": str(exc)[:200]}))
            await ws.close(code=1011, reason="Enterprise proxy error")
        except Exception:
            pass




@app.get("/web/health")
async def web_health() -> dict[str, Any]:
    return {"ok": True, "status": _search_provider_status(), "message": "Configure SHIMS_SEARXNG_URL for best local/private search; otherwise configured APIs or DuckDuckGo fallback are used."}

@app.post("/web/search")
async def web_search(req: SearchRequest) -> dict[str, Any]:
    if not req.query.strip():
        raise HTTPException(400, "Search query required")
    result = await _web_search(req.query.strip(), req.max_results, req.provider)
    store_research_results(req.query.strip(), result.get("provider") or req.provider or "web", result.get("results") or [])
    evidence = evidence_from_search(result)
    action = record_action(
        "web_search",
        f"Web search: {result.get('query') or req.query}"[:220],
        payload={"query": req.query, "max_results": req.max_results, "provider": req.provider, "query_plan": result.get("query_plan")},
        result={"ok": result.get("ok"), "provider": result.get("provider"), "query": result.get("query"), "result_count": len(result.get("results") or [])},
        evidence=evidence,
        requested_level="L3",
        status="completed" if result.get("ok") else "failed",
        summary="Ran planned web search query.",
    )
    trust = build_trust(
        route="web:search",
        evidence=merge_evidence(evidence, evidence_from_action(action.get("action"))),
        missing_evidence=[] if result.get("ok") and evidence else ["No web provider returned verifiable results."],
        action_id=action.get("action_id", ""),
        ledger_hash=action.get("ledger_hash", ""),
        query_plan=result.get("query_plan"),
    )
    result.update(_trust_fields(trust))
    return result

@app.post("/web/plan")
async def web_plan(req: SearchRequest) -> dict[str, Any]:
    return {"ok": True, "plan": plan_search_query(req.query, web_mode=True, force_search=True).to_dict()}

@app.get("/web/plan")
async def web_plan_get(q: str = "", max_results: int = 6, provider: str | None = None) -> dict[str, Any]:
    return {"ok": True, "plan": plan_search_query(q, web_mode=True, force_search=True).to_dict(), "max_results": max_results, "provider": provider or "auto"}



@app.post("/web/deep-research")
async def web_deep_research(req: SearchRequest) -> dict[str, Any]:
    if not req.query.strip():
        raise HTTPException(400, "Query required")
    result = await deep_research(req.query.strip(), _web_search, max_search_results=req.max_results, max_pages=3)
    if result.get("ok"):
        store_research_results(req.query.strip(), "web_crawler", result.get("sources", []))
    return result

@app.post("/web/fetch")
async def web_fetch(req: Request) -> dict[str, Any]:
    body = await req.json()
    url = str(body.get("url", "")).strip()
    if not url:
        raise HTTPException(400, "URL required")
    return await fetch_page(url)

@app.get("/brain/status")
async def brain_status_endpoint() -> dict[str, Any]:
    return {**omni_brain_status(), "background_interval_seconds": BRAIN_BACKGROUND_INTERVAL_SECONDS}


@app.post("/brain/reindex-vectors")
async def brain_reindex_vectors_endpoint() -> dict[str, Any]:
    """One-time migration: build semantic embeddings for all existing memories, knowledge chunks, and research items."""
    return await asyncio.to_thread(reindex_vectors)


@app.post("/api/brain/self-index")
async def brain_self_index_endpoint(force: bool = False) -> dict[str, Any]:
    """Phase 3.1 Self-indexer: ingest allowed SHIMS source roots into the omni-brain."""
    return await asyncio.to_thread(index_shims_source, force=force)


@app.get("/self/status")
@app.get("/brain/self")
async def self_status_endpoint() -> dict[str, Any]:
    data = latest_self_model()
    data["boot_self_awareness_enabled"] = BOOT_SELF_AWARENESS_ENABLED
    data["startup_task_running"] = bool(_self_awareness_task is not None and not _self_awareness_task.done())
    return data


@app.get("/self/notes")
async def self_notes_endpoint() -> dict[str, Any]:
    notes = latest_self_notes()
    return {"ok": bool(notes), "notes": notes, "latest": latest_self_model()}


@app.post("/self/boot-audit")
@app.post("/brain/self/boot-audit")
async def self_boot_audit_endpoint() -> dict[str, Any]:
    return await asyncio.to_thread(run_boot_self_audit, app_name=APP_NAME, app_version=APP_VERSION)

@app.post("/improvement/run")
async def improvement_run_endpoint() -> dict[str, Any]:
    from shared.improvement_loop import run_improvement_cycle
    return await asyncio.to_thread(run_improvement_cycle, _system_prompt())

@app.get("/improvement/runs")
async def improvement_runs_endpoint(limit: int = 20) -> dict[str, Any]:
    from shared.improvement_loop import list_improvement_runs
    return {"ok": True, "runs": list_improvement_runs(limit=limit)}


class PlanCreateRequest(BaseModel):
    goal: str
    steps: list[dict[str, Any]] = []
    context: dict[str, Any] = {}


class PlanIdRequest(BaseModel):
    plan_id: str


class ScheduleCreateRequest(BaseModel):
    title: str
    schedule_type: str
    when: str
    action_type: str
    payload: dict[str, Any] = {}


class ScheduleIdRequest(BaseModel):
    task_id: str


@app.post("/api/plans")
async def plans_create_endpoint(req: PlanCreateRequest) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "plan.create",
                                   {"goal": req.goal, "steps": req.steps, "context": req.context}, allow_gated=True)


@app.get("/api/plans")
async def plans_list_endpoint(status: str | None = None, limit: int = 20) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "plan.list", {"status": status, "limit": limit})


@app.post("/api/plans/get")
async def plans_get_endpoint(req: PlanIdRequest) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "plan.get", {"plan_id": req.plan_id})


@app.post("/api/plans/cancel")
async def plans_cancel_endpoint(req: PlanIdRequest) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "plan.cancel", {"plan_id": req.plan_id})


@app.post("/api/plans/run-wave")
async def plans_run_wave_endpoint(req: PlanIdRequest) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "plan.run_wave", {"plan_id": req.plan_id})


@app.post("/api/plans/run")
async def plans_run_endpoint(req: PlanIdRequest) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "plan.run", {"plan_id": req.plan_id, "max_waves": 20})


@app.post("/api/schedule")
async def schedule_create_endpoint(req: ScheduleCreateRequest) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "schedule.create", {
        "title": req.title,
        "schedule_type": req.schedule_type,
        "when": req.when,
        "action_type": req.action_type,
        "payload": req.payload,
    })


@app.get("/api/schedule")
async def schedule_list_endpoint(enabled_only: bool = False, limit: int = 100) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "schedule.list", {"enabled_only": enabled_only, "limit": limit})


@app.post("/api/schedule/cancel")
async def schedule_cancel_endpoint(req: ScheduleIdRequest) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "schedule.cancel", {"task_id": req.task_id})


class MemorySaveRequest(BaseModel):
    content: str
    key: str = ""
    namespace: str = "agent"
    tags: list[str] = []
    weight: float = 1.0


class MemorySearchRequest(BaseModel):
    query: str
    limit: int = 8


class MemoryIngestMediaRequest(BaseModel):
    path: str
    kind: str  # image | audio | video | screen
    title: str = ""
    tags: list[str] = []
    metadata: dict[str, Any] = {}


@app.post("/api/memory/save")
async def memory_save_endpoint(req: MemorySaveRequest) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "memory.save", {
        "content": req.content,
        "key": req.key,
        "namespace": req.namespace,
        "tags": req.tags,
        "weight": req.weight,
    })


@app.post("/api/memory/search")
async def memory_search_endpoint(req: MemorySearchRequest) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "memory.search", {"query": req.query, "limit": req.limit})


@app.post("/api/memory/ingest-media")
async def memory_ingest_media_endpoint(req: MemoryIngestMediaRequest) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "memory.ingest_media", {
        "path": req.path,
        "kind": req.kind,
        "title": req.title,
        "tags": req.tags,
        "metadata": req.metadata,
    })


@app.post("/brain/context")
@app.post("/rag/search")
async def brain_context(req: BrainContextRequest) -> dict[str, Any]:
    ctx = retrieve_context(req.query, limit=req.limit)
    evidence = evidence_from_brain_context(ctx)
    trust = build_trust(
        route="rag:search",
        evidence=evidence,
        missing_evidence=[] if evidence else ["No matching memory/RAG/research evidence found."],
    )
    ctx.update(_trust_fields(trust))
    return ctx


@app.post("/brain/ingest")
@app.post("/rag/ingest")
async def brain_ingest(req: BrainIngestRequest) -> dict[str, Any]:
    return ingest_knowledge(
        req.title,
        req.text,
        source_type=req.source_type,
        source_uri=req.source_uri,
        tags=req.tags,
        importance=req.importance,
    )


@app.post("/brain/learn")
async def brain_learn(req: BrainLearnRequest) -> dict[str, Any]:
    return run_learning_cycle(limit=req.limit, propose=req.propose)


@app.get("/brain/tasks")
async def brain_tasks(status: str | None = None, limit: int = 50) -> dict[str, Any]:
    return {"ok": True, "tasks": brain_list_tasks(status=status, limit=limit)}


@app.post("/brain/tasks")
async def brain_tasks_enqueue(req: TaskEnqueueRequest) -> dict[str, Any]:
    return brain_schedule_task(req.task_type, req.title or req.task_type, req.payload or {}, priority=req.priority)


@app.post("/brain/tasks/run")
async def brain_tasks_run(max_tasks: int = 25) -> dict[str, Any]:
    """Drain the background task queue now (consolidation, skill extraction, etc.)."""
    return await asyncio.to_thread(brain_drain_tasks, max(1, min(int(max_tasks), 200)))


# ── Generic Background Task API (user-facing) ──
@app.get("/api/tasks")
async def api_tasks_list(status: str = "", limit: int = 50) -> dict[str, Any]:
    from shared.omni_brain import list_tasks
    return {"ok": True, "tasks": list_tasks(status=status or None, limit=limit)}

@app.get("/api/tasks/{task_id}")
async def api_tasks_get(task_id: int) -> dict[str, Any]:
    from shared.omni_brain import get_task
    task = get_task(task_id)
    if not task:
        return {"ok": False, "error": "Task not found"}
    return {"ok": True, "task": task}

@app.post("/api/tasks")
async def api_tasks_enqueue(req: Request) -> dict[str, Any]:
    from shared.omni_brain import schedule_task
    body = await req.json()
    return schedule_task(
        body.get("task_type", ""),
        body.get("title", "Background task"),
        body.get("payload", {}),
        priority=body.get("priority", 5),
    )

@app.post("/api/tasks/{task_id}/cancel")
async def api_tasks_cancel(task_id: int) -> dict[str, Any]:
    from shared.omni_brain import cancel_task
    return cancel_task(task_id)

@app.post("/api/tasks/{task_id}/run")
async def api_tasks_run_now(task_id: int) -> dict[str, Any]:
    from shared.omni_brain import get_task, execute_task, _finish_task
    task = get_task(task_id)
    if not task:
        return {"ok": False, "error": "Task not found"}
    if task["status"] not in ("queued", "running"):
        return {"ok": False, "error": f"Task already {task['status']}"}
    outcome = await asyncio.to_thread(execute_task, task)
    status = "done" if outcome.get("ok") and outcome.get("status") != "failed" else (outcome.get("status") or "failed")
    _finish_task(task_id, status, outcome.get("result") or {"error": outcome.get("error")})
    return {"ok": True, "status": status, "result": outcome.get("result"), "error": outcome.get("error")}


@app.post("/api/settings/auto-evolution")
async def api_settings_auto_evolution(req: Request) -> dict[str, Any]:
    """Toggle auto-evolution mode. When enabled, SHIMS schedules periodic reflection tasks."""
    from shared.config import settings
    from shared.omni_brain import schedule_task
    body = await req.json()
    enabled = bool(body.get("enabled", True))
    # Note: This toggles in-memory only for the current process. For persistence, update .env.
    # We schedule a reflection task immediately if enabled.
    if enabled:
        schedule_task("reflect", "Auto-evolution reflection cycle", {"trigger": "auto_evolution"}, priority=3)
    return {"ok": True, "auto_evolution": enabled, "note": "Auto-evolution " + ("enabled" if enabled else "disabled")}


# ── Browser Agent API (Kimi Claw) ──
@app.post("/api/browser/visit")
async def browser_visit(req: Request) -> dict[str, Any]:
    from shared.browser_agent import visit
    body = await req.json()
    return await visit(body.get("url", ""), body.get("wait_for", ""), body.get("scroll", True))

@app.post("/api/browser/search")
async def browser_search(req: Request) -> dict[str, Any]:
    from shared.browser_agent import search
    body = await req.json()
    return await search(body.get("query", ""), min(int(body.get("max_results") or 8), 12))

@app.post("/api/browser/click")
async def browser_click(req: Request) -> dict[str, Any]:
    from shared.browser_agent import click
    body = await req.json()
    return await click(body.get("url", ""), body.get("selector", ""), body.get("text", ""))

@app.post("/api/browser/extract")
async def browser_extract(req: Request) -> dict[str, Any]:
    from shared.browser_agent import extract
    body = await req.json()
    return await extract(body.get("url", ""), body.get("selector", ""))

@app.post("/api/browser/fill_form")
async def browser_fill_form(req: Request) -> dict[str, Any]:
    from shared.browser_agent import fill_form
    body = await req.json()
    return await fill_form(body.get("url", ""), body.get("fields", {}), body.get("submit_selector", ""))

@app.post("/api/desktop/screenshot")
async def desktop_screenshot() -> dict[str, Any]:
    """Take a desktop screenshot and return the file URL."""
    try:
        from PIL import ImageGrab
        import base64, io, time
        img = ImageGrab.grab()
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
        filename = f"desktop_{int(time.time())}.png"
        out_path = agent_tools.STORAGE_DIR / "screenshots" / filename
        out_path.parent.mkdir(parents=True, exist_ok=True)
        out_path.write_bytes(buf.getvalue())
        url = f"/media/files/screenshot/{filename}"
        return {"ok": True, "filename": filename, "url": url, "base64": b64, "path": str(out_path)}
    except Exception as exc:
        return {"ok": False, "error": str(exc)[:200]}


@app.post("/api/browser/screenshot")
async def browser_screenshot(req: Request) -> dict[str, Any]:
    from shared.browser_agent import screenshot
    body = await req.json()
    return await screenshot(body.get("url", ""), body.get("selector", ""), body.get("full_page", False))



# --- Desktop Bridge integration ---
_BRIDGE_URI = os.environ.get("SHIMS_DESKTOP_BRIDGE_URI", "ws://localhost:9876/bridge")
_BRIDGE_TOKEN = os.environ.get("SHIMS_DESKTOP_BRIDGE_TOKEN", "")

async def _bridge_client():
    from desktop_bridge.bridge_client import DesktopBridge
    if not _BRIDGE_TOKEN:
        return None
    return DesktopBridge(_BRIDGE_URI, _BRIDGE_TOKEN)

@app.post("/api/desktop/bridge/command")
async def desktop_bridge_command(req: Request) -> dict[str, Any]:
    """Execute a command on the connected desktop bridge."""
    body = await req.json()
    bridge = await _bridge_client()
    if bridge is None:
        return {"ok": False, "error": "Desktop bridge not configured. Set SHIMS_DESKTOP_BRIDGE_TOKEN and SHIMS_DESKTOP_BRIDGE_URI."}
    cmd_type = body.get("type", "shell")
    if cmd_type == "shell":
        return await bridge.shell(body.get("command", ""), body.get("cwd"), body.get("timeout", 60))
    if cmd_type == "screenshot":
        return await bridge.screenshot()
    if cmd_type == "system_info":
        return await bridge.system_info()
    if cmd_type == "find_file":
        return await bridge.find_file(body.get("name", ""), body.get("root", "C:\\"))
    if cmd_type == "read_file":
        return await bridge.read_file(body.get("path", ""))
    if cmd_type == "write_file":
        return await bridge.write_file(body.get("path", ""), body.get("content", ""))
    if cmd_type == "ping":
        return await bridge.ping()
    return {"ok": False, "error": f"Unknown bridge command: {cmd_type}"}

@app.get("/api/desktop/bridge/status")
async def desktop_bridge_status() -> dict[str, Any]:
    """Check if desktop bridge is reachable."""
    bridge = await _bridge_client()
    if bridge is None:
        return {"ok": False, "connected": False, "error": "Bridge token not configured"}
    res = await bridge.ping()
    return {"ok": res.get("ok", False), "connected": res.get("ok", False), "detail": res}


@app.post("/api/desktop/bridge/launch")
async def desktop_bridge_launch() -> dict[str, Any]:
    """Launch the desktop bridge process and confirm it is reachable.

    One-click bridge startup for the frontend Desktop Controls panel.
    """
    # If already reachable, just report it.
    bridge = await _bridge_client()
    if bridge is not None:
        res = await bridge.ping()
        if res.get("ok"):
            return {
                "ok": True,
                "started": False,
                "connected": True,
                "uri": _BRIDGE_URI,
                "token": _BRIDGE_TOKEN,
                "message": "Bridge is already running and connected.",
            }

    if not _BRIDGE_TOKEN:
        return {
            "ok": False,
            "started": False,
            "connected": False,
            "error": "Desktop bridge token not configured. Set SHIMS_DESKTOP_BRIDGE_TOKEN in .env.",
        }

    popen_kwargs: dict[str, Any] = {
        "cwd": str(ROOT),
        "stdout": subprocess.DEVNULL,
        "stderr": subprocess.DEVNULL,
    }
    if sys.platform == "win32":
        popen_kwargs["creationflags"] = subprocess.CREATE_NO_WINDOW

    proc = subprocess.Popen(
        [
            sys.executable,
            str(ROOT / "desktop_bridge" / "bridge_server.py"),
            "--host", "0.0.0.0",
            "--port", "9876",
            "--token", _BRIDGE_TOKEN,
        ],
        **popen_kwargs,
    )

    # Wait up to 10 seconds for the bridge to accept connections.
    for _ in range(20):
        await asyncio.sleep(0.5)
        try:
            bridge = await _bridge_client()
            if bridge is not None:
                res = await bridge.ping()
                if res.get("ok"):
                    return {
                        "ok": True,
                        "started": True,
                        "connected": True,
                        "pid": proc.pid,
                        "uri": _BRIDGE_URI,
                        "token": _BRIDGE_TOKEN,
                        "message": "Bridge launched and connected.",
                    }
        except Exception:
            pass

    return {
        "ok": False,
        "started": True,
        "connected": False,
        "pid": proc.pid,
        "error": "Bridge process started but did not respond in time. Check logs or firewall.",
    }
@app.post("/api/browser/scroll")
async def browser_scroll(req: Request) -> dict[str, Any]:
    from shared.browser_agent import scroll
    body = await req.json()
    return await scroll(body.get("url", ""), body.get("direction", "down"), int(body.get("amount") or 800))

@app.post("/api/vision/describe")
async def vision_describe_endpoint(req: Request) -> dict[str, Any]:
    from shared.vision import describe_image
    body = await req.json()
    return await asyncio.to_thread(
        describe_image,
        str(body.get("source", "")).strip(),
        str(body.get("prompt", "Describe this image concisely.")).strip(),
        str(body.get("backend", "auto")).strip(),
    )


class InterpreterRunRequest(BaseModel):
    code: str
    timeout: int = 60


class InterpreterReadRequest(BaseModel):
    workdir: str
    path: str


@app.post("/api/interpreter/run")
async def interpreter_run_endpoint(req: InterpreterRunRequest) -> dict[str, Any]:
    from shared.code_interpreter import run_interpreter
    return await asyncio.to_thread(run_interpreter, req.code, req.timeout)


@app.post("/api/interpreter/read")
async def interpreter_read_endpoint(req: InterpreterReadRequest) -> dict[str, Any]:
    from shared.code_interpreter import read_artifact
    return await asyncio.to_thread(read_artifact, req.workdir, req.path)


@app.get("/media/files/screenshot/{filename}")
async def serve_screenshot(filename: str):
    from fastapi.responses import FileResponse
    from shared.browser_agent import SCREENSHOT_DIR
    try:
        path = safe_relative_path(filename, base=Path(SCREENSHOT_DIR))
    except ValueError as exc:
        raise HTTPException(400, f"Invalid filename: {exc}") from exc
    if not path.exists():
        return {"detail": "Not found"}
    return FileResponse(path)


# ── Agentic core: capabilities, tools, allowed roots, background coder jobs ──
class AgentJobRequest(BaseModel):
    goal: str
    name: str | None = None


class AgentRootRequest(BaseModel):
    path: str


class AgentToolRequest(BaseModel):
    tool: str
    args: dict[str, Any] = {}


class AgentSwarmRequest(BaseModel):
    prompt: str
    agent_ids: list[str] | None = None
    context: dict[str, Any] | None = None
    shared_context: dict[str, Any] | None = None
    use_llm: bool = True  # deterministic path when false
    orchestrate: bool = True  # real meta-orchestrator when true


def _job_events_path(job_id: str | int):
    return agent_tools.STORAGE_DIR / "coder_jobs" / str(job_id) / "events.jsonl"


def _kick_task_drain(n: int = 5) -> None:
    try:
        asyncio.create_task(asyncio.to_thread(brain_drain_tasks, n))
    except RuntimeError:
        pass


@app.get("/agent/capabilities")
async def agent_capabilities() -> dict[str, Any]:
    return agent_tools.capabilities()


@app.get("/agent/tools")
async def agent_tools_list() -> dict[str, Any]:
    return {"ok": True, "specs": agent_tools.tool_specs()}


@app.get("/agent/roots")
async def agent_roots_get() -> dict[str, Any]:
    return {"ok": True, "repo_root": str(agent_tools.REPO_ROOT), "allowed_roots": agent_tools.list_allowed_roots()}


@app.post("/agent/roots")
async def agent_roots_add(req: AgentRootRequest) -> dict[str, Any]:
    return agent_tools.add_allowed_root(req.path)


@app.delete("/agent/roots")
async def agent_roots_remove(path: str) -> dict[str, Any]:
    return agent_tools.remove_allowed_root(path)


@app.post("/agent/tool")
async def agent_tool_run(req: AgentToolRequest) -> dict[str, Any]:
    """Run a single tool directly (safe tools execute; gated tools return needs_approval)."""
    return await asyncio.to_thread(agent_tools.run_tool, req.tool, req.args, allow_gated=False)


@app.post("/agent/run")
async def agent_run(req: ChatRequest) -> StreamingResponse:
    """Force the agentic loop for this message (same stream as chat)."""
    req.agent_mode = True
    if not _agentic_intent(req.message):
        req.message = "/do " + (req.message or "")
    return StreamingResponse(_safe_brain_stream(req), media_type="application/x-ndjson")


@app.post("/agent/swarm")
async def agent_swarm(req: AgentSwarmRequest) -> dict[str, Any]:
    """Dispatch multiple agents in parallel and return a synthesized answer.

    When ``orchestrate`` is true (default), the real meta-orchestrator analyzes
    the prompt, builds a plan, and runs coder/reviewer/tester agents in waves.
    When ``use_llm`` is false, a deterministic offline synthesizer is used.
    """
    if not req.use_llm:
        from shared.swarm import swarm_dict
        return swarm_dict(req.prompt, agent_roles=req.agent_ids or None)
    return await asyncio.to_thread(
        agent_tools.run_tool,
        "agent.swarm",
        {
            "prompt": req.prompt,
            "agent_ids": req.agent_ids,
            "context": req.context or {},
            "shared_context": req.shared_context or {},
            "orchestrate": req.orchestrate,
            "use_llm": req.use_llm,
        },
        allow_gated=True,
    )


@app.post("/agent/jobs")
async def agent_jobs_spawn(req: AgentJobRequest) -> dict[str, Any]:
    res = await asyncio.to_thread(agent_tools.run_tool, "coder.spawn",
                                  {"goal": req.goal, "name": req.name or req.goal[:48]}, allow_gated=True)
    _kick_task_drain()
    return res


@app.get("/agent/jobs")
async def agent_jobs_list(limit: int = 40) -> dict[str, Any]:
    jobs = [t for t in brain_list_tasks(limit=200) if t.get("task_type") == "coder_job"][:limit]
    return {"ok": True, "jobs": jobs}


@app.get("/agent/jobs/{job_id}")
async def agent_job_get(job_id: str) -> dict[str, Any]:
    job = next((t for t in brain_list_tasks(limit=200) if str(t.get("id")) == str(job_id)), None)
    events: list[dict[str, Any]] = []
    ev_path = _job_events_path(job_id)
    if ev_path.exists():
        for line in ev_path.read_text(encoding="utf-8", errors="ignore").splitlines()[-200:]:
            try:
                events.append(json.loads(line))
            except Exception:
                continue
    return {"ok": True, "job": job, "events": events}


@app.post("/agent/jobs/{job_id}/cancel")
async def agent_job_cancel(job_id: str) -> dict[str, Any]:
    try:
        from shared import omni_brain as _ob
        with _ob._connect() as con:  # type: ignore[attr-defined]
            con.execute("UPDATE background_tasks SET status='cancelled', updated_at=? WHERE id=? AND status IN ('queued','running')",
                        (_ob._now(), int(job_id)))
            con.commit()
        return {"ok": True, "cancelled": job_id}
    except Exception as exc:
        return {"ok": False, "error": str(exc)[:200]}


@app.get("/agent/jobs/{job_id}/stream")
async def agent_job_stream(job_id: str) -> StreamingResponse:
    ev_path = _job_events_path(job_id)

    async def gen() -> AsyncGenerator[bytes, None]:
        pos = 0
        idle = 0
        yield b": connected\n\n"
        while idle < 1200:  # end after ~120s of no new events
            if ev_path.exists():
                try:
                    with ev_path.open("r", encoding="utf-8") as f:
                        f.seek(pos)
                        chunk = f.read()
                        pos = f.tell()
                except Exception:
                    chunk = ""
                if chunk:
                    idle = 0
                    for line in chunk.splitlines():
                        if not line.strip():
                            continue
                        yield f"data: {line}\n\n".encode("utf-8")
                        if '"stage": "done"' in line or '"stage":"done"' in line:
                            yield b"event: end\ndata: {}\n\n"
                            return
            idle += 1
            await asyncio.sleep(0.1)
        yield b"event: end\ndata: {}\n\n"

    return StreamingResponse(gen(), media_type="text/event-stream",
                             headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})


# ── Desktop cowork: file organization (confined to the workspace) ───────────
@app.get("/files/workspace")
async def files_workspace_get() -> dict[str, Any]:
    from shared import fileops
    return {"ok": True, "workspace": str(fileops.get_workspace())}


@app.post("/files/workspace")
async def files_workspace_set(req: FileWorkspaceRequest) -> dict[str, Any]:
    from shared import fileops
    try:
        return fileops.set_workspace(req.path)
    except Exception as exc:
        raise HTTPException(400, str(exc))


@app.get("/files/tree")
async def files_tree(subpath: str = "", max_entries: int = 500) -> dict[str, Any]:
    from shared import fileops
    try:
        return fileops.tree(subpath, max_entries=max_entries)
    except ValueError as exc:
        raise HTTPException(400, str(exc))


@app.get("/files/summary")
async def files_summary(subpath: str = "") -> dict[str, Any]:
    from shared import fileops
    try:
        return fileops.summarize_folder(subpath)
    except ValueError as exc:
        raise HTTPException(400, str(exc))


@app.get("/files/read")
async def files_read(relpath: str) -> dict[str, Any]:
    from shared import fileops
    try:
        return fileops.read_text(relpath)
    except ValueError as exc:
        raise HTTPException(400, str(exc))


@app.post("/files/search")
async def files_search(req: FileSearchRequest) -> dict[str, Any]:
    from shared import fileops
    return fileops.search(req.query, in_content=req.in_content)


@app.get("/files/duplicates")
async def files_duplicates() -> dict[str, Any]:
    from shared import fileops
    return await asyncio.to_thread(fileops.find_duplicates)


@app.post("/files/organize/plan")
async def files_organize_plan(subpath: str = "") -> dict[str, Any]:
    from shared import fileops
    try:
        return fileops.propose_organization(subpath)
    except ValueError as exc:
        raise HTTPException(400, str(exc))


@app.post("/files/organize/apply")
async def files_organize_apply(req: FileOrganizeApplyRequest) -> dict[str, Any]:
    from shared import fileops
    return fileops.apply_moves(req.moves)


@app.post("/files/organize/undo")
async def files_organize_undo(req: FileUndoRequest) -> dict[str, Any]:
    from shared import fileops
    return fileops.undo_moves(req.undo_id)


# ── OCR: image → text ───────────────────────────────────────────────────────
@app.get("/ocr/health")
async def ocr_health() -> dict[str, Any]:
    from shared import ocr_service
    return {"ok": True, "available": ocr_service.ocr_available(),
            "engine": ocr_service.engine_name(),
            "note": "Offline OCR ready." if ocr_service.ocr_available()
                    else "OCR disabled on Windows by default (crash risk). Set SHIMS_ENABLE_RAPIDOCR=1 or pull an Ollama vision model."}


@app.post("/ocr")
async def ocr_extract(file: UploadFile = File(...)) -> dict[str, Any]:
    from shared import ocr_service
    data = await file.read()
    if not data:
        raise HTTPException(400, "empty file")
    result = await asyncio.to_thread(ocr_service.ocr_image_bytes, data)
    if not result.get("ok") and result.get("engine") == "none":
        return JSONResponse(status_code=428, content=result)
    return result


# ── Coder workspace (the "separate codex") ──────────────────────────────────
@app.get("/coder/projects")
async def coder_projects() -> dict[str, Any]:
    from shared import coder
    return {"ok": True, "projects": coder.list_projects()}


@app.post("/coder/project")
async def coder_create(req: CoderCreateRequest) -> dict[str, Any]:
    from shared import coder
    return {"ok": True, "project": coder.create_project(req.name, req.goal)}


@app.get("/coder/project/{project_id}")
async def coder_get(project_id: str) -> dict[str, Any]:
    from shared import coder
    return coder.get_project(project_id)


@app.post("/coder/write")
async def coder_write(req: CoderWriteRequest) -> dict[str, Any]:
    from shared import coder
    try:
        return coder.write_file(req.project_id, req.path, req.content)
    except ValueError as exc:
        raise HTTPException(400, str(exc))


@app.post("/coder/run")
async def coder_run(req: CoderRunRequest) -> dict[str, Any]:
    from shared import coder
    return await asyncio.to_thread(coder.run_project, req.project_id, req.entry)


@app.post("/coder/iterate")
async def coder_iterate(req: CoderIterateRequest) -> dict[str, Any]:
    from shared import coder
    return await coder.iterate(req.project_id, req.instruction,
                               provider=req.provider, model=req.model, max_steps=req.max_steps)


# ── Omni Builder: develop the codebase at scale via the configured Anthropic key ──
@app.post("/builder/run")
async def builder_run(req: BuilderRunRequest) -> dict[str, Any]:
    """Anthropic-driven build step applied through the safety harness (compile-gate +
    backup + rollback + git). apply=false previews; apply=true writes + commits."""
    from shared import omni_builder
    return await omni_builder.build_task(
        req.instruction, targets=req.targets, context=req.context,
        provider=req.provider or "anthropic", model=req.model, apply=req.apply,
    )


@app.post("/coder/install")
async def coder_install(req: CoderInstallRequest) -> dict[str, Any]:
    from shared import coder
    return await asyncio.to_thread(coder._install_requirements, req.project_id)


@app.get("/coder/files/{project_id}")
async def coder_files(project_id: str) -> dict[str, Any]:
    from shared import coder
    return {"ok": True, "files": coder.list_files(project_id)}


@app.delete("/coder/file/{project_id}")
async def coder_delete_file(project_id: str, path: str) -> dict[str, Any]:
    from shared import coder
    target = coder._safe_file(project_id, path)
    if target.exists():
        target.unlink()
        return {"ok": True, "deleted": path}
    return {"ok": False, "error": "file not found"}


@app.post("/coder/mkdir/{project_id}")
async def coder_mkdir(project_id: str, path: str) -> dict[str, Any]:
    from shared import coder
    target = coder._safe_file(project_id, path)
    target.mkdir(parents=True, exist_ok=True)
    return {"ok": True, "created": path}


@app.get("/coder/settings")
async def coder_settings_get() -> dict[str, Any]:
    from shared import coder
    return {"ok": True, **coder.get_coder_settings()}


@app.post("/coder/settings")
async def coder_settings_post(req: CoderSettingsRequest) -> dict[str, Any]:
    from shared import coder
    data: dict[str, Any] = {}
    # Explicit null/empty resets to default
    if req.base_dir is None:
        data["base_dir"] = None
    else:
        data["base_dir"] = req.base_dir
    updated = coder.set_coder_settings(data)
    return {"ok": True, **updated}


@app.post("/coder/ai-support")
async def coder_ai_support(req: CoderAiSupportRequest) -> dict[str, Any]:
    """Dedicated AI coding endpoint for non-coders. Generates / refactors / fixes / explains code."""
    from shared import coder, ai as _ai
    model = coder._prefer_coder_model(req.provider, req.model)
    mode = req.mode.lower()

    if mode == "generate" and not req.project_id:
        # Allow quick code generation without an existing project
        system = (
            "You are SHIMS AI Coder. The user cannot code. Write COMPLETE, runnable code. "
            "Respond with STRICT JSON only:\n"
            '{"explanation": "...", "files": {"main.py": "FULL FILE CONTENT", ...}, "requirements": ["dep>=1.0", ...]}\n'
            "Prefer standard library. Add comments."
        )
        prompt = f"USER REQUEST:\n{req.instruction}\n\nWrite the full code."
        result = await _ai.ask_ai(prompt, system=system, provider=req.provider, model=model)
        spec = coder._parse_spec(result.text)
        return {"ok": True, "mode": mode, "explanation": spec.get("explanation", ""),
                "files": spec.get("files", {}), "llm_provider": result.provider, "llm_model": model}

    if not req.project_id:
        return {"ok": False, "error": "project_id required for refactor / fix / explain modes"}

    meta = coder._load_meta(req.project_id)
    if not meta:
        return {"ok": False, "error": "project not found"}

    ctx = coder._project_context(req.project_id, max_chars=12000)

    if mode == "generate":
        system = (
            "You are SHIMS AI Coder. The user cannot code. Write COMPLETE, runnable code. "
            "Respond with STRICT JSON only:\n"
            '{"explanation": "...", "files": {"path/name.py": "FULL FILE CONTENT", ...}}\n'
            "Always return COMPLETE file contents (not diffs)."
        )
        prompt = (
            f"PROJECT GOAL:\n{meta.get('goal') or meta.get('name')}\n\n"
            f"CURRENT FILES:\n{ctx}\n\n"
            f"USER REQUEST (the user cannot code):\n{req.instruction}\n\n"
            "Write the full code changes as JSON."
        )
    elif mode == "refactor":
        system = (
            "You are SHIMS AI Coder. Refactor the code to be cleaner, faster, or more maintainable. "
            "Respond with STRICT JSON only:\n"
            '{"explanation": "...", "files": {"path/name.py": "FULL FILE CONTENT", ...}}\n'
            "Always return COMPLETE file contents (not diffs)."
        )
        prompt = (
            f"PROJECT GOAL:\n{meta.get('goal') or meta.get('name')}\n\n"
            f"CURRENT FILES:\n{ctx}\n\n"
            f"REFACTOR REQUEST:\n{req.instruction}\n\n"
            "Return the complete refactored files as JSON."
        )
    elif mode in ("fix", "debug"):
        system = (
            "You are SHIMS AI Coder. Find and fix bugs. Respond with STRICT JSON only:\n"
            '{"explanation": "...", "files": {"path/name.py": "FULL FILE CONTENT", ...}}\n'
            "Always return COMPLETE file contents (not diffs)."
        )
        prompt = (
            f"PROJECT GOAL:\n{meta.get('goal') or meta.get('name')}\n\n"
            f"CURRENT FILES:\n{ctx}\n\n"
            f"BUG REPORT:\n{req.instruction}\n\n"
            "Return the fixed complete files as JSON."
        )
    elif mode == "explain":
        system = (
            "You are SHIMS AI Coder. Explain code to a non-technical user in plain English. "
            "Be concise but thorough. Use analogies where helpful."
        )
        prompt = (
            f"PROJECT GOAL:\n{meta.get('goal') or meta.get('name')}\n\n"
            f"CURRENT FILES:\n{ctx}\n\n"
            f"USER QUESTION:\n{req.instruction}\n\n"
            "Explain in plain English."
        )
    else:
        return {"ok": False, "error": f"unknown mode: {mode}"}

    result = await _ai.ask_ai(prompt, system=system, provider=req.provider, model=model)

    if mode == "explain":
        return {"ok": True, "mode": mode, "explanation": result.text,
                "llm_provider": result.provider, "llm_model": model}

    spec = coder._parse_spec(result.text)
    files = spec.get("files") or {}
    changed = []
    for path, content in files.items():
        if isinstance(content, str):
            try:
                coder.write_file(req.project_id, path, content)
                changed.append(path)
            except ValueError:
                continue
    if spec.get("run"):
        meta["entry"] = str(spec["run"])
        coder._save_meta(meta)
    return {"ok": True, "mode": mode, "explanation": spec.get("explanation", ""),
            "files_changed": changed, "llm_provider": result.provider,
            "llm_model": model, "project_id": req.project_id}


@app.get("/capture/status")
async def capture_status() -> dict[str, Any]:
    return mailbox_status()


@app.post("/capture/share")
async def capture_share(req: CaptureShareRequest) -> dict[str, Any]:
    if not (req.title or req.text or req.url):
        raise HTTPException(400, "Capture needs a title, text, or URL.")
    result = save_capture(title=req.title, text=req.text, url=req.url, kind=req.kind, source=req.source, metadata=req.metadata)
    evidence = [{"kind": "capture", "title": result.get("item", {}).get("title") or req.title or req.url, "source_uri": result.get("item", {}).get("url") or result.get("item", {}).get("id") or "", "excerpt": result.get("item", {}).get("text") or "", "score": 0.82, "metadata": {"source": req.source, "kind": req.kind}}]
    action = record_action("capture_save", f"Save capture: {req.title or req.url}"[:220], payload=_model_data(req), result={"ok": result.get("ok"), "item_id": result.get("item", {}).get("id")}, evidence=evidence, requested_level="L3", summary="Saved user-provided capture into mailbox and brain memory.")
    trust = build_trust(route="capture:share", evidence=merge_evidence(evidence, evidence_from_action(action.get("action"))), action_id=action.get("action_id", ""), ledger_hash=action.get("ledger_hash", ""))
    result.update(_trust_fields(trust))
    return result


@app.get("/capture/items")
async def capture_items(limit: int = 50, status: str | None = None) -> dict[str, Any]:
    return {"ok": True, "items": list_captures(limit=limit, status=status)}


@app.get("/mailbox/status")
async def mailbox_status_endpoint() -> dict[str, Any]:
    return mailbox_status()


@app.get("/mailbox/policy")
async def mailbox_policy_endpoint() -> dict[str, Any]:
    return mailbox_policy()


@app.get("/mailbox/oauth/start")
async def mailbox_oauth_start(state: str | None = None) -> dict[str, Any]:
    return gmail_auth_url(state=state)


@app.get("/mailbox/oauth/callback", response_class=HTMLResponse)
async def mailbox_oauth_callback(code: str = "", state: str = "", error: str = "") -> HTMLResponse:
    if error:
        return HTMLResponse(f"<h1>SHIMS Gmail authorization failed</h1><p>{error}</p>", status_code=400)
    result = await asyncio.to_thread(exchange_code_for_token, code)
    if result.get("ok"):
        scope = result.get("scope", "")
        refresh = "yes" if result.get("has_refresh_token") else "no"
        return HTMLResponse(
            "<h1>SHIMS Gmail connected ✓</h1>"
            "<p>Tokens stored locally. You can close this tab and return to SHIMS.</p>"
            f"<p><b>Granted scope:</b> {scope}</p>"
            f"<p><b>Refresh token saved:</b> {refresh}</p>"
        )
    return HTMLResponse(
        "<h1>SHIMS Gmail authorization incomplete</h1>"
        f"<p>{result.get('message', 'Token exchange failed.')}</p>"
        f"<p><b>Status:</b> {result.get('status', 'error')}</p>",
        status_code=400 if result.get("status") in {"missing_code", "token_error"} else 200,
    )


@app.post("/mailbox/import")
async def mailbox_import(req: MailboxImportRequest) -> dict[str, Any]:
    if not (req.subject or req.snippet or req.body):
        raise HTTPException(400, "Mailbox import needs a subject, snippet, or body.")
    result = save_mail_message(
        provider=req.provider,
        external_id=req.external_id,
        thread_id=req.thread_id,
        sender=req.sender,
        recipients=req.recipients,
        subject=req.subject,
        snippet=req.snippet,
        body=req.body,
        labels=req.labels,
        received_at=req.received_at,
        source_url=req.source_url,
        metadata=req.metadata,
    )
    evidence = [{"kind": "mailbox", "title": result.get("message", {}).get("subject") or req.subject, "source_uri": result.get("message", {}).get("source_url") or result.get("message", {}).get("id") or "", "excerpt": result.get("message", {}).get("snippet") or req.snippet or req.body, "score": 0.82, "metadata": {"provider": req.provider, "sender": req.sender}}]
    action = record_action("mailbox_import", f"Import mailbox item: {req.subject or 'mail item'}"[:220], payload=_model_data(req), result={"ok": result.get("ok"), "message_id": result.get("message", {}).get("id")}, evidence=evidence, requested_level="L3", summary="Imported user-provided mailbox item into local mailbox and brain memory.")
    trust = build_trust(route="mailbox:import", evidence=merge_evidence(evidence, evidence_from_action(action.get("action"))), action_id=action.get("action_id", ""), ledger_hash=action.get("ledger_hash", ""))
    result.update(_trust_fields(trust))
    return result


@app.get("/mailbox/messages")
async def mailbox_messages(limit: int = 50, provider: str | None = None) -> dict[str, Any]:
    return {"ok": True, "messages": list_mail_messages(limit=limit, provider=provider)}


@app.get("/mailbox/digest")
async def mailbox_digest_endpoint(limit: int = 20) -> dict[str, Any]:
    return mailbox_digest(limit=limit)


@app.post("/mailbox/gmail/sync")
async def mailbox_gmail_sync(req: GmailSyncRequest) -> dict[str, Any]:
    result = await asyncio.to_thread(sync_gmail_metadata, access_token=req.access_token, query=req.query, max_results=req.max_results)
    if not result.get("ok") and result.get("status") == "needs_oauth":
        return JSONResponse(status_code=428, content=result)
    if result.get("ok"):
        evidence = [{"kind": "mailbox", "title": f"Gmail metadata sync ({result.get('stored', 0)} messages)", "source_uri": "gmail:metadata", "excerpt": f"Stored {result.get('stored', 0)} Gmail metadata records.", "score": 0.8, "metadata": {"query": req.query, "provider": "gmail"}}]
        action = record_action("mailbox_import", "Sync Gmail metadata", payload={"query": req.query, "max_results": req.max_results, "scope": "gmail.metadata"}, result={"ok": True, "stored": result.get("stored")}, evidence=evidence, requested_level="L3", summary="Synced Gmail metadata after explicit OAuth token availability.")
        trust = build_trust(route="mailbox:gmail_sync", evidence=merge_evidence(evidence, evidence_from_action(action.get("action"))), action_id=action.get("action_id", ""), ledger_hash=action.get("ledger_hash", ""))
        result.update(_trust_fields(trust))
    return result


@app.post("/mailbox/gmail/send")
async def mailbox_gmail_send(req: GmailSendRequest) -> dict[str, Any]:
    result = await asyncio.to_thread(
        send_gmail_message, req.to, req.subject, req.body,
        cc=req.cc, thread_id=req.thread_id, in_reply_to=req.in_reply_to,
    )
    if not result.get("ok"):
        status_code = 428 if result.get("status") in {"needs_oauth", "scope_required"} else 400
        return JSONResponse(status_code=status_code, content=result)
    evidence = [{"kind": "mailbox", "title": f"Sent mail: {req.subject or '(no subject)'}", "source_uri": f"gmail:{result.get('id', '')}", "excerpt": req.body[:240], "score": 0.85, "metadata": {"to": req.to, "thread_id": result.get("thread_id")}}]
    action = record_action("gmail_send", f"Send email to {req.to}"[:220], payload={"to": req.to, "subject": req.subject}, result={"ok": True, "id": result.get("id")}, evidence=evidence, requested_level="L3", summary="Sent an email through the user-authorized Gmail account.")
    trust = build_trust(route="mailbox:gmail_send", evidence=merge_evidence(evidence, evidence_from_action(action.get("action"))), action_id=action.get("action_id", ""), ledger_hash=action.get("ledger_hash", ""))
    result.update(_trust_fields(trust))
    return result


@app.post("/mailbox/gmail/reply")
async def mailbox_gmail_reply(req: GmailReplyRequest) -> dict[str, Any]:
    result = await asyncio.to_thread(reply_to_gmail_message, req.message_id, req.body)
    if not result.get("ok"):
        status_code = 428 if result.get("status") in {"needs_oauth", "scope_required"} else (404 if result.get("status") == "not_found" else 400)
        return JSONResponse(status_code=status_code, content=result)
    evidence = [{"kind": "mailbox", "title": "Sent Gmail reply", "source_uri": f"gmail:{result.get('id', '')}", "excerpt": req.body[:240], "score": 0.85, "metadata": {"thread_id": result.get("thread_id"), "in_reply_to": req.message_id}}]
    action = record_action("gmail_reply", f"Reply to message {req.message_id}"[:220], payload={"message_id": req.message_id}, result={"ok": True, "id": result.get("id")}, evidence=evidence, requested_level="L3", summary="Replied to a synced mailbox message through the user-authorized Gmail account.")
    trust = build_trust(route="mailbox:gmail_reply", evidence=merge_evidence(evidence, evidence_from_action(action.get("action"))), action_id=action.get("action_id", ""), ledger_hash=action.get("ledger_hash", ""))
    result.update(_trust_fields(trust))
    return result


@app.get("/operator/digest")
async def operator_digest(limit: int = 20, record: bool = False) -> dict[str, Any]:
    return build_operator_digest(limit=limit, record=record)


@app.get("/actions")
async def actions_list(limit: int = 50, status: str | None = None, action_type: str | None = None) -> dict[str, Any]:
    return {"ok": True, "status": action_status(), "actions": list_actions(limit=limit, status=status, action_type=action_type)}


@app.post("/actions")
async def actions_record(req: ActionRecordRequest) -> dict[str, Any]:
    action = record_action(
        req.action_type,
        req.title or req.action_type,
        payload=req.payload,
        result=req.result,
        evidence=req.evidence,
        status=req.status,
        requested_level=req.requested_level,
        summary=req.summary,
    )
    trust = build_trust(
        route="actions:record",
        evidence=evidence_from_action(action.get("action")),
        missing_evidence=[] if req.evidence else ["Action was recorded without external evidence."],
        action_id=action.get("action_id", ""),
        ledger_hash=action.get("ledger_hash", ""),
    )
    return {**action, **_trust_fields(trust)}


@app.get("/actions/{action_id}")
async def actions_get(action_id: str) -> dict[str, Any]:
    action = get_action(action_id)
    if not action:
        raise HTTPException(404, "Action not found")
    trust = build_trust(route="actions:get", evidence=evidence_from_action(action), action_id=action.get("id", ""), ledger_hash=action.get("record_hash", ""))
    return {"ok": True, "action": action, **_trust_fields(trust)}


@app.get("/actions/{action_id}/verify")
async def actions_verify(action_id: str) -> dict[str, Any]:
    result = verify_action(action_id)
    if not result.get("ok") and result.get("reason") == "not_found":
        raise HTTPException(404, "Action not found")
    evidence = [{"kind": "action", "title": f"Action verification {action_id}", "source_uri": action_id, "excerpt": result.get("status") or result.get("reason") or "", "score": 0.95 if result.get("ok") else 0.35, "metadata": result}]
    trust = build_trust(route="actions:verify", evidence=evidence, action_id=action_id, ledger_hash=result.get("ledger_hash", ""))
    return {**result, **_trust_fields(trust)}


@app.get("/actions/pending")
@app.get("/approvals/pending")
async def actions_pending(session_id: str | None = None, include_resolved: bool = False, limit: int = 30) -> dict[str, Any]:
    return {"ok": True, "pending": _list_pending_actions(session_id=session_id, include_resolved=include_resolved, limit=limit)}


@app.post("/actions/approve")
@app.post("/approvals/decide")
async def actions_approve(req: ApprovalDecisionRequest) -> dict[str, Any]:
    pending = _load_pending_action(req.approval_id or "") if req.approval_id else _latest_pending_action()
    if not pending:
        return {"ok": False, "status": "not_found", "message": "No pending approval was found."}
    if not _decision_value(req.decision):
        pending["status"] = "cancelled"
        pending["decision"] = "no"
        pending["note"] = req.note or ""
        pending["resolved_at"] = _utc_now()
        _save_pending_action(pending)
        return {"ok": True, "status": "cancelled", "approval": _public_pending_action(pending)}
    approved_by = req.approved_by or "human-operator"
    result = await _execute_pending_action(pending, approved_by=approved_by)
    pending["status"] = "completed" if result.get("ok") else "failed"
    pending["decision"] = "yes"
    pending["approved_by"] = approved_by
    pending["note"] = req.note or ""
    pending["result"] = result
    pending["resolved_at"] = _utc_now()
    _save_pending_action(pending)
    return {"ok": bool(result.get("ok")), "status": pending["status"], "approval": _public_pending_action(pending), "result": result}


@app.get("/coder/playground/status")
async def coder_playground_status() -> dict[str, Any]:
    apps: list[dict[str, Any]] = []
    for app_index in sorted(GENERATED_APPS_DIR.glob("*/index.html"), key=lambda p: p.stat().st_mtime, reverse=True)[:40]:
        slug = app_index.parent.name
        apps.append({
            "name": slug,
            "relative_path": str(app_index.relative_to(ROOT)).replace("\\", "/"),
            "url": f"/generated-apps/{slug}/index.html",
            "updated_at": datetime.fromtimestamp(app_index.stat().st_mtime).isoformat(timespec="seconds"),
        })
    return {
        "ok": True,
        "roots": sorted(self_evolver.ALLOWED_ROOTS),
        "generated_apps_dir": str(GENERATED_APPS_DIR),
        "apps": apps,
        "pending": _list_pending_actions(limit=12),
        "proposals": list_proposals(limit=12),
    }


@app.post("/coder/playground/propose")
async def coder_playground_propose(req: CoderProposalRequest) -> dict[str, Any]:
    result = _create_coder_proposal(req)
    log_event("coder.propose", route="coder:playground", provider="local", model="self-evolver", ok=bool(result.get("ok")), message=req.relative_path, metadata=result)
    return result


@app.post("/coder/playground/app")
async def coder_playground_app(req: CoderAppRequest) -> dict[str, Any]:
    result = _create_or_propose_coder_app(req)
    log_event("coder.app", route="coder:playground", provider="local", model="self-evolver", ok=bool(result.get("ok")), message=req.name, metadata=result)
    return result


@app.post("/campaigns/plan")
async def campaigns_plan(req: CampaignPlanRequest) -> dict[str, Any]:
    plan = plan_campaign(objective=req.objective, audience=req.audience, offer=req.offer, channels=req.channels, tone=req.tone, due_date=req.due_date)
    evidence = [{"kind": "draft", "title": "Campaign draft plan", "source_uri": "campaign:local", "excerpt": plan.get("brief", {}).get("positioning", ""), "score": 0.64, "metadata": {"channels": plan.get("channels"), "mode": plan.get("mode")}}]
    action = record_action("campaign_draft", f"Draft campaign: {req.objective}"[:220], payload=_model_data(req), result={"ok": True, "channels": plan.get("channels"), "tasks": len(plan.get("tasks") or [])}, evidence=evidence, requested_level="L3", summary="Drafted a local campaign plan. External sends/posts require approval.")
    trust = build_trust(route="campaigns:plan", evidence=merge_evidence(evidence, evidence_from_action(action.get("action"))), action_id=action.get("action_id", ""), ledger_hash=action.get("ledger_hash", ""), requested_level="draft")
    plan.update({"action": action.get("action"), **_trust_fields(trust)})
    return plan


@app.post("/calendar/ics")
async def calendar_ics(req: CalendarIcsRequest) -> dict[str, Any]:
    result = save_ics_event(title=req.title, start=req.start, end=req.end, duration_minutes=req.duration_minutes, description=req.description, location=req.location)
    result = _attach_ledger(result, Path(result["path"]), "calendar_ics")
    evidence = evidence_from_artifact({**result, "type": "calendar", "kind": "ics"})
    action = record_action("calendar_ics_create", f"Create ICS: {req.title}"[:220], payload=_model_data(req), result={k: result.get(k) for k in ("ok", "uid", "title", "start", "end", "url", "sha256", "verified")}, evidence=evidence, requested_level="L3", summary="Created a local ICS calendar draft; no Google Calendar sync was performed.")
    trust = build_trust(route="calendar:ics", evidence=merge_evidence(evidence, evidence_from_action(action.get("action"))), action_id=action.get("action_id", ""), ledger_hash=action.get("ledger_hash", ""), requested_level="draft")
    result.update({"action": action.get("action"), **_trust_fields(trust)})
    return result


# ── Chemistry brain (shims_chem) — exposed to Omni + Android ─────────────────
@app.get("/api/rd/chem/tools")
@app.get("/chem/tools")
async def chem_tools() -> dict[str, Any]:
    from shared import shims_chem_api
    return {"ok": True, "tools": shims_chem_api.get_tool_schemas()}


@app.post("/api/rd/chem/tools/{name}")
@app.post("/chem/tools/{name}")
async def chem_run_tool(name: str, req: ChemToolRequest) -> dict[str, Any]:
    from shared import shims_chem_api
    try:
        return {"ok": True, "result": await asyncio.to_thread(shims_chem_api.run_verifier_tool, name, **req.args)}
    except Exception as exc:
        return {"ok": False, "error": f"chem tool '{name}' failed: {exc}"}


@app.post("/api/rd/chem/verify")
@app.post("/chem/verify")
async def chem_verify(req: ChemVerifyRequest, request: Request) -> dict[str, Any]:
    from shared import shims_chem_api
    smiles = req.smiles.strip()
    if not smiles and request.url.path == "/chem/verify":
        raise HTTPException(400, "smiles is required")
    return {
        "ok": True,
        "smiles": await asyncio.to_thread(shims_chem_api.verify_smiles, smiles),
        "hazards": await asyncio.to_thread(shims_chem_api.verify_hazards, smiles),
    }


@app.post("/api/rd/chem/reaction")
@app.post("/chem/reaction")
async def chem_reaction(req: ChemReactionRequest) -> dict[str, Any]:
    from shared import shims_chem_api
    rxn = (req.reaction or req.rxn_smiles or "").strip()
    if not rxn:
        raise HTTPException(400, "reaction is required")
    return {"ok": True, "result": await asyncio.to_thread(shims_chem_api.verify_reaction, rxn)}


@app.post("/api/rd/chem/retro")
@app.post("/chem/retro")
async def chem_retro(req: ChemRetroRequest) -> dict[str, Any]:
    from shared import shims_chem_api
    target = (req.target or req.target_smiles or "").strip()
    if not target:
        raise HTTPException(400, "target is required")
    max_routes = max(1, min(int(req.max_routes or 5), 10))
    routes = await asyncio.to_thread(shims_chem_api.plan_retro, target, max_routes)
    return {"ok": True, "target": target, "routes": routes}


@app.post("/api/rd/chem/ich")
@app.post("/chem/ich")
async def chem_ich(req: ChemIchRequest) -> dict[str, Any]:
    from shared import shims_chem_api
    result = await asyncio.to_thread(
        shims_chem_api.verify_ich, req.impurity_pct,
        max_daily_dose_g=req.max_daily_dose_g, impurity_name=req.impurity_name,
    )
    return {"ok": True, "result": result}


@app.post("/chem/chemdfm/query")
@app.post("/api/rd/chem/chemdfm/query")
async def chemdfm_query_endpoint(request: Request) -> dict[str, Any]:
    body = await request.json()
    from shared.chemdfm_bridge import chemdfm_query
    return await chemdfm_query(body.get("query", ""), body.get("topic", "general"))


@app.post("/chem/chemdfm/train")
@app.post("/api/rd/chem/chemdfm/train")
async def chemdfm_train_endpoint(request: Request) -> dict[str, Any]:
    body = await request.json()
    from shared.chemdfm_bridge import chemdfm_train
    return chemdfm_train(body.get("fact", ""), body.get("topic", "general"), body.get("validated_by", "human"))


@app.get("/chem/chemdfm/journal")
@app.get("/api/rd/chem/chemdfm/journal")
async def chemdfm_journal() -> dict[str, Any]:
    from shared.chemdfm_bridge import get_journal_summary
    return get_journal_summary()


@app.post("/evals/run")
async def evals_run() -> dict[str, Any]:
    result = run_reliability_evals()
    evidence = [{"kind": "eval", "title": "Reliability eval harness", "source_uri": "evals:run", "excerpt": f"{result.get('passed')}/{result.get('total')} checks passed", "score": 0.9 if result.get("ok") else 0.45, "metadata": {"version": result.get("version")}}]
    action = record_action("action_ledger_record", "Run reliability eval harness", payload={}, result={"ok": result.get("ok"), "passed": result.get("passed"), "total": result.get("total")}, evidence=evidence, requested_level="L3", status="completed" if result.get("ok") else "failed", summary="Ran local reliability eval harness.")
    trust = build_trust(route="evals:run", evidence=merge_evidence(evidence, evidence_from_action(action.get("action"))), action_id=action.get("action_id", ""), ledger_hash=action.get("ledger_hash", ""))
    result.update({"action": action.get("action"), **_trust_fields(trust)})
    return result

@app.get("/web/settings")
async def get_web_settings() -> dict[str, Any]:
    web = _settings.get("web", {})
    return {"ok": True, "settings": {"searxng_url": web.get("searxng_url"), "duckduckgo_fallback": web.get("duckduckgo_fallback"), "tavily_key": _mask_secret(web.get("tavily_key")), "brave_key": _mask_secret(web.get("brave_key")), "serpapi_key": _mask_secret(web.get("serpapi_key"))}}

@app.post("/web/settings")
async def set_web_settings(req: WebSettingsRequest) -> dict[str, Any]:
    web = _settings.setdefault("web", {})
    if req.searxng_url is not None:
        web["searxng_url"] = req.searxng_url.strip().rstrip("/")
        _set_env_persistent("SHIMS_SEARXNG_URL", web["searxng_url"])
    if req.tavily_key is not None:
        web["tavily_key"] = _clean_secret(req.tavily_key)
        _set_env_persistent("TAVILY_API_KEY", web["tavily_key"])
    if req.brave_key is not None:
        web["brave_key"] = _clean_secret(req.brave_key)
        _set_env_persistent("BRAVE_SEARCH_API_KEY", web["brave_key"])
    if req.serpapi_key is not None:
        web["serpapi_key"] = _clean_secret(req.serpapi_key)
        _set_env_persistent("SERPAPI_API_KEY", web["serpapi_key"])
    if req.duckduckgo_fallback is not None:
        web["duckduckgo_fallback"] = bool(req.duckduckgo_fallback)
        _set_env_persistent("SHIMS_DUCKDUCKGO_FALLBACK", "true" if web["duckduckgo_fallback"] else "false")
    return await get_web_settings()

@app.get("/realtime/status")
async def realtime_status() -> dict[str, Any]:
    from shared.realtime_kernel import manifest as realtime_manifest
    m = realtime_manifest()
    return {"ok": True, **m, "transport": {"browser": "web-speech+server-stt", "recommended": "Pipecat WebRTC / websocket transport when pipecat-ai is installed"}, "features": {"half_duplex": True, "barge_in": True, "duplicate_guard": True, "tool_first": True, "streaming_tokens": True}}

@app.get("/agents/list")
async def agents_list() -> dict[str, Any]:
    agents = agent_registry.list_agents()
    return {"ok": True, "agents": [{"id": a["id"], "name": a["name"], "status": a.get("status", "ready"), "role": a["purpose"], "tools": a.get("tools", []), "approval_level": a.get("approval_level", "normal")} for a in agents]}

@app.get("/voice/config")
async def voice_config() -> dict[str, Any]:
    cfg = dict(_settings["voice"])
    cfg["voice_mode"] = settings.voice_mode
    cfg["server_stt_chunk_ms"] = 900
    cfg["stt_correction_enabled"] = False
    return {"ok": True, "config": cfg}

@app.post("/voice/config")
async def set_voice_config(req: VoiceConfigRequest) -> dict[str, Any]:
    cfg = _settings["voice"]
    if req.wake_words:
        cfg["wake_words"] = req.wake_words
    cfg["primary_lang"] = req.primary_lang or "en-IN"
    cfg["secondary_langs"] = req.secondary_langs or ["hi-IN", "en-US"]
    cfg["command_cooldown_seconds"] = max(0.8, float(req.command_cooldown_seconds))
    cfg["silence_timeout_seconds"] = max(0.8, float(req.silence_timeout_seconds))
    cfg["max_auto_replies_without_user"] = max(1, int(req.max_auto_replies_without_user))
    return {"ok": True, "config": cfg}


def _server_stt_available() -> bool:
    try:
        import faster_whisper  # noqa: F401
        return True
    except Exception:
        return False


async def _preload_voice_model() -> None:
    """Warm the faster-whisper model cache during startup so first transcription is fast."""
    if os.environ.get("PYTEST_CURRENT_TEST") or os.getenv("SHIMS_DISABLE_WHISPER", "").lower() in {"1", "true", "yes"}:
        return
    if not _server_stt_available():
        return
    try:
        await asyncio.to_thread(_get_whisper_model)
        log_event("voice.warmup_complete", route="voice:startup", provider="local", model=_active_whisper_model(), ok=True)
    except Exception as exc:
        log_event("voice.warmup_error", route="voice:startup", provider="local", model=_active_whisper_model(), ok=False, message=str(exc)[:200])


_WHISPER_MODEL = None          # cached faster-whisper model instance (load is expensive)
_WHISPER_MODEL_KEY = None       # (model_path, device, compute) tuple the cached instance was built with
_ACTIVE_WHISPER_MODEL = None    # runtime override (in-memory)
_STT_MODEL_STATE = STATE_DIR / "stt_model.json"
_STT_MODELS_DIR = ROOT / "storage" / "models"

# In-flight STT correction tasks keyed by correction_id (usually the session_id).
_STT_CORRECTION_TASKS: dict[str, asyncio.Task] = {}
_STT_CORRECTION_RESULTS: dict[str, dict] = {}


def _active_whisper_model() -> str:
    """Resolve the active Whisper model: runtime override -> saved state -> env -> 'small'."""
    global _ACTIVE_WHISPER_MODEL
    if _ACTIVE_WHISPER_MODEL:
        return _ACTIVE_WHISPER_MODEL
    try:
        if _STT_MODEL_STATE.exists():
            saved = json.loads(_STT_MODEL_STATE.read_text(encoding="utf-8")).get("model")
            if saved:
                _ACTIVE_WHISPER_MODEL = saved
                return saved
    except Exception:
        pass
    return os.getenv("SHIMS_WHISPER_MODEL", "small")


def _list_stt_models() -> list[dict[str, Any]]:
    """Discover locally-available Whisper models under storage/models."""
    models: list[dict[str, Any]] = []
    if _STT_MODELS_DIR.exists():
        for d in sorted(_STT_MODELS_DIR.glob("faster-whisper-*")):
            if (d / "model.bin").exists():
                label = d.name.replace("faster-whisper-", "")
                models.append({"id": str(d), "label": label, "ready": True})
    return models


def _whisper_cached() -> bool:
    """True if the active Whisper model is available locally (no download needed)."""
    if not _server_stt_available():
        return False
    # Instantiating the native WhisperModel can hard-crash (access violation)
    # on some Windows setups — that is uncatchable and kills the process.
    # Never do it inside the test runner, and allow an explicit kill-switch.
    if os.environ.get("PYTEST_CURRENT_TEST") or os.getenv("SHIMS_DISABLE_WHISPER", "").lower() in {"1", "true", "yes"}:
        return False
    try:
        from faster_whisper import WhisperModel  # type: ignore
        WhisperModel(_active_whisper_model(),
                     device=os.getenv("SHIMS_WHISPER_DEVICE", "cpu"),
                     compute_type=os.getenv("SHIMS_WHISPER_COMPUTE", "int8"),
                     local_files_only=True)
        return True
    except Exception:
        return False


def _resolve_device() -> str:
    """Return the effective device — always 'cpu' unless CUDA 12 libs are confirmed present."""
    requested = os.getenv("SHIMS_WHISPER_DEVICE", "cpu").lower()
    if requested in {"auto", "cuda"}:
        # Only use CUDA if cuBLAS 12 is actually loadable; otherwise fall back to CPU silently.
        try:
            import ctypes
            ctypes.CDLL("cublas64_12.dll")
            return "cuda"
        except Exception:
            return "cpu"
    return requested  # explicit 'cpu' or other


def _get_whisper_model():
    """Load the active faster-whisper model, reloading if model/device/compute changed.

    The cache key includes device + compute so a stale GPU instance is never reused
    after settings change or after a cublas error.
    """
    global _WHISPER_MODEL, _WHISPER_MODEL_KEY
    # Native WhisperModel init can hard-crash (access violation) on some
    # Windows setups — uncatchable, kills the process. Never under pytest.
    if os.environ.get("PYTEST_CURRENT_TEST") or os.getenv("SHIMS_DISABLE_WHISPER", "").lower() in {"1", "true", "yes"}:
        raise RuntimeError("Speech-to-text is disabled (SHIMS_DISABLE_WHISPER / test run).")
    device = _resolve_device()
    compute_type = os.getenv("SHIMS_WHISPER_COMPUTE", "int8")
    model_path = _active_whisper_model()
    cache_key = (model_path, device, compute_type)
    if _WHISPER_MODEL is None or _WHISPER_MODEL_KEY != cache_key:
        from faster_whisper import WhisperModel  # type: ignore
        try:
            _WHISPER_MODEL = WhisperModel(model_path, device=device,
                                          compute_type=compute_type, local_files_only=True)
        except Exception:
            _WHISPER_MODEL = WhisperModel(model_path, device=device, compute_type=compute_type)
        _WHISPER_MODEL_KEY = cache_key
    return _WHISPER_MODEL


def _transcribe_sync(path: str, lang: str | None) -> dict[str, Any]:
    model = _get_whisper_model()
    language = None if lang in {None, "auto", ""} else lang.split("-")[0]
    segments, info = model.transcribe(path, language=language, vad_filter=True)
    text = " ".join(seg.text.strip() for seg in segments).strip()
    return {"ok": bool(text), "text": text, "language": getattr(info, "language", lang), "engine": "faster-whisper"}


def _start_stt_correction(raw_text: str, correction_id: str, language: str = "") -> str:
    """Launch an async LLM correction task for a raw transcript.

    Returns the correction_id so callers can await/poll the result.
    """
    raw_text = (raw_text or "").strip()
    if not raw_text or not correction_id:
        return correction_id
    # Cancel any prior in-flight correction for the same id to avoid pile-up.
    old = _STT_CORRECTION_TASKS.pop(correction_id, None)
    if old and not old.done():
        old.cancel()

    async def _run() -> dict[str, Any]:
        try:
            result = await stt_corrector.correct_transcript(raw_text, language=language)
        except Exception as exc:
            result = {"ok": False, "corrected": raw_text, "changed": False, "confidence": 0.0, "explanation": str(exc)}
        _STT_CORRECTION_RESULTS[correction_id] = result
        return result

    _STT_CORRECTION_TASKS[correction_id] = asyncio.create_task(_run())
    return correction_id


async def _await_stt_correction(correction_id: str, timeout: float = 0.6) -> dict[str, Any] | None:
    """Await a correction task with a timeout. Returns None if not ready in time."""
    if not correction_id:
        return None
    task = _STT_CORRECTION_TASKS.get(correction_id)
    if not task:
        return _STT_CORRECTION_RESULTS.get(correction_id)
    if task.done():
        try:
            return task.result()
        except Exception as exc:
            return {"ok": False, "corrected": "", "changed": False, "confidence": 0.0, "explanation": str(exc)}
    try:
        return await asyncio.wait_for(task, timeout=timeout)
    except asyncio.TimeoutError:
        return None


@app.get("/stt/health")
async def stt_health() -> dict[str, Any]:
    installed = _server_stt_available()
    cached = installed and await asyncio.to_thread(_whisper_cached)
    return {
        "ok": True,
        "browser_stt": True,
        "server_stt_installed": installed,
        "server_stt_ready_offline": cached,
        "model": _active_whisper_model(),
        "engine": "faster-whisper" if installed else "browser-only",
        "note": ("Server STT ready." if cached else
                 "Server STT installed but model not downloaded yet. Run: python scripts/download_whisper_model.py"),
        "install": "install_voice_windows.bat",
    }


@app.get("/stt/models")
async def stt_models() -> dict[str, Any]:
    """List locally-available STT models and the active selection (for Settings UI)."""
    models = _list_stt_models()
    active = _active_whisper_model()
    # Surface the active model even if it lives outside storage/models.
    if active not in {m["id"] for m in models}:
        models.insert(0, {"id": active, "label": Path(active).name.replace("faster-whisper-", "") or active, "ready": _whisper_cached()})
    return {"ok": True, "active": active, "models": models}


@app.post("/stt/model")
async def set_stt_model(req: SttModelRequest) -> dict[str, Any]:
    """Switch the active STT model at runtime and persist the choice."""
    global _ACTIVE_WHISPER_MODEL
    choice = (req.model or "").strip()
    allowed = {m["id"] for m in _list_stt_models()} | {"tiny", "base", "small", "medium", "large-v3"}
    if choice not in allowed:
        raise HTTPException(400, f"unknown model '{choice}'. Available: {sorted(allowed)}")
    _ACTIVE_WHISPER_MODEL = choice
    try:
        _STT_MODEL_STATE.write_text(json.dumps({"model": choice}), encoding="utf-8")
    except Exception:
        pass
    ready = await asyncio.to_thread(_whisper_cached)
    return {"ok": True, "active": choice, "ready": ready,
            "note": "Active model updated. It loads on the next voice transcription."}


@app.post("/stt/reset")
async def stt_reset() -> dict[str, Any]:
    """Evict the cached Whisper model so it reloads cleanly on the next transcription.

    Useful after changing SHIMS_WHISPER_DEVICE in .env without restarting the server,
    or to recover from a cublas/CUDA error without a full restart.
    """
    global _WHISPER_MODEL, _WHISPER_MODEL_KEY
    _WHISPER_MODEL = None
    _WHISPER_MODEL_KEY = None
    device = _resolve_device()
    return {"ok": True, "device_will_use": device,
            "note": f"Model cache cleared. Next transcription loads on '{device}'."}


@app.post("/voice/transcribe")
async def voice_transcribe(
    file: UploadFile = File(...),
    lang: str | None = Form("auto"),
    session_id: str | None = Form(None),
    correct: bool = Form(False),
) -> dict[str, Any]:
    suffix = Path(file.filename or "audio.webm").suffix or ".webm"
    path = STT_DIR / _safe_name("voice", suffix.lstrip("."))
    path.write_bytes(await file.read())
    if not _server_stt_available():
        return {"ok": False, "text": "", "reason": "server STT not installed", "install": "install_voice_windows.bat"}
    try:
        result = await asyncio.to_thread(_transcribe_sync, str(path), lang)
    except Exception as exc:
        reason = str(exc)[:260]
        hint = ("Speech model not available offline. Run "
                "'python scripts/download_whisper_model.py' once on a network that allows "
                "huggingface.co, or set SHIMS_WHISPER_MODEL to a local model path.")
        return {"ok": False, "text": "", "reason": reason, "hint": hint}
    if correct and result.get("ok") and result.get("text"):
        correction_id = (session_id or str(uuid.uuid4())).strip()
        _start_stt_correction(result["text"], correction_id, language=result.get("language") or lang or "")
        result["correction_id"] = correction_id
        result["correction_pending"] = True
    elif result.get("ok"):
        result["correction_pending"] = False
    return result


@app.post("/voice/correct")
async def voice_correct(req: Request) -> dict[str, Any]:
    """Start an STT correction for a browser-provided transcript.

    Returns a correction_id the caller can pass to /brain/turn as voice_correction_id.
    """
    data = await req.json()
    raw = str(data.get("text") or "").strip()
    if not raw:
        raise HTTPException(400, "text is required")
    correction_id = str(data.get("session_id") or uuid.uuid4()).strip()
    _start_stt_correction(raw, correction_id, language=str(data.get("language") or ""))
    return {"ok": True, "correction_id": correction_id, "correction_pending": True}


@app.get("/voice/correction")
async def voice_correction_result(correction_id: str) -> dict[str, Any]:
    """Poll the result of an STT correction."""
    result = _STT_CORRECTION_RESULTS.get(correction_id)
    if result is None:
        task = _STT_CORRECTION_TASKS.get(correction_id)
        if task is None:
            raise HTTPException(404, "correction_id not found")
        return {"ok": True, "ready": False, "correction_pending": True}
    return {"ok": True, "ready": True, **result}



def _voice_profiles_file() -> Path:
    return VOICE_PROFILE_DIR / "profiles.json"


def _load_voice_profiles() -> dict[str, Any]:
    path = _voice_profiles_file()
    if path.exists():
        try:
            return json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            pass
    return {"selected": None, "profiles": []}


def _save_voice_profiles(data: dict[str, Any]) -> None:
    _voice_profiles_file().write_text(json.dumps(data, ensure_ascii=False, indent=2, default=str), encoding="utf-8")


@app.get("/voice/profiles")
async def voice_profiles() -> dict[str, Any]:
    data = _load_voice_profiles()
    return {"ok": True, **data, "note": "Voice profile vault stores authorized voice samples locally. OpenVoice/CosyVoice/Piper can be connected as external engines."}

@app.post("/voice/profiles/enroll")
async def voice_profile_enroll(file: UploadFile = File(...), name: str = "owner", consent_phrase: str = "I authorize SHIMS to use this voice profile") -> dict[str, Any]:
    raw = await file.read()
    if len(raw) < 1000:
        raise HTTPException(400, "Voice sample is too small")
    clean_name = re.sub(r"[^A-Za-z0-9_. -]+", "_", name).strip()[:60] or "voice"
    profile_id = hashlib.sha256((clean_name + str(time.time())).encode()).hexdigest()[:16]
    suffix = Path(file.filename or "sample.webm").suffix or ".webm"
    sample_path = VOICE_PROFILE_DIR / f"{profile_id}{suffix}"
    sample_path.write_bytes(raw)
    voiceprint = hashlib.sha256(raw).hexdigest()
    data = _load_voice_profiles()
    profile = {"id": profile_id, "name": clean_name, "sample_file": sample_path.name, "voiceprint_sha256": voiceprint, "consent_phrase": consent_phrase[:200], "created_at": datetime.now().isoformat(timespec="seconds"), "engine": os.getenv("SHIMS_VOICE_CLONE_ENGINE", "profile-vault")}
    data.setdefault("profiles", []).append(profile)
    data["selected"] = profile_id
    _save_voice_profiles(data)
    log_event("voice.profile.enrolled", route="voice:profile", provider="local", model="voiceprint", ok=True, metadata={"profile_id": profile_id, "name": clean_name})
    return {"ok": True, "profile": profile, "selected": profile_id}

@app.post("/voice/profiles/select")
async def voice_profile_select(req: VoiceProfileSelectRequest) -> dict[str, Any]:
    data = _load_voice_profiles()
    ids = {p.get("id") for p in data.get("profiles", [])}
    if req.profile_id not in ids:
        raise HTTPException(404, "Voice profile not found")
    data["selected"] = req.profile_id
    _save_voice_profiles(data)
    return {"ok": True, "selected": req.profile_id}

@app.post("/voice/profiles/delete/{profile_id}")
async def voice_profile_delete(profile_id: str) -> dict[str, Any]:
    data = _load_voice_profiles()
    keep=[]
    removed=None
    for p in data.get("profiles", []):
        if p.get("id") == profile_id:
            removed=p
            try:
                (VOICE_PROFILE_DIR / str(p.get("sample_file"))).unlink(missing_ok=True)
            except Exception:
                pass
        else:
            keep.append(p)
    data["profiles"] = keep
    if data.get("selected") == profile_id:
        data["selected"] = keep[0].get("id") if keep else None
    _save_voice_profiles(data)
    return {"ok": bool(removed), "selected": data.get("selected")}

def _synthesize_pyttsx3_file(text: str, path: Path, rate: int | None = 172, lang: str | None = "en-IN") -> dict[str, Any]:
    """Run pyttsx3 in a worker thread and return a small, explicit TTS contract."""
    try:
        import pyttsx3  # type: ignore

        engine = pyttsx3.init()
        try:
            engine.setProperty("rate", int(rate or 172))
        except Exception:
            pass
        voices = engine.getProperty("voices") or []
        chosen = None
        lang_hint = re.escape(lang or "").replace(r"\-", "[_-]")
        voice_patterns = ["India", "Hindi", r"hi[_-]IN", r"en[_-]IN", "Ravi", "Heera"]
        if lang_hint:
            voice_patterns.append(lang_hint)
        for v in voices:
            hay = f"{getattr(v, 'name', '')} {getattr(v, 'id', '')} {getattr(v, 'languages', '')}"
            if re.search("|".join(voice_patterns), hay, re.I):
                chosen = v.id
                break
        if chosen:
            try:
                engine.setProperty("voice", chosen)
            except Exception:
                pass
        engine.save_to_file(text[:1200], str(path))
        engine.runAndWait()
        try:
            engine.stop()
        except Exception:
            pass
        if path.exists() and path.stat().st_size > 1000:
            return {"ok": True, "engine": "pyttsx3", "spoken": True}
        return {"ok": False, "engine": "pyttsx3", "spoken": False, "error": "pyttsx3 produced no playable audio file"}
    except Exception as exc:
        return {"ok": False, "engine": "pyttsx3", "spoken": False, "error": str(exc)[:240]}

@app.post("/voice/speak")
async def voice_speak(req: SpeakRequest) -> dict[str, Any]:
    text = (req.text or "").strip()
    vp = _load_voice_profiles()
    selected_profile = next((p for p in vp.get("profiles", []) if p.get("id") == vp.get("selected")), None)
    if not text:
        raise HTTPException(400, "No text to speak")
    filename = _safe_name("tts", "wav")
    path = AUDIO_DIR / filename
    backend = (_settings["media"].get("audio_backend") or "auto").lower()
    voice_mode = settings.voice_mode

    # 1. Cloud / OpenAI TTS first when configured or in cloud voice mode.
    if voice_mode == "cloud" or backend in {"auto", "openai", "openai-tts", "cloud"}:
        if os.getenv("OPENAI_API_KEY"):
            try:
                cloud = await asyncio.wait_for(_create_audio(text[:1200]), timeout=12)
                if cloud.get("ok"):
                    cloud["voice_profile"] = selected_profile
                    cloud["spoken"] = True
                    return cloud
            except asyncio.TimeoutError:
                pass
            except Exception as exc:
                pass

    # 2. Local pyttsx3 TTS.
    tts_error = ""
    try:
        timeout = max(3.0, float(os.getenv("SHIMS_TTS_TIMEOUT_SECONDS", "18")))
        synth = await asyncio.wait_for(
            asyncio.to_thread(_synthesize_pyttsx3_file, text, path, req.rate, req.lang),
            timeout=timeout,
        )
        if synth.get("ok"):
            url = f"/media/files/audio/{filename}"
            return {"ok": True, "engine": synth.get("engine", "pyttsx3"), "spoken": True, "file_url": url, "url": url, "type": "audio", "voice_profile": selected_profile}
        tts_error = synth.get("error") or "pyttsx3 did not produce speech audio"
    except asyncio.TimeoutError:
        tts_error = "server TTS timed out before producing speech audio"
    except Exception as exc:
        tts_error = str(exc)[:240]

    # 3. Guaranteed audible fallback tone so frontend still has a playable file and never hangs silently.
    result = await _create_audio(text[:1200])
    result["engine"] = result.get("provider") or "tone-fallback"
    result["spoken"] = False
    result["tts_error"] = tts_error
    result["voice_profile"] = selected_profile
    return result


# ============================================================================
# Wake Word Detection
# ============================================================================

@app.get("/voice/wakeword/status")
async def wakeword_status() -> dict[str, Any]:
    detector = get_detector()
    return {"ok": True, "status": detector.status()}


@app.post("/voice/wakeword/detect")
async def wakeword_detect(file: UploadFile = File(...), transcript: str | None = None) -> dict[str, Any]:
    audio_bytes = await file.read()
    detector = get_detector()
    result = detector.detect(audio_bytes, transcript=transcript)
    if result:
        return {"ok": True, "detected": True, **result}
    # Fallback: run lightweight STT on the chunk and match text wake words
    if not result and _server_stt_available():
        try:
            import tempfile, os
            with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as tmp:
                tmp.write(audio_bytes)
                tmp_path = tmp.name
            try:
                stt_result = await asyncio.to_thread(_transcribe_sync, tmp_path, "auto")
                if stt_result and stt_result.get('text'):
                    text_lower = stt_result['text'].lower().strip()
                    for ww in detector.text_wake_words:
                        if ww in text_lower:
                            return {"ok": True, "detected": True, "label": ww, "score": 1.0, "confidence": 1.0, "backend": "stt"}
            finally:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass
        except Exception:
            pass
    return {"ok": True, "detected": False}


@app.post("/voice/wakeword/enroll")
async def wakeword_enroll(label: str, file: UploadFile = File(...)) -> dict[str, Any]:
    audio_bytes = await file.read()
    trainer = WakeWordTrainer()
    info = trainer.enroll_sample(label, audio_bytes)
    return {"ok": info.get("ok", False), **info}


@app.delete("/voice/wakeword/delete")
async def wakeword_delete(label: str) -> dict[str, Any]:
    trainer = WakeWordTrainer()
    result = trainer.delete_wake_word(label)
    return result


@app.get("/voice/wakeword/list")
async def wakeword_list() -> dict[str, Any]:
    trainer = WakeWordTrainer()
    return {"ok": True, "wake_words": trainer.list_wake_words()}


@app.post("/voice/approve")
async def voice_approve(request: Request) -> dict[str, Any]:
    """Voice-based approval endpoint. Accepts STT transcript and routes to approval logic."""
    data = await request.json()
    transcript = str(data.get("transcript", "")).strip()
    session_id = str(data.get("session_id", "")).strip()
    if not transcript:
        return {"ok": False, "error": "transcript required"}
    decision = _approval_decision_from_text(transcript)
    if decision is None:
        return {"ok": False, "error": "Could not parse approval decision from transcript", "transcript": transcript, "hints": "Say 'approve', 'yes', 'go ahead', 'cancel', or 'no'"}
    pending = _latest_pending_action(session_id or None)
    if not pending:
        return {"ok": False, "error": "No pending action to approve", "transcript": transcript, "decision": decision}
    if decision is False:
        pending["status"] = "cancelled"
        pending["decision"] = "no"
        pending["resolved_at"] = _utc_now()
        _save_pending_action(pending)
        return {"ok": True, "decision": "cancelled", "approval_id": pending.get("approval_id"), "title": pending.get("title")}
    result = await _execute_pending_action(pending, approved_by="voice-human")
    pending["status"] = "completed" if result.get("ok") else "failed"
    pending["decision"] = "yes"
    pending["approved_by"] = "voice-human"
    pending["resolved_at"] = _utc_now()
    pending["result"] = result
    _save_pending_action(pending)
    return {"ok": True, "decision": "approved", "approval_id": pending.get("approval_id"), "title": pending.get("title"), "result": result}


@app.get("/media/settings")
async def get_media_settings() -> dict[str, Any]:
    settings = dict(_settings["media"])
    settings["audio_api_key"] = _mask_secret(os.getenv("SHIMS_AUDIO_API_KEY"))
    settings["video_api_key"] = _mask_secret(os.getenv("SHIMS_VIDEO_API_KEY"))
    settings["openai_key"] = _mask_secret(os.getenv("OPENAI_API_KEY"))
    diffusers_device = "unavailable"
    diffusers_warning = ""
    if settings.get("diffusers_enabled"):
        try:
            import torch  # type: ignore
            diffusers_device = "cuda" if torch.cuda.is_available() else "cpu"
            diffusers_warning = _diffusers_cpu_guard_reason(settings.get("diffusers_model") or "", diffusers_device)
        except Exception as exc:
            diffusers_warning = str(exc)[:180]
    settings["diffusers_device"] = diffusers_device
    settings["diffusers_warning"] = diffusers_warning
    return {
        "ok": True,
        "settings": settings,
        "providers": {
            "image": ["auto", "stable-diffusion", "openai", "diffusers", "local"],
            "audio": ["auto", "openai", "generic", "local"],
            "video": ["auto", "openai", "generic", "local"],
        },
    }

@app.post("/media/settings")
async def set_media_settings(req: MediaSettingsRequest) -> dict[str, Any]:
    media = _settings["media"]
    if req.image_backend:
        media["image_backend"] = req.image_backend.strip().lower()
        _set_env_persistent("SHIMS_IMAGE_BACKEND", media["image_backend"])
    if req.audio_backend:
        media["audio_backend"] = req.audio_backend.strip().lower()
        _set_env_persistent("SHIMS_AUDIO_BACKEND", media["audio_backend"])
    if req.video_backend:
        media["video_backend"] = req.video_backend.strip().lower()
        _set_env_persistent("SHIMS_VIDEO_BACKEND", media["video_backend"])
    if req.stable_diffusion_url is not None:
        media["stable_diffusion_url"] = req.stable_diffusion_url.strip().rstrip("/")
        _set_env_persistent("STABLE_DIFFUSION_URL", media["stable_diffusion_url"])
    if req.comfyui_url is not None:
        media["comfyui_url"] = req.comfyui_url.strip().rstrip("/")
        _set_env_persistent("COMFYUI_URL", media["comfyui_url"])
    if req.diffusers_enabled is not None:
        media["diffusers_enabled"] = bool(req.diffusers_enabled)
        _set_env_persistent("SHIMS_ENABLE_DIFFUSERS", "true" if media["diffusers_enabled"] else "false")
    if req.diffusers_model:
        media["diffusers_model"] = req.diffusers_model.strip()
        _set_env_persistent("SHIMS_DIFFUSERS_MODEL", media["diffusers_model"])
    if req.openai_tts_model is not None:
        media["openai_tts_model"] = req.openai_tts_model.strip() or "gpt-4o-mini-tts"
        _set_env_persistent("OPENAI_TTS_MODEL", media["openai_tts_model"])
    if req.openai_tts_voice is not None:
        media["openai_tts_voice"] = req.openai_tts_voice.strip() or "alloy"
        _set_env_persistent("OPENAI_TTS_VOICE", media["openai_tts_voice"])
    if req.openai_video_model is not None:
        media["openai_video_model"] = req.openai_video_model.strip() or "sora-2"
        _set_env_persistent("OPENAI_VIDEO_MODEL", media["openai_video_model"])
    if req.openai_video_size is not None:
        media["openai_video_size"] = req.openai_video_size.strip() or "1280x720"
        _set_env_persistent("OPENAI_VIDEO_SIZE", media["openai_video_size"])
    if req.openai_video_seconds is not None:
        media["openai_video_seconds"] = max(1, min(int(req.openai_video_seconds), 20))
        _set_env_persistent("OPENAI_VIDEO_SECONDS", str(media["openai_video_seconds"]))
    if req.audio_api_url is not None:
        media["audio_api_url"] = req.audio_api_url.strip().rstrip("/")
        _set_env_persistent("SHIMS_AUDIO_API_URL", media["audio_api_url"])
    if req.video_api_url is not None:
        media["video_api_url"] = req.video_api_url.strip().rstrip("/")
        _set_env_persistent("SHIMS_VIDEO_API_URL", media["video_api_url"])
    if req.audio_api_key is not None:
        _set_env_persistent("SHIMS_AUDIO_API_KEY", _clean_secret(req.audio_api_key))
    if req.video_api_key is not None:
        _set_env_persistent("SHIMS_VIDEO_API_KEY", _clean_secret(req.video_api_key))
    return await get_media_settings()

@app.post("/media/generate")
async def media_generate(req: MediaRequest) -> dict[str, Any]:
    result = await _create_media(req.kind, req.prompt, req.theme, req.quality, req.provider, privacy_mode=req.privacy_mode)
    evidence = evidence_from_artifact(result)
    action = record_action(
        "artifact_generate",
        f"Generate {req.kind}: {req.prompt}"[:220],
        payload=_model_data(req),
        result={k: result.get(k) for k in ("ok", "type", "kind", "title", "filename", "url", "file_url", "sha256", "verified")},
        evidence=evidence,
        requested_level="L3",
        status="completed" if result.get("ok", True) else "failed",
        summary=f"Generated local {req.kind} artifact.",
    )
    trust = build_trust(
        route="media:generate",
        evidence=merge_evidence(evidence, evidence_from_action(action.get("action"))),
        missing_evidence=[] if result.get("verified") else ["Artifact was created without a verified document ledger hash."],
        action_id=action.get("action_id", ""),
        ledger_hash=action.get("ledger_hash", ""),
    )
    result.update(_trust_fields(trust))
    return result

@app.get("/media/library")
async def media_library() -> dict[str, Any]:
    items: list[dict[str, Any]] = []
    for kind, folder in [("image", IMAGE_DIR), ("audio", AUDIO_DIR), ("video", VIDEO_DIR), ("pdf", PDF_DIR), ("ppt", PPT_DIR), ("document", DOC_DIR)]:
        for p in sorted(folder.glob("*"), key=lambda x: x.stat().st_mtime, reverse=True)[:30]:
            rel_folder = {"image":"images","audio":"audio","video":"video","pdf":"pdf","ppt":"ppt","document":"documents"}[kind]
            items.append({"type": kind, "kind": kind, "filename": p.name, "title": p.stem, "url": f"/media/files/{rel_folder}/{p.name}", "file_url": f"/media/files/{rel_folder}/{p.name}"})
    items.sort(key=lambda it: (MEDIA_DIR / ({"image":"images","audio":"audio","video":"video","pdf":"pdf","ppt":"ppt","document":"documents"}[it["kind"]]) / it["filename"]).stat().st_mtime, reverse=True)
    return {"ok": True, "items": items[:80]}

@app.get("/documents")
async def documents() -> dict[str, Any]:
    files = [p.name for p in sorted(PDF_DIR.glob("*.pdf"), key=lambda x: x.stat().st_mtime, reverse=True)]
    return {"ok": True, "documents": files}


@app.get("/documents/profiles")
async def document_profiles() -> dict[str, Any]:
    from shared.document_engine import available_profiles
    return {"ok": True, "profiles": available_profiles()}


@app.post("/documents/rich-docx")
async def documents_rich_docx(req: RichDocxRequest) -> dict[str, Any]:
    """Generate a richly-formatted Word document from structured content blocks.

    Blocks support headings, sub-headings, bullet/numbered lists, indenting,
    per-run bold/italic/size/color, callout notes and tables, styled by a named
    profile (corporate / regulatory / modern / minimal).
    """
    from shared.document_engine import available_profiles, build_docx
    if req.profile not in available_profiles():
        raise HTTPException(400, f"unknown profile '{req.profile}'; choose from {available_profiles()}")
    if not req.blocks:
        raise HTTPException(400, "blocks are required")
    from uuid import uuid4
    slug = re.sub(r"[^A-Za-z0-9_-]+", "_", req.title.strip())[:50] or "document"
    filename = f"{slug}_{uuid4().hex[:8]}.docx"
    out = DOC_DIR / filename
    try:
        await asyncio.to_thread(
            build_docx, req.title, req.blocks,
            profile=req.profile, output_path=out,
            letterhead=req.letterhead, subtitle=req.subtitle,
        )
    except Exception as exc:
        raise HTTPException(400, f"document generation failed: {exc}")
    url = f"/media/files/documents/{filename}"
    return {"ok": True, "type": "docx", "kind": "docx", "title": req.title,
            "profile": req.profile, "filename": filename, "url": url,
            "file_url": url, "download_url": url}

@app.post("/documents/generate")
async def documents_generate(request: Request) -> dict[str, Any]:
    data = await request.json()
    kind = str(data.get("kind") or "pdf")
    title = str(data.get("title") or "SHIMS Document")
    content = data.get("content") or title
    if isinstance(content, dict):
        body = "\n".join(f"{k}: {v}" for k, v in content.items())
    else:
        body = str(content)
    if kind.lower() in {"ppt", "pptx", "powerpoint"}:
        return await _create_ppt(title + "\n" + body)
    return await _create_pdf(body, title=title)

@app.get("/memory")
async def memory(namespace: str | None = None, q: str | None = None, limit: int = 100) -> dict[str, Any]:
    return {"ok": True, "memories": brain_list_memories(namespace=namespace, query=q, limit=limit)}

@app.get("/memory/search")
async def memory_search(q: str = "", limit: int = 20) -> dict[str, Any]:
    return {"ok": True, "memories": brain_list_memories(query=q, limit=limit)}

@app.post("/memory/save")
async def memory_save(req: MemorySaveRequest) -> dict[str, Any]:
    return {"ok": True, "memory": brain_remember(req.namespace, req.key, req.value, tags=req.tags, pinned=req.pinned, weight=req.weight, source=req.source)}

@app.delete("/memory/{memory_id}")
async def memory_forget(memory_id: int) -> dict[str, Any]:
    return {"ok": brain_forget_memory(memory_id)}

_BUILTIN_SKILLS = [
    {"name": "Tool-first generation", "description": "Image/PDF/PPT/audio/video/docx tools execute before the LLM."},
    {"name": "Desktop file cowork", "description": "Organize, search, de-duplicate and summarize files in your workspace (with undo)."},
    {"name": "OCR", "description": "Extract text from images and screenshots, fully offline."},
    {"name": "Coder workspace", "description": "A separate codex: plan → write → run → fix multi-file projects."},
    {"name": "Long-term memory", "description": "SQLite-backed memory and RAG context persist across sessions."},
    {"name": "Background learning", "description": "Episodes become consolidated memories and learned skills."},
    {"name": "Mail read/reply", "description": "Read and reply to Gmail through your own OAuth consent."},
    {"name": "Enterprise bridge", "description": "Query and create enterprise records when paired."},
]


@app.get("/skills")
async def skills() -> dict[str, Any]:
    from shared import skills as skill_store
    learned = [
        {"id": s["id"], "name": s["name"], "description": s.get("summary", ""),
         "learned": True, "pinned": s.get("pinned", False), "tags": s.get("tags", [])}
        for s in skill_store.list_skills(limit=200)
    ]
    return {"ok": True, "builtin": _BUILTIN_SKILLS, "learned": learned,
            "skills": _BUILTIN_SKILLS + learned}


@app.post("/skills/save")
async def skills_save(req: SkillSaveRequest) -> dict[str, Any]:
    from shared import skills as skill_store
    return {"ok": True, "skill": skill_store.save_skill(
        req.name, req.summary, body=req.body, tags=req.tags or [], pinned=req.pinned,
        source="user", skill_id=req.skill_id)}


@app.delete("/skills/{skill_id}")
async def skills_forget(skill_id: str) -> dict[str, Any]:
    from shared import skills as skill_store
    return {"ok": skill_store.forget_skill(skill_id)}

@app.get("/tasks")
async def tasks() -> dict[str, Any]:
    return {"ok": True, "tasks": [{"id":"v16-reliability-core", "title":"Reliability Core", "status":"implemented", "diff":"trust envelopes + evidence + action ledger + operator/campaign/calendar/evals"}, {"id":"v15-omni-brain", "title":"Omni cognitive kernel", "status":"implemented", "diff":"durable memory + RAG + research capture + background learning"}, {"id":"v11-unified-brain", "title":"Unified brain/router", "status":"implemented", "diff":"single turn pipeline + provider registry + tool-first routing"}, {"id":"v11-enterprise", "title":"Enterprise restoration", "status":"implemented", "diff":"dashboards + GST/doc studio + bridge"}, *brain_list_tasks(limit=20)]}

@app.get("/sessions")
async def sessions() -> list[dict[str, Any]]:
    return [
        {
            "id": sid,
            "title": (msgs[0]["content"][:60] if msgs else "New chat"),
            "message_count": len(msgs),
            "updated_index": idx,
        }
        for idx, (sid, msgs) in enumerate(reversed(list(_sessions.items())))
    ]

@app.post("/sessions/new")
async def sessions_new() -> dict[str, Any]:
    session_id = str(uuid.uuid4())
    _sessions[session_id] = []
    return {"ok": True, "session_id": session_id, "title": "New chat", "messages": []}

@app.get("/sessions/{session_id}")
async def session_detail(session_id: str) -> dict[str, Any]:
    if session_id not in _sessions:
        raise HTTPException(404, "Session not found")
    messages = _sessions.get(session_id) or []
    return {
        "ok": True,
        "session_id": session_id,
        "title": (messages[0]["content"][:60] if messages else "New chat"),
        "message_count": len(messages),
        "messages": messages,
    }

@app.delete("/sessions/{session_id}")
async def session_delete(session_id: str) -> dict[str, Any]:
    existed = session_id in _sessions
    _sessions.pop(session_id, None)
    return {"ok": existed, "deleted": session_id}

@app.get("/system/providers")
async def providers() -> dict[str, Any]:
    installed = await _ollama_models_raw()
    hf_names = await _hf_names()
    hf_ready = bool(hf_names)
    return {"providers": [
        {"id":"ollama", "label":"Ollama Local", "configured": True, "status":"ready" if installed else "offline", "model": _preferred_local_model([m["name"] for m in installed])},
        {"id":"huggingface", "label":"Hugging Face Local", "configured": True, "status":"ready" if hf_ready else "offline", "model": hf_names[0] if hf_ready else DEFAULT_HUGGINGFACE_MODEL},
        *[{ "id":p, "label": p.title() if p != "anthropic" else "Anthropic / Claude", "configured": _provider_configured(p), "status":"ready" if _provider_configured(p) else "missing key", "model": PROVIDER_DEFAULTS[p]} for p in ["openai", "anthropic", "gemini", "kimi", "deepseek", "qwen"]]
    ]}

@app.post("/system/providers/{pid}/test")
async def provider_test(pid: str, request: Request) -> dict[str, Any]:
    pid = pid.lower().strip()
    try: body = await request.json()
    except Exception: body = {}
    if isinstance(body, dict):
        key = _clean_secret(body.get("api_key") or body.get(f"{pid}_api_key"))
        model = str(body.get("model") or "").strip()
        if key and pid in PROVIDER_ENV:
            _set_env_persistent(PROVIDER_ENV[pid], key)
        if model and pid in PROVIDER_DEFAULTS:
            PROVIDER_DEFAULTS[pid] = model
            _set_env_persistent(f"{pid.upper()}_MODEL", model)
    if pid == "ollama":
        names = await _ollama_names()
        return {"ok": bool(names), "reply": "Ollama online" if names else f"Ollama offline at {OLLAMA_HOST}", "models": names}
    if pid not in PROVIDER_DEFAULTS:
        raise HTTPException(400, "Unknown provider")
    if not _provider_configured(pid):
        return {"ok": False, "reply": f"missing {PROVIDER_ENV.get(pid, 'API key')}", "model": PROVIDER_DEFAULTS.get(pid, "")}
    if pid in {"openai", "gemini", "kimi", "deepseek", "qwen", "huggingface"}:
        try:
            answer, route = await _run_llm(pid, PROVIDER_DEFAULTS[pid], [{"role":"user", "content":"Reply with exactly: SHIMS key ok"}])
            ok = "SHIMS key ok" in answer or "key ok" in answer.lower()
            return {"ok": ok, "reply": answer[:260], "model": PROVIDER_DEFAULTS[pid], "configured": True, "route": route}
        except Exception as exc:
            return {"ok": False, "reply": str(exc)[:260], "model": PROVIDER_DEFAULTS[pid], "configured": True, "error": str(exc)[:260]}
    if pid == "anthropic":
        try:
            text = await _anthropic_chat(PROVIDER_DEFAULTS[pid], [{"role":"user", "content":"Reply with exactly: SHIMS key ok"}])
            ok = "SHIMS key ok" in text or "key ok" in text.lower()
            return {"ok": ok, "reply": text[:260], "model": PROVIDER_DEFAULTS[pid], "configured": True}
        except Exception as exc:
            return {"ok": False, "reply": str(exc)[:260], "model": PROVIDER_DEFAULTS[pid], "configured": True, "error": str(exc)[:260]}
    return {"ok": True, "reply": "configured", "model": PROVIDER_DEFAULTS.get(pid, "")}

@app.get("/system/provider-keys")
async def get_provider_keys() -> dict[str, Any]:
    providers_obj = {p: {"configured": _provider_configured(p), "model": PROVIDER_DEFAULTS[p], "masked": _mask_secret(os.getenv(PROVIDER_ENV.get(p, "")) if p in PROVIDER_ENV else "")} for p in PROVIDER_DEFAULTS}
    return {"ok": True, "providers": providers_obj, "provider_list": [{"id": p, **info} for p, info in providers_obj.items()]}

@app.post("/system/provider-keys")
async def set_provider_key(req: ProviderKeyRequest) -> dict[str, Any]:
    p = req.provider.lower().strip()
    if p not in PROVIDER_DEFAULTS:
        raise HTTPException(400, "Unknown provider")
    if req.model:
        PROVIDER_DEFAULTS[p] = req.model.strip()
        _set_env_persistent(f"{p.upper()}_MODEL", PROVIDER_DEFAULTS[p])
    if req.action == "clear" and p in PROVIDER_ENV:
        os.environ.pop(PROVIDER_ENV[p], None)
        _set_env_persistent(PROVIDER_ENV[p], "")
        return {"ok": True, "provider": p, "configured": False, "model": PROVIDER_DEFAULTS[p]}
    cleaned = _clean_secret(req.api_key)
    if cleaned and p in PROVIDER_ENV:
        _set_env_persistent(PROVIDER_ENV[p], cleaned)
    return {"ok": True, "provider": p, "configured": _provider_configured(p), "masked": _mask_secret(os.getenv(PROVIDER_ENV.get(p, "")) if p in PROVIDER_ENV else ""), "model": PROVIDER_DEFAULTS[p]}

@app.post("/system/settings")
async def set_settings(req: SettingsRequest) -> dict[str, Any]:
    mapping = {"gemini": req.gemini_api_key, "openai": req.openai_api_key, "anthropic": req.anthropic_api_key, "kimi": req.kimi_api_key, "deepseek": req.deepseek_api_key, "qwen": req.qwen_api_key}
    saved = []
    for p, key in mapping.items():
        cleaned = _clean_secret(key)
        if cleaned and p in PROVIDER_ENV:
            _set_env_persistent(PROVIDER_ENV[p], cleaned); saved.append(p)
    # HuggingFace endpoint settings (local OpenAI-compatible server)
    if req.huggingface_base_url is not None:
        url = str(req.huggingface_base_url).strip().rstrip("/")
        _set_env_persistent("HUGGINGFACE_BASE_URL", url)
        global HUGGINGFACE_HOST
        HUGGINGFACE_HOST = url or HUGGINGFACE_HOST
    if req.huggingface_api_key is not None:
        key = _clean_secret(req.huggingface_api_key)
        _set_env_persistent("HUGGINGFACE_API_KEY", key)
    if req.huggingface_model is not None:
        model = str(req.huggingface_model).strip()
        _set_env_persistent("HUGGINGFACE_MODEL", model)
        global DEFAULT_HUGGINGFACE_MODEL, PROVIDER_DEFAULTS
        DEFAULT_HUGGINGFACE_MODEL = model
        PROVIDER_DEFAULTS["huggingface"] = model
        saved.append("huggingface")
    return {"ok": True, "saved": saved, "providers": {p: {"configured": _provider_configured(p), "model": PROVIDER_DEFAULTS[p]} for p in PROVIDER_DEFAULTS}}

@app.post("/system/reset-local")
async def reset_local() -> dict[str, Any]:
    names = await _ollama_names()
    model = _preferred_local_model(names)
    _set_env_persistent("SHIMS_PROVIDER", "ollama")
    _set_env_persistent("SHIMS_OLLAMA_MODEL", model)
    PROVIDER_DEFAULTS["ollama"] = model
    return {"ok": True, "provider": "ollama", "model": model}

ENTERPRISE_NOT_CONFIGURED = {"ok": False, "enabled": False, "message": "Enterprise integration is not configured. Set SHIMS_ENTERPRISE_URL and SHIMS_ENTERPRISE_PAIRING_ENABLED=true to enable."}

@app.get("/enterprise/status")
async def enterprise_status() -> dict[str, Any]:
    if not ENTERPRISE_ENABLED:
        return ENTERPRISE_NOT_CONFIGURED
    try:
        async with httpx.AsyncClient(timeout=2.5) as client:
            r = await client.get(f"{ENTERPRISE_URL}/health")
            return {"ok": r.status_code < 400, "enabled": True, "enterprise": r.json() if r.headers.get("content-type", "").startswith("application/json") else r.text[:200], "url": ENTERPRISE_URL}
    except Exception as exc:
        return {"ok": False, "enabled": True, "url": ENTERPRISE_URL, "detail": str(exc)[:220]}

@app.post("/enterprise/task")
async def enterprise_task(request: Request) -> dict[str, Any]:
    if not ENTERPRISE_ENABLED:
        return ENTERPRISE_NOT_CONFIGURED
    payload = await request.json()
    try:
        async with httpx.AsyncClient(timeout=20) as client:
            r = await client.post(f"{ENTERPRISE_URL}/task", json=payload)
            r.raise_for_status()
            return r.json()
    except Exception as exc:
        return {"ok": False, "detail": str(exc)[:260]}

@app.get("/enterprise/commands")
async def enterprise_commands() -> dict[str, Any]:
    """List all bridge commands Enterprise supports."""
    if not ENTERPRISE_ENABLED:
        return ENTERPRISE_NOT_CONFIGURED
    return {
        "ok": True,
        "commands": [
            {"cmd": "summary", "desc": "Factory-wide harmonized summary"},
            {"cmd": "list_dashboard", "desc": "Department dashboard data (rd, qc, warehouse, production, procurement)"},
            {"cmd": "create_experiment", "desc": "Create R&D experiment record"},
            {"cmd": "create_procurement_request", "desc": "Create procurement request"},
            {"cmd": "harmonize", "desc": "Cross-department harmonization analysis"},
            {"cmd": "create_gst_invoice", "desc": "Generate GST e-invoice"},
            {"cmd": "create_ewaybill", "desc": "Generate e-waybill"},
            {"cmd": "create_qms_record", "desc": "Create QMS deviation/CAPA/change control"},
            {"cmd": "create_lims_sample", "desc": "Create LIMS sample record"},
            {"cmd": "create_ebr_step", "desc": "Create eBR batch step"},
            {"cmd": "create_document", "desc": "Create quotation, PO, SOP, lab notebook, etc."},
            {"cmd": "run_ai_lab", "desc": "Run AI lab process design or document formatting"},
        ],
    }

@app.post("/enterprise/command")
async def enterprise_command_proxy(request: Request) -> dict[str, Any]:
    """Proxy any command to Enterprise /api/bridge/command."""
    if not ENTERPRISE_ENABLED:
        return ENTERPRISE_NOT_CONFIGURED
    payload = await request.json()
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            headers = {"X-Bridge-Token": settings.bridge_token or ""}
            r = await client.post(f"{ENTERPRISE_URL}/api/bridge/command", json=payload, headers=headers)
            r.raise_for_status()
            return r.json()
    except Exception as exc:
        return {"ok": False, "detail": str(exc)[:260]}

@app.get("/enterprise/dashboard")
async def enterprise_dashboard() -> dict[str, Any]:
    """Fetch Enterprise dashboard summary."""
    if not ENTERPRISE_ENABLED:
        return ENTERPRISE_NOT_CONFIGURED
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            headers = {"X-Bridge-Token": settings.bridge_token or ""}
            r = await client.post(
                f"{ENTERPRISE_URL}/api/bridge/command",
                json={"command": "summary", "payload": {}},
                headers=headers,
            )
            r.raise_for_status()
            return {"ok": True, "enterprise": r.json(), "url": ENTERPRISE_URL}
    except Exception as exc:
        return {"ok": False, "url": ENTERPRISE_URL, "detail": str(exc)[:220]}

@app.get("/tasks/scheduled")
async def scheduled_tasks() -> dict[str, Any]:
    """Return upcoming scheduled tasks from the Omni brain task queue."""
    try:
        from shared.omni_brain import list_tasks
        tasks = list_tasks(limit=50)
        return {"ok": True, "tasks": tasks}
    except Exception as exc:
        return {"ok": False, "detail": str(exc)[:260]}

@app.get("/evolution/status")
async def evolution_status() -> dict[str, Any]:
    lessons = build_daily_lessons(limit=500)
    return {
        "ok": True,
        "version": APP_VERSION,
        "daily_lessons": lessons,
        "recent_events": recent_events(25),
        "proposals": list_proposals(limit=10),
        "mode": "observe -> analyze -> propose -> sandbox -> validate -> human approve -> apply -> rollback if validation fails",
        "safety": "self-evolver/security/config are immutable; patches are archived and audited",
        "capability_check_endpoint": "/evolution/capability-check",
    }


def _evolution_capability_targets(revision: str) -> list[dict[str, Any]]:
    revision = re.sub(r"[^A-Za-z0-9_.:-]+", "-", revision or "").strip("-") or datetime.now().strftime("%Y%m%d-%H%M%S")
    stamp = datetime.now().isoformat()
    return [
        {
            "id": "backend",
            "relative_path": "backend/generated_features/omni_backend_probe.py",
            "scope": "backend_feature",
            "tests": [[sys.executable, "-m", "py_compile", "backend/generated_features/omni_backend_probe.py"]],
            "new_content": (
                "from __future__ import annotations\n\n"
                f"REVISION = {revision!r}\n"
                f"GENERATED_AT = {stamp!r}\n\n"
                "def backend_probe() -> dict[str, str]:\n"
                "    return {\n"
                "        'surface': 'backend',\n"
                "        'capability': 'self-evolution can add or update backend code through the guarded pipeline',\n"
                "        'revision': REVISION,\n"
                "        'generated_at': GENERATED_AT,\n"
                "    }\n"
            ),
        },
        {
            "id": "frontend",
            "relative_path": "frontend/self_evolution_probe.js",
            "scope": "frontend_feature",
            "tests": [[sys.executable, "-c", "from pathlib import Path; p=Path('frontend/self_evolution_probe.js'); s=p.read_text(encoding='utf-8'); assert 'shimsOmniEvolutionProbe' in s and 'self-evolution' in s; print('frontend-probe-readable')"]],
            "new_content": (
                "// Generated by SHIMS Omni guarded self-evolution capability check.\n"
                "export const shimsOmniEvolutionProbe = Object.freeze({\n"
                "  surface: 'frontend',\n"
                "  capability: 'self-evolution can add or update frontend code through the guarded pipeline',\n"
                "  feature: 'self-evolution',\n"
                f"  revision: {json.dumps(revision)},\n"
                f"  generatedAt: {json.dumps(stamp)}\n"
                "});\n\n"
                "if (typeof window !== 'undefined') {\n"
                "  window.shimsOmniEvolutionProbe = shimsOmniEvolutionProbe;\n"
                "}\n"
            ),
        },
        {
            "id": "feature",
            "relative_path": "shared/generated_skills/omni_feature_probe.py",
            "scope": "generated_feature",
            "tests": [[sys.executable, "-m", "py_compile", "shared/generated_skills/omni_feature_probe.py"]],
            "new_content": (
                "from __future__ import annotations\n\n"
                f"REVISION = {revision!r}\n"
                f"GENERATED_AT = {stamp!r}\n\n"
                "def feature_probe() -> dict[str, str]:\n"
                "    return {\n"
                "        'surface': 'generated_feature',\n"
                "        'capability': 'self-evolution can create and update feature modules through the guarded pipeline',\n"
                "        'revision': REVISION,\n"
                "        'generated_at': GENERATED_AT,\n"
                "    }\n"
            ),
        },
    ]


@app.post("/evolution/capability-check")
@app.post("/api/v13/evolution/capability-check")
async def evolution_capability_check(req: EvolutionCapabilityCheckRequest) -> dict[str, Any]:
    """Exercise the real guarded patch pipeline across backend, frontend, and feature files."""
    requested = {str(x).strip().lower() for x in (req.targets or []) if str(x).strip()}
    targets = [t for t in _evolution_capability_targets(req.revision or "") if not requested or t["id"] in requested]
    if not targets:
        return {"ok": False, "status": "no_targets", "message": "No matching capability-check targets requested.", "available_targets": ["backend", "frontend", "feature"]}
    apply_requested = bool(req.apply)
    if apply_requested and (req.approval_phrase or "").strip() != "I_APPROVE_SHIMS_PATCH":
        return {
            "ok": False,
            "status": "approval_required",
            "message": "Applying source changes requires approval_phrase='I_APPROVE_SHIMS_PATCH'. Dry-run validation is available with apply=false.",
            "mode": "guarded-self-evolution",
            "targets": [t["id"] for t in targets],
        }
    approved_by = (req.approved_by or "human-operator").strip() or "human-operator"
    if approved_by.lower() in {"shims", "ai", "assistant"}:
        approved_by = "human-operator"
    results: list[dict[str, Any]] = []
    for target in targets:
        proposal = create_proposal(
            target["relative_path"],
            target["new_content"],
            reason=f"Capability check: prove SHIMS Omni can evolve {target['id']} files through the guarded pipeline.",
            author=approved_by,
            scope=target["scope"],
            tests=target["tests"],
        )
        item: dict[str, Any] = {"target": target["id"], "relative_path": target["relative_path"], "proposal": proposal}
        if not proposal.get("ok"):
            item.update({"ok": False, "status": proposal.get("status") or "proposal_failed", "message": proposal.get("message")})
            results.append(item)
            continue
        proposal_id = proposal["proposal_id"]
        validation = validate_proposal(proposal_id, validation=target["tests"])
        item["validation"] = {"ok": validation.status == "validated", "status": validation.status, "message": validation.message, **validation.details}
        if validation.status != "validated":
            item.update({"ok": False, "status": validation.status, "message": validation.message})
            results.append(item)
            continue
        if apply_requested:
            approval = approve_proposal(proposal_id, approved_by=approved_by, note="Capability check approved by explicit operator phrase.")
            item["approval"] = {"ok": approval.status == "approved", "status": approval.status, "message": approval.message, **approval.details}
            if approval.status != "approved":
                item.update({"ok": False, "status": approval.status, "message": approval.message})
                results.append(item)
                continue
            applied = apply_proposal(proposal_id, approved_by=approved_by, validation=target["tests"])
            item["apply"] = {"ok": applied.status == "applied", "status": applied.status, "message": applied.message, **applied.details}
            item.update({"ok": applied.status == "applied", "status": applied.status, "message": applied.message})
        else:
            item.update({"ok": True, "status": "validated", "message": "Dry-run proposal validated in sandbox; live files were not changed."})
        results.append(item)
    ok = all(item.get("ok") for item in results)
    log_event("evolution.capability_check", route="evolution:capability-check", provider="local", model="self-evolver", ok=ok, metadata={"apply": apply_requested, "results": results})
    return {
        "ok": ok,
        "mode": "guarded-self-evolution",
        "applied": apply_requested,
        "approval_gate": "approval_phrase='I_APPROVE_SHIMS_PATCH' required for live source changes",
        "targets": results,
    }


@app.post("/evolution/self-check")
@app.post("/api/v13/evolution/self-check")
async def evolution_self_check(req: EvolutionSelfCheckRequest) -> dict[str, Any]:
    """Inspect SHIMS code and create a real, validated patch proposal.

    This does NOT apply patches. It only proposes them for human review.
    """
    from shared.self_check import run_self_check
    return await run_self_check(
        scope=req.scope,
        relative_path=req.relative_path,
        goal=req.goal,
        test_path=req.test_path,
    )


@app.post("/evolution/reflect")
async def evolution_reflect() -> dict[str, Any]:
    lessons = build_daily_lessons(limit=1000)
    log_event("evolution.reflect", route="evolution", provider="local", model="telemetry", ok=True, metadata=lessons)
    return {"ok": True, "daily_lessons": lessons}

@app.get("/evolution/proposals")
@app.get("/api/v13/evolution/proposals")
async def evolution_proposals(limit: int = 100) -> dict[str, Any]:
    return {"ok": True, "proposals": list_proposals(limit=limit)}

@app.post("/evolution/propose")
@app.post("/api/v13/evolution/propose")
async def evolution_propose(req: EvolutionProposalRequest) -> dict[str, Any]:
    proposal = create_proposal(
        req.relative_path,
        req.new_content,
        reason=req.reason or "",
        author=req.author or "user",
        scope=req.scope or "code",
        tests=req.tests,
    )
    log_event("evolution.propose", route="evolution", provider="local", model="self_evolver", ok=bool(proposal.get("ok")), message=req.reason or "", metadata={"proposal": proposal})
    return proposal

@app.post("/evolution/validate/{proposal_id}")
@app.post("/api/v13/evolution/validate/{proposal_id}")
async def evolution_validate(proposal_id: str) -> dict[str, Any]:
    result = validate_proposal(proposal_id)
    log_event("evolution.validate", route="evolution", provider="local", model="self_evolver", ok=result.status == "validated", message=proposal_id, metadata=result.details)
    return {"ok": result.status == "validated", "status": result.status, "message": result.message, **result.details}

@app.post("/evolution/approve/{proposal_id}")
@app.post("/api/v13/evolution/approve/{proposal_id}")
async def evolution_approve(proposal_id: str, req: EvolutionApprovalRequest | None = None) -> dict[str, Any]:
    req = req or EvolutionApprovalRequest()
    result = approve_proposal(proposal_id, approved_by=req.approved_by or "human", note=req.note or "")
    log_event("evolution.approve", route="evolution", provider="local", model="self_evolver", ok=result.status == "approved", message=proposal_id, metadata=result.details)
    return {"ok": result.status == "approved", "status": result.status, "message": result.message, **result.details}

@app.post("/evolution/apply/{proposal_id}")
@app.post("/api/v13/evolution/apply/{proposal_id}")
async def evolution_apply(proposal_id: str, req: EvolutionApplyRequest | None = None) -> dict[str, Any]:
    req = req or EvolutionApplyRequest()
    if getattr(req, "auto_approve_after_validation", False):
        validation = validate_proposal(proposal_id)
        if validation.status != "validated":
            return {"ok": False, "status": validation.status, "message": validation.message, **validation.details}
        approval = approve_proposal(proposal_id, approved_by=getattr(req, "approved_by", "human") or "human", note="auto approval requested by API caller after validation")
        if approval.status != "approved":
            return {"ok": False, "status": approval.status, "message": approval.message, **approval.details}
    result = apply_proposal(proposal_id, approved_by=getattr(req, "approved_by", "human") or "human", approval_phrase=getattr(req, "approval_phrase", "") or "", validation=getattr(req, "validation", None))
    log_event("evolution.apply", route="evolution", provider="local", model="self_evolver", ok=result.status == "applied", message=proposal_id, metadata=result.details)
    return {"ok": result.status == "applied", "status": result.status, "message": result.message, **result.details}

@app.post("/evolution/apply-quick")
@app.post("/api/v13/evolution/apply-quick")
async def evolution_apply_quick(req: EvolutionProposalRequest) -> dict[str, Any]:
    proposal = create_proposal(req.relative_path, req.new_content, reason=req.reason or "quick apply", author=req.author or "user", scope=req.scope or "code", tests=req.tests)
    if not proposal.get("ok"):
        return proposal
    validation = validate_proposal(proposal["id"])
    if validation.status != "validated":
        return {"ok": False, "status": validation.status, "message": validation.message, **validation.details}
    approval = approve_proposal(proposal["id"], approved_by=req.author or "user", note="quick apply after sandbox validation")
    if approval.status != "approved":
        return {"ok": False, "status": approval.status, "message": approval.message, **approval.details}
    applied = apply_proposal(proposal["id"], approved_by=req.author or "user")
    return {"ok": applied.status == "applied", "status": applied.status, "message": applied.message, **applied.details}


@app.get("/evolution/proposals/{proposal_id}/approval-card")
@app.get("/api/v13/evolution/proposals/{proposal_id}/approval-card")
async def evolution_approval_card(proposal_id: str) -> dict[str, Any]:
    """Return a one-tap approve/discard card for the UI."""
    try:
        return {"ok": True, "card": approval_card(proposal_id)}
    except Exception as exc:
        return {"ok": False, "message": str(exc)[:220]}


@app.post("/evolution/proposals/{proposal_id}/undo")
@app.post("/api/v13/evolution/proposals/{proposal_id}/undo")
async def evolution_undo(proposal_id: str) -> dict[str, Any]:
    """Revert an applied patch within the undo window."""
    result = undo_apply(proposal_id)
    log_event("evolution.undo", route="evolution", provider="local", model="self-evolver", ok=result.status == "undone", message=proposal_id, metadata=result.details)
    return {"ok": result.status == "undone", "status": result.status, "message": result.message, **result.details}


@app.post("/api/improvement/run")
async def improvement_run(system_prompt_text: str = "") -> dict[str, Any]:
    """Run one eval → reflect → propose improvement cycle."""
    try:
        result = run_improvement_cycle(system_prompt_text=system_prompt_text)
        return {"ok": True, **result}
    except Exception as exc:
        return {"ok": False, "message": str(exc)[:260]}


@app.get("/api/improvement/runs")
async def improvement_runs(limit: int = 20) -> dict[str, Any]:
    try:
        return {"ok": True, "runs": list_improvement_runs(limit=limit)}
    except Exception as exc:
        return {"ok": False, "message": str(exc)[:260]}


@app.post("/api/plans/learn")
async def plans_learn(min_steps: int = 2, limit: int = 20) -> dict[str, Any]:
    """Scan completed plans and turn them into reusable skills."""
    try:
        return learn_from_completed_plans(min_steps=min_steps, limit=limit)
    except Exception as exc:
        return {"ok": False, "message": str(exc)[:260]}


@app.post("/api/plans/{plan_id}/learn")
async def plan_learn_one(plan_id: str) -> dict[str, Any]:
    """Turn one completed plan into a skill."""
    try:
        return plan_to_skill(plan_id)
    except Exception as exc:
        return {"ok": False, "message": str(exc)[:260]}


@app.post("/api/plans/{plan_id}/fail")
async def plan_record_failure(plan_id: str, reason: str = "") -> dict[str, Any]:
    """Record a plan failure pattern for planner improvement."""
    try:
        return record_plan_failure(plan_id, reason=reason or "unspecified")
    except Exception as exc:
        return {"ok": False, "message": str(exc)[:260]}


@app.get("/api/plans/suggest")
async def plans_suggest(goal: str) -> dict[str, Any]:
    """Suggest a learned plan for a goal."""
    try:
        return suggest_plan_for_goal(goal)
    except Exception as exc:
        return {"ok": False, "message": str(exc)[:260]}


@app.get("/api/v13/health")
async def api_v13_health() -> dict[str, Any]:
    return {"ok": True, "name": APP_NAME, "version": APP_VERSION, "brain": "unified-v13", "features": ["tool-first-routing", "mcp-style-manifest", "guarded-self-evolution", "autonomy-policy", "android-bridge", "enterprise-six-stack-foundation"]}

@app.post("/api/v13/chat/turn")
async def api_v13_chat_turn(req: ChatRequest) -> StreamingResponse:
    return StreamingResponse(_safe_brain_stream(req), media_type="application/x-ndjson")

@app.get("/api/v13/mcp/manifest")
async def api_v13_mcp_manifest() -> dict[str, Any]:
    return mcp_manifest()

@app.get("/api/v13/autonomy/policy")
async def api_v13_autonomy_policy() -> dict[str, Any]:
    return autonomy_policy()

@app.post("/api/v13/autonomy/check")
async def api_v13_autonomy_check(req: AutonomyCheckRequest) -> dict[str, Any]:
    return check_autonomy(req.action, req.requested_level)

@app.get("/ledger/verify")
async def ledger_verify(path: str) -> dict[str, Any]:
    return verify_document(path)

@app.post("/ledger/register")
async def ledger_register(request: Request) -> dict[str, Any]:
    data = await request.json()
    path = data.get("path")
    if not path:
        raise HTTPException(400, "path required")
    return ledger_document(path, data.get("document_type") or "document", data.get("metadata") or {})

@app.exception_handler(HTTPException)
async def http_exception_handler(request: Request, exc: HTTPException):
    return JSONResponse({"detail": exc.detail}, status_code=exc.status_code)

# Backward-compatible helpers used by older smoke tests.
def _infer_provider(model: str | None, provider: str | None = None) -> str:
    if provider and provider.lower().strip() == 'ollama':
        return 'ollama'
    if model and _looks_local_model(model):
        return 'ollama'
    cloud = _cloud_provider_from_model(model)
    if cloud:
        return cloud
    if provider and provider.lower().strip() not in {'', 'auto'}:
        return provider.lower().strip()
    return 'ollama'


def _resolve_model(provider: str, model: str | None) -> str:
    provider = (provider or 'ollama').lower().strip()
    if provider == 'ollama':
        return model or DEFAULT_OLLAMA_MODEL
    if model and _looks_local_model(model):
        return model
    return model or PROVIDER_DEFAULTS.get(provider, DEFAULT_OLLAMA_MODEL)


# ============================================================================
# Support & Abuse Reporting
# ============================================================================

SUPPORT_DIR = ROOT / 'storage' / 'support'
SUPPORT_DIR.mkdir(parents=True, exist_ok=True)


@app.post("/api/v15/support/abuse-report")
async def api_abuse_report(req: Request) -> dict[str, Any]:
    body = await req.json()
    report = {
        "timestamp": datetime.utcnow().isoformat(),
        "category": body.get("category", "other"),
        "description": body.get("description", ""),
        "platform": body.get("platform", "unknown"),
        "app_version": body.get("app_version", "unknown"),
    }
    path = SUPPORT_DIR / f"abuse_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.json"
    path.write_text(json.dumps(report, indent=2), encoding="utf-8")
    return {"ok": True, "message": "Report received. We will review within 48 hours."}


@app.post("/api/v15/support/ticket")
async def api_support_ticket(req: Request) -> dict[str, Any]:
    body = await req.json()
    ticket = {
        "timestamp": datetime.utcnow().isoformat(),
        "subject": body.get("subject", ""),
        "body": body.get("body", ""),
        "platform": body.get("from_platform", "unknown"),
        "app_version": body.get("app_version", "unknown"),
    }
    path = SUPPORT_DIR / f"ticket_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.json"
    path.write_text(json.dumps(ticket, indent=2), encoding="utf-8")
    return {"ok": True, "message": "Ticket created. Support will respond via email within 24 hours."}


# ============================================================================
# Subscription
# ============================================================================

SUBSCRIPTION_FILE = ROOT / 'storage' / 'subscriptions.json'


def _load_subscriptions() -> dict:
    if SUBSCRIPTION_FILE.exists():
        return json.loads(SUBSCRIPTION_FILE.read_text(encoding="utf-8"))
    return {}


@app.get("/api/v15/subscription/status")
async def api_subscription_status(req: Request) -> dict[str, Any]:
    subs = _load_subscriptions()
    device_id = req.headers.get("X-Device-Id", "anonymous")
    user_sub = subs.get(device_id, {})
    tier = user_sub.get("tier", "free")
    expires = user_sub.get("expires_at")
    return {
        "ok": True,
        "tier": tier,
        "is_premium": tier == "premium",
        "expires_at": expires,
        "features": {
            "unlimited_chats": tier == "premium",
            "cloud_sync": tier == "premium",
            "advanced_voice": tier == "premium",
            "priority_support": tier == "premium",
        }
    }


@app.post("/api/v15/subscription/activate")
async def api_subscription_activate(req: Request) -> dict[str, Any]:
    body = await req.json()
    device_id = body.get("device_id", "")
    tier = body.get("tier", "premium")
    duration_days = body.get("duration_days", 30)
    if not device_id:
        return {"ok": False, "error": "device_id required"}
    subs = _load_subscriptions()
    subs[device_id] = {
        "tier": tier,
        "activated_at": datetime.utcnow().isoformat(),
        "expires_at": (datetime.utcnow() + timedelta(days=duration_days)).isoformat(),
    }
    SUBSCRIPTION_FILE.write_text(json.dumps(subs, indent=2), encoding="utf-8")
    return {"ok": True, "tier": tier, "expires_at": subs[device_id]["expires_at"]}


# ============================================================================
# Analytics (privacy-respecting, opt-in)
# ============================================================================

ANALYTICS_FILE = ROOT / 'storage' / 'analytics_events.jsonl'
ANALYTICS_FILE.parent.mkdir(parents=True, exist_ok=True)


@app.post("/api/v15/analytics/event")
async def api_analytics_event(req: Request) -> dict[str, Any]:
    body = await req.json()
    event = {
        "timestamp": datetime.utcnow().isoformat(),
        "event_name": body.get("event_name", "unknown"),
        "platform": body.get("platform", "unknown"),
        "app_version": body.get("app_version", "unknown"),
        "session_id": body.get("session_id", ""),
    }
    with ANALYTICS_FILE.open("a", encoding="utf-8") as f:
        f.write(json.dumps(event) + "\n")
    return {"ok": True}


# ============================================================================
# Model / Provider Settings
# ============================================================================

@app.get("/api/v15/settings/models")
async def api_settings_models() -> dict[str, Any]:
    # Dynamically query Ollama so installed local models appear in the Settings dropdown.
    # Only tool-capable models are shown unless SHIMS_SHOW_ALL_MODELS is set.
    ollama_models = []
    try:
        raw = await _ollama_models_raw(timeout=1.5)
        for m in raw:
            size_gb = round(m.get("size", 0) / (1024**3), 1)
            ollama_models.append({
                "id": m["name"],
                "name": m["name"],
                "description": f"{m.get('family', 'local')} · {m.get('quantization', 'GGUF')} · {size_gb}GB",
                "ram_mb": int(size_gb * 1024),
                "tool_capable": m.get('tool_capable', is_tool_capable(m['name'])),
            })
    except Exception:
        pass
    # Fallback hardcoded list if Ollama is unreachable
    if not ollama_models:
        ollama_models = [
            {"id": "llama3.2:latest", "name": "Llama 3.2", "description": "Fast, balanced", "ram_mb": 4000},
            {"id": "qwen2.5:7b", "name": "Qwen 2.5", "description": "Smart, detailed", "ram_mb": 6000},
            {"id": "phi3:latest", "name": "Phi-3", "description": "Lightweight", "ram_mb": 2500},
        ]
    return {
        "ok": True,
        "providers": [
            {"id": "ollama", "name": "Local AI (Ollama)", "requires_internet": False, "privacy": "on-device"},
            {"id": "openai", "name": "OpenAI", "requires_internet": True, "privacy": "cloud"},
            {"id": "anthropic", "name": "Anthropic", "requires_internet": True, "privacy": "cloud"},
            {"id": "gemini", "name": "Gemini", "requires_internet": True, "privacy": "cloud"},
            {"id": "deepseek", "name": "DeepSeek", "requires_internet": True, "privacy": "cloud"},
            {"id": "kimi", "name": "Kimi", "requires_internet": True, "privacy": "cloud"},
            {"id": "huggingface", "name": "HuggingFace Local", "requires_internet": False, "privacy": "on-device"},
            {"id": "qwen", "name": "Qwen / Alibaba", "requires_internet": True, "privacy": "cloud"},
            {"id": "auto", "name": "Auto-Select", "requires_internet": False, "privacy": "mixed"},
        ],
        "models": _build_settings_models(ollama_models)
    }


def _build_settings_models(ollama_models: list[dict[str, Any]]) -> dict[str, list[dict[str, Any]]]:
    """Group curated RECOMMENDED_MODELS by provider for the Settings dropdowns.

    This keeps the Settings model pickers in sync with /chat/models so users can
    choose any known latest or older model for each provider.
    """
    models: dict[str, list[dict[str, Any]]] = {"ollama": ollama_models}
    for pid in ["openai", "anthropic", "gemini", "deepseek", "kimi", "qwen", "huggingface"]:
        entries: list[dict[str, Any]] = []
        for m in RECOMMENDED_MODELS:
            if m.get("provider") != pid:
                continue
            desc = m.get("notes") or m.get("role") or "Cloud model"
            entries.append({
                "id": m["name"],
                "name": m["name"],
                "description": desc,
                "ram_mb": 0,
                "tool_capable": bool(m.get("tool_capable")),
            })
        # De-duplicate while preserving order
        seen: set[str] = set()
        deduped: list[dict[str, Any]] = []
        for e in entries:
            if e["id"] not in seen:
                seen.add(e["id"])
                deduped.append(e)
        models[pid] = deduped
    models["auto"] = [{"id": "auto", "name": "Auto", "description": "Best available", "ram_mb": 0}]
    return models


@app.post("/api/v15/settings/agent-models")
async def api_settings_agent_models(req: AgentModelsRequest) -> dict[str, Any]:
    """Persist per-role specialist model env vars."""
    mapping = {
        "SHIMS_ROUTER_MODEL": req.router_model,
        "SHIMS_FAST_MODEL": req.fast_model,
        "SHIMS_SMART_MODEL": req.smart_model,
        "SHIMS_CODER_MODEL": req.coder_model,
        "SHIMS_CREATIVE_MODEL": req.creative_model,
        "SHIMS_CHEMISTRY_MODEL": req.chemistry_model,
        "SHIMS_RESEARCH_MODEL": req.research_model,
    }
    saved = []
    for env_var, value in mapping.items():
        if value is not None:
            cleaned = str(value).strip()
            _set_env_persistent(env_var, cleaned)
            if cleaned:
                saved.append(env_var)
    return {"ok": True, "saved": saved}


# ============================================================================
# AI Document Ingestion
# ============================================================================

INGESTED_DOCS_DIR = ROOT / "storage" / "ingested_documents"
INGESTED_DOCS_DIR.mkdir(parents=True, exist_ok=True)


def _extract_text_from_file(path: Path) -> tuple[str, str | None]:
    """Extract text from PDF, DOCX, TXT, etc. Returns (text, error_or_None)."""
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf" and PdfReader is not None:
            reader = PdfReader(str(path))
            parts = []
            for page in reader.pages:
                try:
                    txt = page.extract_text()
                    if txt:
                        parts.append(txt)
                except Exception:
                    pass
            return "\n".join(parts), None
        if suffix in {".docx", ".doc"} and DocxDocument is not None:
            doc = DocxDocument(str(path))
            parts: list[str] = []
            # Body paragraphs
            for p in doc.paragraphs:
                if p.text:
                    parts.append(p.text)
            # Tables (common place for structured document content)
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells if cell.text)
                    if row_text:
                        parts.append(row_text)
            # Headers/footers across sections
            for section in doc.sections:
                for p in section.header.paragraphs:
                    if p.text:
                        parts.append(p.text)
                for p in section.footer.paragraphs:
                    if p.text:
                        parts.append(p.text)
            text = "\n".join(parts)
            if not text.strip():
                return "", f"No readable text found in {suffix} file (empty paragraphs/tables)."
            return text, None
        if suffix in {".txt", ".md", ".csv", ".json"}:
            return path.read_text(encoding="utf-8", errors="ignore"), None
        # Fallback: try to read as text
        return path.read_text(encoding="utf-8", errors="ignore"), None
    except Exception as exc:
        return "", f"Extraction failed for {suffix}: {exc}"


async def _ai_summarize_document(text: str, filename: str, provider: str | None = None, model: str | None = None) -> dict[str, Any]:
    provider, model, _ = await _resolve_provider_model(provider, model)
    if not await _provider_ready_for_llm(provider, model):
        return {"ok": False, "error": f"Provider {provider} not ready", "summary": text[:2000], "key_points": [], "entities": []}
    prompt = (
        f"Analyze the following document ({filename}) and return ONLY a JSON object with these keys:\n"
        "summary: a concise 3-5 sentence summary\n"
        "key_points: array of 5-8 key bullet points\n"
        "entities: array of important named entities (people, companies, products, regulations, etc.)\n"
        "category: best single category (e.g., SOP, COA, Report, Email, Contract, Research, Invoice, Manual, Other)\n"
        "sentiment: overall tone (neutral, positive, negative, urgent, cautionary)\n\n"
        "Document text (first 12000 chars):\n" + text[:12000] + "\n\nReturn only valid JSON."
    )
    messages = [
        {"role": "system", "content": "You are a precise document analyst. Return only JSON."},
        {"role": "user", "content": prompt},
    ]
    try:
        raw, route = await _run_llm(provider, model, messages)
        cleaned = raw.strip()
        if cleaned.startswith("```"):
            cleaned = cleaned.split("\n", 1)[1] if "\n" in cleaned else ""
        if cleaned.endswith("```"):
            cleaned = cleaned.rsplit("\n", 1)[0]
        cleaned = cleaned.strip()
        parsed = json.loads(cleaned)
        return {
            "ok": True,
            "summary": parsed.get("summary", ""),
            "key_points": parsed.get("key_points", []),
            "entities": parsed.get("entities", []),
            "category": parsed.get("category", "Other"),
            "sentiment": parsed.get("sentiment", "neutral"),
            "route": route,
        }
    except Exception as exc:
        return {"ok": False, "error": str(exc)[:260], "summary": text[:2000], "key_points": [], "entities": [], "category": "Other", "sentiment": "neutral"}


@app.post("/api/v15/documents/ingest")
async def api_document_ingest(file: UploadFile = File(...), provider: str | None = None, model: str | None = None) -> dict[str, Any]:
    suffix = Path(file.filename or "document.txt").suffix or ".txt"
    safe_name = _safe_name("ingest", suffix.lstrip("."))
    file_path = INGESTED_DOCS_DIR / safe_name
    file_path.write_bytes(await file.read())
    raw_text, extract_error = _extract_text_from_file(file_path)
    if not raw_text.strip():
        return {"ok": False, "error": extract_error or "Could not extract text from document"}
    ai_result = await _ai_summarize_document(raw_text, file.filename or "document", provider=provider, model=model)
    full_text_for_ingest = f"DOCUMENT: {file.filename or 'uploaded'}\nCATEGORY: {ai_result.get('category', 'Other')}\nSUMMARY: {ai_result.get('summary', '')}\nKEY POINTS: {chr(10).join(ai_result.get('key_points', []))}\nENTITIES: {', '.join(ai_result.get('entities', []))}\nSENTIMENT: {ai_result.get('sentiment', 'neutral')}\n\nFULL TEXT:\n{raw_text[:8000]}"
    ingest = ingest_knowledge(
        title=file.filename or "Uploaded document",
        text=full_text_for_ingest,
        source_type="document",
        source_uri=str(file_path),
        tags=["document", ai_result.get("category", "Other").lower(), "upload"],
        importance=1.2,
    )
    ingest.update({
        "ai": ai_result,
        "filename": file.filename,
        "saved_path": str(file_path),
        "char_count": len(raw_text),
    })
    return ingest


@app.get("/api/v15/documents/ingested")
async def api_document_ingested_list() -> dict[str, Any]:
    docs = []
    for f in sorted(INGESTED_DOCS_DIR.iterdir(), key=lambda x: x.stat().st_mtime, reverse=True):
        try:
            docs.append({"name": f.name, "size": f.stat().st_size, "modified": f.stat().st_mtime})
        except Exception:
            pass
    return {"ok": True, "documents": docs}


# ============================================================================
# Content Policy
# ============================================================================

@app.get("/api/v15/content-policy")
async def api_content_policy() -> dict[str, Any]:
    return {
        "ok": True,
        "policy_url": "https://shims.jklifecare.com/content-policy",
        "categories": ["hate_speech", "harassment", "violence", "csam", "illegal_acts", "malware", "self_harm", "spam"],
        "report_email": "support@jklifecare.com",
        "response_time_hours": 48
    }


# ============================================================================
# Business Card Scanner (SHIMS CardScan)
# ============================================================================

SCAN_CONTACTS_DIR = ROOT / "storage" / "contacts"
SCAN_CONTACTS_DIR.mkdir(parents=True, exist_ok=True)


class ScanParseRequest(BaseModel):
    ocr_text: str
    provider: str | None = None
    model: str | None = None


class ScanCompanyRequest(BaseModel):
    company_name: str
    provider: str | None = None
    model: str | None = None


class ScanVcfRequest(BaseModel):
    name: str = ""
    title: str = ""
    company: str = ""
    phone: str = ""
    email: str = ""
    website: str = ""
    address: str = ""
    social: str = ""
    notes: str = ""


@app.post("/api/v15/scan/parse")
async def api_scan_parse(req: ScanParseRequest) -> dict[str, Any]:
    """Use AI to parse raw OCR text from a business card into structured fields."""
    provider, model, _ = await _resolve_provider_model(req.provider, req.model)
    if not await _provider_ready_for_llm(provider, model):
        return {"ok": False, "error": f"Provider {provider} not ready. Choose an installed Ollama model or configure a cloud key."}
    prompt = (
        "Extract business card information from this OCR text and return ONLY a JSON object with these keys: "
        "name, title, company, phone, email, website, address, social. "
        "If a field is not found, use an empty string. "
        "OCR Text:\n" + req.ocr_text + "\n\nReturn only valid JSON, nothing else."
    )
    messages = [
        {"role": "system", "content": "You are a precise business card parser. Return only JSON."},
        {"role": "user", "content": prompt},
    ]
    try:
        text, route = await _run_llm(provider, model, messages)
        # Strip markdown fences if the model wrapped JSON in them
        cleaned = text.strip()
        if cleaned.startswith("```"):
            cleaned = cleaned.split("\n", 1)[1] if "\n" in cleaned else ""
        if cleaned.endswith("```"):
            cleaned = cleaned.rsplit("\n", 1)[0]
        cleaned = cleaned.strip()
        parsed = json.loads(cleaned)
        return {
            "ok": True,
            "parsed": {
                "name": parsed.get("name", ""),
                "title": parsed.get("title", ""),
                "company": parsed.get("company", ""),
                "phone": parsed.get("phone", ""),
                "email": parsed.get("email", ""),
                "website": parsed.get("website", ""),
                "address": parsed.get("address", ""),
                "social": parsed.get("social", ""),
            },
            "route": route,
        }
    except Exception as exc:
        return {"ok": False, "error": str(exc)[:260]}


@app.post("/api/v15/scan/company")
async def api_scan_company(req: ScanCompanyRequest) -> dict[str, Any]:
    """Use AI to look up company info by name."""
    provider, model, _ = await _resolve_provider_model(req.provider, req.model)
    if not await _provider_ready_for_llm(provider, model):
        return {"ok": False, "error": f"Provider {provider} not ready. Choose an installed Ollama model or configure a cloud key."}
    prompt = (
        f'You are a business intelligence assistant. For the company "{req.company_name}", provide:\n'
        "1. A brief description (2-3 sentences)\n"
        "2. Their main products or services (list 4-6 items)\n"
        "3. Industry category\n"
        "4. Estimated company size (startup/SME/enterprise)\n\n"
        'Return ONLY a JSON object with keys: description, products (array), industry, size.\n'
        "Return only valid JSON, nothing else."
    )
    messages = [
        {"role": "system", "content": "You are a concise business intelligence assistant. Return only JSON."},
        {"role": "user", "content": prompt},
    ]
    try:
        text, route = await _run_llm(provider, model, messages)
        cleaned = text.strip()
        if cleaned.startswith("```"):
            cleaned = cleaned.split("\n", 1)[1] if "\n" in cleaned else ""
        if cleaned.endswith("```"):
            cleaned = cleaned.rsplit("\n", 1)[0]
        cleaned = cleaned.strip()
        parsed = json.loads(cleaned)
        return {
            "ok": True,
            "company": {
                "name": req.company_name,
                "description": parsed.get("description", ""),
                "products": parsed.get("products", []),
                "industry": parsed.get("industry", ""),
                "size": parsed.get("size", ""),
            },
            "route": route,
        }
    except Exception as exc:
        return {"ok": False, "error": str(exc)[:260]}


@app.post("/api/v15/scan/vcf")
async def api_scan_vcf(req: ScanVcfRequest) -> Response:
    """Generate a downloadable VCF contact card from parsed business card data."""
    vcf_lines = ["BEGIN:VCARD", "VERSION:3.0"]
    if req.name:
        vcf_lines.append(f"FN:{req.name}")
        parts = req.name.strip().split()
        if len(parts) >= 2:
            vcf_lines.append(f"N:{parts[-1]};{' '.join(parts[:-1])};;;")
        else:
            vcf_lines.append(f"N:;{req.name};;;")
    if req.title:
        vcf_lines.append(f"TITLE:{req.title}")
    if req.company:
        vcf_lines.append(f"ORG:{req.company}")
    if req.phone:
        vcf_lines.append(f"TEL;TYPE=CELL:{req.phone}")
    if req.email:
        vcf_lines.append(f"EMAIL;TYPE=WORK:{req.email}")
    if req.website:
        vcf_lines.append(f"URL:{req.website}")
    if req.address:
        vcf_lines.append(f"ADR;TYPE=WORK:;;{req.address};;;;")
    if req.social:
        vcf_lines.append(f"X-SOCIALPROFILE:{req.social}")
    if req.notes:
        vcf_lines.append(f"NOTE:{req.notes}")
    vcf_lines.append("END:VCARD")
    vcf_body = "\r\n".join(vcf_lines) + "\r\n"
    filename = f"{req.name or 'contact'}_shims.vcf"
    return Response(
        content=vcf_body,
        media_type="text/vcard",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@app.post("/api/v15/scan/save")
async def api_scan_save(req: Request) -> dict[str, Any]:
    """Save a scanned contact to local storage."""
    body = await req.json()
    contact = body.get("contact", {})
    contact_id = body.get("id") or str(uuid.uuid4())
    path = SCAN_CONTACTS_DIR / f"{contact_id}.json"
    path.write_text(json.dumps(contact, indent=2, default=str), encoding="utf-8")
    return {"ok": True, "id": contact_id}


@app.get("/api/v15/scan/contacts")
async def api_scan_contacts() -> dict[str, Any]:
    """List saved scanned contacts."""
    contacts = []
    for f in sorted(SCAN_CONTACTS_DIR.glob("*.json"), reverse=True):
        try:
            contacts.append(json.loads(f.read_text(encoding="utf-8")))
        except Exception:
            continue
    return {"ok": True, "contacts": contacts}


# ── Neural Governor (Omni) ───────────────────────────────────────────────────

@app.get("/neural-governor", include_in_schema=False)
async def neural_governor_page() -> HTMLResponse:
    """Omni Neural Governor dashboard."""
    html = '''<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Neural Governor — SHIMS Omni</title>
<style>
body{background:#0b1220;color:#e2e8f0;font-family:system-ui,-apple-system,BlinkMacSystemFont,Segoe UI,Roboto,sans-serif;margin:0;padding:1.5rem}
.card{background:#0f172a;border:1px solid #1e293b;border-radius:12px;padding:1.25rem;margin-bottom:1rem}
.card h3{margin:0 0 .75rem 0;font-size:1.1rem;color:#e2e8f0}
.grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(280px,1fr));gap:1rem}
.metric{background:#1e293b;border-radius:8px;padding:.75rem 1rem;display:flex;justify-content:space-between;align-items:center}
.metric label{color:#94a3b8;font-size:.85rem}
.metric value{color:#f8fafc;font-weight:600;font-size:.95rem}
.chat-box{background:#0b1220;border:1px solid #1e293b;border-radius:12px;padding:1rem;min-height:180px;max-height:320px;overflow-y:auto;margin-bottom:.75rem}
.chat-input{display:flex;gap:.5rem}
.chat-input input{flex:1;background:#1e293b;border:1px solid #334155;border-radius:8px;padding:.6rem .9rem;color:#f8fafc;font-size:.95rem}
.chat-input button{background:#3b82f6;border:none;border-radius:8px;padding:.6rem 1.2rem;color:#fff;font-weight:600;cursor:pointer}
.chat-input button:hover{background:#2563eb}
.message-user{color:#60a5fa;margin:.5rem 0}
.message-assistant{color:#e2e8f0;margin:.5rem 0}
.trust-card{background:#1e293b;border-left:4px solid #3b82f6;border-radius:6px;padding:.75rem;margin-top:.5rem;font-size:.85rem;color:#94a3b8}
.status-dot{display:inline-block;width:10px;height:10px;border-radius:50%;margin-right:6px}
.status-closed{background:#22c55e}.status-open{background:#ef4444}.status-half_open{background:#f59e0b}
a{color:#60a5fa}
</style>
</head>
<body>
<div style="display:flex;align-items:center;gap:12px;margin-bottom:8px">
  <a href="/" style="color:#60a5fa;text-decoration:none;font-size:13px">← Back to Omni</a>
  <a href="/neural-agent" style="color:#60a5fa;text-decoration:none;font-size:13px">🧠 Neural Agent</a>
</div>
<h1>🧬 Neural Governor</h1>
<p style="color:#94a3b8">Cognitive governance layer for SHIMS Omni. Monitors drift, arbitrates corrections, and learns your style.</p>
<div class="grid">
<div class="card"><h3>System Health</h3><div id="health">Loading...</div></div>
<div class="card"><h3>Provider Circuits</h3><div id="circuits">Loading...</div></div>
<div class="card"><h3>Personal Profile</h3><div id="profile">Loading...</div></div>
<div class="card" style="grid-column:1/-1"><h3>Governed Chat</h3>
<div id="chat-box" class="chat-box"><div style="color:#64748b;font-size:.9rem">Start a conversation. Every response is drift-checked and lineage-tracked.</div></div>
<div class="chat-input"><input id="chat-prompt" type="text" placeholder="Ask anything..." onkeydown="if(event.key==='Enter')sendChat()"><button onclick="sendChat()">Send</button></div>
<div id="chat-status" style="color:#64748b;font-size:.75rem;margin-top:.5rem"></div>
</div>
</div>
<script>
async function loadDiagnostics(){
  try{
    const r=await fetch('/api/neural-governor/diagnostics');
    const d=await r.json();
    const h=d.hardware||{};
    document.getElementById('health').innerHTML=
      '<div class=metric><label>Platform</label><value>'+h.platform+'</value></div>'+
      '<div class=metric><label>RAM</label><value>'+h.total_ram_gb+' GB</value></div>'+
      '<div class=metric><label>VRAM</label><value>'+h.vram_gb+' GB</value></div>'+
      '<div class=metric><label>CUDA</label><value>'+(h.cuda_available?'Yes':'No')+'</value></div>'+
      '<div class=metric><label>Internet</label><value>'+(h.internet_available?'Online':'Offline')+'</value></div>';
    let chtml='';
    (d.circuits||[]).forEach(c=>{
      chtml+='<div class=metric><label>'+c.provider+'</label><value><span class="status-dot status-'+c.status+'"></span>'+c.status.toUpperCase()+' <small>(F:'+c.failures+' S:'+c.successes+')</small></value></div>';
    });
    document.getElementById('circuits').innerHTML=chtml||'<span style=color:#64748b>No circuit data yet.</span>';
  }catch(e){document.getElementById('health').textContent='Error: '+e.message;}
}
async function loadProfile(){
  try{
    const r=await fetch('/api/neural-governor/personal/profile');
    const d=await r.json();
    const p=d.profile||{};
    document.getElementById('profile').innerHTML=
      '<div class=metric><label>Writing Style</label><value>'+(p.writing_style||'formal')+'</value></div>'+
      '<div class=metric><label>Technical Depth</label><value>'+(p.technical_depth||3)+'/5</value></div>'+
      '<div class=metric><label>Tone</label><value>'+(p.communication_tone||'professional')+'</value></div>'+
      '<div class=metric><label>Learning</label><value>'+(p.learning_enabled?'Enabled':'Paused')+'</value></div>';
  }catch(e){document.getElementById('profile').textContent='Error: '+e.message;}
}
async function sendChat(){
  const input=document.getElementById('chat-prompt');
  const box=document.getElementById('chat-box');
  const status=document.getElementById('chat-status');
  const prompt=input.value.trim();
  if(!prompt)return;
  input.value='';
  box.innerHTML+='<div class=message-user><strong>You:</strong> '+escapeHtml(prompt)+'</div>';
  status.textContent='Thinking...';
  try{
    const res=await fetch('/api/neural-governor/chat',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({prompt,session_id:'gov_dashboard'})});
    const data=await res.json();
    if(data.output){
      box.innerHTML+='<div class=message-assistant><strong>SHIMS:</strong> '+escapeHtml(data.output)+'</div>';
      if(data.drift_report&&data.drift_report.triggered){
        box.innerHTML+='<div class=trust-card>⚠️ Drift detected: '+data.drift_report.signals_triggered.join(', ')+' | Score: '+data.drift_report.composite+'</div>';
      }
      status.textContent='Model: '+data.routing.model+' | Latency: '+data.latency_ms+'ms | Trust: '+data.trust_score;
    }else{
      status.textContent='Error: '+(data.detail||'Unknown error');
    }
  }catch(e){status.textContent='Network error: '+e.message;}
  box.scrollTop=box.scrollHeight;
}
function escapeHtml(text){const div=document.createElement('div');div.textContent=text;return div.innerHTML;}
loadDiagnostics();loadProfile();
</script>
</body>
</html>''';
    return HTMLResponse(html, headers={"Cache-Control": "no-cache, no-store, must-revalidate"})


@app.post("/api/neural-governor/chat")
async def api_neural_governor_chat(req: ChatRequest) -> dict[str, Any]:
    """Governed chat via Neural Governor."""
    if not req.message:
        raise HTTPException(400, "Prompt required")
    session_id = req.session_id or str(uuid.uuid4())
    gov = NeuralGovernor(user_id=0, session_id=session_id)
    result = await gov.chat(
        prompt=req.message,
        provider_preference=req.provider,
        model_preference=req.model,
    )
    return result


@app.get("/api/neural-governor/lineage/{lineage_id}")
async def api_neural_governor_lineage(lineage_id: str) -> dict[str, Any]:
    record = get_lineage(lineage_id)
    if not record:
        raise HTTPException(404, "Lineage not found")
    return record


@app.post("/api/neural-governor/feedback")
async def api_neural_governor_feedback(req: Request) -> dict[str, Any]:
    body = await req.json()
    lineage_id = str(body.get("lineage_id", ""))
    rating = int(body.get("rating", 0))
    notes = str(body.get("notes", ""))
    if rating not in (-1, 1):
        raise HTTPException(400, "Rating must be 1 or -1")
    gov = NeuralGovernor(user_id=0)
    ok = await gov.feedback(lineage_id, rating, notes)
    return {"ok": ok}


@app.get("/api/neural-governor/models")
async def api_neural_governor_models() -> dict[str, Any]:
    return {"ok": True, "models": model_registry_list()}


@app.get("/api/neural-governor/drift/summary")
async def api_neural_governor_drift_summary(days: int = 7) -> dict[str, Any]:
    return get_drift_summary(0, days)


@app.get("/api/neural-governor/diagnostics")
async def api_neural_governor_diagnostics() -> dict[str, Any]:
    return {
        "ok": True,
        "hardware": quick_profile(),
        "router_status": get_router_status(),
        "circuits": get_all_circuits(),
        "recent_resources": get_recent_snapshots(10),
    }


@app.get("/api/neural-governor/evolution/proposals")
async def api_neural_governor_proposals(status: str | None = None) -> dict[str, Any]:
    return {"ok": True, "proposals": governor_list_proposals(status)}



@app.get("/api/neural-governor/evolution/proposals/{proposal_uuid}")
async def api_neural_governor_get_proposal(proposal_uuid: str) -> dict[str, Any]:
    proposal = governor_get_proposal(proposal_uuid)
    if not proposal:
        raise HTTPException(404, "Proposal not found")
    return {"ok": True, "proposal": proposal}

@app.post("/api/neural-governor/evolution/proposals/{proposal_uuid}/review")
async def api_neural_governor_review_proposal(proposal_uuid: str, req: Request) -> dict[str, Any]:
    body = await req.json()
    approved = bool(body.get("approved"))
    notes = str(body.get("notes", ""))
    return review_proposal(proposal_uuid, 0, approved, notes)


@app.get("/api/neural-governor/personal/profile")
async def api_neural_governor_personal_profile() -> dict[str, Any]:
    profile = get_profile(0)
    if not profile:
        profile = ensure_profile(0)
    return {"ok": True, "profile": profile.to_dict()}


@app.post("/api/neural-governor/personal/profile")
async def api_neural_governor_update_profile(req: Request) -> dict[str, Any]:
    body = await req.json()
    profile = ensure_profile(0)
    if "writing_style" in body:
        profile.writing_style = str(body["writing_style"])
    if "technical_depth" in body:
        profile.technical_depth = max(1, min(5, int(body["technical_depth"])))
    if "communication_tone" in body:
        profile.communication_tone = str(body["communication_tone"])
    if "learning_enabled" in body:
        profile.learning_enabled = bool(body["learning_enabled"])
    save_profile(profile)
    return {"ok": True, "profile": profile.to_dict()}


@app.get("/api/neural-governor/patent/generate")
async def api_neural_governor_patent() -> dict[str, Any]:
    return generate_patent_spec()


# ── Enhanced Coder Playground v2 ─────────────────────────────────────────────

@app.get("/coder/v2/templates")
async def coder_v2_templates() -> dict[str, Any]:
    from shared.coder_v2 import list_templates
    return {"ok": True, "templates": list_templates()}


@app.post("/coder/v2/project")
async def coder_v2_create_project(req: Request) -> dict[str, Any]:
    from shared.coder_v2 import create_project
    body = await req.json()
    name = str(body.get("name", "Untitled")).strip()
    template = body.get("template")
    return create_project(name, template)


@app.get("/coder/v2/projects")
async def coder_v2_list_projects() -> dict[str, Any]:
    from shared.coder_v2 import list_projects
    return {"ok": True, "projects": list_projects()}


@app.get("/coder/v2/project/{project_id}")
async def coder_v2_get_project(project_id: str) -> dict[str, Any]:
    from shared.coder_v2 import get_project
    project = get_project(project_id)
    if not project:
        raise HTTPException(404, "Project not found")
    return {"ok": True, "project": project}


@app.get("/coder/v2/project/{project_id}/files")
async def coder_v2_list_files(project_id: str, subdir: str = "") -> dict[str, Any]:
    from shared.coder_v2 import list_files
    return {"ok": True, "files": list_files(project_id, subdir)}


@app.get("/coder/v2/project/{project_id}/file")
async def coder_v2_read_file(project_id: str, path: str = "") -> dict[str, Any]:
    from shared.coder_v2 import read_file
    return read_file(project_id, path)


@app.post("/coder/v2/project/{project_id}/file")
async def coder_v2_write_file(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v2 import write_file
    body = await req.json()
    file_path = str(body.get("path", ""))
    content = str(body.get("content", ""))
    return write_file(project_id, file_path, content)


@app.delete("/coder/v2/project/{project_id}/file")
async def coder_v2_delete_file(project_id: str, path: str = "") -> dict[str, Any]:
    from shared.coder_v2 import delete_file
    return delete_file(project_id, path)


@app.post("/coder/v2/project/{project_id}/mkdir")
async def coder_v2_mkdir(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v2 import mkdir
    body = await req.json()
    dir_path = str(body.get("path", ""))
    return mkdir(project_id, dir_path)


@app.post("/coder/v2/project/{project_id}/run")
async def coder_v2_run(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v2 import run_project
    body = await req.json()
    entry_file = body.get("entry_file")
    return run_project(project_id, entry_file)


@app.post("/coder/v2/project/{project_id}/upload-folder")
async def coder_v2_upload_folder(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v2 import upload_folder
    body = await req.json()
    files = body.get("files", {})
    return upload_folder(project_id, files)


@app.get("/coder/v2/project/{project_id}/export")
async def coder_v2_export(project_id: str) -> FileResponse:
    from shared.coder_v2 import export_project
    result = export_project(project_id)
    if not result.get("ok"):
        raise HTTPException(400, result.get("error", "Export failed"))
    path = Path(result["zip_path"])
    return FileResponse(path, filename=f"{project_id}.zip", media_type="application/zip")


@app.post("/coder/v2/import")
async def coder_v2_import(req: Request) -> dict[str, Any]:
    from shared.coder_v2 import import_project
    body = await req.json()
    zip_b64 = body.get("zip_b64", "")
    name = str(body.get("name", "Imported"))
    if not zip_b64:
        raise HTTPException(400, "zip_b64 required")
    import base64
    try:
        zip_data = base64.b64decode(zip_b64)
    except Exception:
        raise HTTPException(400, "Invalid base64")
    return import_project(zip_data, name)


# ── Git Endpoints ────────────────────────────────────────────────────────────

@app.post("/coder/v2/project/{project_id}/git/init")
async def coder_v2_git_init(project_id: str) -> dict[str, Any]:
    from shared.coder_v2 import git_init
    return git_init(project_id)


@app.get("/coder/v2/project/{project_id}/git/status")
async def coder_v2_git_status(project_id: str) -> dict[str, Any]:
    from shared.coder_v2 import git_status
    return git_status(project_id)


@app.get("/coder/v2/project/{project_id}/git/log")
async def coder_v2_git_log(project_id: str, n: int = 10) -> dict[str, Any]:
    from shared.coder_v2 import git_log
    return git_log(project_id, n)


@app.get("/coder/v2/project/{project_id}/git/diff")
async def coder_v2_git_diff(project_id: str) -> dict[str, Any]:
    from shared.coder_v2 import git_diff
    return git_diff(project_id)


@app.post("/coder/v2/project/{project_id}/git/commit")
async def coder_v2_git_commit(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v2 import git_commit
    body = await req.json()
    message = str(body.get("message", "Update"))
    return git_commit(project_id, message)


@app.get("/coder/v2/project/{project_id}/git/branch")
async def coder_v2_git_branch(project_id: str) -> dict[str, Any]:
    from shared.coder_v2 import git_branch
    return git_branch(project_id)


# ── AI Context Endpoint ──────────────────────────────────────────────────────

@app.post("/coder/v2/project/{project_id}/ai-context")
async def coder_v2_ai_context(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v2 import build_ai_context
    body = await req.json()
    mentions = body.get("mention_files", [])
    return build_ai_context(project_id, mentions)


# ── Terminal Endpoints ───────────────────────────────────────────────────────

@app.post("/coder/v2/project/{project_id}/terminal/start")
async def coder_v2_terminal_start(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v2 import terminal_start
    body = await req.json()
    shell = body.get("shell")
    return terminal_start(project_id, shell)


@app.post("/coder/v2/terminal/{term_id}/read")
async def coder_v2_terminal_read(term_id: str) -> dict[str, Any]:
    from shared.coder_v2 import terminal_read
    return terminal_read(term_id)


@app.post("/coder/v2/terminal/{term_id}/write")
async def coder_v2_terminal_write(term_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v2 import terminal_write
    body = await req.json()
    data = str(body.get("data", ""))
    return terminal_write(term_id, data)


@app.post("/coder/v2/terminal/{term_id}/kill")
async def coder_v2_terminal_kill(term_id: str) -> dict[str, Any]:
    from shared.coder_v2 import terminal_kill
    return terminal_kill(term_id)


# ═══════════════════════════════════════════════════════════════════════════════
# CODER PLAYGROUND v3 — Full Power IDE Endpoints
# ═══════════════════════════════════════════════════════════════════════════════

@app.get("/coder-v3", response_class=HTMLResponse)
async def coder_v3_page() -> str:
    path = Path(__file__).resolve().parents[2] / "frontend" / "coder_v3.html"
    if path.exists():
        return path.read_text(encoding="utf-8")
    return "<h1>Coder v3 frontend not found</h1>"


# ── Shell & Process Management ────────────────────────────────────────────────

@app.post("/coder/v3/project/{project_id}/shell")
async def coder_v3_shell(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import run_shell_command
    body = await req.json()
    return run_shell_command(
        project_id,
        command=body.get("command", ""),
        timeout=body.get("timeout", 300),
        background=body.get("background", False),
        env_vars=body.get("env_vars"),
    )


@app.get("/coder/v3/processes")
async def coder_v3_processes(project_id: str = "") -> dict[str, Any]:
    from shared.coder_v3 import list_processes
    return {"ok": True, "processes": list_processes(project_id or None)}


@app.get("/coder/v3/process/{proc_id}/output")
async def coder_v3_process_output(proc_id: str, clear: bool = True) -> dict[str, Any]:
    from shared.coder_v3 import read_process_output
    return read_process_output(proc_id, clear=clear)


@app.post("/coder/v3/process/{proc_id}/kill")
async def coder_v3_process_kill(proc_id: str) -> dict[str, Any]:
    from shared.coder_v3 import kill_process
    return kill_process(proc_id)


# ── Port Management ───────────────────────────────────────────────────────────

@app.post("/coder/v3/project/{project_id}/port/allocate")
async def coder_v3_allocate_port(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import allocate_port
    body = await req.json()
    return allocate_port(project_id, purpose=body.get("purpose", "app"))


@app.get("/coder/v3/project/{project_id}/ports")
async def coder_v3_list_ports(project_id: str) -> dict[str, Any]:
    from shared.coder_v3 import list_allocated_ports
    return {"ok": True, "ports": list_allocated_ports(project_id)}


@app.post("/coder/v3/port/{port}/release")
async def coder_v3_release_port(port: int) -> dict[str, Any]:
    from shared.coder_v3 import release_port
    return release_port(port)


# ── Environment Variables ─────────────────────────────────────────────────────

@app.get("/coder/v3/project/{project_id}/env")
async def coder_v3_get_env(project_id: str) -> dict[str, Any]:
    from shared.coder_v3 import get_env_vars
    return get_env_vars(project_id)


@app.post("/coder/v3/project/{project_id}/env")
async def coder_v3_set_env(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import set_env_vars
    body = await req.json()
    return set_env_vars(project_id, body.get("env", {}), merge=body.get("merge", True))


# ── Search & Replace ──────────────────────────────────────────────────────────

@app.get("/coder/v3/project/{project_id}/search")
async def coder_v3_search(
    project_id: str,
    query: str,
    regex: bool = False,
    case_sensitive: bool = False,
    file_pattern: str = "*",
) -> dict[str, Any]:
    from shared.coder_v3 import search_in_project
    return search_in_project(project_id, query, regex, case_sensitive, file_pattern)


@app.post("/coder/v3/project/{project_id}/find-replace")
async def coder_v3_find_replace(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import find_replace_all
    body = await req.json()
    return find_replace_all(
        project_id,
        find=body.get("find", ""),
        replace=body.get("replace", ""),
        regex=body.get("regex", False),
        case_sensitive=body.get("case_sensitive", False),
        file_pattern=body.get("file_pattern", "*"),
    )


# ── Bulk Operations ───────────────────────────────────────────────────────────

@app.post("/coder/v3/project/{project_id}/bulk-delete")
async def coder_v3_bulk_delete(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import bulk_delete
    body = await req.json()
    return bulk_delete(project_id, body.get("paths", []))


@app.post("/coder/v3/project/{project_id}/bulk-move")
async def coder_v3_bulk_move(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import bulk_move
    body = await req.json()
    return bulk_move(project_id, body.get("moves", []))


@app.post("/coder/v3/project/{project_id}/bulk-copy")
async def coder_v3_bulk_copy(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import bulk_copy
    body = await req.json()
    return bulk_copy(project_id, body.get("copies", []))


# ── Package Management ────────────────────────────────────────────────────────

@app.post("/coder/v3/project/{project_id}/install")
async def coder_v3_install(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import install_packages
    body = await req.json()
    return install_packages(project_id, body.get("packages", []), manager=body.get("manager"))


# ── Code Formatting ───────────────────────────────────────────────────────────

@app.post("/coder/v3/project/{project_id}/format")
async def coder_v3_format(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import format_code
    body = await req.json()
    return format_code(project_id, file_path=body.get("file_path"))


# ── Test Runner ───────────────────────────────────────────────────────────────

@app.post("/coder/v3/project/{project_id}/test")
async def coder_v3_test(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import run_tests
    body = await req.json()
    return run_tests(project_id, target=body.get("target"))


# ── File Watcher ──────────────────────────────────────────────────────────────

@app.post("/coder/v3/project/{project_id}/watcher/start")
async def coder_v3_watcher_start(project_id: str) -> dict[str, Any]:
    from shared.coder_v3 import start_file_watcher
    return start_file_watcher(project_id)


@app.post("/coder/v3/project/{project_id}/watcher/stop")
async def coder_v3_watcher_stop(project_id: str) -> dict[str, Any]:
    from shared.coder_v3 import stop_file_watcher
    return stop_file_watcher(project_id)


# ── Git Enhancements ──────────────────────────────────────────────────────────

@app.get("/coder/v3/project/{project_id}/git/blame/{file_path:path}")
async def coder_v3_git_blame(project_id: str, file_path: str) -> dict[str, Any]:
    from shared.coder_v3 import git_blame
    return git_blame(project_id, file_path)


@app.get("/coder/v3/project/{project_id}/git/history/{file_path:path}")
async def coder_v3_git_history(project_id: str, file_path: str, n: int = 20) -> dict[str, Any]:
    from shared.coder_v3 import git_file_history
    return git_file_history(project_id, file_path, n=n)


@app.post("/coder/v3/project/{project_id}/git/stash")
async def coder_v3_git_stash(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import git_stash
    body = await req.json()
    return git_stash(project_id, message=body.get("message"))


@app.get("/coder/v3/project/{project_id}/git/stash")
async def coder_v3_git_stash_list(project_id: str) -> dict[str, Any]:
    from shared.coder_v3 import git_stash_list
    return git_stash_list(project_id)


@app.post("/coder/v3/project/{project_id}/git/stash/pop")
async def coder_v3_git_stash_pop(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import git_stash_pop
    body = await req.json()
    return git_stash_pop(project_id, stash=body.get("stash", "stash@{0}"))


@app.get("/coder/v3/project/{project_id}/git/remote")
async def coder_v3_git_remote(project_id: str) -> dict[str, Any]:
    from shared.coder_v3 import git_remote_list
    return git_remote_list(project_id)


@app.post("/coder/v3/project/{project_id}/git/remote")
async def coder_v3_git_remote_add(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import git_remote_add
    body = await req.json()
    return git_remote_add(project_id, body.get("name", "origin"), body.get("url", ""))


@app.get("/coder/v3/project/{project_id}/git/tag")
async def coder_v3_git_tag(project_id: str) -> dict[str, Any]:
    from shared.coder_v3 import git_tag_list
    return git_tag_list(project_id)


@app.post("/coder/v3/project/{project_id}/git/tag")
async def coder_v3_git_tag_create(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import git_create_tag
    body = await req.json()
    return git_create_tag(project_id, body.get("tag", ""), message=body.get("message", ""))


@app.get("/coder/v3/project/{project_id}/git/diff-branches")
async def coder_v3_git_diff_branches(project_id: str, branch_a: str = "main", branch_b: str = "HEAD") -> dict[str, Any]:
    from shared.coder_v3 import git_diff_branches
    return git_diff_branches(project_id, branch_a, branch_b)


# ── Diff Engine ───────────────────────────────────────────────────────────────

@app.get("/coder/v3/project/{project_id}/diff")
async def coder_v3_diff(project_id: str, file_a: str, file_b: str) -> dict[str, Any]:
    from shared.coder_v3 import diff_files
    return diff_files(project_id, file_a, file_b)


@app.get("/coder/v3/project/{project_id}/diff-versions")
async def coder_v3_diff_versions(project_id: str, file_path: str, ref_a: str = "HEAD", ref_b: str = "") -> dict[str, Any]:
    from shared.coder_v3 import diff_file_versions
    return diff_file_versions(project_id, file_path, ref_a, ref_b)


# ── Symbols & Outline ─────────────────────────────────────────────────────────

@app.get("/coder/v3/project/{project_id}/symbols/{file_path:path}")
async def coder_v3_symbols(project_id: str, file_path: str) -> dict[str, Any]:
    from shared.coder_v3 import extract_symbols
    return extract_symbols(project_id, file_path)


@app.get("/coder/v3/project/{project_id}/outline")
async def coder_v3_outline(project_id: str) -> dict[str, Any]:
    from shared.coder_v3 import project_outline
    return project_outline(project_id)


# ── Health Check ──────────────────────────────────────────────────────────────

@app.get("/coder/v3/project/{project_id}/health")
async def coder_v3_health(project_id: str) -> dict[str, Any]:
    from shared.coder_v3 import project_health
    return project_health(project_id)


# ── Snapshots ─────────────────────────────────────────────────────────────────

@app.post("/coder/v3/project/{project_id}/snapshot")
async def coder_v3_snapshot(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import create_snapshot
    body = await req.json()
    return create_snapshot(project_id, name=body.get("name"))


@app.get("/coder/v3/snapshots")
async def coder_v3_snapshots(project_id: str = "") -> dict[str, Any]:
    from shared.coder_v3 import list_snapshots
    return {"ok": True, "snapshots": list_snapshots(project_id or None)}


@app.post("/coder/v3/snapshot/{snap_id}/restore")
async def coder_v3_restore_snapshot(snap_id: str) -> dict[str, Any]:
    from shared.coder_v3 import restore_snapshot
    return restore_snapshot(snap_id)


# ── Permission System ─────────────────────────────────────────────────────────

@app.post("/coder/v3/permission/request")
async def coder_v3_permission_request(req: Request) -> dict[str, Any]:
    from shared.coder_v3 import request_permission
    body = await req.json()
    return request_permission(
        body.get("project_id", ""),
        body.get("operation", ""),
        body.get("reason", ""),
        body.get("requested_by", "user"),
    )


@app.get("/coder/v3/permissions")
async def coder_v3_permissions(project_id: str = "", status: str = "") -> dict[str, Any]:
    from shared.coder_v3 import list_permission_requests
    return {"ok": True, "requests": list_permission_requests(project_id or None, status or None)}


@app.post("/coder/v3/permission/{req_id}/approve")
async def coder_v3_permission_approve(req_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import approve_permission
    body = await req.json()
    return approve_permission(req_id, approved_by=body.get("approved_by", "admin"))


@app.post("/coder/v3/permission/{req_id}/deny")
async def coder_v3_permission_deny(req_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import deny_permission
    body = await req.json()
    return deny_permission(req_id, denied_by=body.get("denied_by", "admin"), reason=body.get("reason", ""))


# ── AI Inline Assistance ──────────────────────────────────────────────────────

@app.post("/coder/v3/project/{project_id}/ai/assist")
async def coder_v3_ai_assist(project_id: str, req: Request) -> dict[str, Any]:
    from shared.coder_v3 import ai_apply, ai_assist
    body = await req.json()
    # Phase 1.3 — Soul, Brain & Swarm: support action='apply' on the assist endpoint.
    if body.get("action", "explain") == "apply":
        return ai_apply(
            project_id,
            file_path=body.get("file_path", ""),
            ai_response=body.get("ai_response", ""),
        )
    return await ai_assist(
        project_id,
        action=body.get("action", "explain"),
        file_path=body.get("file_path", ""),
        selection=body.get("selection", ""),
        context=body.get("context", ""),
        model=body.get("model", ""),
    )


@app.post("/coder/v3/project/{project_id}/ai/apply")
async def coder_v3_ai_apply(project_id: str, req: Request) -> dict[str, Any]:
    # Phase 1.3 — Soul, Brain & Swarm: dedicated apply endpoint for AI-generated code.
    from shared.coder_v3 import ai_apply
    body = await req.json()
    return ai_apply(
        project_id,
        file_path=body.get("file_path", ""),
        ai_response=body.get("ai_response", ""),
    )


# ── Command Palette ───────────────────────────────────────────────────────────

@app.get("/coder/v3/command-palette")
async def coder_v3_command_palette(query: str = "") -> dict[str, Any]:
    from shared.coder_v3 import get_command_palette
    return {"ok": True, "commands": get_command_palette(query)}


# ── Lint & Dependencies ───────────────────────────────────────────────────────

@app.get("/coder/v3/project/{project_id}/lint/{file_path:path}")
async def coder_v3_lint(project_id: str, file_path: str) -> dict[str, Any]:
    from shared.coder_v3 import lint_file
    return lint_file(project_id, file_path)


@app.get("/coder/v3/project/{project_id}/dependencies")
async def coder_v3_dependencies(project_id: str) -> dict[str, Any]:
    from shared.coder_v3 import analyze_dependencies
    return analyze_dependencies(project_id)


# ── Git Import ────────────────────────────────────────────────────────────────

@app.post("/coder/v3/import-git")
async def coder_v3_import_git(req: Request) -> dict[str, Any]:
    from shared.coder_v3 import import_from_git
    body = await req.json()
    return import_from_git(body.get("url", ""), name=body.get("name"), branch=body.get("branch", "main"))


# ── Auto-complete ─────────────────────────────────────────────────────────────

@app.get("/coder/v3/project/{project_id}/completions")
async def coder_v3_completions(project_id: str, file_path: str, line: int = 1, column: int = 0) -> dict[str, Any]:
    from shared.coder_v3 import get_completions
    return get_completions(project_id, file_path, line, column)


# ── Audit Log ─────────────────────────────────────────────────────────────────

@app.get("/coder/v3/audit-log")
async def coder_v3_audit_log(project_id: str = "", limit: int = 100) -> dict[str, Any]:
    from shared.coder_v3 import get_audit_log
    return {"ok": True, "entries": get_audit_log(project_id or None, limit)}


# ═══════════════════════════════════════════════════════════════════════════════
# NEURAL AGENT — Self-Evolution Control Endpoints
# ═══════════════════════════════════════════════════════════════════════════════

@app.get("/neural-agent", response_class=HTMLResponse)
async def neural_agent_page() -> str:
    path = Path(__file__).resolve().parents[2] / "frontend" / "neural_agent.html"
    if path.exists():
        return path.read_text(encoding="utf-8")
    return "<h1>Neural Agent frontend not found</h1>"


@app.get("/api/neural-agent/model-status")
async def neural_agent_model_status() -> dict[str, Any]:
    from shared.neural_agent import get_model_status
    return await asyncio.to_thread(get_model_status)


@app.get("/api/neural-agent/proposals")
async def neural_agent_proposals(status: str = "", limit: int = 50) -> dict[str, Any]:
    from shared.neural_agent import list_all_proposals
    proposals = await asyncio.to_thread(list_all_proposals, status=status or None, limit=limit)
    return {"ok": True, "proposals": proposals}


@app.get("/api/neural-agent/proposals/{proposal_id}")
async def neural_agent_proposal_detail(proposal_id: str) -> dict[str, Any]:
    from shared.neural_agent import get_proposal_full
    proposal = await asyncio.to_thread(get_proposal_full, proposal_id)
    if not proposal:
        return {"ok": False, "error": "Proposal not found"}
    return {"ok": True, "proposal": proposal}


@app.post("/api/neural-agent/proposals/{proposal_id}/test")
async def neural_agent_proposal_test(proposal_id: str) -> dict[str, Any]:
    from shared.neural_agent import test_proposal
    return await asyncio.to_thread(test_proposal, proposal_id)


@app.post("/api/neural-agent/proposals/{proposal_id}/accept")
async def neural_agent_proposal_accept(proposal_id: str, req: Request) -> dict[str, Any]:
    from shared.neural_agent import accept_proposal
    body = await req.json()
    return await asyncio.to_thread(accept_proposal, proposal_id, body.get("reviewer", "user"), body.get("notes", ""))


@app.post("/api/neural-agent/proposals/{proposal_id}/reject")
async def neural_agent_proposal_reject(proposal_id: str, req: Request) -> dict[str, Any]:
    from shared.neural_agent import reject_proposal
    body = await req.json()
    return await asyncio.to_thread(reject_proposal, proposal_id, body.get("reviewer", "user"), body.get("notes", ""))


@app.post("/api/neural-agent/proposals/{proposal_id}/apply")
async def neural_agent_proposal_apply(proposal_id: str, req: Request) -> dict[str, Any]:
    from shared.neural_agent import apply_proposal
    body = await req.json()
    return await asyncio.to_thread(apply_proposal, proposal_id, body.get("approved_by", "user"), body.get("approval_phrase", ""))


@app.post("/api/neural-agent/generate")
async def neural_agent_generate(req: Request) -> dict[str, Any]:
    from shared.neural_agent import generate_proposal
    body = await req.json()
    return await asyncio.to_thread(generate_proposal,
        intent=body.get("intent", ""),
        file_path=body.get("file_path", ""),
        instructions=body.get("instructions", ""),
    )


@app.post("/api/neural-agent/reflect")
async def neural_agent_reflect() -> dict[str, Any]:
    from shared.neural_agent import run_reflection
    return await asyncio.to_thread(run_reflection)


class ChemQueryRequest(BaseModel):
    query: str
    topic: str = "general"


class ChemTrainRequest(BaseModel):
    fact: str
    topic: str = "general"
    validated_by: str = "human"


@app.post("/api/chem/chemdfm/query")
async def chem_chemdfm_query_endpoint(req: ChemQueryRequest) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "chem.chemdfm_query", {"query": req.query, "topic": req.topic})


@app.post("/api/chem/chemdfm/train")
async def chem_chemdfm_train_endpoint(req: ChemTrainRequest) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "chem.chemdfm_train", {"fact": req.fact, "topic": req.topic, "validated_by": req.validated_by})


@app.get("/api/chem/chemdfm/journal")
async def chem_chemdfm_journal_endpoint(mode: str = "summary", limit: int = 100) -> dict[str, Any]:
    return await asyncio.to_thread(agent_tools.run_tool, "chem.chemdfm_journal", {"mode": mode, "limit": limit})


# ═══════════════════════════════════════════════════════════════════════════════
# Phase 1.1 — Soul, Brain & Swarm: coder v3 file/run/iterate, swarm, self-index
# ═══════════════════════════════════════════════════════════════════════════════

@app.get("/coder/v3/project/{project_id}/file")
async def coder_v3_read_file(project_id: str, path: str = "") -> dict[str, Any]:
    """Read a single file from a coder v2/v3 project."""
    from shared.coder_v2 import read_file
    return read_file(project_id, path)


@app.post("/coder/v3/project/{project_id}/file")
async def coder_v3_write_file(project_id: str, req: Request) -> dict[str, Any]:
    """Write a single file in a coder v2/v3 project."""
    from shared.coder_v2 import write_file
    body = await req.json()
    file_path = str(body.get("path", ""))
    content = str(body.get("content", ""))
    return write_file(project_id, file_path, content)


@app.delete("/coder/v3/project/{project_id}/file")
async def coder_v3_delete_file(project_id: str, path: str = "") -> dict[str, Any]:
    """Delete a single file from a coder v2/v3 project."""
    from shared.coder_v2 import delete_file
    return delete_file(project_id, path)


@app.post("/coder/v3/project/{project_id}/run")
async def coder_v3_run_project(project_id: str, req: Request) -> dict[str, Any]:
    """Run a coder v2/v3 project."""
    from shared.coder_v2 import run_project
    body = await req.json()
    entry_file = body.get("entry_file")
    return run_project(project_id, entry_file)


class CoderV3IterateRequest(BaseModel):
    instruction: str
    provider: str | None = None
    model: str | None = None
    max_steps: int = 2


@app.post("/coder/v3/project/{project_id}/ai/iterate")
async def coder_v3_ai_iterate(project_id: str, req: CoderV3IterateRequest) -> dict[str, Any]:
    """Run the legacy coder iterate loop against a v2/v3 project."""
    from shared import coder
    return await coder.iterate(
        project_id,
        req.instruction,
        provider=req.provider,
        model=req.model,
        max_steps=req.max_steps,
    )

