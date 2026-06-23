from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Dict
from .settings import BASE_DIR, settings


def safe_name(name: str) -> str:
    return ''.join(c if c.isalnum() or c in ['-', '_'] else '_' for c in name)[:80] or 'shims_document'


def generated_path(filename: str) -> Path:
    out = settings.generated_dir if settings.generated_dir.is_absolute() else BASE_DIR / settings.generated_dir
    out.mkdir(parents=True, exist_ok=True)
    return out / filename


def generate_document(title: str, body: str, output_type: str = 'docx') -> Path:
    ext = output_type.lower().strip().lstrip('.')
    path = generated_path(f'{safe_name(title)}.{ext}')
    if ext == 'txt':
        path.write_text(title + '\n\n' + body, encoding='utf-8')
    elif ext == 'docx':
        from docx import Document
        doc = Document()
        doc.add_heading(title, 0)
        doc.add_paragraph(settings.company_name)
        for para in body.split('\n\n'):
            doc.add_paragraph(para)
        doc.save(path)
    elif ext == 'pdf':
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(str(path), pagesize=A4)
        w, h = A4
        y = h - 50
        c.setFont('Helvetica-Bold', 15)
        c.drawString(50, y, title[:80])
        y -= 25
        c.setFont('Helvetica', 10)
        for line in (body or '').splitlines():
            for chunk in [line[i:i+92] for i in range(0, len(line), 92)] or ['']:
                if y < 50:
                    c.showPage(); y = h - 50
                c.drawString(50, y, chunk)
                y -= 14
        c.save()
    elif ext == 'xlsx':
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = 'SHIMS'
        ws.append(['Title', title])
        for i, line in enumerate(body.splitlines(), 1):
            ws.append([i, line])
        wb.save(path)
    elif ext == 'pptx':
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body[:1000]
        prs.save(path)
    else:
        path = generated_path(f'{safe_name(title)}.txt')
        path.write_text(title + '\n\n' + body, encoding='utf-8')
    return path


def generate_coa_docx(record, template) -> Path:
    from docx import Document
    schema = json.loads(template.schema_json or '{"fields":[]}')
    values = json.loads(record.results_json or '{}')
    path = generated_path(f'COA_{safe_name(record.batch_no)}_{record.id}.docx')
    doc = Document()
    doc.add_heading('Certificate of Analysis', 0)
    doc.add_paragraph(settings.company_name)
    doc.add_paragraph(settings.company_address)
    doc.add_paragraph(f'GST: {settings.company_gst} | Email: {settings.company_email} | Phone: {settings.company_phone}')
    doc.add_paragraph(f'Product: {record.product_name}')
    doc.add_paragraph(f'Batch No: {record.batch_no}')
    doc.add_paragraph(f'Sample ID: {record.sample_id}')
    table = doc.add_table(rows=1, cols=4)
    headers = table.rows[0].cells
    headers[0].text = 'Parameter'; headers[1].text = 'Specification'; headers[2].text = 'Result'; headers[3].text = 'Method'
    for field in schema.get('fields', []):
        row = table.add_row().cells
        name = field.get('name')
        result = values.get(name, '')
        if field.get('unit') and result != '':
            result = f'{result} {field.get("unit")}'
        row[0].text = str(field.get('label', name))
        row[1].text = str(field.get('spec', ''))
        row[2].text = str(result)
        row[3].text = str(field.get('method', ''))
    doc.add_paragraph('Prepared by QC. Human approval required before regulated release.')
    doc.save(path)
    return path
