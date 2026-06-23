from __future__ import annotations

import asyncio
import json
import os
import re
import time
import wave
import math
from pathlib import Path
from typing import Any, Dict

import httpx
from fastapi import FastAPI, Request
from fastapi.responses import StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

ROOT = Path.home() / "shims_mobile_runtime"
MEDIA = ROOT / "media"
for sub in ["images", "pdf", "ppt", "audio", "video"]:
    (MEDIA / sub).mkdir(parents=True, exist_ok=True)

app = FastAPI(title="SHIMS Android Offline-Lite Runtime", version="0.11.0")
app.mount("/media", StaticFiles(directory=str(MEDIA)), name="media")

class Turn(BaseModel):
    message: str | None = None
    input: str | None = None
    provider: str = "auto"
    model: str = "local-gguf"
    locale: str = "en-IN"
    source: str = "android"

class MediaReq(BaseModel):
    kind: str | None = None
    type: str | None = None
    prompt: str = "SHIMS generated artifact"

@app.get("/health")
def health() -> Dict[str, Any]:
    return {"ok": True, "name": "SHIMS Android Offline-Lite Runtime", "mode": "termux", "media": str(MEDIA)}

@app.get("/api/v11/models")
def models() -> Dict[str, Any]:
    return {
        "ok": True,
        "mode": "termux-offline-lite",
        "installed": ["local-gguf via llama.cpp server at 127.0.0.1:8080"],
        "recommended": [
            {"name": "Qwen2.5-0.5B/1.5B GGUF", "role": "phone offline fast"},
            {"name": "Llama-3.2-1B/3B GGUF", "role": "phone offline fallback"},
            {"name": "qwen2.5:7b", "role": "Predator/Ollama host"},
        ],
        "note": "Official Ollama is not packaged here; phone-only local inference uses llama.cpp GGUF. Desktop SHIMS uses Ollama pull/list/select."
    }

@app.post("/api/v11/models/pull")
async def pull(req: Request):
    body = await req.json()
    model = body.get("model", "")
    async def gen():
        yield json.dumps({"status": "received", "model": model}) + "\n"
        yield json.dumps({"status": "android_note", "message": "Ollama pull runs on Predator/Desktop SHIMS. For phone-only offline, download a GGUF model and run llama.cpp server."}) + "\n"
    return StreamingResponse(gen(), media_type="application/x-ndjson")

def _intent(text: str) -> str | None:
    t = text.lower()
    if any(w in t for w in ["image", "tasveer", "photo", "picture", "चित्र"]): return "image"
    if any(w in t for w in ["pdf", "document", "invoice", "coa", "report"]): return "pdf"
    if any(w in t for w in ["ppt", "powerpoint", "presentation", "deck"]): return "ppt"
    if any(w in t for w in ["audio", "sound", "voice note", "speak file"]): return "audio"
    if any(w in t for w in ["video", "movie", "clip"]): return "video"
    return None

@app.post("/api/v11/chat/turn")
async def chat_turn(turn: Turn):
    text = (turn.message or turn.input or "").strip()
    kind = _intent(text)
    async def stream():
        if not text:
            yield json.dumps({"type":"done","answer":"I did not receive text. Speak again or type your command.","provider":"termux","model":turn.model}) + "\n"
            return
        if kind:
            artifact = await generate_artifact(kind, text)
            yield json.dumps({"type":"done","answer":f"Created {kind} locally on Android offline-lite runtime.","provider":"tool","model":"termux","media_result":artifact}) + "\n"
            return
        llm = await call_llama_cpp(text)
        yield json.dumps({"type":"done","answer":llm,"provider":"termux-llama.cpp" if llm.startswith("LLM:") else "offline-fallback","model":turn.model,"emotion":"calm"}) + "\n"
    return StreamingResponse(stream(), media_type="application/x-ndjson")

@app.post("/media/generate")
async def media_generate(req: MediaReq):
    kind = (req.kind or req.type or "image").lower()
    return await generate_artifact(kind, req.prompt)

async def call_llama_cpp(prompt: str) -> str:
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            r = await client.post("http://127.0.0.1:8080/completion", json={"prompt": prompt, "n_predict": 256, "temperature": 0.4})
            if r.status_code == 200:
                data = r.json()
                return "LLM: " + (data.get("content") or data.get("response") or "").strip()
    except Exception:
        pass
    return "I am running in Android offline-lite mode. Connect to Predator SHIMS for full Omni brain, or start llama.cpp server in Termux for phone-local LLM chat."

async def generate_artifact(kind: str, prompt: str) -> Dict[str, Any]:
    slug = re.sub(r"[^a-zA-Z0-9]+", "_", prompt.lower()).strip("_")[:48] or "shims"
    ts = int(time.time())
    if kind == "pdf": return create_pdf(slug, ts, prompt)
    if kind == "ppt": return create_ppt(slug, ts, prompt)
    if kind == "audio": return create_wav(slug, ts, prompt)
    if kind == "video": return create_video_note(slug, ts, prompt)
    return create_image(slug, ts, prompt)

def create_image(slug: str, ts: int, prompt: str) -> Dict[str, Any]:
    from PIL import Image, ImageDraw, ImageFont
    path = MEDIA / "images" / f"{slug}_{ts}.png"
    img = Image.new("RGB", (1024, 1024), (8, 14, 28))
    d = ImageDraw.Draw(img)
    d.rectangle([40,40,984,984], outline=(77,227,255), width=4)
    d.text((70,80), "SHIMS Android Image", fill=(255,255,255))
    d.text((70,140), prompt[:220], fill=(220,240,255))
    d.text((70,930), "Offline-lite placeholder. Use Predator Stable Diffusion for real AI images.", fill=(255,200,100))
    img.save(path)
    return {"ok": True, "kind": "image", "url": f"/media/images/{path.name}", "note": "local placeholder image"}

def create_pdf(slug: str, ts: int, prompt: str) -> Dict[str, Any]:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
    path = MEDIA / "pdf" / f"{slug}_{ts}.pdf"
    c = canvas.Canvas(str(path), pagesize=A4)
    w,h = A4
    c.setFont("Helvetica-Bold", 18); c.drawString(42, h-60, "SHIMS Android PDF")
    c.setFont("Helvetica", 11); c.drawString(42, h-85, "J.K. Lifecare Centers Pvt. Ltd.")
    y=h-125
    for line in [prompt[i:i+92] for i in range(0,len(prompt),92)][:36]:
        c.drawString(42, y, line); y-=16
    c.save()
    return {"ok": True, "kind": "pdf", "url": f"/media/pdf/{path.name}"}

def create_ppt(slug: str, ts: int, prompt: str) -> Dict[str, Any]:
    from pptx import Presentation
    path = MEDIA / "ppt" / f"{slug}_{ts}.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "SHIMS Android PPT"
    slide.placeholders[1].text = prompt
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Generated by SHIMS"
    slide2.placeholders[1].text = "Offline-lite PPT output. Use Enterprise Document Studio for branded production output."
    prs.save(path)
    return {"ok": True, "kind": "ppt", "url": f"/media/ppt/{path.name}"}

def create_wav(slug: str, ts: int, prompt: str) -> Dict[str, Any]:
    path = MEDIA / "audio" / f"{slug}_{ts}.wav"
    rate=22050; dur=1.5; freq=440
    with wave.open(str(path), 'w') as wf:
        wf.setnchannels(1); wf.setsampwidth(2); wf.setframerate(rate)
        for i in range(int(rate*dur)):
            val=int(32767*0.18*math.sin(2*math.pi*freq*i/rate))
            wf.writeframesraw(val.to_bytes(2,'little',signed=True))
    return {"ok": True, "kind": "audio", "url": f"/media/audio/{path.name}", "note":"tone placeholder; Android TTS speaks responses natively"}

def create_video_note(slug: str, ts: int, prompt: str) -> Dict[str, Any]:
    path = MEDIA / "video" / f"{slug}_{ts}.txt"
    path.write_text("SHIMS Android video storyboard\n\n" + prompt + "\n\nInstall/use FFmpeg or Predator SHIMS for real MP4 generation.\n", encoding='utf-8')
    return {"ok": True, "kind": "video", "url": f"/media/video/{path.name}", "note":"storyboard fallback"}
